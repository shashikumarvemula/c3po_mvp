"""
PowerPoint Chart Legend Fix Script

This script fixes the issue where PowerPoint charts show a legend entry for every
week instead of just showing the series names (mx_claims, median, lower_bound, upper_bound).

Usage:
    python fix_chart_legend.py <path_to_powerpoint_file> <slide_number> <chart_title>

Example:
    python fix_chart_legend.py output.pptx 5 "Weekly Claims Outlier Detection"
"""

import sys
from pptx import Presentation
from pptx.chart.data import CategoryChartData
import pandas as pd

def fix_chart_legend(pptx_path, slide_number, chart_title):
    """
    Fix the chart legend by ensuring proper series/category structure.
    
    Args:
        pptx_path: Path to the PowerPoint file
        slide_number: Slide number (1-indexed)
        chart_title: Title of the chart to fix
    """
    print(f"Opening PowerPoint file: {pptx_path}")
    prs = Presentation(pptx_path)
    
    # Get the slide
    slide_idx = slide_number - 1
    if slide_idx >= len(prs.slides):
        print(f"❌ Error: Slide {slide_number} does not exist")
        return False
    
    slide = prs.slides[slide_idx]
    print(f"✅ Found slide {slide_number}")
    
    # Find the chart
    chart_shape = None
    for shape in slide.shapes:
        if shape.has_chart:
            try:
                if shape.chart.chart_title.text_frame.text == chart_title:
                    chart_shape = shape
                    break
            except:
                pass
    
    if chart_shape is None:
        print(f"❌ Error: Chart with title '{chart_title}' not found on slide {slide_number}")
        return False
    
    chart = chart_shape.chart
    print(f"✅ Found chart: {chart_title}")
    
    # Get the embedded Excel data
    try:
        import io
        from openpyxl import load_workbook
        
        chart_part = chart.part
        xlsx_part = None
        
        for rel_id in list(chart_part.rels.keys()):
            try:
                rel = chart_part.rels[rel_id]
                if not rel.is_external and hasattr(rel, 'target_part'):
                    part = rel.target_part
                    if hasattr(part, 'blob'):
                        try:
                            test_wb = load_workbook(io.BytesIO(part.blob))
                            xlsx_part = part
                            test_wb.close()
                            break
                        except:
                            continue
            except:
                continue
        
        if not xlsx_part:
            print("❌ Error: Could not find embedded Excel data")
            return False
        
        # Read the embedded Excel
        wb = load_workbook(io.BytesIO(xlsx_part.blob))
        ws = wb.active
        
        print(f"📊 Current Excel structure:")
        print(f"   Rows: {ws.max_row}, Columns: {ws.max_column}")
        
        # Read headers
        headers = []
        for col in range(1, ws.max_column + 1):
            cell_value = ws.cell(1, col).value
            headers.append(str(cell_value) if cell_value else f"Col{col}")
        
        print(f"   Headers: {headers}")
        
        # Check if data needs to be transposed
        # If we have many columns (one per week), we need to transpose
        if ws.max_column > 10:  # Likely has week columns
            print("⚠️  Detected transposed data structure - fixing...")
            
            # Read all data
            data = []
            for row in range(1, ws.max_row + 1):
                row_data = []
                for col in range(1, ws.max_column + 1):
                    row_data.append(ws.cell(row, col).value)
                data.append(row_data)
            
            # Clear worksheet
            ws.delete_rows(1, ws.max_row)
            
            # Write transposed data
            # New structure: Column A = weeks, Column B = mx_claims, Column C = median, etc.
            # Old structure: Row 1 = headers, Row 2 = mx_claims, Row 3 = median, etc.
            
            # Assuming old structure had:
            # Row 1: Week headers
            # Row 2: mx_claims values
            # Row 3: median values
            # Row 4: lower_bound values
            # Row 5: upper_bound values
            
            num_weeks = len(data[0]) - 1  # Exclude first column (series name)
            
            # Write new headers
            ws.cell(1, 1, "week_end_date")
            ws.cell(1, 2, data[1][0] if len(data) > 1 else "mx_claims")
            ws.cell(1, 3, data[2][0] if len(data) > 2 else "median")
            ws.cell(1, 4, data[3][0] if len(data) > 3 else "lower_bound")
            ws.cell(1, 5, data[4][0] if len(data) > 4 else "upper_bound")
            
            # Write transposed data
            for week_idx in range(num_weeks):
                row_num = week_idx + 2
                ws.cell(row_num, 1, data[0][week_idx + 1])  # Week
                if len(data) > 1:
                    ws.cell(row_num, 2, data[1][week_idx + 1])  # mx_claims
                if len(data) > 2:
                    ws.cell(row_num, 3, data[2][week_idx + 1])  # median
                if len(data) > 3:
                    ws.cell(row_num, 4, data[3][week_idx + 1])  # lower_bound
                if len(data) > 4:
                    ws.cell(row_num, 5, data[4][week_idx + 1])  # upper_bound
            
            # Save back to embedded Excel
            output = io.BytesIO()
            wb.save(output)
            output.seek(0)
            xlsx_part._blob = output.read()
            wb.close()
            
            print(f"✅ Transposed data: {num_weeks} weeks, 4 series")
            
            # Now recreate the chart data
            chart_data = CategoryChartData()
            
            # Read the weeks for categories
            weeks = []
            for row in range(2, num_weeks + 2):
                week_val = ws.cell(row, 1).value
                weeks.append(str(week_val) if week_val else f"Week{row-1}")
            
            chart_data.categories = weeks
            
            # Add series (read from the transposed data we just wrote)
            series_names = ["mx_claims", "median", "lower_bound", "upper_bound"]
            for col_idx, series_name in enumerate(series_names, start=2):
                values = []
                for row in range(2, num_weeks + 2):
                    val = ws.cell(row, col_idx).value
                    values.append(float(val) if val is not None else None)
                chart_data.add_series(series_name, values)
                print(f"   Added series: {series_name}")
            
            # Replace chart data
            chart.replace_data(chart_data)
            print("✅ Chart data structure fixed")
        else:
            print("✅ Data structure looks correct (no transpose needed)")
        
        # Save the presentation
        output_path = pptx_path.replace('.pptx', '_fixed.pptx')
        prs.save(output_path)
        print(f"✅ Saved fixed presentation to: {output_path}")
        
        return True
        
    except Exception as e:
        print(f"❌ Error: {e}")
        import traceback
        traceback.print_exc()
        return False

if __name__ == "__main__":
    if len(sys.argv) < 4:
        print("Usage: python fix_chart_legend.py <pptx_file> <slide_number> <chart_title>")
        print('Example: python fix_chart_legend.py output.pptx 5 "Weekly Claims Outlier Detection"')
        sys.exit(1)
    
    pptx_file = sys.argv[1]
    slide_num = int(sys.argv[2])
    chart_title = sys.argv[3]
    
    success = fix_chart_legend(pptx_file, slide_num, chart_title)
    sys.exit(0 if success else 1)
