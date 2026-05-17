import os
import io
import chainlit as cl
from PyPDF2 import PdfReader
from pptx import Presentation
from docx import Document
import pandas as pd
import openpyxl
import csv


class FileReader:
    """Class to read various file types and extract text content."""
    def __init__(self):
        print("FileReader initialized")
        self.file_path = None
        self.extension = None
        self.file_type = None
        self.file_content = None
        self.file_name = None
        self.file_size = None
        self.file_type = None

    async def extract_text_from_file(self, file_path, extension):
        """Extract text content from various file types"""
        print(f"Extracting text from file with extension {extension}...")
        
        # Text files
        if extension in ['.txt', '.md', '.json', '.html', '.xml', '.log']:
            try:
                print("Extracting text from text file...")
                with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                    return f.read()
            except Exception as e:
                cl.error(f"Error reading text file: {str(e)}")
                return f"Error extracting text from {extension} file: {str(e)}"
        
        # CSV files - separate handling for better error management
        elif extension == '.csv':
            try:
                print("Extracting text from CSV file...")
                df = pd.read_csv(file_path)
                return df.to_string()
            except pd.errors.EmptyDataError:
                return "The CSV file is empty."
            except pd.errors.ParserError:
                try:
                    # Fallback to basic CSV reading
                    text = ""
                    with open(file_path, 'r', encoding='utf-8', errors='replace') as csvfile:
                        reader = csv.reader(csvfile)
                        for row in reader:
                            text += ",".join(row) + "\n"
                    return text
                except Exception as e:
                    cl.error(f"Error reading CSV file with basic reader: {str(e)}")
                    return f"Error extracting text from CSV file: {str(e)}"
            except Exception as e:
                cl.error(f"Error reading CSV file: {str(e)}")
                return f"Error extracting text from CSV file: {str(e)}"
        
        # PDF files
        elif extension == '.pdf':
            try:
                print("Extracting text from PDF file...")
                text = ""
                with open(file_path, "rb") as f:
                    pdf = PdfReader(f)
                    for page in pdf.pages:
                        text += page.extract_text() + "\n"
                return text if text.strip() else "PDF contained no extractable text (may be scanned images)"
            except Exception as e:
                cl.error(f"Error reading PDF file: {str(e)}")
                return f"Error extracting text from PDF file: {str(e)}"

        # Word documents
        elif extension in ['.docx', '.doc']:
            try:
                print("Extracting text from Word document...")
                doc = Document(file_path)
                return "\n".join([paragraph.text for paragraph in doc.paragraphs])
            except Exception as e:
                cl.error(f"Error reading Word document: {str(e)}")
                return f"Error extracting text from Word document: {str(e)}"
        
        # PowerPoint presentations
        elif extension in ['.pptx', '.ppt']:
            try:
                print("Extracting text from PowerPoint presentation...")
                pres = Presentation(file_path)
                text = ""
                for slide in pres.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            text += shape.text + "\n"
                return text
            except Exception as e:
                cl.error(f"Error reading PowerPoint presentation: {str(e)}")
                return f"Error extracting text from PowerPoint presentation: {str(e)}"
        
        # Excel files
        elif extension in ['.xlsx', '.xls']:
            try:
                print("Extracting text from Excel file...")
                df = pd.read_excel(file_path)
                return df.to_string()
            except Exception as e:
                try:
                    # Fallback for complex Excel files
                    print("Trying alternate Excel reading method...")
                    workbook = openpyxl.load_workbook(file_path, data_only=True)
                    text = ""
                    for sheet_name in workbook.sheetnames:
                        sheet = workbook[sheet_name]
                        text += f"\n--- Sheet: {sheet_name} ---\n"
                        for row in sheet.iter_rows(values_only=True):
                            text += "\t".join([str(cell) if cell is not None else "" for cell in row]) + "\n"
                    return text
                except Exception as inner_e:
                    cl.error(f"Error reading Excel file: {str(e)} and fallback failed: {str(inner_e)}")
                    return f"Error extracting text from Excel file: {str(e)}"
        
        # Return message for unsupported file types
        else:
            message = f"Unsupported file type: {extension}"
            cl.error(message)
            return message
    
    async def read(self, file_path, extension):
        """Read the file and return its content as a string."""
        try:
            content = await self.extract_text_from_file(file_path, extension)
            if not content or content.strip() == "":
                return "The file appears to be empty or contains no extractable text."
            return content
        except Exception as e:
            error_message = f"Error reading file: {str(e)}"
            cl.error(error_message)
            return error_message