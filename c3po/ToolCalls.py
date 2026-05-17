from Semi_Structured_Bot.opensearch_execution import OpensearchExecutionManager
from Structured_Bot.agent_tools import SalesPerformanceMetrics
import asyncio
import json
import time
import pandas as pd
import datetime
import io
import chainlit as cl
from Structured_Bot.helper import FileNameExtractor
import ast
from collections import Counter
import re
import os
from Structured_Bot.llm_requests_new  import LLM
from Files_Images_Handling import ElementsHandling
files_images_handler=ElementsHandling()
from openpyxl import load_workbook


# Abstract base class for all tools
class BaseTool:
    """Base class that defines the interface for all tool classes"""
    def __init__(self):
        """Initialize the base tool"""
        self.result = None
    
    async def run(self, *args, **kwargs):
        """Execute the tool's functionality"""
        pass


class RunHybridSearch(BaseTool):
    """Tool for executing hybrid search operations"""
    
    def __init__(self, opensearch_manager):
        """Initialize with an OpensearchExecutionManager instance"""
        self.opensearch = opensearch_manager

        
    async def run(self, tool_params):
        arguments = tool_params.get("arguments")
        user_message = tool_params.get("user_message")
        """Execute a hybrid search operation"""
        print("Inside RunHybridSearch.run function")
        print("arguments: ", arguments)
        search_query = arguments.get("query")
        print("search_query: ", search_query)
        if cl.user_session.get("chat_profile", "") == os.getenv("BOT_TYPE_PMR"):
            user_message=cl.user_session.get('tool_params').get('user_question','')
            function_response=await self.opensearch.pmr_hybrid_search(user_message, search_query)
            print("function_response: ", function_response)
            dataframe=pd.DataFrame(function_response)
            if dataframe.empty:
                query_to_ui=search_query
                query_to_ui['query']['bool']['must'][1]['knn']['pmr_embedding_pii_removal_profile_cosine_updated']['vector'] = 'question_embedding'
                return query_to_ui, "EMPTY DATA Retrieved. can you please call 'run_hybrid_search' again without any changes", dataframe.info()
            else:
                for index in dataframe.index:
                    dataframe.loc[index,"combined_text"]=f""" Interviewer Question: {dataframe.loc[index,'questions']} \n Respondent Answer: {dataframe.loc[index,'answers']} """
                    tool_params["result"]=dataframe
                    print("dataframe at pmr: ",dataframe.head())
                query_to_ui=search_query
                query_to_ui['query']['bool']['must'][1]['knn']['pmr_embedding_pii_removal_profile_cosine_updated']['vector'] = 'question_embedding'
                return query_to_ui, "DATA RETRIEVED SUCCESSFULLY. Use text_processor tool to extract insights from the data.", dataframe
        elif cl.user_session.get("chat_profile", "") == os.getenv("BOT_TYPE_MARKET_MAP"):
            function_response=await self.opensearch.market_map_hybrid_search(user_message, search_query)
            print("function_response: ", function_response)
            dataframe=pd.DataFrame(function_response)
            for index in dataframe.index:
                dataframe.loc[index,"combined_text"]=f""" Interviewer Question: {dataframe.loc[index,'questions']} \n Respondent Answer: {dataframe.loc[index,'answers']} """
                tool_params["result"]=dataframe
                print("dataframe at pmr: ",dataframe.head())
            search_query['query']['bool']['must'][1]['knn']['mrkt_map_embedding']['vector'] = 'question_embedding'
            return search_query, "DATA RETRIEVED SUCCESSFULLY. Use text_processor tool to extract insights from the data.", dataframe

        elif cl.user_session.get("chat_profile", "") == os.getenv("BOT_TYPE_EARLY_EXP"):
            function_response=await self.opensearch.early_exp_hybrid_search(user_message, search_query)
            print("function_response: ", function_response)
            dataframe=pd.DataFrame(function_response)
            for index in dataframe.index:
                dataframe.loc[index,"combined_text"]=f""" Interviewer Question: {dataframe.loc[index,'questions']} \n Respondent Answer: {dataframe.loc[index,'answers']} """
                tool_params["result"]=dataframe
                print("dataframe at pmr: ",dataframe.head())
            search_query['query']['bool']['must'][1]['knn']['init_expr_embedding_v2_ingest']['vector'] = 'question_embedding'
            return search_query, "DATA RETRIEVED SUCCESSFULLY. Use text_processor tool to extract insights from the data.", dataframe

        elif cl.user_session.get("chat_profile", "") == os.getenv("BOT_TYPE_PBC_MARKET_RESEARCH"):
            function_response=await self.opensearch.pbc_market_research_hybrid_search(user_message, search_query)
            print("function_response: ", function_response)
            dataframe=pd.DataFrame(function_response)
            for index in dataframe.index:
                dataframe.loc[index,"combined_text"]=f""" Interviewer Question: {dataframe.loc[index,'questions']} \n Respondent Answer: {dataframe.loc[index,'answers']} """
                tool_params["result"]=dataframe
                print("dataframe at pbc market research: ",dataframe.head())
            search_query['query']['bool']['must'][1]['knn']['pbc_market_research_embedding']['vector'] = 'question_embedding'
            return search_query, "DATA RETRIEVED SUCCESSFULLY. Use text_processor tool to extract insights from the data.", dataframe

        else:
            function_response = await self.opensearch.execute_hybrid_search(user_message, search_query)
            print("function_response: ", function_response)
            return search_query, function_response, None


class RunSqlOnOpensearch(BaseTool):
    """Tool for executing SQL queries on OpenSearch"""
    
    def __init__(self, opensearch_manager):
        """Initialize with an OpensearchExecutionManager instance"""
        self.opensearch = opensearch_manager
        
    async def run(self, tool_params):
        arguments = tool_params.get("arguments")
        user_message = tool_params.get("user_message")
        """Execute an SQL query on OpenSearch"""
        print("Inside RunSqlOnOpensearch.run function")
        print("arguments: ", arguments)
        sql_query = arguments.get("query")
        print("search_query: ", sql_query)
        function_response = await self.opensearch.execute_sql_on_opensearch(user_message, sql_query)
        print("function_response: ", function_response)
        return sql_query, function_response, None


class CalculateDateRange(BaseTool):
    """Tool for calculating date ranges for different metrics"""
    
    def __init__(self):
        """Initialize the CalculateDateRange tool"""
        pass
        
    async def run(self, tool_params):
        arguments = tool_params.get("arguments")
        """Calculate date ranges based on the provided weekend date"""
        print("Inside CalculateDateRange.run called...")
        formatted_latest_weekend_date = arguments.get("query")
        date_ranges = SalesPerformanceMetrics.calculate_date_ranges(formatted_latest_weekend_date)

        m12_dates_dic = {
            "r12m_start_date": date_ranges["r12m_start_date"],
            "r12m_end_date": date_ranges["r12m_end_date"],
            "p12m_start_date": date_ranges["p12m_start_date"],
            "p12m_end_date": date_ranges["p12m_end_date"],
        }

        w4_dates_dic = {
            "start_r4w": date_ranges["start_r4w"],
            "end_r4w": date_ranges["end_r4w"],
            "start_p4w": date_ranges["start_p4w"],
            "end_p4w": date_ranges["end_p4w"],
        }

        w13_dates_dic = {
            "start_r13w": date_ranges["start_r13w"],
            "end_r13w": date_ranges["end_r13w"],
            "start_p13w": date_ranges["start_p13w"],
            "end_p13w": date_ranges["end_p13w"],
        }

        qtd_dates_dic = {key: value for key, value in date_ranges.items() if key.startswith("start_q") or key.startswith("end_q")}

        function_response = f"""Here are the different Date Ranges for different metrics:
        {m12_dates_dic} \n
        {w4_dates_dic} \n
        {w13_dates_dic} \n
        {qtd_dates_dic} \n

        Remember that R1M(Recent 1 month) is same as R4W and P1M is similar to P4W
        AND R3M is similar to R13M and P3M is similar to P13M and Q1TD(quarter 1 to TO(latest) Date) is same as Q2TD,Q3TD,Q4TD.
        """
        return None, function_response, None

class GenerateSQLCode(BaseTool):
    """Tool for generating SQL code"""
    def __init__(self):
        """Initialize the GenerateSQLCode tool"""
        pass
    async def run(self, tool_params):
        cl.user_session.set("coding_language", "Python")
        arguments = tool_params.get("arguments")
        dataware_house = tool_params.get("dataware_house")
        sql_code = arguments.get("query")
        cl.user_session.set("sql_code", sql_code)
        print("sql_code is",sql_code)
        if sql_code:
            try:
                warehouse = cl.make_async(dataware_house.get_data_from_wareshouse)
                result = await warehouse(sql_code)
                tool_params = cl.user_session.get("tool_params")
                tool_params["result"] = result
                cl.user_session.set("tool_params",tool_params)
                print("Data from sql warehouse")
                if result.empty:
                    raise Exception("Empty DataFrame")
                print(result.head())
                print(result.tail())
                buffer = io.StringIO()
                result.info(buf=buffer)
                info_string = buffer.getvalue()
                print("result information:", info_string)
                #result.to_csv("Structured_Bot/output_files/result.csv")
                
                function_response = f"""SQL EXECUTED SUCCESSFULLY.
                    Dataset: 'result' dataframe with {len(result)} rows
                    - result.head(): {result.head()}
                    - result.tail() : {result.tail()}
                    - result.info(): {info_string}
                    - Date format: Already in string format (use directly)
                    Use python_after_sql tool. NEVER use any return statement in the python code.
                    """ 
                           
                return sql_code, function_response, result
            except Exception as e:
                function_response=f"""SQL QUERY RESULT : {e}, PLEASE RETRY. Modifying the sql query"""
                print("Error while executing SQL code:", e)
                cl.user_session.set("coding_language", "SQL")
                return sql_code, function_response, ''
                      

class PythonAfterSQL(BaseTool): 
    """Tool for executing Python code after SQL execution"""
    def __init__(self):
        """Initialize the PythonAfterSQL tool"""
        pass
    async def run(self, tool_params):
        cl.user_session.set("coding_language", "SQL")
        namespace = tool_params["namespace"]
        namespace["result"] = tool_params["result"]
        result = tool_params["result"]
        print("result is",result)
        print(result)
        arguments = tool_params.get("arguments")
        python_code = arguments.get("query")
        print("python_code is",python_code)
        tool_params["python_code"] = python_code
        cl.user_session.set("tool_params",tool_params)
        if python_code:
            try:
                exec(python_code,namespace,namespace)
                if cl.user_session.get("chat_profile")==os.getenv("BOT_TYPE_PMR"):
                    return python_code, "Task completed successfully,Now call insights_response to generate insights", result
                return python_code, "Task completed successfully", result
            except Exception as e:
                function_response=f"Here is the error : {e}"
                cl.user_session.set("coding_language", "Python")
                return python_code,function_response,""
            
class TextProcessor(BaseTool):

    
    def extract_insights_simple(self,insight_draft):
        try:
            # Clean the string first
            cleaned = insight_draft.replace('\n', ' ').replace('\r', ' ').replace('\t', ' ')
            
            # Extract JSON part
            start = cleaned.find("{")
            end = cleaned.rfind("}") + 1
            
            if start != -1 and end != 0:
                json_part = cleaned[start:end]
                print("insights draft after processing", json_part)
                return json.loads(json_part)
        except:
            pass
        
        # Fallback to manual extraction
        start_marker = '"insights_draft": "'
        if start_marker in insight_draft:
            start_idx = insight_draft.find(start_marker) + len(start_marker)
            end_idx = insight_draft.rfind('"}')
            if end_idx != -1:
                content = insight_draft[start_idx:end_idx]
                return {"insights_draft": content}
        print("insights draft after processing", insight_draft)
        return {"insights_draft": insight_draft}

    
    async def get_insight_draft(self):
        system_prompt=self.insight_draft_propmt
        relevant_chunks=self.data.to_json(orient='index')
        user_content = f"Here are the user question : {self.question}\n and here are the relevant chunks : {relevant_chunks}"
        insight_draft,email,total_tokens_conv,input_tokens_conv,output_tokens_conv=await self.request_llm.send_request_non_streaming(system_prompt,user_content,self.email,self.total_tokens_conv,self.input_tokens_conv,self.output_tokens_conv)
        # insight_draft = insight_draft[insight_draft.find("{"):insight_draft.rfind('}')+1].replace("\n"," ")

        print("insights draft before processing: ",insight_draft)
        return self.extract_insights_simple(insight_draft)
        # print("insights : ",insight_draft)
        # print("\n\n\n\n")
        # return insight_draft

    async def get_topics(self):
        # for index in list(self.data.index):
        #     full_text=self.data.loc[index,self.column]
        #     print("index: ",index)
        #     if not pd.isna(self.data.loc[index,self.column]):
        #         user_content=f"Here are the user question : {self.question }\n and text : {full_text}"
        #         print("content before sending to llm for topic: ",user_content)
        #         rel_topic,email,total_tokens_conv,input_tokens_conv,output_tokens_conv = self.request_llm.send_request_non_streaming(self.topic_prompt, user_content,self.email,self.total_tokens_conv,self.input_tokens_conv,self.output_tokens_conv)
        #         print("type :",type(rel_topic))
        #         rel_topic=rel_topic[rel_topic.find("{"):rel_topic.find('}')+1]
        #         rel_topic=ast.literal_eval(rel_topic)
        #         self.data.loc[index,'Topic of Text']=str(rel_topic['topic_of_text'])
        #     else:
        #         pass
        # print("Picking topics completed")
        # return self.data 

        llm_tasks = []
        valid_indices = []
        self.insights_draft=""
        if cl.user_session.get("chat_profile", "") in [os.getenv("BOT_TYPE_PMR"),os.getenv("BOT_TYPE_MARKET_MAP"),os.getenv("BOT_TYPE_EARLY_EXP"),os.getenv("BOT_TYPE_PBC_MARKET_RESEARCH")]:
            self.insights_draft=await self.get_insight_draft()
            print("insights draft : ",self.insights_draft['insights_draft'])

        for index in list(self.data.index):
            full_text = self.data.loc[index, self.column]
            print("index: ", index)
            if not pd.isna(self.data.loc[index, self.column]):

                if self.insights_draft:
                    columns=[column if column not in ['questions','answers','source'] else '' for column in self.data.columns]
                    # data_tro_user_non_user.loc[0,['']]
                    columns=list(set(columns))
                    columns.remove('')
                    respondent_metadata=self.data.loc[index,columns].to_dict()
                    user_content = f"Here are the ** user question ** : {self.question}\n and ** text ** : {full_text} and the ** insight draft for reference ** : {self.insights_draft} and the ** metadata of the respondent **: {respondent_metadata}"
                
                else:
                    user_content = f"Here are the user question : {self.question}\n and text : {full_text}"

                print(f"content before sending to llm for topic (index {index}): ", user_content)
                
                # Create LLM task directly
                llm_task = self.request_llm.send_request_non_streaming(
                    self.topic_prompt, 
                    user_content, 
                    self.email, 
                    self.total_tokens_conv, 
                    self.input_tokens_conv, 
                    self.output_tokens_conv
                )
                llm_tasks.append(llm_task)
                valid_indices.append(index)
        
        if llm_tasks:
            print(f"Processing {len(llm_tasks)} LLM requests in parallel...")
            llm_results = await asyncio.gather(*llm_tasks, return_exceptions=True)
            
            # Process results and update dataframe
            for i, result in enumerate(llm_results):
                index = valid_indices[i]
                
                if isinstance(result, Exception):
                    print(f"Exception occurred for index {index}: {result}")
                    continue
                
                try:
                    rel_topic, email, total_tokens_conv, input_tokens_conv, output_tokens_conv = result
                    print(f"type for index {index}:", type(rel_topic))
                    if cl.user_session.get("chat_profile", "") == os.getenv("BOT_TYPE_PHYSICIAN_OPINIONS_V2"):
                    # Extract JSON-like payload from response
                        rel_topic_str = rel_topic[rel_topic.find("{"):rel_topic.find('}')+1]
                        rel_topic_obj = ast.literal_eval(rel_topic_str)

                        topics_list = []
                        # Support original format: {"topic_of_text": ["TOPIC1", ...]}
                        if isinstance(rel_topic_obj, dict) and 'topic_of_text' in rel_topic_obj:
                            topics_list = rel_topic_obj.get('topic_of_text', [])
                        # Support new format: {"topic1": "positive", "topic2": "mixed", ...}
                        elif isinstance(rel_topic_obj, dict):
                            topics_list = list(rel_topic_obj.keys())
                        # Defensive: if model returned list directly
                        elif isinstance(rel_topic_obj, list):
                            topics_list = rel_topic_obj
                        # Defensive: if single string
                        elif isinstance(rel_topic_obj, str):
                            topics_list = [rel_topic_obj]

                        # Normalize topics_list to list[str]
                        if isinstance(topics_list, str):
                            topics_list = [topics_list]
                        if topics_list is None:
                            topics_list = []

                        # Remove empties and ensure all elements are strings
                        topics_list = [str(t).strip() for t in topics_list if str(t).strip() != '']

                        self.data.loc[index, 'Topic of Text'] = str(topics_list)
                    else:
                        rel_topic = rel_topic[rel_topic.find("{"):rel_topic.rfind('}')+1]
                        if cl.user_session.get('chat_profile') in [os.getenv("BOT_TYPE_PMR"),os.getenv("BOT_TYPE_EARLY_EXP"),os.getenv("BOT_TYPE_PBC_MARKET_RESEARCH")]:
                            rel_topic = ast.literal_eval(rel_topic)['topics']
                            print('rel_topic after pmr extraction: ',rel_topic)
                        else:    
                            rel_topic = ast.literal_eval(rel_topic)['topic_of_text']
                        if 'Topic of Text' not in self.data.columns:
                            self.data['Topic of Text'] = '[]'
                            print('empty list initialised')
                        initial_topics=ast.literal_eval(self.data.loc[index,'Topic of Text'])
                        while '' in initial_topics:
                            initial_topics.remove('')
                        print('initial_topics before adding new topics :',initial_topics)
                        while '' in rel_topic:
                            rel_topic.remove('')
                        print('rel_topic after removing empty strings :',rel_topic)
                        print('rel_topic before initialize 1 :',rel_topic)
                        if rel_topic:
                            print('type of initial_topics :',type(initial_topics))
                            print('type of rel_topic :',type(rel_topic))
                            total_topics=initial_topics+rel_topic
                            # total_topics=list(set(total_topics))
                            print('rel_topic before initialize :',total_topics)
                            self.data.loc[index, 'Topic of Text'] = str(total_topics)
                except Exception as e:
                    print(f"Error processing result for index {index}: {e}")
                    continue
        
        print("Picking topics completed")
        return self.data 

    async def extract_json_simple(self,response_text):
        """Extract JSON from LLM response - handles extra text before/after JSON"""
        
        # Find JSON by counting braces (handles your specific case)
        start = response_text.find('{')
        if start == -1:
            return {}
            
        brace_count = 0
        for i, char in enumerate(response_text[start:], start):
            if char == '{':
                brace_count += 1
            elif char == '}':
                brace_count -= 1
                if brace_count == 0:
                    json_text = response_text[start:i + 1]
                    try:
                        return json.loads(json_text)
                    except:
                        break
        
        # Fallback to your current method
        try:
            start = response_text.find("{")
            end = response_text.rfind('}') + 1
            return json.loads(response_text[start:end])
        except:
            return {}


    async def get_standard_topic(self):
        topics=[]
        for index in list(self.data.index):
            cell = self.data.loc[index,'Topic of Text'] if 'Topic of Text' in self.data.columns else None
            try:
                if cell is None or (isinstance(cell, float) and pd.isna(cell)):
                    continue
                if isinstance(cell, list):
                    topics += cell
                else:
                    cell_str = str(cell).strip()
                    if cell_str in ('', 'nan', 'None'):
                        continue
                    parsed = ast.literal_eval(cell_str)
                    if isinstance(parsed, list):
                        topics += parsed
                    elif isinstance(parsed, str):
                        topics.append(parsed)
            except Exception as e:
                print(f"Skipping index {index} while aggregating topics due to parse issue: {e}")
                continue
        user_content=f"Here are the user question : {self.question }\n and topics : {topics} and the **insights draft**:{self.insights_draft}"
        standardised,email,total_tokens_conv,input_tokens_conv,output_tokens_conv=await self.request_llm.send_request_non_streaming(self.standard_topic_prompt, user_content,self.email,self.total_tokens_conv,self.input_tokens_conv,self.output_tokens_conv)
        # standardised=standardised[standardised.find("{"):standardised.rfind('}')+1]
        print("standardising output: ",standardised)
        standardised=await self.extract_json_simple(standardised)
        if standardised.get('context_confirmation',False):
            relevant_chunks=self.data.to_json(orient='index')
            user_content=f"Here are the user question : {self.question }\n and topics : {topics} and the **complete raw data**:{relevant_chunks} "
            standardised,email,total_tokens_conv,input_tokens_conv,output_tokens_conv=await self.request_llm.send_request_non_streaming(self.standard_topic_prompt, user_content,self.email,self.total_tokens_conv,self.input_tokens_conv,self.output_tokens_conv)
            print("standardising output after LLM: ",standardised)
            standardised=await self.extract_json_simple(standardised)
        # standardised=ast.literal_eval(standardised)
        # standardised = standardised.replace("'", '"')
        print("standardizing topics before literal eval: ",standardised)
        # standardised = json.loads(standardised)
        # standardised=ast.literal_eval(standardised)

        print("standardizing topics after literal eval: ",standardised)
        # print(standardised)
        if cl.user_session.get("chat_profile", "") == (os.getenv("BOT_TYPE_PHYSICIAN_OPINIONS_V2")):
        
            for index in list(self.data.index):
                ans=[]
                cell = self.data.loc[index,'Topic of Text'] if 'Topic of Text' in self.data.columns else None
                try:
                    if cell is None or (isinstance(cell, float) and pd.isna(cell)):
                        self.data.loc[index,"Standardized"] = str(ans)
                        continue
                    if isinstance(cell, list):
                        iterable = cell
                    else:
                        cell_str = str(cell).strip()
                        if cell_str in ('', 'nan', 'None'):
                            self.data.loc[index,"Standardized"] = str(ans)
                            continue
                        parsed = ast.literal_eval(cell_str)
                        iterable = parsed if isinstance(parsed, list) else [parsed]
                    for elem in iterable:
                        ans.append(standardised.get(elem,''))
                        print(f"key : {elem} and value : {standardised.get(elem,None)}")
                    self.data.loc[index,"Standardized"]=str(ans)
                except Exception as e:
                    print(f"Error standardizing topics for index {index}: {e}")
                    self.data.loc[index,"Standardized"] = str(ans)
        else:
            for index in self.data.index:
                standard_topic_mapping=[]
                for topic in ast.literal_eval(self.data.loc[index,'Topic of Text']):
                    for key in standardised:
                        if topic['label'] in standardised[key]['original_topics']:
                            standard_topic_mapping.append(key)
                # print(f"index: {index} and topics: {self.data.loc[index,'Topic of Text']} and grouped : {standard}")
                self.data.loc[index,"Standardized"]=str(standard_topic_mapping)
                print("\n\n")
            print("Standardizing topics completed")
        return self.data
    
    async def get_grouped_topics(self,indices,grouped_topics=[]):
        print("in get_grouped_topics")
        # prompt=system_prompt.format(topics=list(data.loc[:,'Topic of Text']),question=question)
        # prompt=system_prompt.format(topics=prep_topics,question=question)
        topics=[]
        for index in indices:
            
            try:
                if self.data.loc[index,'Standardized']!='':
                    sample=ast.literal_eval(self.data.loc[index,'Standardized'])
                    while '' in sample:
                        sample.remove('')
                    topics+=sample
            except Exception as e:
                topics+=['None']
        # print("topics : ",topics)
        user_content=f"Here are the user question : {self.question}\n and topics : {topics} and so far grouped names: {grouped_topics}"
        # Ensure credentials are refreshed before each call
        standardised=await self.request_llm.send_request_non_streaming(self.topic_grouping_prompt,user_content,self.email,self.total_tokens_conv,self.input_tokens_conv,self.output_tokens_conv)
        standardised=standardised[0]
        print("standardised before preprocess: ",standardised)
        # standardised=str(standardised)
        # standardised=standardised.replace("\n","")
        standardised=standardised[standardised.find("{"):standardised.find('}')+1]
        # standardised = standardised.replace("'", '"')
        if standardised!="":
            print("standardizing topics: ",standardised)
            print("before literal eval sending")
            standardised=ast.literal_eval(standardised)
            print("after literal eval")
            print("standardizing topics: ",standardised)
            print(standardised)
            for index in indices:
                print("before literal in loop for index:",index)
                if self.data.loc[index,'Standardized']!='':
                    ast.literal_eval(self.data.loc[index,'Standardized'])
                    print("after literal in loop for index:",index)

            for index in indices:
                ans=[]
                if self.data.loc[index,'Standardized']!='':
                    for elem in ast.literal_eval(self.data.loc[index,'Standardized']):
                        ans.append(standardised.get(elem,''))
                    # print(f"key : {elem} and value : {standardised.get(elem,None)}")
                    self.data.loc[index,"Standardized_Grouped"]=str(ans)
                else:
                    pass
                    # print("Hi in else")
            print("Grouping Standardised topics completed")
            return self.data,standardised
        else:
            return self.data,{}
        

    async def summarize(self):
        standard=[]
    
        for index in list(self.data.index):
            standard=standard+ast.literal_eval(self.data.loc[index,'Standardized_Grouped'])
        standard=dict(Counter(standard))
    
        df=pd.Series(standard)
        df.sort_values(inplace=True)
        return df
    
    async def run(self,tool_params):
        self.data=tool_params["result"]
        self.topic_prompt=tool_params["topic_picking_prompt"]
        self.standard_topic_prompt=tool_params["standard_topic_prompt"]
        self.topic_grouping_prompt=tool_params["topic_grouping_prompt"]
        self.insight_draft_propmt=tool_params.get("insight_draft_prompt","")
        # Sentiment prompt is optional; when provided, we will run sentiment extraction
        self.sentiment_prompt=tool_params.get("sentiment_prompt", "")
        self.question=tool_params["user_question"]
        self.request_llm=LLM()
        print("column for processing is ",tool_params["text_column"])
        self.column=tool_params["text_column"]
        self.email=tool_params['email']
        self.total_tokens_conv=tool_params['total_tokens_conv']
        self.input_tokens_conv=tool_params['input_tokens_conv']
        self.output_tokens_conv=tool_params['output_tokens_conv']
        for i in range(1):
            self.data=await self.get_topics()
        self.data=await self.get_standard_topic()
        if cl.user_session.get("chat_profile")==os.getenv("BOT_TYPE_PHYSICIAN_OPINIONS_V2"):
            grouped_topics_mapp={}
            grouped_topics=[]
            for i in range(0,len(self.data),25):
                self.data,output_process=await self.get_grouped_topics(list(self.data.index)[i:i+25],grouped_topics)
                for key in list(output_process.keys()):
                    value=output_process[key]
                    
                    grouped_topics_mapp[value]=grouped_topics_mapp.get(value,0)+1
                    if grouped_topics_mapp[value]==1:
                        grouped_topics.append(value)
            
            print("grouped topics mapping: ",grouped_topics_mapp)
        
        # Run sentiment extraction if prompt is available and chatbot is Physician Opinions or PMR
            chatbot_name=cl.user_session.get("chat_profile", "")
            try:
                should_run_sentiment = bool(self.sentiment_prompt and chatbot_name == (os.getenv("BOT_TYPE_PHYSICIAN_OPINIONS_V2")))
            except Exception:
                should_run_sentiment = bool(self.sentiment_prompt)
            if should_run_sentiment:
                print("Starting sentiment extraction using sentiment_prompt")
                self.data = await self.get_sentiments_for_grouped_topics()
                print("Sentiment extraction completed")   
        # print("grouped topics mapping: ",grouped_topics_mapp)
        # self.results=await self.summarize()
        print("Summarized results")
        print(self.data.head())
        tool_params["result"]=self.data

        buffer = io.StringIO()
        self.data.info(buf=buffer)
        info_string = buffer.getvalue()
        print("result information:", info_string)
        
        function_response=f"""
        Text extraction Task completed successfully
        result head : {self.data.head()}
        result tail : {self.data.tail()}
        result info : {info_string}
        """
        cl.user_session.set("tool_params",tool_params)
        return "Topic Extraction Completed Successfully,NOW go forward for visualization part python_after_sql tool",function_response,self.data

    async def get_sentiments_for_grouped_topics(self):
        """Extract sentiment per standardized grouped topic using combined text per row.
        Results are stored in a new column 'Topic_Sentiment_Map' as a JSON-stringified dict.
        """
        if not self.sentiment_prompt:
            return self.data
        llm_tasks = []
        valid_indices = []
        for index in list(self.data.index):
            try:
                combined_text = str(self.data.loc[index, self.column]) if not pd.isna(self.data.loc[index, self.column]) else ""
            except Exception:
                combined_text = ""
            try:
                grouped_list_raw = self.data.loc[index, 'Standardized_Grouped'] if 'Standardized_Grouped' in self.data.columns else "[]"
                grouped_list = []
                if isinstance(grouped_list_raw, list):
                    grouped_list = grouped_list_raw
                else:
                    grouped_list = ast.literal_eval(str(grouped_list_raw)) if str(grouped_list_raw).strip() not in ("", "nan", "None") else []
                # Clean topics
                grouped_list = [str(t).strip() for t in grouped_list if str(t).strip() not in ('', 'uncategorized', 'uncategorized topic')]
            except Exception as e:
                print(f"Error parsing Standardized_Grouped for index {index}: {e}")
                grouped_list = []

            if combined_text and grouped_list:
                user_content = (
                    f"User Question: {self.question}\n"
                    f"Combined Text: {combined_text}\n"
                    f"Standardized_Grouped Topics: {grouped_list}\n"
                    "Return a JSON object mapping each topic to one of these sentiments: positive, negative, neutral, no sentiment."
                )
                task = self.request_llm.send_request_non_streaming(
                    self.sentiment_prompt,
                    user_content,
                    self.email,
                    self.total_tokens_conv,
                    self.input_tokens_conv,
                    self.output_tokens_conv
                )
                llm_tasks.append(task)
                valid_indices.append(index)
            else:
                self.data.loc[index, 'Topic_Sentiment_Map'] = str({})

        if llm_tasks:
            print(f"Processing {len(llm_tasks)} sentiment LLM requests in parallel...")
            llm_results = await asyncio.gather(*llm_tasks, return_exceptions=True)
            for i, result in enumerate(llm_results):
                index = valid_indices[i]
                if isinstance(result, Exception):
                    print(f"Exception occurred during sentiment for index {index}: {result}")
                    self.data.loc[index, 'Topic_Sentiment_Map'] = str({})
                    continue
                try:
                    resp_text, email, total_tokens_conv, input_tokens_conv, output_tokens_conv = result
                    resp_str = resp_text[resp_text.find("{"):resp_text.find('}')+1]
                    sentiment_map = {}
                    try:
                        sentiment_map = ast.literal_eval(resp_str)
                    except Exception:
                        sentiment_map = {}
                    # Normalize values and filter to allowed set
                    allowed = {"positive", "negative", "neutral", "no sentiment"}
                    clean_map = {}
                    for k, v in (sentiment_map.items() if isinstance(sentiment_map, dict) else []):
                        key = str(k).strip()
                        val = str(v).strip().lower()
                        if val not in allowed:
                            continue
                        clean_map[key] = val
                    self.data.loc[index, 'Topic_Sentiment_Map'] = str(clean_map)
                except Exception as e:
                    print(f"Error processing sentiment result for index {index}: {e}")
                    self.data.loc[index, 'Topic_Sentiment_Map'] = str({})

        return self.data
    
    
class InsightsResponse(BaseTool):
    def __init__(self):
        pass
    async def run(self,tool_params):
        data_list = []
        saved_files_lst=tool_params.get("excel_files_list", [])
        print("saved files list in tool call :",saved_files_lst)

        for file in saved_files_lst:
            try:
                if file.endswith('.csv'):
                    df = pd.read_csv(file)
                elif file.endswith('.xlsx') or file.endswith('.xls'):
                    df = pd.read_excel(file)
                else:
                    continue
                data_list.append({
                    "File Name": file,
                    "File's Content": df.to_dict(orient='records')
                })
            except Exception as e:
                print(f"Error processing file {file}: {e}")

        json_data_combined = json.dumps(data_list, indent=4)
        list_data = json.loads(json_data_combined)
        print("json_data_combined (Python Object):", json_data_combined)
        # Function to check if the first key in the first dictionary has a value
        def check_if_data_exist(data):
            try:
                # Validate if data is a non-empty list
                if not isinstance(data, list) or not data:
                    print(f"Invalid Data Structure (Not a list or empty): {data}")
                    return False
                file_entry = data[0]  # Get the first dictionary entry
                if not isinstance(file_entry, dict):
                    print(f"Invalid Data Structure (First entry not a dictionary): {file_entry}")
                    return False
                # Ensure "File's Content" exists, is a list, and is non-empty
                file_content = file_entry.get("File's Content")
                if not isinstance(file_content, list) or not file_content:
                    print("Missing or empty 'File's Content'")
                    return False
                first_dict = file_content[0]  # Get the first dictionary inside "File's Content"
                if not isinstance(first_dict, dict) or not first_dict:
                    print("First dictionary inside 'File's Content' is empty or not a dictionary")
                    return False
                # Get the first key in the dictionary
                first_key = next(iter(first_dict), None)
                if first_key is None:  # If there's no key, return False
                    print("No key found in the first dictionary")
                    return False
                first_value = first_dict.get(first_key)
                print(f"First Key: {first_key}, First Value: {first_value}")
                # Return True if the first key exists and has a non-empty, non-null value
                return bool(first_value)
            except Exception as e:
                print(f"Error while checking JSON structure: {e}")
                return False
        # Measure execution time
        start_time = time.time()
        # Run the function with corrected data
        result = check_if_data_exist(list_data)
        # Print result
        insights_resp={}
        print(f"Result for JSON Data: {result}")
        if not result:
            print("Json Data Does not Exist")
            insights_resp["content"]=""
            return None,"",None
        else:
            print("Json Data Exists")
        print(f"Execution Time For Insights: {time.time() - start_time}")
        # Parse the JSON string back into a Python object
        parsed_data = json.loads(json_data_combined)
        json_data_combined = [file_data["File's Content"] for file_data in parsed_data]
        print("Here is final json data combined:", json_data_combined)
        pmr_topic_length=0
        analysis_results_pmr=None
        for elem in list_data:
            if 'analysis_results' in elem['File Name']:
                pmr_topic_length=len(elem["File's Content"])
                analysis_results_pmr=elem["File's Content"]
        print("pmr topic length:",pmr_topic_length)
        formatted_response = [
            {
                "type": "text",
                "text": f"Here is final json data combined: {json.dumps(json_data_combined)}"
            }
        ]
        if cl.user_session.get("chat_profile") in [os.getenv("BOT_TYPE_PMR"),os.getenv("BOT_TYPE_PHYSICIAN_OPINIONS_V2"),os.getenv("BOT_TYPE_MARKET_MAP"),os.getenv("BOT_TYPE_EARLY_EXP"),os.getenv("BOT_TYPE_PBC_MARKET_RESEARCH")]:
            insights_prompt=tool_params['insights_prompt']
        #     for elem in json_data_combined:
            print("list data :",list_data)
        #         if elem["File Name"]=="analysis_results.csv":
            formatted_response = [
            {
                "type": "text",
                "text": f"Now provide Insights for the data attached below,through following this instructions and must generate insights for {pmr_topic_length} topics mentioned : "+insights_prompt
            },
            {
                'type':'text',
                "text": f"Must generate insights for {pmr_topic_length} topics mentioned as follows: {analysis_results_pmr}"
            },
            {
                "type":"json",
                "Json data from opensearch for user question":json.dumps(json_data_combined)
            }
        ]
        
        if cl.user_session.get("chat_profile")==os.getenv("BOT_TYPE_DSG"):
            insights_prompt = tool_params['insights_prompt']
            formatted_response = [
                {
                "type": "text",
                "text": insights_prompt.format(user_question = tool_params["user_question"],result_data_json=json.dumps(json_data_combined))
            }
            ]

        tool_params['insights_call']=True
        return None,formatted_response,None
    
class CalculateNSPCaptureRatio(BaseTool):
    """
    Universal tool for NSP analysis supporting 4 analysis types:
    1. capture_ratio - Time series capture ratio analysis
    2. projection_oral - R12M projection factor for oral drugs
    3. projection_iv - R12M projection factor for IV drugs (with optional drug_name filtering)
    4. projection_adc_trend - 6-month trend for ADC projection factor
    """
    
    # Drug classification mappings
    ORAL_DRUGS = {
        'ABEMACICLIB', 'PALBOCICLIB', 'RIBOCICLIB', 'CDK',
        'OLAPARIB', 'TALAZOPARIB',
        'ALPELISIB', 'ITOVEBI', 'INAVOLISIB', 'TRUQAP',
        'EVEROLIMUS', 'CAPIVASERTIB',
        'ANASTROZOLE', 'LETROZOLE', 'EXEMESTANE',
        'TAMOXIFEN', 'TOREMIFENE', 'OTHER HORMONAL',
        'CAPECITABINE', 'XELODA',
        'OTHER TARGETED'
    }
    
    IV_DRUGS = {
        'TRODELVY', 'SACITUZUMAB', 'SACITUZUMAB GOVITECAN',
        'ENHERTU', 'TRASTUZUMAB', 'FAM-TRASTUZUMAB',
        'DATROWAY',
        'GEMZAR', 'HALAVEN', 'ERIBULIN', 'VINORELBINE',
        'OTHER CHEMO', 'TAXANE', 'EPIRUBICIN', 'IXABEPILONE',
        'FULVESTRANT', 'ORSERDU'
    }
    
    def __init__(self):
        super().__init__()
        
    def _classify_drug(self, drug_name):
        """Classify drug as oral or IV"""
        if not drug_name:
            return None
            
        drug_upper = drug_name.upper().strip()
        
        # Check exact matches
        if drug_upper in self.ORAL_DRUGS:
            return 'oral'
        if drug_upper in self.IV_DRUGS:
            return 'iv'
        
        # Check partial matches
        for oral_drug in self.ORAL_DRUGS:
            if oral_drug in drug_upper or drug_upper in oral_drug:
                return 'oral'
        
        for iv_drug in self.IV_DRUGS:
            if iv_drug in drug_upper or drug_upper in iv_drug:
                return 'iv'
        
        # Default to oral
        return 'oral'
    
    def _generate_capture_ratio_sql(self, drug_name, period_type, drug_type, batch_month):
        """Generate SQL for capture ratio analysis"""
        function_name = 'nsp_capture_ratio_oral' if drug_type == 'oral' else 'nsp_capture_ratio_iv'
        batch_param = f"'{batch_month}'" if batch_month else "NULL"
        
        sql_query = f"""
        SELECT * 
        FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.{function_name}(
            '{drug_name.upper().strip()}',
            '{period_type.upper().strip()}',
            {batch_param}
        )
        ORDER BY period ASC;
        """
        
        return sql_query
    
    def _generate_projection_oral_sql(self):
        """Generate SQL for oral projection factor"""
        sql_query = """
        SELECT * 
        FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.nsp_projection_factor_oral();
        """
        return sql_query
    
    def _generate_projection_iv_sql(self, drug_name=None):
        """Generate SQL for IV projection factor with optional drug filtering"""
        if drug_name and drug_name.strip():
            drug_param = f"'{drug_name.upper().strip()}'"
        else:
            drug_param = "NULL"
        
        sql_query = f"""
        SELECT * 
        FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.nsp_projection_factor_iv({drug_param})
        ORDER BY 
            CASE CATEGORY
                WHEN 'ADC' THEN 1
                WHEN 'CHEMO' THEN 2
                WHEN 'FULVESTRANT' THEN 3
            END;
        """
        return sql_query
    
    def _generate_projection_adc_trend_sql(self):
        """Generate SQL for ADC projection factor trend"""
        sql_query = """
        SELECT * 
        FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.nsp_projection_factor_adc_trend()
        ORDER BY batch_month ASC;
        """
        return sql_query
    
    async def run(self, tool_params):
        """Execute NSP analysis based on analysis_type"""
        print("NSP Analysis Tool Started")
        
        # Extract parameters
        arguments = tool_params.get("arguments", {})
        analysis_type = arguments.get("analysis_type", "").strip().lower()
        drug_name = arguments.get("drug_name", "").strip() if arguments.get("drug_name") else None
        period_type = arguments.get("period_type", "QUARTERLY").upper()
        batch_month = arguments.get("batch_month", "").strip() if arguments.get("batch_month") else None
        
        print(f"Parameters: analysis_type={analysis_type}, drug_name={drug_name}, period_type={period_type}, batch_month={batch_month}")
        
        # Validate analysis_type
        valid_types = ["capture_ratio", "projection_oral", "projection_iv", "projection_adc_trend"]
        if analysis_type not in valid_types:
            error_msg = f"Error: analysis_type must be one of {valid_types}. Got: {analysis_type}"
            print(error_msg)
            return None, error_msg, None
        
        try:
            # Route to appropriate SQL generation
            if analysis_type == "capture_ratio":
                if not drug_name:
                    error_msg = "Error: drug_name is required for capture_ratio analysis"
                    print(error_msg)
                    return None, error_msg, None
                
                if period_type not in ["MONTHLY", "QUARTERLY", "YEARLY"]:
                    error_msg = f"Error: period_type must be MONTHLY, QUARTERLY, or YEARLY. Got: {period_type}"
                    print(error_msg)
                    return None, error_msg, None
                
                drug_type = self._classify_drug(drug_name)
                print(f"Drug '{drug_name}' classified as {drug_type.upper()}")
                sql_query = self._generate_capture_ratio_sql(drug_name, period_type, drug_type, batch_month)
                
            elif analysis_type == "projection_oral":
                print("Generating oral projection factor SQL")
                sql_query = self._generate_projection_oral_sql()
                
            elif analysis_type == "projection_iv":
                print(f"Generating IV projection factor SQL{f' for {drug_name}' if drug_name else ' (all categories)'}")
                sql_query = self._generate_projection_iv_sql(drug_name)
                
            elif analysis_type == "projection_adc_trend":
                print("Generating ADC projection factor trend SQL")
                sql_query = self._generate_projection_adc_trend_sql()
            
            print(f"\nGenerated SQL:\n{sql_query}")
            
            # Store context
            cl.user_session.set("nsp_sql_query", sql_query)
            cl.user_session.set("nsp_analysis_type", analysis_type)
            cl.user_session.set("nsp_drug_name", drug_name)
            
            # Execute SQL
            dataware_house = tool_params.get("dataware_house")
            if not dataware_house:
                error_msg = "Error: dataware_house connection not available"
                print(error_msg)
                return sql_query, error_msg, None
            
            print("Executing SQL query...")
            warehouse = cl.make_async(dataware_house.get_data_from_wareshouse)
            result = await warehouse(sql_query)
            
            # Validate result
            if result is None or result.empty:
                error_msg = f"No data returned. Please verify parameters."
                print(error_msg)
                return sql_query, error_msg, pd.DataFrame()
            
            print(f"Success: Retrieved {len(result)} rows")
            
            # Store result
            tool_params["result"] = result
            cl.user_session.set("tool_params", tool_params)
            cl.user_session.set("nsp_result", result)
            
            # Simple response - just head() and instruction
            if result.empty:
                    raise Exception("Empty DataFrame")
            print(result.head())
            print(result.tail())
            buffer = io.StringIO()
            result.info(buf=buffer)
            info_string = buffer.getvalue()
            print("result information:", info_string)
            #result.to_csv("Structured_Bot/output_files/result.csv")
            
            function_response = f"""SQL EXECUTED SUCCESSFULLY.
                Dataset: 'result' dataframe with {len(result)} rows
                - result.head(): {result.head()}
                - result.tail() : {result.tail()}
                - result.info(): {info_string}
                - Date format: Already in string format (use directly)
                Use python_after_sql tool. NEVER use any return statement in the python code.
                """ 
            
            print("Analysis complete")
            return sql_query, function_response, result
            
        except Exception as e:
            error_msg = f"Error in NSP analysis: {str(e)}"
            print(error_msg)
            import traceback
            traceback.print_exc()
            return sql_query , error_msg, None
        


        
class CreatePresentationDeck(BaseTool):
    """Tool for creating PowerPoint presentations with multiple charts"""
    
    def __init__(self):
        super().__init__()
        self.config_files = {
            'monthly': 'DS&G/Input_Files/monthly_deck.xlsx',
            'weekly': 'DS&G/Input_Files/weekly_deck.xlsx'
        }
        self.template_files = {
            'monthly': 'DS&G/monthly_template.pptx',
            'weekly': 'DS&G/weekly_template.pptx'
        }

    def _debug_print_chart_titles(self, prs):
        """Debug: Print all chart titles in the presentation"""
        print("\n" + "="*60)
        print("DEBUG: Chart Titles Found in Template")
        print("="*60)
        
        for slide_idx, slide in enumerate(prs.slides, start=1):
            print(f"\n📍 Slide {slide_idx}:")
            
            chart_count = 0
            for shape_idx, shape in enumerate(slide.shapes):
                if shape.has_chart:
                    chart_count += 1
                    try:
                        title = shape.chart.chart_title.text_frame.text
                        print(f"   Chart #{chart_count}: '{title}'")
                    except:
                        print(f"   Chart #{chart_count}: [No Title]")
            
            if chart_count == 0:
                print(f"   ⚠️  No charts found on this slide")
        
        print("="*60 + "\n")
    
    def _update_chart_from_config(self, prs, slide_number: int, chart_title: str, df: pd.DataFrame, 
                               x_axis_column: str, y_axis_column: str, 
                               lower_threshold: str, upper_threshold: str, 
                               median_threshold: str = None,
                               chart_type: str = None,
                               key_event_column: str = None,
                               python_code: str = None,
                               insight_text: str = None,
                               latest_date_column: str = None):
        """
        Update chart - handles both single series and multi-batch charts
        """
        import io
        from openpyxl import load_workbook
        from pptx.chart.data import CategoryChartData
        from pptx.util import Pt
        from pptx.dml.color import RGBColor

        is_batch_chart = False
        embedded_excel_updated = False

        
        # Get the slide
        slide_idx = slide_number - 1
        if slide_idx >= len(prs.slides):
            raise ValueError(f"Slide {slide_number} does not exist")
        
        slide = prs.slides[slide_idx]
        
        if chart_title and pd.notna(chart_title) and str(chart_title).strip().lower() not in ['nan', 'none', '']:
            # Find chart by title
            chart_shape = None
            for shape in slide.shapes:
                if shape.has_chart:
                    try:
                        if shape.chart.chart_title.text_frame.text == chart_title:
                            chart_shape = shape
                            break
                    except:
                        pass
            
            if chart_shape is None:
                print(f"   ⚠️  Chart with title '{chart_title}' not found on slide {slide_number} - Skipping chart update")
                chart = None
            else:
                chart = chart_shape.chart
                print(f"📊 Updating chart '{chart_title}' on slide {slide_number}")
        else:
            print(f"   ℹ️  No chart title provided - Skipping chart update")
            chart = None
        
        
        # ========== CHECK IF THIS IS A BATCH CHART (Only if chart found) ==========
        if chart:
            is_batch_chart = 'batch_label' in df.columns and 'batch_rank' in df.columns
        
            if is_batch_chart:
                print(f"   🔄 Detected batch chart - creating separate series...")
                print(f"   Original shape: {df.shape}")
                print(f"   Batches found: {df['batch_label'].unique().tolist()}")
                
                # ========== CREATE SEPARATE DATAFRAMES FOR EACH BATCH ==========
                try:
                    # Use actual batch labels from the data, sorted by batch_rank
                    batch_order = df.sort_values('batch_rank')['batch_label'].unique().tolist()
                    
                    # Create a list to hold data for each batch
                    all_batch_data = []
                    max_length = 0
                    
                    for batch_name in batch_order:
                        # Filter data for this batch
                        batch_df = df[df['batch_label'] == batch_name].copy()
                        batch_df = batch_df.sort_values(x_axis_column)
                        
                        # Get x and y values (drop NaN)
                        batch_df_clean = batch_df[[x_axis_column, y_axis_column]].dropna()
                        
                        all_batch_data.append({
                            'name': batch_name,
                            'x': batch_df_clean[x_axis_column].tolist(),
                            'y': batch_df_clean[y_axis_column].tolist()
                        })
                        
                        max_length = max(max_length, len(batch_df_clean))
                        
                        print(f"      {batch_name}: {len(batch_df_clean)} data points")
                    
                    # Check if there's a data_label column to add as an additional series
                    if 'data_label' in df.columns:
                        # Use the first batch's x-axis values and data_label from that batch
                        first_batch_name = batch_order[0]
                        first_batch_df = df[df['batch_label'] == first_batch_name].copy()
                        first_batch_df = first_batch_df.sort_values(x_axis_column)
                        
                        # Convert data_label to numeric, keeping NaN values
                        first_batch_df['data_label'] = pd.to_numeric(
                            first_batch_df['data_label'].replace('null', pd.NA).replace('NULL', pd.NA), 
                            errors='coerce'
                        )
                        
                        # Check if there are any non-null values
                        non_null_count = first_batch_df['data_label'].notna().sum()
                        if non_null_count > 0:
                            # Keep nulls as None so they appear at the correct positions
                            # DON'T use dropna() - this preserves the position of data points
                            y_values = [float(v) if pd.notna(v) else None for v in first_batch_df['data_label']]
                            x_values = first_batch_df[x_axis_column].tolist()
                            
                            all_batch_data.append({
                                'name': 'Data Label',
                                'x': x_values,
                                'y': y_values
                            })
                            
                            print(f"      Data Label: {non_null_count} non-null values out of {len(y_values)} points")
                    
                    print(f"   Max series length: {max_length}")
                    
                except Exception as e:
                    print(f"   ⚠️  Batch processing failed: {str(e)}")
                    raise
        
        # ========== UPDATE EMBEDDED EXCEL WITH SEPARATE SERIES ==========
        if chart:
            try:
                chart_part = chart.part
                xlsx_part = None
                
                rel_ids = list(chart_part.rels.keys())
                
                for rel_id in rel_ids:
                    try:
                        rel = chart_part.rels[rel_id]
                        
                        if not rel.is_external and hasattr(rel, 'target_part'):
                            part = rel.target_part
                            
                            if hasattr(part, 'blob'):
                                try:
                                    test_wb = load_workbook(io.BytesIO(part.blob))
                                    xlsx_part = part
                                    test_wb.close()
                                    print(f"   ✅ Found embedded Excel")
                                    break
                                except:
                                    continue
                    except:
                        continue
                
                if xlsx_part and is_batch_chart:
                    # Update embedded Excel with separate series structure
                    wb = load_workbook(io.BytesIO(xlsx_part.blob))
                    ws = wb.active
                    
                    # Clear all existing data
                    ws.delete_rows(1, ws.max_row)
                    
                    # ========== WRITE HEADERS ==========
                    col_idx = 1
                    for batch_data in all_batch_data:
                        ws.cell(1, col_idx, f"{batch_data['name']} - Week")
                        ws.cell(1, col_idx + 1, batch_data['name'])
                        col_idx += 2
                    
                    print(f"   Created {len(all_batch_data)} separate series")
                    
                    # ========== WRITE DATA FOR EACH SERIES ==========
                    for batch_idx, batch_data in enumerate(all_batch_data):
                        start_col = batch_idx * 2 + 1
                        
                        for row_idx, (x_val, y_val) in enumerate(zip(batch_data['x'], batch_data['y']), start=2):
                            ws.cell(row_idx, start_col, str(x_val))      # X value
                            ws.cell(row_idx, start_col + 1, float(y_val))  # Y value
                        
                        print(f"      {batch_data['name']}: {len(batch_data['x'])} rows written")
                    
                    # Save back
                    output = io.BytesIO()
                    wb.save(output)
                    output.seek(0)
                    xlsx_part._blob = output.read()
                    wb.close()
                    
                    embedded_excel_updated = True
                    print(f"   ✅ Updated embedded Excel with separate series")
                    
                elif xlsx_part and not is_batch_chart:
                    # Regular chart (non-batch) - use friendly names for outlier charts
                    wb = load_workbook(io.BytesIO(xlsx_part.blob))
                    ws = wb.active
                    
                    if ws.max_row > 1:
                        ws.delete_rows(2, ws.max_row)
                    
                    is_outlier_chart = chart_type and str(chart_type).lower() == 'outlier'
                    col_mapping = [(x_axis_column, x_axis_column)]
                    
                    if y_axis_column and pd.notna(y_axis_column) and y_axis_column in df.columns:
                        name = 'Weekly Claims' if is_outlier_chart else y_axis_column
                        col_mapping.append((y_axis_column, name))
                    
                    if median_threshold and pd.notna(median_threshold) and str(median_threshold).strip() and median_threshold in df.columns:
                        name = 'Median' if is_outlier_chart else median_threshold
                        col_mapping.append((median_threshold, name))
                    
                    if upper_threshold and pd.notna(upper_threshold) and str(upper_threshold).strip() and upper_threshold in df.columns:
                        name = 'Upper Bound' if is_outlier_chart else upper_threshold
                        col_mapping.append((upper_threshold, name))
                        
                        name = 'Lower Bound' if is_outlier_chart else lower_threshold
                        col_mapping.append((lower_threshold, name))
                    
                    if 'data_label' in df.columns:
                        col_mapping.append(('data_label', 'Data Label'))
                    
                    for col_idx, (_, header) in enumerate(col_mapping, start=1):
                        ws.cell(row=1, column=col_idx, value=str(header))
                    
                    threshold_vals = {}
                    if is_outlier_chart:
                        if median_threshold in df.columns:
                            threshold_vals[median_threshold] = float(df[median_threshold].iloc[0]) if pd.notna(df[median_threshold].iloc[0]) else None
                        if upper_threshold in df.columns:
                            threshold_vals[upper_threshold] = float(df[upper_threshold].iloc[0]) if pd.notna(df[upper_threshold].iloc[0]) else None
                        if lower_threshold in df.columns:
                            threshold_vals[lower_threshold] = float(df[lower_threshold].iloc[0]) if pd.notna(df[lower_threshold].iloc[0]) else None
                    
                    for row_idx, (_, row) in enumerate(df.iterrows(), start=2):
                        for col_idx, (df_col, _) in enumerate(col_mapping, start=1):
                            if col_idx == 1:
                                ws.cell(row=row_idx, column=col_idx, value=str(row[df_col]))
                            else:
                                val = threshold_vals.get(df_col, row[df_col]) if is_outlier_chart and df_col in threshold_vals else row[df_col]
                                if pd.notna(val):
                                    ws.cell(row=row_idx, column=col_idx, value=float(val))
                    
                    output = io.BytesIO()
                    wb.save(output)
                    output.seek(0)
                    xlsx_part._blob = output.read()
                    wb.close()
                    
                    embedded_excel_updated = True
                    print(f"   ✅ Updated embedded Excel: {len(df)} rows, {len(col_mapping)} columns")
            
            except Exception as e:
                print(f"   ⚠️  Embedded Excel update failed: {str(e)}")
                import traceback
                traceback.print_exc()
        
        # ========== FALLBACK METHOD ==========
        if chart and not embedded_excel_updated:
            print(f"   🔄 Using fallback: chart.replace_data()")
            
            try:
                if is_batch_chart:
                    # For batch charts, use the separate series data
                    chart_data = CategoryChartData()
                    
                    # Use the longest series for categories (usually the first batch)
                    # But ensure we use a series that has the x-axis values
                    main_series_data = next((b for b in all_batch_data if b['name'] != 'Data Label'), all_batch_data[0])
                    chart_data.categories = [str(x) for x in main_series_data['x']]
                    
                    # Add each series
                    for batch_data in all_batch_data:
                        # If x-values match main series, use y-values directly
                        if batch_data['x'] == main_series_data['x']:
                            values = batch_data['y']
                        else:
                            # Align values to categories if x-values differ (shouldn't happen for batches but good for safety)
                            values = []
                            batch_map = dict(zip(batch_data['x'], batch_data['y']))
                            for cat in main_series_data['x']:
                                values.append(batch_map.get(cat, None))
                        
                        chart_data.add_series(batch_data['name'], values)
                    
                    chart.replace_data(chart_data)
                else:
                    # Regular chart - Plot y_axis_column and optional thresholds
                    chart_data = CategoryChartData()
                    chart_data.categories = df[x_axis_column].astype(str).tolist()
                    
                    # Check if this is an outlier chart using chart_type column
                    is_outlier_chart = chart_type and str(chart_type).lower() == 'outlier'
                    
                    # Add the main y-axis series (if specified)
                    if y_axis_column and pd.notna(y_axis_column) and y_axis_column in df.columns:
                        # Try to convert to numeric if it's not already
                        if not pd.api.types.is_numeric_dtype(df[y_axis_column]):
                            try:
                                # Try converting (handles strings like "5.2%", "1,234", etc.)
                                df[y_axis_column] = pd.to_numeric(df[y_axis_column].astype(str).str.replace('%', '').str.replace(',', ''), errors='coerce')
                                print(f"      Converted '{y_axis_column}' to numeric")
                            except:
                                pass
                        
                        if pd.api.types.is_numeric_dtype(df[y_axis_column]):
                            values = [float(v) if pd.notna(v) else None for v in df[y_axis_column]]
                            # Use friendly name for outlier chart
                            series_name = 'Weekly Claims' if is_outlier_chart else y_axis_column
                            chart_data.add_series(series_name, values)
                            print(f"      Added main series: {series_name}")
                        else:
                            print(f"      ⚠️  y_axis_column '{y_axis_column}' is not numeric and could not be converted")
                    
                    # Add median threshold if specified (as constant line)
                    if is_outlier_chart and median_threshold and pd.notna(median_threshold) and str(median_threshold).strip() and median_threshold in df.columns:
                        if pd.api.types.is_numeric_dtype(df[median_threshold]):
                            # Use first value and repeat for all categories (horizontal line)
                            constant_value = float(df[median_threshold].iloc[0]) if pd.notna(df[median_threshold].iloc[0]) else None
                            values = [constant_value] * len(df)
                            chart_data.add_series('Median', values)
                            print(f"      Added median threshold: Median (constant: {constant_value})")
                    
                    # Add upper threshold if specified (as constant line)
                    if upper_threshold and pd.notna(upper_threshold) and str(upper_threshold).strip() and upper_threshold in df.columns:
                        # Try to convert to numeric if it's not already
                        if not pd.api.types.is_numeric_dtype(df[upper_threshold]):
                            try:
                                df[upper_threshold] = pd.to_numeric(df[upper_threshold].astype(str).str.replace('%', '').str.replace(',', ''), errors='coerce')
                                print(f"      Converted '{upper_threshold}' to numeric")
                            except:
                                pass
                        
                        if pd.api.types.is_numeric_dtype(df[upper_threshold]):
                            # Use first value and repeat for all categories (horizontal line)
                            constant_value = float(df[upper_threshold].iloc[0]) if pd.notna(df[upper_threshold].iloc[0]) else None
                            values = [constant_value] * len(df)
                            series_name = 'Upper Bound' if is_outlier_chart else upper_threshold
                            chart_data.add_series(series_name, values)
                            print(f"      Added upper threshold: {series_name} (constant: {constant_value})")
                    
                    # Add lower threshold if specified (as constant line)
                    if lower_threshold and pd.notna(lower_threshold) and str(lower_threshold).strip() and lower_threshold in df.columns:
                        # Try to convert to numeric if it's not already
                        if not pd.api.types.is_numeric_dtype(df[lower_threshold]):
                            try:
                                df[lower_threshold] = pd.to_numeric(df[lower_threshold].astype(str).str.replace('%', '').str.replace(',', ''), errors='coerce')
                                print(f"      Converted '{lower_threshold}' to numeric")
                            except:
                                pass
                        
                        if pd.api.types.is_numeric_dtype(df[lower_threshold]):
                            # Use first value and repeat for all categories (horizontal line)
                            constant_value = float(df[lower_threshold].iloc[0]) if pd.notna(df[lower_threshold].iloc[0]) else None
                            values = [constant_value] * len(df)
                            series_name = 'Lower Bound' if is_outlier_chart else lower_threshold
                            chart_data.add_series(series_name, values)
                            print(f"      Added lower threshold: {series_name} (constant: {constant_value})")
                    
                    # Add Data Label series if specified
                    if 'data_label' in df.columns:
                        # Convert to numeric if needed
                        if not pd.api.types.is_numeric_dtype(df['data_label']):
                            try:
                                df['data_label'] = pd.to_numeric(df['data_label'].replace('null', pd.NA).replace('NULL', pd.NA), errors='coerce')
                            except:
                                pass
                        
                        values = [float(v) if pd.notna(v) else None for v in df['data_label']]
                        # Only add if there's at least one non-null value
                        if any(v is not None for v in values):
                            chart_data.add_series('Data Label', values)
                            print(f"      Added Data Label series")
                    
                    # Log the chart data structure before replacing
                    print(f"   📊 Chart data structure:")
                    print(f"      Categories: {len(chart_data.categories)} items")
                    print(f"      First 5 categories: {chart_data.categories[:5]}")
                    print(f"      Last 5 categories: {chart_data.categories[-5:]}")
                    
                    chart.replace_data(chart_data)
                
                print(f"   ✅ Chart data replaced using fallback")
                
            except ValueError as ve:
                if "target-mode is external" in str(ve):
                    print(f"   ⚠️  Chart has external relationship issue (broken template chart)")
                    print(f"   💡 Fix: Open template, delete & recreate chart on slide {slide_number}")
                    print(f"   ⏭️  Skipping this chart, continuing with others...")
                    return  # Skip this chart, don't fail entire deck
                else:
                    print(f"   ❌ Fallback failed: {str(ve)}")
                    import traceback
                    traceback.print_exc()
                    raise
            except Exception as e:
                print(f"   ❌ Fallback failed: {str(e)}")
                import traceback
                traceback.print_exc()
                raise
        
        # ========== APPLY COLORS ==========
        if chart:
            try:
                print(f"   🎨 Applying series colors...")
                
                batch_colors = [
                    RGBColor(31, 119, 180),   # Blue
                    RGBColor(255, 127, 14),   # Orange
                    RGBColor(44, 160, 44)     # Green
                ]
                
                plot = chart.plots[0]
                
                for idx, series in enumerate(plot.series):
                    try:
                        line = series.format.line
                        if idx < len(batch_colors):
                            line.color.rgb = batch_colors[idx]
                            line.width = Pt(2.5)
                    except:
                        pass
                
                print(f"   ✅ All series set to line chart type")
                print(f"   ✅ Colors applied")
                
            except Exception as e:
                print(f"   ⚠️  Color application failed: {e}")
        
        # ========== MARK OUTLIERS ==========
        if chart:
            try:
                # Check if this is an outlier chart and has outlier_flag column
                if chart_type and str(chart_type).lower() == 'outlier' and 'outlier_flag' in df.columns:
                    print(f"   🔴 Marking outliers based on outlier_flag column...")
                    
                    plot = chart.plots[0]
                    if len(plot.series) > 0:
                        main_series = plot.series[0]  # First series is "Weekly Claims"
                        
                        outlier_count = 0
                        # Iterate through dataframe to find outliers
                        for idx, row in df.iterrows():
                            try:
                                # Get position in series
                                position = df.index.get_loc(idx)
                                point = main_series.points[position]
                                
                                if row['outlier_flag'] != 'Normal':
                                    # This is an outlier - add RED diamond marker
                                    point.marker.style = 2  # Diamond shape
                                    point.marker.format.fill.solid()
                                    point.marker.format.fill.fore_color.rgb = RGBColor(255, 0, 0)  # Red
                                    point.marker.size = 8
                                    outlier_count += 1
                                else:
                                    # Normal point - no marker
                                    point.marker.style = -1  # No marker
                            except Exception as e:
                                print(f"      ⚠️  Could not style point at position {position}: {e}")
                        
                        print(f"   ✅ Marked {outlier_count} outliers with red markers")
            except Exception as e:
                print(f"   ⚠️  Outlier marking failed: {e}")
        
        # ========== ADD KEY EVENT ANNOTATIONS ==========
        if chart:
            try:
                if key_event_column and pd.notna(key_event_column) and str(key_event_column).strip() and key_event_column in df.columns:
                    print(f"   🏷️  Adding key event annotations from column '{key_event_column}'...")
                    
                    plot = chart.plots[0]
                    if len(plot.series) > 0:
                        # Determine which series to annotate
                        target_series = None
                        target_series_index = 0
                        
                        # If data_label column exists, try to find the "Data Label" series
                        if 'data_label' in df.columns:
                            # First try to find by name
                            for i, ser in enumerate(plot.series):
                                print(f"      Series {i}: '{ser.name}'")
                                if ser.name == "Data Label":
                                    target_series = ser
                                    target_series_index = i
                                    print(f"   🎯 Targeting 'Data Label' series (index {i}) for annotations")
                                    break
                            
                            # If not found by name, but we know we added it (it's usually the last one)
                            if target_series is None and len(plot.series) > 1:
                                target_series = plot.series[-1]
                                target_series_index = len(plot.series) - 1
                                print(f"   🎯 Targeting last series (index {target_series_index}) as fallback for 'Data Label'")
                        
                        # Fallback to first series if still not found
                        if target_series is None:
                            target_series = plot.series[0]
                            print(f"   🎯 Targeting main series (index 0) for annotations")
                        
                        num_points = len(target_series.points)
                        
                        # For batch charts, we need to align with the series data
                        # The "Data Label" series (and main series) corresponds to the first batch's x-values
                        if is_batch_chart and 'batch_rank' in df.columns:
                            # Get the first batch (rank 1)
                            first_batch_label = df.sort_values('batch_rank')['batch_label'].iloc[0]
                            df_for_events = df[df['batch_label'] == first_batch_label].copy().reset_index(drop=True)
                            print(f"      Filtering to batch '{first_batch_label}': {len(df_for_events)} rows")
                        else:
                            df_for_events = df.reset_index(drop=True)
                        
                        event_count = 0
                        for idx, row in df_for_events.iterrows():
                            try:
                                event_text = row[key_event_column]
                                
                                # Skip if no event for this data point
                                if pd.isna(event_text) or str(event_text).strip() == '' or str(event_text).lower() == 'null':
                                    continue
                                
                                # Position is the row index within this filtered dataframe
                                position = idx
                                
                                if position >= num_points:
                                    print(f"      ⚠️  Position {position} out of range (series has {num_points} points)")
                                    continue
                                    
                                point = target_series.points[position]
                                
                                # Enable data label for this point
                                dl = point.data_label
                                dl.has_text_frame = True
                                
                                # Clear existing text and add new
                                tf = dl.text_frame
                                p = tf.paragraphs[0]
                                p.clear()
                                run = p.add_run()
                                run.text = str(event_text)
                                
                                # Style the data label (Standard PowerPoint styling)
                                try:
                                    from pptx.util import Pt
                                    from pptx.dml.color import RGBColor
                                    from pptx.enum.text import PP_ALIGN
                                    
                                    run.font.size = Pt(9)
                                    run.font.bold = True
                                    run.font.color.rgb = RGBColor(0, 0, 0)  # Black text
                                    
                                    # Add light yellow background fill to the data label
                                    try:
                                        fill = dl.format.fill
                                        fill.solid()
                                        fill.fore_color.rgb = RGBColor(255, 255, 200)  # Light yellow background
                                        
                                        # Add border/line around the label
                                        line = dl.format.line
                                        line.color.rgb = RGBColor(180, 180, 150)  # Light olive border
                                        line.width = Pt(1)
                                    except Exception as fill_err:
                                        print(f"      ⚠️  Could not add background: {fill_err}")
                                    
                                    # Position above the point
                                    dl.position = 0  # 0 = ABOVE
                                except:
                                    pass
                                
                                event_count += 1
                                print(f"      🏷️  Added event at position {position}: '{event_text}'")
                                
                            except Exception as e:
                                print(f"      ⚠️  Could not add event at position {position}: {e}")
                        
                        print(f"   ✅ Added {event_count} key event annotations")
            except Exception as e:
                print(f"   ⚠️  Key event annotation failed: {e}")
        
        if chart:
            print(f"✅ Chart '{chart_title}' updated successfully")

        # ========== UPDATE INSIGHT BOX (insight_box_1) ==========
        try:
            if insight_text and str(insight_text).strip() and insight_text in df.columns and not df.empty:
                actual_insight_text = str(df[insight_text].iloc[0])
                
                if actual_insight_text and str(actual_insight_text).strip().lower() not in ['nan', 'none', '']:
                    print(f"   📝 Updating insight box from column '{insight_text}'...")
                    
                    # Find shape named "insight_box_1"
                    insight_shape = None
                    for shape in slide.shapes:
                        if shape.name == "insight_box_1":
                            insight_shape = shape
                            break
                    
                    if insight_shape and insight_shape.has_text_frame:
                        text_frame = insight_shape.text_frame
                        
                        # Split input text by newlines to get bullets
                        new_bullets = [line for line in str(actual_insight_text).split('\n') if line.strip()]
                        
                        if new_bullets:
                            # Get existing paragraphs
                            paragraphs = text_frame.paragraphs
                            
                            # Update existing paragraphs
                            min_len = min(len(paragraphs), len(new_bullets))
                            for i in range(min_len):
                                p = paragraphs[i]
                                new_text = new_bullets[i]
                                
                                # Update first run to preserve its style (font, size, etc.)
                                if p.runs:
                                    p.runs[0].text = new_text
                                    # Clear subsequent runs in same paragraph to avoid duplicates/confusion
                                    for r in p.runs[1:]:
                                        r.text = ''
                                else:
                                    p.add_run().text = new_text
                            
                            # Handle extra new bullets (Add them)
                            if len(new_bullets) > len(paragraphs):
                                # Capture style from first paragraph to apply to new ones
                                ref_font_name = None
                                ref_font_size = None
                                if paragraphs and paragraphs[0].runs:
                                    ref_run = paragraphs[0].runs[0]
                                    ref_font_name = ref_run.font.name
                                    ref_font_size = ref_run.font.size

                                for i in range(len(paragraphs), len(new_bullets)):
                                    new_p = text_frame.add_paragraph()
                                    new_run = new_p.add_run()
                                    new_run.text = new_bullets[i]
                                    
                                    # Apply captured style if available
                                    if ref_font_name:
                                        new_run.font.name = ref_font_name
                                    if ref_font_size:
                                        new_run.font.size = ref_font_size

                            # Handle excess existing paragraphs (Delete them)
                            elif len(paragraphs) > len(new_bullets):
                                # We need to access the element to delete
                                to_remove = []
                                for i in range(len(new_bullets), len(paragraphs)):
                                    to_remove.append(paragraphs[i]._element)
                                
                                for elem in to_remove:
                                    elem.getparent().remove(elem)
                                    
                            print(f"   ✅ Updated 'insight_box_1' with {len(new_bullets)} bullets")
                    else:
                        print(f"   ⚠️  Shape 'insight_box_1' not found or has no text frame")
        except Exception as e:
            print(f"   ⚠️  Insight box update failed: {str(e)}")

        # ========== UPDATE LATEST DATE (latest_date) ==========
        try:
            if latest_date_column and str(latest_date_column).strip() and latest_date_column in df.columns:
                print(f"   📅 Updating latest date from column '{latest_date_column}'...")
                
                # Get the value from the first row of date column
                if not df.empty:
                    date_val = str(df[latest_date_column].iloc[0])
                    
                    found_date_shape = False
                    for shape in slide.shapes:
                        if shape.name == "latest_date":
                            found_date_shape = True
                            if shape.has_text_frame:
                                text_frame = shape.text_frame
                                if text_frame.paragraphs:
                                    p = text_frame.paragraphs[0]
                                    if p.runs:
                                        # Update text preserving formatting
                                        p.runs[0].text = date_val
                                        # Clear subsequent runs
                                        for r in p.runs[1:]:
                                            r.text = ''
                                    else:
                                        p.add_run().text = date_val
                                    
                                    print(f"   ✅ Updated 'latest_date' to '{date_val}'")
                            break
                    
                    if not found_date_shape:
                        print(f"   ⚠️  Shape 'latest_date' not found")
        except Exception as e:
            print(f"   ⚠️  Latest date update failed: {str(e)}")

        
    async def run(self, tool_params):
        """Execute deck creation with UI loader and timeout handling"""
        print("🎨 Deck Creation Tool Started")
        
        arguments = tool_params.get("arguments", {})
        arguments["timeout_minutes"] = 60
        dataware_house = tool_params.get("dataware_house")
        
        # Extract parameters
        deck_type = arguments.get("deck_type", "monthly").lower()
        template_name = arguments.get("template_name", "")
        max_workers = arguments.get("max_parallel_workers", 21)
        timeout_minutes = arguments.get("timeout_minutes", 60)
        
        # Validate deck type
        if deck_type not in ['monthly', 'weekly']:
            error_msg = f"Error: deck_type must be 'monthly' or 'weekly'. Got: {deck_type}"
            return None, error_msg, None
        
        # # Validate max_workers range
        # if max_workers < 15 or max_workers > 20:
        #     print(f"⚠️  max_workers {max_workers} out of range, using default 20")
        #     max_workers = 20

        timeout_seconds = timeout_minutes * 60
        
        # Get configuration and template paths
        config_path = self.config_files.get(deck_type)
        template_path = template_name if template_name else self.template_files.get(deck_type)
        
        # Check files exist
        if not os.path.exists(config_path):
            error_msg = f"Error: Configuration file not found: {config_path}"
            return None, error_msg, None
        
        if not os.path.exists(template_path):
            error_msg = f"Error: Template file not found: {template_path}"
            return None, error_msg, None
        
        # Generate output filename
        current_date = datetime.datetime.now().strftime("%Y_%m_%d")
        output_filename = f"Data Strategy & Governance – Oncology GenAI Set-up_{deck_type}_{current_date}.pptx"
        output_path = f"Structured_Bot/output_files/{output_filename}"
        
        loader_msg = None
        
        try:
            # ========== START LOADER ==========
            cl.user_session.set("is_executing", True)
            loader_msg = cl.Message(
                content="",
                elements=[cl.CustomElement(name="FunctionExecutionLoader", props={
                    "isExecuting": True,
                    "functionName": f"Creating {deck_type.upper()} Deck with {max_workers} parallel workers...",
                    "message": f"Please wait while we execute queries and generate your presentation... (Timeout: {timeout_minutes} mins)"
                })]
            )
            await loader_msg.send()
            
            # Step 1: Read configuration
            print(f"📖 Reading configuration from {config_path}")
            df_config = pd.read_excel(config_path)
            print(f"Found {len(df_config)} charts to create")
            
            # Validate required columns
            required_cols = ['slide_number', 'title', 'sql_code']
            missing_cols = [col for col in required_cols if col not in df_config.columns]
            if missing_cols:
                raise ValueError(f"Missing required columns in config: {missing_cols}")
            
            # Optional columns with defaults
            if 'x_axis_column' not in df_config.columns:
                df_config['x_axis_column'] = None
            if 'y_axis_column' not in df_config.columns:
                df_config['y_axis_column'] = None
            if 'lower_threshold' not in df_config.columns:
                df_config['lower_threshold'] = None
            if 'upper_threshold' not in df_config.columns:
                df_config['upper_threshold'] = None
            if 'median_threshold' not in df_config.columns:
                df_config['median_threshold'] = None
            if 'chart_type' not in df_config.columns:
                df_config['chart_type'] = None
            if 'key_event_column' not in df_config.columns:
                df_config['key_event_column'] = None
            if 'python_code' not in df_config.columns:
                df_config['python_code'] = None
            
            # Step 2: Execute all queries in parallel
            print(f"🚀 Executing {len(df_config)} SQL queries in parallel (workers: {max_workers}, timeout: {timeout_minutes} mins)...")
            
            queries = {
                idx: row['sql_code'] 
                for idx, row in df_config.iterrows()
            }
            
            query_results = {}
            start_time = time.time()
            
            # Create async warehouse function
            warehouse = cl.make_async(dataware_house.get_data_from_wareshouse)
            
            async def execute_single_query(query_id, sql_query):
                """Execute a single query using warehouse"""
                try:
                    print(f"[Query {query_id}] Executing...")
                    query_start = time.time()
                    
                    result = await warehouse(sql_query)
                    
                    elapsed = time.time() - query_start
                    print(f"[Query {query_id}] ✅ Completed in {elapsed:.2f}s - {len(result)} rows")
                    
                    return query_id, result, "SUCCESS"
                except Exception as e:
                    error_msg = str(e)
                    print(f"[Query {query_id}] ❌ Failed: {error_msg}")
                    return query_id, None, error_msg
            
            # Execute all queries with timeout
            print(f"Executing queries using asyncio.gather with {timeout_seconds}s timeout...")
            query_tasks = [
                execute_single_query(query_id, sql_query) 
                for query_id, sql_query in queries.items()
            ]
            
            try:
                results = await asyncio.wait_for(
                    asyncio.gather(*query_tasks, return_exceptions=True),
                    timeout=timeout_seconds
                )
                
                # Process results
                completed = 0
                for result in results:
                    if isinstance(result, Exception):
                        print(f"⚠️  Query exception: {str(result)}")
                        continue
                    
                    query_id, df, status = result
                    if status == "SUCCESS" and df is not None and not df.empty:
                        query_results[query_id] = df
                        completed += 1
                        print(f"✅ Progress: {completed}/{len(queries)} queries completed ({int((completed/len(queries))*100)}%)")
                
                total_time = time.time() - start_time
                
                print(f"\n{'='*60}")
                print(f"Parallel execution completed in {total_time:.2f}s")
                print(f"Successful queries: {len(query_results)}/{len(queries)}")
                print(f"{'='*60}\n")
                
                if len(query_results) == 0:
                    raise Exception("No queries executed successfully!")
                
            except asyncio.TimeoutError:
                elapsed_time = time.time() - start_time
                timeout_msg = (
                    f"⏱️ **Databricks Query Timeout**\n\n"
                    f"The queries took longer than {timeout_minutes} minutes to execute (elapsed: {elapsed_time/60:.1f} mins).\n\n"
                    f"**Please try again in a few minutes.**"
                )
                
                print(f"\n{'='*60}")
                print(f"❌ TIMEOUT: Queries exceeded {timeout_minutes} minute limit")
                print(f"{'='*60}\n")
                
                cl.user_session.set("is_executing", False)
                await loader_msg.remove()
                
                await cl.Message(content=timeout_msg).send()
                
                return timeout_msg, timeout_msg + " DO NOT TRIGGER THE DECK CREATION TOOL AGAIN!", None
            
            # Step 3: Load template
            print(f"📄 Loading template: {template_path}")
            from pptx import Presentation
            prs = Presentation(template_path)

            self._debug_print_chart_titles(prs)  # Debug: Print all chart titles

            
            print(f"🎨 Updating {len(query_results)} charts in template...")
            
            successful_charts = 0
            failed_charts = []
            
            # Step 4: Update each chart using config
            for idx, row in df_config.iterrows():
                slide_num = int(row['slide_number'])
                title = row['title']
                
                if idx not in query_results:
                    print(f"⚠️  Skipping slide {slide_num} - Query failed")
                    failed_charts.append((slide_num, title))
                    continue
                
                df_data = query_results[idx]
                
                try:
                    # Get all configuration parameters
                    x_axis_col = row.get('x_axis_column')
                    y_axis_col = row.get('y_axis_column')
                    lower_thresh = row.get('lower_threshold')
                    upper_thresh = row.get('upper_threshold')
                    median_thresh = row.get('median_threshold')
                    chart_type_val = row.get('chart_type')
                    key_event_col = row.get('key_event_column')
                    python_code = row.get('python_code')
                    insight_text = row.get('insight_text')
                    latest_date_col = row.get('latest_date')
                    
                    # DEBUG: Print configuration values
                    print(f"   🔍 Configuration for '{title}':")
                    print(f"      x_axis_column: {x_axis_col}")
                    print(f"      y_axis_column: {y_axis_col}")
                    print(f"      median_threshold: {median_thresh}")
                    print(f"      lower_threshold: {lower_thresh}")
                    print(f"      upper_threshold: {upper_thresh}")
                    print(f"      chart_type: {chart_type_val}")
                    print(f"      key_event_column: {key_event_col}")
                    print(f"      insight_text_column: {insight_text}")
                    print(f"      latest_date: {latest_date_col}")
                    print(f"   📊 Available columns in SQL result: {list(df_data.columns)}")
                    
                    # Update chart using new method
                    self._update_chart_from_config(
                        prs=prs,
                        slide_number=slide_num,
                        chart_title=title,
                        df=df_data,
                        x_axis_column=x_axis_col,
                        y_axis_column=y_axis_col,
                        lower_threshold=lower_thresh,
                        upper_threshold=upper_thresh,
                        median_threshold=median_thresh,
                        chart_type=chart_type_val,
                        key_event_column=key_event_col,
                        python_code=python_code,
                        insight_text=insight_text,
                        latest_date_column=latest_date_col
                    )
                    
                    successful_charts += 1
                    print(f"✅ Chart updated on slide {slide_num}: {title}")
                    
                except Exception as e:
                    print(f"❌ Failed to update chart on slide {slide_num}: {str(e)}")
                    import traceback
                    traceback.print_exc()
                    failed_charts.append((slide_num, title))
            
            # Step 5: Save presentation
            print(f"💾 Saving presentation to {output_path}")
            prs.save(output_filename)
            
            # ========== STOP LOADER ==========
            cl.user_session.set("is_executing", False)
            await loader_msg.remove()
            
            # Send completion message
            completion_message = f"Deck creation completed successfully!\n\n✅ {successful_charts} slides updated"
            if failed_charts:
                completion_message += f"\n⚠️ {len(failed_charts)} charts failed"
            completion_message += "\n\nThanks for your patience!"
            
            await cl.Message(content=completion_message).send()
            
            # Display and delete
            await files_images_handler.pptx_file_handler([output_filename])
            os.remove(output_filename)
            
            tool_params["result"] = f"Deck created successfully: {output_filename}"
            cl.user_session.set("tool_params", tool_params)
            
            return f"Deck created successfully: {output_filename}", f"Deck created successfully: {output_filename}", output_path
            
        except Exception as e:
            error_msg = f"❌ Error creating deck: {str(e)}"
            print(error_msg)
            import traceback
            traceback.print_exc()
            
            if loader_msg:
                cl.user_session.set("is_executing", False)
                await loader_msg.remove()
                await cl.Message(content=f"Deck creation failed.\n\n{error_msg}").send()
            
            return error_msg, error_msg, None

# Update ToolCallsExecutor to include the updated tool
class ToolCallsExecutor:
    """
    This class handles the execution of tool calls in the chatbot.
    """
    def __init__(self):
        """Initialize with shared resources."""
        self.opensearch = OpensearchExecutionManager()
        
        self.tool_map = {
            "run_hybrid_search": lambda: RunHybridSearch(self.opensearch),
            "run_sql_on_opensearch": lambda: RunSqlOnOpensearch(self.opensearch),
            "calculate_date_ranges": lambda: CalculateDateRange(),
            "generate_sql_code": lambda: GenerateSQLCode(),
            "python_after_sql": lambda: PythonAfterSQL(),
            "insights_response": lambda: InsightsResponse(),
            "text_processor": lambda: TextProcessor(),
            "calculate_nsp_capture_ratio": lambda: CalculateNSPCaptureRatio(),  
            "create_presentation_deck": lambda: CreatePresentationDeck(),  # NEW TOOL
# UPDATED TOOL
        }
    
    def get_tool(self, function_name):
        """
        Returns the appropriate tool class instance based on function name.
        
        Args:
            function_name (str): Name of the function/tool to get
            
        Returns:
            BaseTool: An instance of the tool class
            
        Raises:
            ValueError: If function name is not recognized
        """
        if function_name in self.tool_map:
            return self.tool_map[function_name]()
        else:
            raise ValueError(f"Unknown function: {function_name}")





















    

        
        