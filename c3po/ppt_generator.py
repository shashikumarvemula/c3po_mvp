""" 
ENHANCED LLM-POWERED CODE-TO-PPT GENERATOR
==========================================
Complete implementation with support for:
- Structured bot (original functionality with pie charts, single-value charts)
- RM_insights bot (topic-sentiment analysis with summarized insights)

Key Features:
- LLM-powered code analysis for better plot recognition
- Editable charts in PowerPoint with proper data mapping
- Smart insights fitting with optimized layout
- Support for complex multi-series visualizations
- Enhanced error handling and debugging
- Pie chart support for any data structure
- Single-value/hardcoded category charts
- RM_insights bot with topic-wise sentiment breakdown
"""

import pandas as pd
import numpy as np
import re
import os
import json
import asyncio
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.chart import XL_CHART_TYPE
from pptx.chart.data import CategoryChartData, XyChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_THEME_COLOR
from datetime import datetime
import warnings
warnings.filterwarnings('ignore')
import chainlit as cl

class GeneratePPT:
    """Enhanced LLM-powered converter with intelligent plot recognition"""
    
    def __init__(self):
        self.prs = None
        self.detected_plots = []
        self.insights = ""
        self.llm_client = ""
        self.model_id = ""
        
        # Enhanced professional color palette
        self.chart_colors = [
            RGBColor(31, 119, 180),   # Blue
            RGBColor(255, 127, 14),   # Orange  
            RGBColor(44, 160, 44),    # Green
            RGBColor(214, 39, 40),    # Red
            RGBColor(148, 103, 189),  # Purple
            RGBColor(140, 86, 75),    # Brown
            RGBColor(227, 119, 194),  # Pink
            RGBColor(127, 127, 127),  # Gray
            RGBColor(188, 189, 34),   # Olive
            RGBColor(23, 190, 207),   # Cyan
        ]

    # ==================== BOT TYPE DETECTION ====================
    
    def is_structured_bot(self):
        """Check if current bot is Structured bot"""
        try:
            return cl.user_session.get("chat_profile", "") == os.getenv("BOT_TYPE_STRUCTURED")
        except:
            return False
    
    def is_rm_insights_bot(self):
        """Check if current bot is RM_insights"""
        try:
            return cl.user_session.get("chat_profile", "") == os.getenv("BOT_TYPE_RM_INSIGHTS")
        except:
            return False

    # ==================== RM_INSIGHTS BOT METHODS ====================
    
    async def summarize_insights_to_main_points(self, insights_text):
        """Summarize insights to 2-3 main points using LLM"""
        try:
            if not self.llm_client or not insights_text or not insights_text.strip():
                return insights_text
            
            system_content = """You are an expert at summarizing insights into concise main points.
Your task:
1. Read the provided insights carefully
2. Extract the 2-3 MOST IMPORTANT main points
3. Each point should be ONE clear sentence (max 25 words)
4. Focus on actionable insights and key findings
5. Use bullet point format

Return ONLY the bullet points, nothing else. Format:
- [First main point]
- [Second main point]
- [Third main point if applicable]"""

            user_content = f"""Summarize these insights into 2-3 main points:

{insights_text}

Return only the bullet points."""

            response = await self.llm_client.chat.completions.create(
                messages=[
                    {"role": "system", "content": system_content},
                    {"role": "user", "content": user_content}
                ],
                model=self.model_id,
                max_tokens=300,
                temperature=0.3,
                stream=False
            )
            
            summarized = response.choices[0].message.content.strip()
            print(f"✅ Summarized insights: {summarized}")
            return summarized
            
        except Exception as e:
            print(f"⚠️ Error summarizing insights: {e}")
            # Fallback: simple truncation to first 3 sentences
            sentences = insights_text.split('. ')
            if len(sentences) > 3:
                return '. '.join(sentences[:3]) + '.'
            return insights_text

    def get_rm_insights_topic_data(self, df):
        """Process data to create topic-wise sentiment breakdown"""
        try:
            topic_sentiment_data = {}
            topic_insight_ids = {}
            
            for index in df.index:
                try:
                    sentiment_data = df.loc[index, 'Sentiment']
                    insight_id = df.loc[index, 'Insight id']
                    
                    if pd.notna(sentiment_data) and sentiment_data != '':
                        if isinstance(sentiment_data, str):
                            sentiment_dict = json.loads(sentiment_data.replace("'", '"'))
                        elif isinstance(sentiment_data, dict):
                            sentiment_dict = sentiment_data
                        else:
                            continue
                        
                        for topic, sentiment in sentiment_dict.items():
                            if topic.strip() != '' and not topic.startswith('Excluded:'):
                                if topic not in topic_sentiment_data:
                                    topic_sentiment_data[topic] = {
                                        'Positive': 0, 'Negative': 0, 'Neutral': 0, 'No sentiment': 0
                                    }
                                    topic_insight_ids[topic] = set()
                                
                                if sentiment in ['Positive', 'Negative', 'Neutral']:
                                    topic_sentiment_data[topic][sentiment] += 1
                                else:
                                    topic_sentiment_data[topic]['No sentiment'] += 1
                                
                                if pd.notna(insight_id):
                                    topic_insight_ids[topic].add(insight_id)
                except:
                    continue
            
            structured_data = []
            for topic, sentiments in topic_sentiment_data.items():
                if sum(sentiments.values()) > 0:
                    structured_data.append({
                        'Topic': topic,
                        'Positive': sentiments['Positive'],
                        'Negative': sentiments['Negative'],
                        'Neutral': sentiments['Neutral'],
                        'No sentiment': sentiments['No sentiment'],
                        'Total': sum(sentiments.values()),
                        'Unique_IDs': len(topic_insight_ids.get(topic, set()))
                    })
            
            structured_data.sort(key=lambda x: x['Unique_IDs'], reverse=True)
            if len(structured_data) > 10:
                structured_data = structured_data[:10]
            
            return structured_data
        except Exception as e:
            print(f"⚠️ Error processing topic data: {e}")
            return []

    def get_rm_insights_stats(self, df):
        """Extract statistics for RM_insights"""
        try:
            stats = {}
            total_unique_ids = set()
            
            for index in df.index:
                insight_id = df.loc[index, 'Insight id']
                if pd.notna(insight_id):
                    total_unique_ids.add(insight_id)
            
            stats['total_unique_ids'] = len(total_unique_ids)
            
            # Question
            stats['question'] = ""
            for col in ['Question', 'question']:
                if col in df.columns:
                    q = df[col].dropna()
                    if len(q) > 0:
                        stats['question'] = str(q.iloc[0])
                        break
            
            if not stats['question']:
                try:
                    result_df = cl.user_session.get('tool_params', {}).get('result')
                    if result_df is not None:
                        for col in ['Question', 'question']:
                            if col in result_df.columns:
                                q = result_df[col].dropna()
                                if len(q) > 0:
                                    stats['question'] = str(q.iloc[0])
                                    break
                except:
                    pass
            
            # Insight
            stats['insight'] = ""
            try:
                session_insights = cl.user_session.get("insights", "")
                if session_insights and session_insights.strip():
                    stats['insight'] = session_insights.strip()
            except:
                pass
            
            if not stats['insight']:
                for col in ['Insight', 'insight']:
                    if col in df.columns:
                        ins = df[col].dropna()
                        if len(ins) > 0:
                            stats['insight'] = str(ins.iloc[0])
                            break
            
            if not stats['insight']:
                try:
                    result_df = cl.user_session.get('tool_params', {}).get('result')
                    if result_df is not None:
                        for col in ['Insight', 'insight']:
                            if col in result_df.columns:
                                ins = result_df[col].dropna()
                                if len(ins) > 0:
                                    stats['insight'] = str(ins.iloc[0])
                                    break
                except:
                    pass
            
            return stats
        except Exception as e:
            return {'total_unique_ids': 0, 'question': '', 'insight': ''}

    def create_rm_insights_chart_data(self, topic_data):
        """Create chart data with topics as categories"""
        try:
            chart_data = CategoryChartData()
            
            categories = []
            for item in topic_data:
                topic = item['Topic']
                if len(topic) > 50:
                    topic = topic[:47] + "..."
                categories.append(topic)
            
            chart_data.categories = categories
            
            no_sent = [item['No sentiment'] for item in topic_data]
            negative = [item['Negative'] for item in topic_data]
            neutral = [item['Neutral'] for item in topic_data]
            positive = [item['Positive'] for item in topic_data]
            
            chart_data.add_series('No sentiment', no_sent)
            chart_data.add_series('Negative', negative)
            chart_data.add_series('Neutral', neutral)
            chart_data.add_series('Positive', positive)
            
            return chart_data, topic_data
        except Exception as e:
            return None, None

    async def create_rm_insights_slide(self, df, plot_spec):
        """Create RM_insights slide with topic-wise visualization and summarized insights"""
        try:
            self.prs = Presentation()
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])
            
            topic_data = self.get_rm_insights_topic_data(df)
            if not topic_data:
                return None
            
            stats = self.get_rm_insights_stats(df)
            
            # Title
            title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.2), Inches(9), Inches(0.4))
            title_frame = title_box.text_frame
            title_frame.text = "Topic-Sentiment Distribution Analysis"
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(16)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(44, 62, 80)
            title_para.alignment = PP_ALIGN.CENTER
            
            # Counts
            text_y = 0.65
            count_box = slide.shapes.add_textbox(Inches(0.5), Inches(text_y), Inches(9), Inches(0.25))
            count_frame = count_box.text_frame
            count_frame.text = f"Total Insights: {len(df)}"
            count_para = count_frame.paragraphs[0]
            count_para.font.size = Pt(11)
            count_para.font.bold = True
            count_para.font.color.rgb = RGBColor(44, 62, 80)
            count_para.alignment = PP_ALIGN.CENTER
            text_y += 0.3
            
            # Question
            if stats.get('question'):
                q_box = slide.shapes.add_textbox(Inches(6.7), Inches(0.95), Inches(2.8), Inches(0.5))
                q_frame = q_box.text_frame
                q_frame.word_wrap = True
                q_para = q_frame.paragraphs[0]
                q_para.text = f"Question: {stats['question']} "
                q_para.font.size = Pt(10)
                q_para.font.name = 'Trebuchet MS'
                q_para.font.bold = False
                q_para.font.italic = False
                q_para.font.color.rgb = RGBColor(192, 192, 192)
                
                # Add insight count in italic
                run = q_para.add_run()
                run.text = f"({len(df)} insights)"
                run.font.size = Pt(10)
                run.font.name = 'Trebuchet MS'
                run.font.italic = True
                run.font.color.rgb = RGBColor(192, 192, 192)
                
                text_y += 0.45
            
            # Chart
            chart_data, _ = self.create_rm_insights_chart_data(topic_data)
            if not chart_data:
                return None
            
            chart = slide.shapes.add_chart(
                XL_CHART_TYPE.BAR_STACKED,
                Inches(0.5), Inches(text_y), Inches(6), Inches(4.5),
                chart_data
            ).chart
            
            chart.has_title = False
            
            if hasattr(chart, 'value_axis') and chart.value_axis:
                chart.value_axis.visible = False
                chart.value_axis.has_major_gridlines = False
                chart.value_axis.has_minor_gridlines = False
            
            if hasattr(chart, 'category_axis') and chart.category_axis:
                chart.category_axis.tick_labels.font.size = Pt(10)
                chart.category_axis.tick_labels.font.name = 'Trebuchet MS'
                chart.category_axis.tick_labels.font.color.rgb = RGBColor(64, 64, 64)
            
            colors = [
                RGBColor(177, 200, 220),  # No sentiment - Light cement blue (#B1C8DC)
                RGBColor(255, 0, 0),       # Negative - Red (#FF0000)
                RGBColor(255, 255, 0),     # Neutral - Yellow (#FFFF00)
                RGBColor(0, 176, 80)       # Positive - Green (#00B050)
            ]
            
            if hasattr(chart, 'series') and chart.series:
                for i, series in enumerate(chart.series):
                    if i < len(colors):
                        # Apply color to series fill
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = colors[i]
                        # Also apply to each data point
                        for point in series.points:
                            point.format.fill.solid()
                            point.format.fill.fore_color.rgb = colors[i]                    
                    if hasattr(series, 'has_data_labels'):
                        series.has_data_labels = True
                        if hasattr(series, 'data_labels'):
                            series.data_labels.show_value = True
                            series.data_labels.show_category_name = False
                            series.data_labels.show_legend_key = False
                            series.data_labels.position = 1
                            series.data_labels.font.size = Pt(10)
                            series.data_labels.font.bold = True
                            series.data_labels.font.color.rgb = RGBColor(0, 0, 0)
            
            if hasattr(chart, 'has_legend'):
                chart.has_legend = True
                from pptx.enum.chart import XL_LEGEND_POSITION
                chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                chart.legend.include_in_layout = False
                if hasattr(chart.legend, 'font'):
                    chart.legend.font.size = Pt(10)
            
            # Key Insights - SUMMARIZED TO 2-3 POINTS
            insights_box = slide.shapes.add_textbox(Inches(6.7), Inches(text_y), Inches(2.8), Inches(3.5))
            insights_frame = insights_box.text_frame
            insights_frame.word_wrap = True
            insights_frame.clear()
            insights_frame.margin_left = Inches(0.1)
            insights_frame.margin_right = Inches(0.1)
            insights_frame.margin_top = Inches(0.1)
            
            # Get and summarize insights
            original_insights = stats.get('insight', "Analysis of safety data and adverse events.")
            print(f"📝 Original insights length: {len(original_insights)} chars")
            
            # Summarize to 2-3 main points
            summarized_insights = await self.summarize_insights_to_main_points(original_insights)
            print(f"✅ Summarized insights: {summarized_insights}")
            
            # Header for Key Insights
            p = insights_frame.paragraphs[0]
            p.text = "Key Insight:"
            p.font.size = Pt(12)
            p.font.name = 'Trebuchet MS'
            p.font.bold = True
            p.font.color.rgb = RGBColor(44, 62, 80)
            p.space_after = Pt(8)
            
            # Add summarized points
            lines = summarized_insights.split('\n')
            for line in lines:
                line = line.strip()
                if line:
                    # Remove existing bullet if present
                    if line.startswith('•'):
                        line = line[1:].strip()
                    elif line.startswith('-'):
                        line = line[1:].strip()
                    elif line.startswith('*'):
                        line = line[1:].strip()
                    
                    p = insights_frame.add_paragraph()
                    p.text = f"• {line}"
                    p.font.size = Pt(10)
                    p.font.name = 'Trebuchet MS'
                    p.font.color.rgb = RGBColor(64, 64, 64)
                    p.space_after = Pt(6)
                    p.level = 0
            
            insights_box.fill.solid()
            insights_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
            insights_box.line.color.rgb = RGBColor(180, 180, 180)
            insights_box.line.width = Pt(1)
            
            return self.prs
        except Exception as e:
            print(f"❌ Error: {e}")
            import traceback
            traceback.print_exc()
            return None

    # ==================== STRUCTURED BOT - LLM ANALYSIS METHODS ====================
    
    async def analyze_code_with_llm(self, python_code, df_columns, df_sample):
        """Use LLM to intelligently analyze Python visualization code"""
        if not self.llm_client:
            print("⚠️ No LLM client provided, falling back to regex parsing")
            return await self.fallback_plot_detection(python_code)
       
        system_content = """You are an expert Python data visualization analyst. Analyze the given Python code and extract detailed plot specifications.
Your task:
1. Identify ALL plots/charts in the code
2. Extract exact column names, plot types, styling, and parameters, especially HUE parameters
3. Handle multi-series plots, subplots, hue-based groupings, and complex visualizations
4. Return a structured JSON response
5. Pay special attention to HUE parameters - when hue is specified, it creates multiple series
6. You have to grasp exact python code and extract the specifications, so that it will be in deck exact as expected.

IMPORTANT HUE HANDLING:
- When you detect 'hue=' parameter in seaborn plots, set "multi_series": true
- The hue column creates separate series for each unique value
- This is different from explicit multi-series with multiple y columns

SPECIAL PIE CHART HANDLING:
- For plt.pie() or ax.pie() calls, set type="pie"
- Extract the data source (sizes parameter) and labels parameter
- For pie charts, x and y can be null since they use different data structure
- Look for calculated values like "100 - percentage" patterns
- Extract series_names from labels parameter or variable assignments
- Pay attention to custom colors in the colors parameter


Return JSON format:
{
    "plots": [
        {
            "type": "line|bar|scatter|hist|box|heatmap|pie",
            "x": "column_name or expression the acutal names",
            "y": "column_name or expression or list for multi-series" the actual names,
            "hue": "column_name for grouping (VERY IMPORTANT - creates multiple series)",
            "title": "chart title",
            "xlabel": "x-axis label",
            "ylabel": "y-axis label",
            "multi_series": true/false (TRUE if hue is present OR multiple y columns),
            "series_names": ["series1", "series2"] (if multi-series),
            "style_params": {
                "color": "color_specification",
                "marker": "marker_style",
                "linestyle": "line_style",
                "alpha": 0.7,
                "figsize": [width, height]
            },
            "subplot_info": {
                "is_subplot": true/false,
                "position": [row, col] (if subplot)
            },
           
        }
    ],
    "analysis_confidence": 0.95,
    "detected_libraries": ["matplotlib", "seaborn", "plotly"],
    "complexity_level": "simple|moderate|complex"
}

Critical Rules:
- If you see 'hue=' in seaborn code, ALWAYS set "multi_series": true
- Extract the exact hue column name from hue='column_name'
- Be precise with column names - extract exact strings
- Handle both matplotlib and seaborn syntax
- Extract titles from plt.title(), ax.set_title(), or title= parameters
- Look for styling parameters like colors, markers, line styles"""


        # Prepare sample data info for context
        sample_data_info = f"""
Available DataFrame columns: {list(df_columns)}
Sample data (first 3 rows):
{df_sample.to_string()}
Data types: {df_sample.dtypes.to_dict()}
"""

        user_content = f"""Analyze this Python visualization code and extract plot specifications:

PYTHON CODE:
```python
{python_code}
DATAFRAME CONTEXT:
{sample_data_info}
Please provide a detailed JSON analysis of all plots in the code."""
        try:
            response = await self.llm_client.chat.completions.create(
                messages=[
                    {"role": "system", "content": system_content},
                    {"role": "user", "content": user_content}
                ],
                model=self.model_id,
                max_tokens=2048,
                temperature=0.1,
                top_p=1,
                stream=False
            )
           
            response_text = response.choices[0].message.content
            print(f"🤖 LLM Analysis Response: {response_text[:200]}...")
           
            # Extract JSON from response
            try:
                # Find JSON block in response
                json_start = response_text.find('{')
                json_end = response_text.rfind('}') + 1
                if json_start != -1 and json_end > json_start:
                    json_str = response_text[json_start:json_end]
                    analysis = json.loads(json_str)
                    return analysis
                else:
                    print("⚠️ No valid JSON found in LLM response")
                    return await self.fallback_plot_detection(python_code)
            except json.JSONDecodeError as e:
                print(f"⚠️ JSON parsing error: {e}")
                return await self.fallback_plot_detection(python_code)
           
        except Exception as e:
            print(f"⚠️ LLM analysis failed: {e}")
            return await self.fallback_plot_detection(python_code)

    async def fallback_plot_detection(self, python_code):
        """Fallback regex-based plot detection"""
        plots = []
   
        # Enhanced regex patterns
        patterns = [
        (r'sns\.lineplot\s*\([^)]*x\s*=\s*["\']([^"\']+)["\'][^)]*y\s*=\s*["\']([^"\']+)["\'][^)]*hue\s*=\s*["\']([^"\']+)["\']', 'line_hue'),
        (r'sns\.scatterplot\s*\([^)]*x\s*=\s*["\']([^"\']+)["\'][^)]*y\s*=\s*["\']([^"\']+)["\'][^)]*hue\s*=\s*["\']([^"\']+)["\']', 'scatter_hue'),
        (r'sns\.barplot\s*\([^)]*x\s*=\s*["\']([^"\']+)["\'][^)]*y\s*=\s*["\']([^"\']+)["\'][^)]*hue\s*=\s*["\']([^"\']+)["\']', 'bar_hue'),
   
        # Matplotlib patterns
        (r'plt\.plot\s*\(\s*([^,\)]+)\s*,\s*([^,\)]+)[^)]*label\s*=\s*["\']([^"\']+)["\']', 'line_labeled'),
        (r'plt\.plot\s*\(\s*([^,\)]+)(?:\s*,\s*([^,\)]+))?', 'line'),
        (r'plt\.scatter\s*\(\s*([^,\)]+)\s*,\s*([^,\)]+)', 'scatter'),
        (r'plt\.bar\s*\(\s*([^,\)]+)\s*,\s*([^,\)]+)', 'bar'),
        (r'plt\.hist\s*\(\s*([^,\)]+)', 'hist'),
   
        # Regular seaborn patterns
        (r'sns\.lineplot\s*\([^)]*x\s*=\s*["\']([^"\']+)["\'][^)]*y\s*=\s*["\']([^"\']+)["\']', 'line'),
        (r'sns\.scatterplot\s*\([^)]*x\s*=\s*["\']([^"\']+)["\'][^)]*y\s*=\s*["\']([^"\']+)["\']', 'scatter'),
        (r'sns\.barplot\s*\([^)]*x\s*=\s*["\']([^"\']+)["\'][^)]*y\s*=\s*["\']([^"\']+)["\']', 'bar'),


        (r'plt\.pie\s*\(\s*([^,\)]+)(?:\s*,\s*labels\s*=\s*([^,\)]+))?', 'pie'),
        (r'ax\.pie\s*\(\s*([^,\)]+)(?:\s*,\s*labels\s*=\s*([^,\)]+))?', 'pie'),

        ]
   
        # Extract title
        title_match = re.search(r'plt\.title\s*\(\s*["\']([^"\']+)["\']', python_code, re.IGNORECASE)
        title = title_match.group(1) if title_match else "Data Visualization"
   
        # Process patterns
        for pattern, plot_type in patterns:
            matches = re.findall(pattern, python_code, re.IGNORECASE)
            for match in matches:
                if plot_type.endswith('_hue'):
                    # Hue-based plot
                    base_type = plot_type.replace('_hue', '')
                    plots.append({
                        "type": base_type,
                        "x": self.clean_column_name(match[0]),
                        "y": self.clean_column_name(match[1]),
                        "hue": self.clean_column_name(match[2]),
                        "title": title,
                        "multi_series": True
                    })
                elif plot_type == 'line_labeled':
                    plots.append({
                        "type": "line",
                        "x": self.clean_column_name(match[0]),
                        "y": [self.clean_column_name(match[1])],
                        "title": title,
                        "multi_series": True,
                        "series_names": [match[2]]
                    })
                elif len(match) >= 2:
                    plots.append({
                        "type": plot_type.replace('_hue', ''),
                        "x": self.clean_column_name(match[0]),
                        "y": self.clean_column_name(match[1]),
                        "hue": self.clean_column_name(match[2]) if len(match) > 2 else None,
                        "title": title,
                        "multi_series": len(match) > 2
                    })

   
        return {
            "plots": plots,
            "analysis_confidence": 0.7,
            "detected_libraries": ["matplotlib", "seaborn"],
            "complexity_level": "simple"
        }

    def clean_column_name(self, expr):
        """Clean and extract column name from expression"""
        expr = str(expr).strip()
   
        # Remove common patterns
        patterns = [
            r'df\[["\']([^"\']+)["\']\]',
            r'data\[["\']([^"\']+)["\']\]',
            r'^["\']([^"\']+)["\']$'
        ]
   
        for pattern in patterns:
            match = re.search(pattern, expr)
            if match:
                return match.group(1)
   
        # Clean prefixes
        for prefix in ['df.', 'data.']:
            if expr.startswith(prefix):
                expr = expr[len(prefix):]
   
        return expr.replace('"', '').replace("'", '')

    # ==================== STRUCTURED BOT - SINGLE VALUE CHART METHODS ====================
   
    def create_single_value_chart_data(self, df, plot_spec):
        """Handle single value charts where x is a hardcoded category - GENERALIZED VERSION"""
        try:
            x_col = plot_spec['x']
            y_col = plot_spec['y']
           
            print(f"🔧 Processing single value chart: x={x_col} (type: {type(x_col)}), y={y_col}")
           
            # Normalize the category name
            category_name = self.normalize_single_value_category(x_col)
            print(f"📋 Normalized category: '{category_name}'")
           
            # Create chart data
            chart_data = CategoryChartData()
            chart_data.categories = [category_name]
           
            # Get the value from the DataFrame
            if y_col in df.columns:
                value = df[y_col].iloc[0]
                values = [float(value) if pd.notna(value) else 0.0]
                chart_data.add_series('Value', values)
                print(f"✅ Created single value chart: '{category_name}' = {values[0]}")
                return chart_data, 'single_value_bar'
            else:
                print(f"❌ Y column '{y_col}' not found in DataFrame")
                return None, None
           
        except Exception as e:
            print(f"❌ Error creating single value chart: {e}")
            import traceback
            traceback.print_exc()
            return None, None
   
    def is_single_value_chart(self, x_col, df):
        """Detect if this is a single-value chart with hardcoded categories"""
        try:
            # Case 1: x_col is a list (like ['Key Accounts'])
            if isinstance(x_col, list):
                return True
           
            # Case 2: x_col is a string representation of a list (like "['Key Accounts']")
            if isinstance(x_col, str) and (x_col.startswith('[') and x_col.endswith(']')):
                return True
           
            # Case 3: x_col is a quoted string (like "'Key Accounts'")
            if isinstance(x_col, str) and (x_col.startswith("'") or x_col.startswith('"')):
                return True
           
            # Case 4: x_col is not in DataFrame columns but looks like a category
            if isinstance(x_col, str) and x_col not in df.columns:
                # Check if it contains common single-value chart patterns
                single_value_patterns = [
                    'average', 'total', 'count', 'sum', 'key accounts',
                    'overall', 'summary', 'aggregate', 'combined'
                ]
                if any(pattern in x_col.lower() for pattern in single_value_patterns):
                    return True
           
            return False
        except Exception as e:
            print(f"⚠️ Error detecting single value chart: {e}")
            return False
           
    def normalize_single_value_category(self, x_col):
        """Convert various x_col formats to a clean category name"""
        try:
            # Case 1: Already a list
            if isinstance(x_col, list):
                return x_col[0] if x_col else "Category"
           
            # Case 2: String representation of a list
            if isinstance(x_col, str) and x_col.startswith('[') and x_col.endswith(']'):
                import ast
                try:
                    parsed = ast.literal_eval(x_col)
                    if isinstance(parsed, list) and parsed:
                        return parsed[0]
                except:
                    # Fallback: extract content between brackets
                    content = x_col.strip('[]').strip().strip("'\"")
                    return content if content else "Category"
           
            # Case 3: Quoted string
            if isinstance(x_col, str) and (x_col.startswith("'") or x_col.startswith('"')):
                return x_col.strip("'\"")
           
            # Case 4: Plain string
            return str(x_col)
           
        except Exception as e:
            print(f"⚠️ Error normalizing category: {e}")
            return "Category"

    # ==================== STRUCTURED BOT - PIE CHART METHODS ====================

    def create_generalized_pie_chart_data(self, df, plot_spec):
        """
        Create pie chart data in a generalized way that works with any pie chart code
        Analyzes the DataFrame to find the best columns for pie chart representation
        """
        try:
            print(f"🥧 Creating generalized pie chart data for: {plot_spec.get('title', 'Pie Chart')}")
           
            from pptx.chart.data import CategoryChartData
            chart_data = CategoryChartData()
           
            # Strategy 1: Look for percentage columns (most common in pie charts)
            percentage_cols = [col for col in df.columns if 'percentage' in col.lower() or 'percent' in col.lower() or '%' in col]
           
            if percentage_cols:
                print(f"📊 Found percentage column: {percentage_cols[0]}")
                return self._create_pie_from_percentage_column(df, plot_spec, percentage_cols[0], chart_data)
           
            # Strategy 2: Look for count/frequency columns
            count_cols = [col for col in df.columns if any(word in col.lower() for word in ['count', 'freq', 'total', 'sum', 'amount', 'sales', 'revenue'])]
            text_cols = df.select_dtypes(include=['object', 'string']).columns.tolist()
           
            if count_cols and text_cols:
                print(f"📊 Found count column: {count_cols[0]} with categories: {text_cols[0]}")
                return self._create_pie_from_count_columns(df, plot_spec, text_cols[0], count_cols[0], chart_data)
           
            # Strategy 3: Use any numeric column with categorical data
            numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
            if numeric_cols and text_cols:
                print(f"📊 Using numeric column: {numeric_cols[0]} with categories: {text_cols[0]}")
                return self._create_pie_from_count_columns(df, plot_spec, text_cols[0], numeric_cols[0], chart_data)
           
            # Strategy 4: Create pie chart from value counts of categorical column
            if text_cols:
                print(f"📊 Creating pie from value counts of: {text_cols[0]}")
                return self._create_pie_from_value_counts(df, plot_spec, text_cols[0], chart_data)
           
            # Strategy 5: If only numeric data, create distribution pie
            if numeric_cols:
                print(f"📊 Creating distribution pie from: {numeric_cols[0]}")
                return self._create_pie_from_numeric_distribution(df, plot_spec, numeric_cols[0], chart_data)
           
            print("❌ No suitable data found for pie chart")
            return None, None
           
        except Exception as e:
            print(f"❌ Error creating generalized pie chart: {e}")
            import traceback
            traceback.print_exc()
            return None, None
   
    def _create_pie_from_percentage_column(self, df, plot_spec, percentage_col, chart_data):
        """Create pie chart from a percentage column"""
        try:
            # Get the main percentage value
            main_percentage = float(df[percentage_col].iloc[0])
            other_percentage = 100 - main_percentage
           
            # Try to get meaningful labels from plot title or column names
            title = plot_spec.get('title', '')
            series_names = plot_spec.get('series_names', [])
           
            if series_names and len(series_names) >= 2:
                labels = series_names[:2]
            elif 'texas' in title.lower():
                labels = ['Texas Oncology', 'Other Accounts']
            elif any(word in title.lower() for word in ['vs', 'versus', ':']):
                # Try to extract labels from title
                parts = title.replace(':', ' vs ').split(' vs ')
                if len(parts) >= 2:
                    labels = [parts[0].strip(), parts[1].strip()]
                else:
                    labels = ['Primary', 'Others']
            else:
                # Generic labels
                labels = ['Main Category', 'Others']
           
            values = [main_percentage, other_percentage]
           
            chart_data.categories = labels
            chart_data.add_series('Distribution', values)
           
            print(f"✅ Percentage-based pie chart: {labels[0]}={values[0]:.1f}%, {labels[1]}={values[1]:.1f}%")
            return chart_data, 'pie'
           
        except Exception as e:
            print(f"❌ Error in percentage pie chart: {e}")
            return None, None
   
    def _create_pie_from_count_columns(self, df, plot_spec, category_col, value_col, chart_data):
        """Create pie chart from category and value columns"""
        try:
            # Sort by value and take top categories (limit to 8 for readability)
            df_sorted = df.sort_values(value_col, ascending=False).head(8)
           
            # Get labels and values
            labels = []
            for label in df_sorted[category_col]:
                label_str = str(label)
                # Truncate long labels
                if len(label_str) > 25:
                    label_str = label_str[:22] + "..."
                labels.append(label_str)
           
            values = df_sorted[value_col].fillna(0).tolist()
           
            # If we have many small values, group them as "Others"
            if len(df) > len(df_sorted):
                others_value = df[~df.index.isin(df_sorted.index)][value_col].sum()
                if others_value > 0:
                    labels.append("Others")
                    values.append(others_value)
           
            chart_data.categories = labels
            chart_data.add_series('Distribution', values)
           
            print(f"✅ Category-value pie chart with {len(labels)} segments")
            return chart_data, 'pie'
           
        except Exception as e:
            print(f"❌ Error in category-value pie chart: {e}")
            return None, None
   
    def _create_pie_from_value_counts(self, df, plot_spec, category_col, chart_data):
        """Create pie chart from value counts of a categorical column"""
        try:
            # Get value counts and limit to top 8
            value_counts = df[category_col].value_counts().head(8)
           
            labels = [str(label)[:25] for label in value_counts.index]
            values = value_counts.values.tolist()
           
            chart_data.categories = labels
            chart_data.add_series('Count', values)
           
            print(f"✅ Value counts pie chart with {len(labels)} categories")
            return chart_data, 'pie'
           
        except Exception as e:
            print(f"❌ Error in value counts pie chart: {e}")
            return None, None
   
    def _create_pie_from_numeric_distribution(self, df, plot_spec, numeric_col, chart_data):
        """Create pie chart from numeric data distribution (quintiles)"""
        try:
            # Create quintiles for distribution
            data = df[numeric_col].dropna()
            quintiles = pd.qcut(data, q=5, labels=['Very Low', 'Low', 'Medium', 'High', 'Very High'])
            value_counts = quintiles.value_counts()
           
            labels = [str(label) for label in value_counts.index]
            values = value_counts.values.tolist()
           
            chart_data.categories = labels
            chart_data.add_series('Distribution', values)
           
            print(f"✅ Numeric distribution pie chart with {len(labels)} quintiles")
            return chart_data, 'pie'
           
        except Exception as e:
            print(f"❌ Error in numeric distribution pie chart: {e}")
            return None, None
   
    def is_pie_chart_valid(self, plot_spec, df):
        """
        Generalized validation for pie charts
        Pie charts don't need traditional x/y validation
        """
        try:
            plot_type = plot_spec.get('type', '')
            if plot_type != 'pie':
                return False
           
            # Pie charts are valid if we have any data to work with
            if len(df) == 0 or len(df.columns) == 0:
                return False
           
            # Check if we have at least one column we can use
            has_numeric = len(df.select_dtypes(include=[np.number]).columns) > 0
            has_categorical = len(df.select_dtypes(include=['object', 'string']).columns) > 0
            has_percentage = any('percentage' in col.lower() or 'percent' in col.lower() for col in df.columns)
           
            if has_numeric or has_categorical or has_percentage:
                print(f"✅ Pie chart validation passed: numeric={has_numeric}, categorical={has_categorical}, percentage={has_percentage}")
                return True
           
            print(f"❌ Pie chart validation failed: no suitable data columns")
            return False
           
        except Exception as e:
            print(f"❌ Error validating pie chart: {e}")
            return False

    # ==================== STRUCTURED BOT - CHART DATA CREATION ====================
   
    def create_editable_chart_data(self, df, plot_spec):
        """Create chart data that maps properly to PowerPoint's data model with enhanced stacked bar support"""
        try:
            x_col = plot_spec['x']
            y_col = plot_spec['y']
            plot_type = plot_spec['type']
            hue_col = plot_spec.get('hue')
           
            print(f"📊 Creating chart data: type={plot_type}, x={x_col}, y={y_col}")

            if self.is_single_value_chart(x_col, df):
                print(f"🎯 Detected single value chart")
                return self.create_single_value_chart_data(df, plot_spec)
            
            if hue_col and hue_col in df.columns and isinstance(y_col, str) and y_col in df.columns:
                print(f"🎨 Processing HUE-based chart: hue column = {hue_col}")
                chart_data = CategoryChartData()
               
                # Get unique hue values (series)
                hue_values = df[hue_col].unique()
                print(f"🔍 Found {len(hue_values)} unique hue values: {list(hue_values)}")
               
                # Get unique x values and sort them properly
                if pd.api.types.is_datetime64_any_dtype(df[x_col]):
                    # For datetime, sort chronologically
                    x_values = sorted(df[x_col].unique())
                elif pd.api.types.is_numeric_dtype(df[x_col]):
                    # For numeric, sort numerically
                    x_values = sorted(df[x_col].unique())
                else:
                    # For categorical, maintain original order or sort alphabetically
                    x_values = df[x_col].unique()
               
                # Limit x values for readability (take recent/relevant data)
                if len(x_values) > 20:
                    x_values = x_values[-20:]  # Take last 20 for time series
               
                # Create category labels
                category_labels = []
                for x_val in x_values:
                    if pd.api.types.is_datetime64_any_dtype(df[x_col]):
                        # Format dates nicely
                        label = pd.to_datetime(x_val).strftime('%m/%d')
                    else:
                        label = str(x_val)
                        if len(label) > 15:
                            label = label[:12] + "..."
                    category_labels.append(label)
               
                chart_data.categories = category_labels
               
                # Create a series for each hue value
                for hue_value in hue_values:
                    # Filter data for this hue value
                    hue_data = df[df[hue_col] == hue_value]
                   
                    # Create series data aligned with x_values
                    series_data = []
                    for x_val in x_values:
                        # Find corresponding y value for this x and hue combination
                        matching_rows = hue_data[hue_data[x_col] == x_val]
                        if len(matching_rows) > 0:
                            y_val = matching_rows[y_col].iloc[0]  # Take first match
                            series_data.append(float(y_val) if pd.notna(y_val) else 0.0)
                        else:
                            series_data.append(0.0)  # No data point for this x value
                   
                    # Create readable series name
                    series_name = str(hue_value)
                    if len(series_name) > 25:
                        series_name = series_name[:22] + "..."
                   
                    chart_data.add_series(series_name, series_data)
                    print(f"✅ Added hue series: {series_name} with {len(series_data)} data points")
               
                return chart_data, 'hue_multi_series'

               
            elif plot_spec.get('multi_series', False) and isinstance(y_col, list):
                chart_data = CategoryChartData()
               
                # Get unique x values and sort them
                if x_col in df.columns:
                    # For parent_name and similar, limit to top entries for readability
                    if plot_type in ['bar', 'stacked_bar'] and 'percentage' in str(y_col).lower():
                        # Sort by sum of first y column and take top 10
                        first_y = y_col[0] if y_col else x_col
                        if first_y in df.columns:
                            df_sorted = df.sort_values(first_y, ascending=False).head(15)
                        else:
                            df_sorted = df.head(15)
                    else:
                        df_sorted = df.head(20)
                   
                    x_values = df_sorted[x_col].tolist()
                   
                    # Truncate long names for better display
                    category_labels = []
                    for x_val in x_values:
                        label = str(x_val)
                        if len(label) > 25:
                            # Smart truncation - keep important parts
                            words = label.split()
                            if len(words) > 3:
                                label = f"{words[0]} {words[1]}... {words[-1]}"
                            else:
                                label = label[:22] + "..."
                        category_labels.append(label)
                   
                    chart_data.categories = category_labels
                   
                    # Add each series
                    for i, y_column in enumerate(y_col):
                        if y_column in df.columns:
                            # Get data for this series
                            series_data = []
                            for _, row in df_sorted.iterrows():
                                y_val = row[y_column]
                                series_data.append(float(y_val) if pd.notna(y_val) else 0.0)
                           
                            # Smart series naming
                            if plot_spec.get('series_names') and i < len(plot_spec['series_names']):
                                series_name = plot_spec['series_names'][i]
                            else:
                                # Create readable series name from column name
                                series_name = y_column.replace('_', ' ').replace('bc', 'BC').title()
                                # Special handling for common patterns
                                series_name = series_name.replace('Bc', 'BC').replace('Percentage', '%')
                           
                            chart_data.add_series(series_name[:30], series_data)
                            print(f"✅ Added series: {series_name} with {len(series_data)} data points")
                        else:
                            print(f"⚠️ Column {y_column} not found, skipping series")
               
                return chart_data, 'multi_series'
           
            # Single series charts
            elif x_col in df.columns and (y_col in df.columns or plot_type == 'hist'):
               
                if plot_type == 'hist':
                    # Histogram data
                    chart_data = CategoryChartData()
                    if pd.api.types.is_numeric_dtype(df[x_col]):
                        hist_data, bins = np.histogram(df[x_col].dropna(), bins=min(15, len(df[x_col].unique())))
                        bin_labels = [f"{bins[i]:.1f}-{bins[i+1]:.1f}" for i in range(len(bins)-1)]
                    else:
                        value_counts = df[x_col].value_counts().head(15)
                        hist_data = value_counts.values
                        bin_labels = [str(label)[:15] for label in value_counts.index]
                   
                    chart_data.categories = bin_labels
                    chart_data.add_series('Frequency', hist_data.tolist())
                    return chart_data, 'histogram'
               
                elif plot_type in ['bar', 'barh']:
                    # Bar chart data - improved handling
                    chart_data = CategoryChartData()
                   
                    # Sort and limit data for better visualization
                    if pd.api.types.is_numeric_dtype(df[y_col]):
                        df_sorted = df.sort_values(y_col, ascending=False).head(15)
                    else:
                        df_sorted = df.head(15)
                   
                    # Create categories with smart truncation
                    categories = []
                    for cat in df_sorted[x_col]:
                        cat_str = str(cat)
                        if len(cat_str) > 20:
                            cat_str = cat_str[:17] + "..."
                        categories.append(cat_str)
                   
                    chart_data.categories = categories
                    values = df_sorted[y_col].fillna(0).tolist()
                    chart_data.add_series(y_col.replace('_', ' ').title(), values)
                    return chart_data, 'bar'
               
                elif plot_type == 'scatter':
                    # Scatter plot data
                    chart_data = XyChartData()
                    series = chart_data.add_series('Data Points')
                   
                    # Sample data for performance
                    sample_size = min(100, len(df))
                    df_sample = df.sample(n=sample_size, random_state=42) if len(df) > sample_size else df
                   
                    for _, row in df_sample.iterrows():
                        x_val = row[x_col]
                        y_val = row[y_col]
                        if pd.notna(x_val) and pd.notna(y_val):
                            try:
                                series.add_data_point(float(x_val), float(y_val))
                            except (ValueError, TypeError):
                                continue
                   
                    return chart_data, 'scatter'
               
                else:  # Default to line chart
                    chart_data = CategoryChartData()
                   
                    # Sort data by x column for proper line chart
                    df_sorted = df.sort_values(x_col).head(50)  # Limit points
                   
                    categories = [str(x)[:15] for x in df_sorted[x_col].tolist()]
                    chart_data.categories = categories
                    chart_data.add_series('Values', df_sorted[y_col].tolist())
                    return chart_data, 'line'
           
            else:
                print(f"⚠️ Invalid columns or configuration: x={x_col}, y={y_col}, type={plot_type}")
                return None, None
               
        except Exception as e:
            print(f"❌ Error creating chart data: {e}")
            import traceback
            traceback.print_exc()
            return None, None

    # ==================== STRUCTURED BOT - VALIDATION ====================
   
    def enhanced_plot_validation(self, plots, df):
        """Enhanced plot validation that handles all chart types including pie charts"""
       
        for plot in plots:
            plot_type = plot.get('type', '')
            x_col = plot.get('x')
            y_col = plot.get('y')
           
            print(f"🔍 Validating plot: type={plot_type}, x={x_col}, y={y_col}")
           
            # Special handling for pie charts
            if plot_type == 'pie':
                if self.is_pie_chart_valid(plot, df):
                    return plot
                else:
                    continue
           
            # Handle single value charts
            if self.is_single_value_chart(x_col, df) and isinstance(y_col, str) and y_col in df.columns:
                print(f"✅ Single value chart validated: category={x_col}, value_col={y_col}")
                return plot
           
            # Validate x column exists (for non-pie charts)
            try:
                x_col_valid = isinstance(x_col, str) and x_col in df.columns
            except TypeError:
                x_col_valid = False
           
            if not x_col_valid:
                print(f"❌ X column '{x_col}' not found in dataframe or is invalid type")
                continue
           
            # Handle hue-based charts
            hue_col = plot.get('hue')
            if hue_col and hue_col in df.columns and isinstance(y_col, str) and y_col in df.columns:
                print(f"✅ HUE-based plot validated: x={x_col}, y={y_col}, hue={hue_col}")
                plot['multi_series'] = True
                return plot
           
            # Handle multi-series charts
            if plot.get('multi_series', False) and isinstance(y_col, list):
                missing_cols = [col for col in y_col if col not in df.columns]
                if not missing_cols:
                    print(f"✅ Multi-series plot validated: {len(y_col)} series")
                    return plot
                else:
                    print(f"❌ Missing columns for multi-series: {missing_cols}")
                    continue
           
            # Handle single series charts
            elif isinstance(y_col, str):
                if plot_type == 'hist':
                    print(f"✅ Histogram plot validated")
                    return plot
                elif y_col in df.columns:
                    print(f"✅ Single-series plot validated")
                    return plot
                else:
                    print(f"❌ Y column '{y_col}' not found in dataframe")
                    continue
           
            print(f"❌ Plot validation failed for type={plot_type}")
       
        return None

    # ==================== STRUCTURED BOT - CHART CREATION ====================
   
    def create_chart_from_spec(self, slide, df, plot_spec, color_index=0):
        """Create PowerPoint chart from LLM-analyzed plot specification"""
        try:
            print(f"🎨 Creating {plot_spec['type']} chart: {plot_spec.get('title', 'Untitled')}")
           
            # Special handling for pie charts
            if plot_spec['type'] == 'pie':
                chart_data, chart_subtype = self.create_generalized_pie_chart_data(df, plot_spec)
            else:
                chart_data, chart_subtype = self.create_editable_chart_data(df, plot_spec)
                
            if not chart_data:
                return None
           
            # Map plot types to PowerPoint chart types
            chart_type_mapping = {
                'line': XL_CHART_TYPE.LINE_MARKERS,
                'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'barh': XL_CHART_TYPE.BAR_CLUSTERED,
                'scatter': XL_CHART_TYPE.XY_SCATTER,
                'histogram': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'pie': XL_CHART_TYPE.PIE
            }
           
            ppt_chart_type = chart_type_mapping.get(plot_spec['type'], XL_CHART_TYPE.COLUMN_CLUSTERED)
            if plot_spec['type'] == 'barh':
                ppt_chart_type = XL_CHART_TYPE.BAR_CLUSTERED
           
            # Create chart with better positioning and spacing for neat layout
            chart = slide.shapes.add_chart(
                ppt_chart_type,
                Inches(0.4), Inches(1.4), Inches(5.7), Inches(4.2),  # Better dimensions and positioning
                chart_data
            ).chart
           
            # Set title
            chart.has_title = True
            title_text = plot_spec.get('title', f"{plot_spec['type'].title()} Chart")
            chart.chart_title.text_frame.text = title_text[:60]  # Truncate long titles
           
            # Apply enhanced styling
            self.apply_enhanced_chart_styling(chart, plot_spec, color_index)
            
            # Special styling for pie charts
            if plot_spec['type'] == 'pie':
                self.apply_pie_chart_styling(chart, plot_spec)
           
            return chart
           
        except Exception as e:
            print(f"❌ Error creating chart: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def create_chart_from_spec_enhanced(self, slide, df, plot_spec, color_index=0):
        """Enhanced chart creation that handles pie charts"""
        try:
            plot_type = plot_spec['type']
            print(f"🎨 Creating {plot_type} chart: {plot_spec.get('title', 'Untitled')}")
           
            # Special handling for pie charts
            if plot_type == 'pie':
                chart_data, chart_subtype = self.create_generalized_pie_chart_data(df, plot_spec)
            else:
                chart_data, chart_subtype = self.create_editable_chart_data(df, plot_spec)
           
            if not chart_data:
                print(f"❌ Failed to create chart data for {plot_type}")
                return None
           
            # Map plot types to PowerPoint chart types
            chart_type_mapping = {
                'line': XL_CHART_TYPE.LINE_MARKERS,
                'bar': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'barh': XL_CHART_TYPE.BAR_CLUSTERED,
                'scatter': XL_CHART_TYPE.XY_SCATTER,
                'histogram': XL_CHART_TYPE.COLUMN_CLUSTERED,
                'pie': XL_CHART_TYPE.PIE
            }
           
            ppt_chart_type = chart_type_mapping.get(plot_type, XL_CHART_TYPE.COLUMN_CLUSTERED)
           
            # Create chart with proper positioning
            chart = slide.shapes.add_chart(
                ppt_chart_type,
                Inches(0.4), Inches(1.4), Inches(5.7), Inches(4.2),
                chart_data
            ).chart
           
            # Set title
            chart.has_title = True
            title_text = plot_spec.get('title', f"{plot_type.title()} Chart")
            chart.chart_title.text_frame.text = title_text[:60]
           
            # Apply styling
            self.apply_enhanced_chart_styling(chart, plot_spec, color_index)
            # Special styling for pie charts
            if plot_type == 'pie':
                self.apply_pie_chart_styling(chart, plot_spec)
           
            return chart
           
        except Exception as e:
            print(f"❌ Error creating chart: {e}")
            import traceback
            traceback.print_exc()
            return None

    # ==================== STRUCTURED BOT - STYLING ====================
   
    def apply_enhanced_chart_styling(self, chart, plot_spec, color_index=0):
        """Apply enhanced styling with improved x-axis label positioning"""
        try:
            # Title styling
            if chart.has_title and chart.chart_title:
                title_font = chart.chart_title.text_frame.paragraphs[0].font
                title_font.size = Pt(14)
                title_font.bold = True
                title_font.color.rgb = RGBColor(44, 62, 80)
           
            # Series styling with proper colors
            if hasattr(chart, 'series') and chart.series:
                for i, series in enumerate(chart.series):
                    color = self.chart_colors[i % len(self.chart_colors)]
                   
                    # Apply fill color
                    if hasattr(series, 'format') and hasattr(series.format, 'fill'):
                        series.format.fill.solid()
                        series.format.fill.fore_color.rgb = color
                   
                    # Apply line color for line charts
                    if hasattr(series, 'format') and hasattr(series.format, 'line'):
                        series.format.line.color.rgb = color
                        series.format.line.width = Pt(2.5)  # Slightly thicker lines
           
            # Enhanced axis styling
            if hasattr(chart, 'value_axis') and chart.value_axis:
                if hasattr(chart.value_axis, 'tick_labels'):
                    chart.value_axis.tick_labels.font.size = Pt(9)
                    chart.value_axis.tick_labels.font.color.rgb = RGBColor(64, 64, 64)
               
                # Set axis title if specified
                if plot_spec.get('ylabel') and chart.value_axis.has_title:
                    chart.value_axis.axis_title.text_frame.text = plot_spec['ylabel'][:30]
                    chart.value_axis.axis_title.text_frame.paragraphs[0].font.size = Pt(10)
                    chart.value_axis.axis_title.text_frame.paragraphs[0].font.bold = True
           
            # Enhanced category axis (X-axis) styling with better label positioning
            if hasattr(chart, 'category_axis') and chart.category_axis:
                if hasattr(chart.category_axis, 'tick_labels'):
                    # Improved font styling
                    chart.category_axis.tick_labels.font.size = Pt(8)  # Slightly smaller for better fit
                    chart.category_axis.tick_labels.font.color.rgb = RGBColor(64, 64, 64)
                    chart.category_axis.tick_labels.font.name = 'Calibri'
                   
                    # Better label rotation for readability
                    try:
                        # Rotate labels to prevent overlap
                        chart.category_axis.tick_labels.rotation = -45
                        # Position labels below the axis
                        chart.category_axis.tick_labels.offset = 100  # Move labels down
                    except Exception as e:
                        print(f"X-axis label positioning warning: {e}")
               
                # Enhanced axis title positioning
                if plot_spec.get('xlabel') and chart.category_axis.has_title:
                    chart.category_axis.axis_title.text_frame.text = plot_spec['xlabel'][:30]
                    title_para = chart.category_axis.axis_title.text_frame.paragraphs[0]
                    title_para.font.size = Pt(10)
                    title_para.font.bold = True
                    title_para.font.color.rgb = RGBColor(44, 62, 80)
               
                # Improve axis line styling
                try:
                    if hasattr(chart.category_axis, 'format') and hasattr(chart.category_axis.format, 'line'):
                        chart.category_axis.format.line.color.rgb = RGBColor(128, 128, 128)
                        chart.category_axis.format.line.width = Pt(1)
                except:
                    pass
           
            # Enhanced legend positioning for multi-series charts
            if (plot_spec.get('multi_series', False) or plot_spec.get('hue') or (hasattr(chart, 'series') and len(chart.series) > 1)):
                try:
                    chart.has_legend = True
                    if chart.legend:
                        from pptx.enum.chart import XL_LEGEND_POSITION
                        chart.legend.position = XL_LEGEND_POSITION.RIGHT  # Better for hue-based charts
                        chart.legend.include_in_layout = False
                       
                        if hasattr(chart.legend, 'font'):
                            chart.legend.font.size = Pt(8)
                            chart.legend.font.color.rgb = RGBColor(64, 64, 64)
                            chart.legend.font.name = 'Calibri'
                except Exception as e:
                    print(f"Legend styling warning: {e}")
           
            # Add gridlines for better readability
            try:
                if hasattr(chart, 'value_axis') and chart.value_axis:
                    chart.value_axis.has_major_gridlines = True
                    if hasattr(chart.value_axis, 'major_gridlines'):
                        chart.value_axis.major_gridlines.format.line.color.rgb = RGBColor(220, 220, 220)
                        chart.value_axis.major_gridlines.format.line.width = Pt(0.5)
            except Exception as e:
                print(f"Gridlines styling warning: {e}")
           
        except Exception as e:
            print(f"Warning in chart styling: {e}")
            import traceback
            traceback.print_exc()
   
    def apply_pie_chart_styling(self, chart, plot_spec):
        """Apply special styling for pie charts"""
        try:
            # Enable data labels with percentages
            if hasattr(chart, 'plots') and chart.plots:
                plot = chart.plots[0]
                if hasattr(plot, 'has_data_labels'):
                    plot.has_data_labels = True
                    if hasattr(plot, 'data_labels'):
                        plot.data_labels.show_percentage = True
                        plot.data_labels.show_value = False
                        plot.data_labels.position = 2  # Outside end
           
            # Apply custom colors if specified in style_params
            style_params = plot_spec.get('style_params', {})
            colors = style_params.get('color', [])
           
            if colors and hasattr(chart, 'series') and chart.series:
                series = chart.series[0]
                if hasattr(series, 'points'):
                    for i, point in enumerate(series.points):
                        if i < len(colors):
                            color_hex = colors[i]
                            if color_hex.startswith('#'):
                                # Convert hex to RGB
                                r = int(color_hex[1:3], 16)
                                g = int(color_hex[3:5], 16)
                                b = int(color_hex[5:7], 16)
                                point.format.fill.solid()
                                point.format.fill.fore_color.rgb = RGBColor(r, g, b)
                        else:
                            # Use default colors
                            color = self.chart_colors[i % len(self.chart_colors)]
                            point.format.fill.solid()
                            point.format.fill.fore_color.rgb = color
           
            print("✅ Pie chart styling applied successfully")
           
        except Exception as e:
            print(f"⚠️ Warning in pie chart styling: {e}")

    # ==================== STRUCTURED BOT - INSIGHTS TEXTBOX ====================
           
    def create_optimized_insights_textbox(self, slide, insights_text):
        """Create optimized insights text box that displays full content without trimming"""
        print(f"🔍 Creating insights textbox with content: {insights_text[:100]}...")
       
        # Optimized positioning for single slide layout - larger box for full content
        insights_box = slide.shapes.add_textbox(
            Inches(6.3), Inches(1.2), Inches(3.2), Inches(4.5)  # Larger dimensions
        )
       
        text_frame = insights_box.text_frame
        text_frame.word_wrap = True
        text_frame.margin_left = Inches(0.08)
        text_frame.margin_right = Inches(0.08)
        text_frame.margin_top = Inches(0.08)
        text_frame.margin_bottom = Inches(0.08)
       
        # Clear existing content
        text_frame.clear()
       
        # Ensure we have content to display
        if not insights_text or not insights_text.strip():
            insights_text = "📊 Analysis Insights\n\n• Data visualization created\n• Chart ready for presentation\n• Editable format available"
       
        print(f"📝 Processing full insights text: {len(insights_text)} characters")
       
        # NO TRIMMING - Display full content with smart formatting
        # Split into manageable sections but don't limit content
        sections = insights_text.split('\n\n')
       
        # Add first paragraph directly to existing frame
        first_section = True
       
        for section_idx, section in enumerate(sections):  # Process ALL sections
            if not section.strip():
                continue
               
            lines = section.split('\n')
           
            for line_idx, line in enumerate(lines):
                if not line.strip():
                    continue
               
                # Use existing paragraph for first line, create new ones for others
                if first_section and line_idx == 0:
                    p = text_frame.paragraphs[0]
                    first_section = False
                else:
                    p = text_frame.add_paragraph()
               
                # Keep full line content - NO TRUNCATION
                p.text = line.strip()
               
                # Font styling
                font = p.font
                font.size = Pt(8)  # Smaller font to fit more content
                font.name = 'Calibri'
               
                # Style based on content
                if line.strip().startswith(('📊', '📈', '🔍', '💡', '•', '-', '*')) or ':' in line:
                    font.bold = True
                    font.color.rgb = RGBColor(44, 62, 80)
                else:
                    font.color.rgb = RGBColor(64, 64, 64)
               
                # Tight spacing to fit more content
                p.space_after = Pt(2)
           
            # Add minimal space between sections
            if section_idx < len(sections) - 1:
                p = text_frame.add_paragraph()
                p.text = ""
                p.space_after = Pt(4)
       
        # Enhanced visual styling
        insights_box.fill.solid()
        insights_box.fill.fore_color.rgb = RGBColor(248, 249, 250)
       
        # More visible border
        insights_box.line.color.rgb = RGBColor(180, 180, 180)
        insights_box.line.width = Pt(1)
       
        print(f"✅ Full insights textbox created with {len(text_frame.paragraphs)} paragraphs")
        return insights_box

    # ==================== MAIN CONVERSION METHOD ====================
   
    async def convert(self, df, python_code):
        """Main conversion function - checks bot type FIRST, then routes accordingly"""
        try:
            # PRIORITY 1: Check for RM_insights bot
            if self.is_rm_insights_bot():
                print(f"🎯 RM_insights bot detected - using topic-sentiment formatting")
                full_df = cl.user_session.get('tool_params', {}).get("result", df)
                if full_df is None or len(full_df) == 0:
                    full_df = df
                plot_spec = {'type': 'bar', 'title': 'Topic-Sentiment Analysis'}
                return await self.create_rm_insights_slide(full_df, plot_spec)
            
            # PRIORITY 2: Check for Structured bot (default to this if not RM_insights)
            print(f"📊 Processing structured data: {len(df)} rows, {len(df.columns)} columns")
            print(f"📋 Columns: {list(df.columns)}")
           
            # Prepare sample data for LLM
            df_sample = df.head(3)
           
            # Use LLM to analyze code
            analysis = await self.analyze_code_with_llm(python_code, df.columns, df_sample)
           
            if not analysis or not analysis.get('plots'):
                print("❌ No plots detected in code analysis")
                return None
           
            plots = analysis['plots']
            print(f"🎯 LLM detected {len(plots)} plots with {analysis.get('analysis_confidence', 0):.2f} confidence")
           
            # Find best valid plot with enhanced validation
            valid_plot = self.enhanced_plot_validation(plots, df)
           
            # If no valid plot found, create a fallback plot
            if not valid_plot:
                print("⚠️ No valid plot found, creating fallback visualization")
                # Find numeric columns for fallback
                numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
                text_cols = df.select_dtypes(include=['object']).columns.tolist()
               
                if len(numeric_cols) >= 2 and len(text_cols) >= 1:
                    valid_plot = {
                        'type': 'bar',
                        'x': text_cols[0],
                        'y': numeric_cols[:3],  # Use first 3 numeric columns
                        'title': 'Data Analysis Dashboard',
                        'xlabel': text_cols[0].replace('_', ' ').title(),
                        'ylabel': 'Values',
                        'multi_series': True,
                        'series_names': [col.replace('_', ' ').title() for col in numeric_cols[:3]]
                    }
                    print(f"✅ Created fallback plot: {valid_plot['x']} vs {valid_plot['y']}")
                elif len(numeric_cols) >= 1 and len(text_cols) >= 1:
                    valid_plot = {
                        'type': 'bar',
                        'x': text_cols[0],
                        'y': numeric_cols[0],
                        'title': 'Data Analysis Dashboard',
                        'xlabel': text_cols[0].replace('_', ' ').title(),
                        'ylabel': numeric_cols[0].replace('_', ' ').title(),
                        'multi_series': False
                    }
                    print(f"✅ Created simple fallback plot: {valid_plot['x']} vs {valid_plot['y']}")
                else:
                    print("❌ Cannot create fallback plot - insufficient columns")
           
            if not valid_plot:
                print("❌ No valid plot found with available columns")
                return None
           
            # Create presentation
            self.prs = Presentation()
            slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])  # Blank layout
           
            # Add main title
            title_box = slide.shapes.add_textbox(
                Inches(0.5), Inches(0.2), Inches(9), Inches(0.8)
            )
            title_frame = title_box.text_frame
            title_frame.text = valid_plot.get('title', 'Data Analysis Dashboard')
            title_para = title_frame.paragraphs[0]
            title_para.font.size = Pt(18)
            title_para.font.bold = True
            title_para.font.color.rgb = RGBColor(44, 62, 80)
            title_para.alignment = PP_ALIGN.CENTER
           
            # Create chart
            chart = self.create_chart_from_spec_enhanced(slide, df, valid_plot, 0)
           
            if chart:
                print(f"✅ Created editable chart: {valid_plot['type']}")
               
                # Generate insights content
                insights_title = "📊 Key Insights"
               
                print(f"🔍 Processing insights - Custom: {bool(self.insights)}, Length: {len(self.insights) if self.insights else 0}")
               
                if not self.insights or not self.insights.strip():
                    print("📝 Generating default insights...")
                    # Generate smart default insights
                    if valid_plot.get('multi_series', False):
                        y_cols = valid_plot['y']
                        series_count = len(y_cols) if isinstance(y_cols, list) else 1
                        default_insights = f"""📊 Dataset Overview
- Records: {len(df):,}
- Variables: {len(df.columns)}
- Chart: Multi-Series {valid_plot['type'].title()}

📈 Analysis Focus
- X-axis: {valid_plot['x'].replace('_', ' ').title()}
- Series: {series_count} data series
- Comparison visualization

🔍 Key Observations
- Multi-dimensional data analysis
- Clear trend visualization
- Professional formatting applied

💡 Business Impact
- Performance comparison enabled
- Trend analysis simplified
- Decision-ready insights"""
                    else:
                        # Single series insights
                        x_col = valid_plot['x']
                        y_col = valid_plot['y']
                       
                        try:
                            if x_col in df.columns and y_col in df.columns:
                                if pd.api.types.is_numeric_dtype(df[y_col]):
                                    y_mean = df[y_col].mean()
                                    y_range = df[y_col].max() - df[y_col].min()
                                    stat_text = f"• Mean: {y_mean:.2f}\n• Range: {y_range:.2f}"
                                else:
                                    unique_count = df[y_col].nunique()
                                    stat_text = f"• Unique values: {unique_count}"
                            else:
                                stat_text = "• Data patterns identified"
                        except:
                            stat_text = "• Statistical analysis complete"
                       
                        default_insights = f"""📊 Dataset Summary
- Records: {len(df):,}
- Analysis: {valid_plot['type'].title()} Chart

📈 Variables
- X: {x_col.replace('_', ' ').title()}
- Y: {y_col.replace('_', ' ').title()}

🔍 Statistics
{stat_text}

💡 Insights
- Clear data visualization
- Editable chart format
- Ready for presentation

🎯 Next Steps
- Review data trends
- Analyze patterns
- Make data-driven decisions"""
                   
                    insights_content = default_insights
                    print(f"✅ Generated default insights: {len(insights_content)} characters")
                else:
                    insights_content = self.insights.strip()
                    print(f"✅ Using custom insights: {len(insights_content)} characters")
               
                full_insights = f"{insights_title}\n\n{insights_content}"
                print(f"📋 Final insights text length: {len(full_insights)}")
               
                # Add optimized insights textbox
                insights_box = self.create_optimized_insights_textbox(slide, full_insights)
                print(f"📍 Insights box position: {insights_box.left}, {insights_box.top}, {insights_box.width}, {insights_box.height}")
               
                # Add footer with better positioning to avoid overlap with x-labels
                footer_box = slide.shapes.add_textbox(
                    Inches(0.5), Inches(6.1), Inches(9), Inches(0.4)  # Moved down to avoid x-label overlap
                )
                footer_frame = footer_box.text_frame
                chart_info = f"Chart: {valid_plot['type'].title()}"
                if analysis.get('detected_libraries'):
                    chart_info += f" • Libraries: {', '.join(analysis['detected_libraries'])}"
                footer_text = f"LLM-Enhanced Analysis • {chart_info} • Confidence: {analysis.get('analysis_confidence', 0.9):.1%}"
                footer_frame.text = footer_text
                footer_para = footer_frame.paragraphs[0]
                footer_para.font.size = Pt(7)
                footer_para.font.italic = True
                footer_para.font.color.rgb = RGBColor(128, 128, 128)
                footer_para.alignment = PP_ALIGN.CENTER
               
                print(f"\n🎉 SUCCESS! Created LLM-enhanced editable slide")
                print(f"📄 Chart Type: {valid_plot['type']} with proper data mapping")
                print(f"🤖 LLM Confidence: {analysis.get('analysis_confidence', 0.9):.1%}")
                return self.prs
            else:
                print("❌ Failed to create chart from LLM analysis")
                return None
           
        except Exception as e:
            print(f"❌ Error in conversion: {e}")
            import traceback
            traceback.print_exc()
            return None

    # ==================== INSIGHT SUMMARIZATION ====================
       
    async def insight_summarize(self, insights):
        """Summarize insights - used for structured bot"""
        system_prompt = """
You are a good summarizer of the given text.
You are given with insights of a data, your task is to summarize the provided insights and summarize so that it covers all points,
Always ensure that the whole insights should fit in the 10 points and summarize, always ensure that 10 points covers all the insights.
Just provide the insights no another points

Examples

    insights : "some large paragraph"
    output : "10 points no heading"
"""
        user_content = """
Here are the insights of data : {insights}
"""
        try:
            user_content = user_content.format(insights=insights)
            response = await self.llm_client.chat.completions.create(
                messages=[
                    {"role": "system", "content": system_prompt},
                    {"role": "user", "content": user_content}
                ],
                model=self.model_id,
                max_tokens=2048,
                temperature=0.1,
                top_p=1,
                stream=False
            )
            response_text = response.choices[0].message.content
            print(f"🤖 LLM Summarization Response: {response_text[:200]}...")
            return response_text
        except Exception as e:
            print(f"❌ Error in insight summarization: {e}")
            raise e

    # ==================== PUBLIC API METHODS ====================
   
    async def create_ppt(self, df, python_code, filename="analysis.pptx", insights=""):
        """
        Enhanced LLM-powered function: DataFrame + Python code → Professional PPT
       
        Key Features:
        - Automatic bot type detection (Structured vs RM_insights)
        - LLM-powered code analysis for accurate plot detection
        - Editable charts in PowerPoint with proper data mapping
        - Optimized single-slide layout with insights
        - Enhanced error handling and debugging
        - Support for complex multi-series visualizations
        - Pie chart support for any data structure
        - Single-value/hardcoded category charts
        - RM_insights with topic-sentiment breakdown and summarized insights
       
        Args:
            df: pandas DataFrame
            python_code: String containing visualization code
            filename: Output filename (optional)
            insights: Custom insights text (optional)
       
        Returns:
            String: Filename of created PPT, or None if failed
        """
        
        # Summarize insights if not RM_insights bot
        if not self.is_rm_insights_bot():
            insights = await self.insight_summarize(insights)
            print("📝 Insights summarized for structured bot")
        
        insights = insights.replace("\n\n", "\n")
        self.insights = insights
       
        try:
            # Run async conversion
            if hasattr(asyncio, 'get_running_loop'):
                try:
                    loop = asyncio.get_running_loop()
                    prs = await self.convert(df, python_code)
                except RuntimeError:
                    # No running loop, create new one
                    prs = asyncio.run(self.convert(df, python_code))
            else:
                # Fallback for older Python versions
                prs = asyncio.run(self.convert(df, python_code))
           
            if prs:
                prs.save(filename)
                print(f"📄 LLM-Enhanced PPT saved as: {filename}")
               
                # Enhanced Excel export with LLM analysis
                excel_filename = filename.replace('.pptx', '_analysis.xlsx')
                try:
                    with pd.ExcelWriter(excel_filename, engine='openpyxl') as writer:
                        # Original data
                        df.to_excel(writer, sheet_name='Data', index=False)
                       
                        # LLM analysis summary
                        if self.detected_plots:
                            analysis_summary = []
                            for i, plot in enumerate(self.detected_plots):
                                analysis_summary.append({
                                    'Plot_ID': i+1,
                                    'Chart_Type': plot.get('type', 'Unknown'),
                                    'X_Column': plot.get('x', 'N/A'),
                                    'Y_Column': str(plot.get('y', 'N/A')),
                                    'Title': plot.get('title', 'Untitled'),
                                    'Multi_Series': plot.get('multi_series', False),
                                    'Has_Hue': plot.get('hue') is not None,
                                    'Confidence': plot.get('confidence', 'N/A')
                                })
                           
                            analysis_df = pd.DataFrame(analysis_summary)
                            analysis_df.to_excel(writer, sheet_name='LLM_Analysis', index=False)
                       
                        # Data statistics
                        stats_summary = {
                            'Metric': ['Total_Rows', 'Total_Columns', 'Numeric_Columns', 'Text_Columns', 'Missing_Values'],
                            'Value': [
                                len(df),
                                len(df.columns),
                                len(df.select_dtypes(include=[np.number]).columns),
                                len(df.select_dtypes(include=['object']).columns),
                                df.isnull().sum().sum()
                            ]
                        }
                        stats_df = pd.DataFrame(stats_summary)
                        stats_df.to_excel(writer, sheet_name='Data_Stats', index=False)
                   
                    print(f"📈 Enhanced analysis exported to: {excel_filename}")
                except Exception as e:
                    print(f"⚠️ Excel export warning: {e}")
               
                return filename
            else:
                return None
               
        except Exception as e:
            print(f"❌ Error in create_ppt: {e}")
            import traceback
            traceback.print_exc()
            return None
    
    def create_ppt_sync(self, df, python_code, filename="analysis.pptx", insights=""):
        """
        Synchronous wrapper for create_ppt - useful when you can't use async
        """
        return asyncio.run(self.create_ppt(df, python_code, filename, insights))
    
    def update_deck(self):
        """
        Update deck from Chainlit session parameters
        Works with both Structured and RM_insights bots
        """
        self.llm_client = cl.user_session.get("llm_client")
        self.model_id = cl.user_session.get("llm_client_model_id")
        df = cl.user_session.get('tool_params').get("result")
        python_code = cl.user_session.get('tool_params').get("python_code")
        filename = cl.user_session.get("deck_filename")
        insights = cl.user_session.get("insights")
        
        self.create_ppt_sync(df, python_code, filename, insights)
        
        # Clean up Excel file
        files = [filename.replace('.pptx', '_analysis.xlsx')]
        for file in files:
            if os.path.exists(file):
                try:
                    os.remove(file)
                    print(f"🗑️ {file} has been deleted.")
                except Exception as e:
                    print(f"⚠️ Error while deleting {file}: {str(e)}")


# ==================== END OF GeneratePPT CLASS ====================