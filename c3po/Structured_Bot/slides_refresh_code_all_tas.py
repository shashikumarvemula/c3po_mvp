#!/usr/bin/env python
# coding: utf-8

# In[ ]:

import numpy as np
import time
import zipfile
import xml.etree.ElementTree as ET
import openpyxl
import os
import tempfile
import shutil
import pandas as pd
import matplotlib.pyplot as plt
from datetime import datetime, timedelta
import pytz
from functools import reduce
import dask.dataframe as dd
from pptx.util import Pt
from databricks import sql
import os
import pandas as pd
from datetime import datetime, timedelta
import pytz
from decimal import Decimal
import pandas as pd
from datetime import datetime
from dateutil.relativedelta import relativedelta
import concurrent.futures



import pandas as pd
import numpy as np
import datetime
from dateutil.relativedelta import relativedelta
from pptx import Presentation
from pptx.chart.data import ChartData, CategoryChartData
from databricks import sql
import os
from datetime import datetime
from dateutil.relativedelta import relativedelta
import time

from databricks.sql.exc import (
    MaxRetryDurationError,
    RequestError,
    SessionAlreadyClosedError,
)
from chainlit.oauth_providers import providers,get_oauth_provider


from sql_warehouse_handling import DataWarehouse





class Update_Deck:
    def __init__(self):
        
        self.DATABRICKS_SERVER_HOSTNAME =os.environ["DATABRICKS_SERVER_HOSTNAME"]
        self.DATABRICKS_HTTP_PATH = os.environ["DATABRICKS_HTTP_PATH"]
        self.access_token = os.getenv("INITIAL_ACCESS_TOKEN")
        print("init of update deck access token", self.access_token)
        self.warehouse = DataWarehouse(
            host_name=self.DATABRICKS_SERVER_HOSTNAME,
            http_path=self.DATABRICKS_HTTP_PATH,
            access_token=self.access_token
        )
        self.connection = None
        self.cursor = None
        print("In update deck init method connection created ")

    def connect_to_databricks(self):

        try:
            print("refrest token in update deck")
            self.access_token = self.warehouse.refresh_token()
            print("new access token in update deck from warehouse", self.access_token)
            self.connection = sql.connect(
                server_hostname=self.DATABRICKS_SERVER_HOSTNAME,
                http_path=self.DATABRICKS_HTTP_PATH,
                access_token=self.access_token,
                 _tls_no_verify = True
            )
            self.cursor = self.connection.cursor()
            print("Databricks Connected Successfully")
        except Exception as e:
            print(f"Database connection error: {str(e)}")
            raise

    def update_cell_text(self,cell, value, type_= ''):
        print("update_cell_text!!!!", self.update_cell_text)
        if type_ == "chart_title_":
            print("in update chart title")
            cell.text_frame.text = value
        else:
            text_frame = cell.text_frame
            if text_frame.paragraphs:
                paragraph = text_frame.paragraphs[0]
                if paragraph.runs:
                    para_len = len(paragraph.runs)
                    if para_len > 1:
                        print("len> 1", para_len)
                        for run_idx, run in enumerate(paragraph.runs):
                            run.text = ''
                    print("type of value!!!!!", type(value))
                    if "." not in str(value):
                        try:
                            value = int(value)
                            value = f"{value:,}"
                        except:
                            pass
                    paragraph.runs[0].text = str(value)


class Update_Onc_Claims_Deck(Update_Deck):
    def __init__(self):
        super().__init__()

    def update_deck(self): # whole_data
        tnbc_tps_data = None
        perst_data_old_batch = None
        print("In update onc claims deck before connection to databricks")
        self.connect_to_databricks()
        print("Connection Established")

        query = "SELECT MAX(Year_Month) FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`onc_claims_data`"
        # latest_year_month = pd.to_datetime(str(pd.DataFrame(self.cursor.execute(query).fetchall(), columns=[elem[0] for elem in self.cursor.description]).iloc[0,0]), format='%Y%m')
        latest_year_month =pd.DataFrame(self.cursor.execute(query).fetchall(), columns=[elem[0] for elem in self.cursor.description])
        # Access the value from the first row and first column
        latest_year_month_value = latest_year_month.iloc[0, 0]

        # latest_year_month_value=202412
        # latest_year_month = pd.to_datetime(str(latest_year_month_value), format='%Y%m')
        latest_year_month_dt = pd.to_datetime(str(latest_year_month_value), format='%Y%m')
        date_obj = pd.to_datetime(latest_year_month_dt)
        formatted_str = date_obj.strftime("%b '%y")
        print("formatted in desired form",formatted_str)
        prev_year_same_month = date_obj - pd.DateOffset(years=1)

        formatted_str_prev = prev_year_same_month.strftime("%b '%y")

        print("formatted pre year",formatted_str_prev)

        # latest_year_month = pd.to_datetime(str(tnbc_nps_data_apr['Year_Month'].max()), format='%Y%m')

        print("here is the latest year in date time format",latest_year_month)

        # def recent_22_months_range(max_year_month):
        #     # Convert integer or string YYYYMM to date (1st day of month)
        #     date = datetime.strptime(str(max_year_month), "%Y%m")

        #     # Calculate start date: 21 months before the max (for 22 months total including max)
        #     start_date = date - relativedelta(months=21)

        #     # Format back to YYYYMM integer
        #     print("start date in date format",start_date)
        #     start_yyyymm = str(start_date.strftime("%Y%m"))
        #     end_yyyymm = str(date.strftime("%Y%m"))

        #     return start_yyyymm, end_yyyymm

        # # print("here is the range of dates",recent_22_months_range(latest_year_month_value))

        # start_ym, end_ym = recent_22_months_range(latest_year_month_value)

        # # Step 3: Use these in your SQL query to filter data
        # start_time=time.time()
        # data_query = f"""
        #     SELECT *
        #     FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`onc_claims_data`
        #     WHERE Year_Month >= {start_ym} AND Year_Month<={end_ym}
        # """
        # print("excuting the query to get tnbc tps data apr : ",data_query)
        # tnbc_tps_data_apr = self.warehouse.get_data_from_wareshouse(data_query)


        # Function to calculate start and end date for 22 months
        def recent_22_months_range(max_year_month):
            date = datetime.strptime(str(max_year_month), "%Y%m")
            start_date = date - relativedelta(months=21)
            
            start_yyyymm = str(start_date.strftime("%Y%m"))
            end_yyyymm = str(date.strftime("%Y%m"))
            
            return start_yyyymm, end_yyyymm

        # Function to load data for each month
        def load_month_data(year_month, new_patient_flag="",line_of_therapy=""):
            # Construct the query for the given month
            print("retreiving the data for year month", year_month)
            print("new patient flag added to the query")
            data_query = ""
            if new_patient_flag == '1':
                data_query = f"""
                    SELECT *
                    FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`onc_claims_data`
                    WHERE Year_Month = {year_month} and new_patient_flag = '1'
                """
            elif line_of_therapy=='2':
                data_query = f"""
                    SELECT *
                    FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`onc_claims_data`
                    WHERE Year_Month = {year_month} and Line_of_Therapy_LOT >= '2'
                """
            elif line_of_therapy=='3+':
                data_query = f"""
                    SELECT *
                    FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`onc_claims_data`
                    WHERE Year_Month = {year_month} and Line_of_Therapy_LOT >= '3'
                """
                
            # Execute the query and return the result
            print("data query",data_query)
            try:
                result = self.warehouse.get_data_from_wareshouse(data_query)
                if result.empty:
                    print(f"No data found for year month {year_month}")
                print("retrieved the data for year month", year_month , "with result ", result.head())
                return result
            except Exception as e:
                print("error in loading data in warehouse",e)
                # return pd.DataFrame()

        # Split into months and load data concurrently
        def load_data_concurrently(latest_year_month_value,new_patient_flag="",line_of_therapy=""):
            # Get the start and end year_month values
            start_ym, end_ym = recent_22_months_range(latest_year_month_value)
            print("latest_year_month_value in 22 months range",latest_year_month_value)
            print("start_ym",start_ym,"end_ym",end_ym)
            
            # Generate all months in the range
            months_to_load = []
            new_patient_flag_list=[]
            line_of_therapy_list=[]
            current_month = datetime.strptime(start_ym, "%Y%m")
            end_month = datetime.strptime(end_ym, "%Y%m")
            
            while current_month <= end_month:
                months_to_load.append(current_month.strftime("%Y%m"))
                current_month += relativedelta(months=1)
                if new_patient_flag:
                    new_patient_flag_list.append('1')
                else:
                    new_patient_flag_list.append('')
                if line_of_therapy:
                    line_of_therapy_list.append(line_of_therapy)
                else:
                    line_of_therapy_list.append('')

            # Use ThreadPoolExecutor to load data concurrently
            all_data = []
            print("before concurrent loading of data")
            with concurrent.futures.ThreadPoolExecutor() as executor:
                # Map the load_month_data function to the months
                print("while loading data concurrently")
                results = executor.map(load_month_data,months_to_load,new_patient_flag_list,line_of_therapy_list)
                
                # Concatenate all the results
                for result in results:
                    all_data.append(result)
            
            # Concatenate all data into a single DataFrame
            final_data = pd.concat(all_data, ignore_index=True)
            
            return final_data

        # Example usage
        start_time = time.time()
        # tnbc_tps_data_apr = load_data_concurrently(latest_year_month_value)
        # print("Time taken to load data:", time.time() - start_time)

        # print("query executed to get tnbc tps data apr")
        # end_time=time.time()
        # print("tnbc tps data apr Husna ",tnbc_tps_data_apr.head())
        # print("time take to load",end_time - start_time)

        # print("columns the tps data have",tnbc_tps_data_apr.columns)
        # tnbc_tps_data_apr['new_patient_flag'] = tnbc_tps_data_apr['new_patient_flag'].astype(int)


        # tnbc_nps_data_apr = tnbc_tps_data_apr[tnbc_tps_data_apr['new_patient_flag'] == 1]
        try:
            tnbc_nps_data_apr=load_data_concurrently(latest_year_month_value,new_patient_flag='1')
            size_mb = tnbc_nps_data_apr.memory_usage(deep=True).sum() / (1024 * 1024)
            print(f"DataFrame size nps: {size_mb:.2f} MB")
        except Exception as e:
            print("error in loading data",e)
            #tnbc_nps_data_apr=load_data_concurrently(latest_year_month_value,line_of_therapy='2')



        print("nps data claims",tnbc_nps_data_apr.head())


        def get_adopter_dabbler_shares(orig_data):
            data = orig_data.copy()
            data['Year_Month'] = pd.to_datetime(data['Year_Month'].astype(str), format='%Y%m')
            # Convert 'Line_of_Therapy_LOT' to integer before comparison
            data['Line_of_Therapy_LOT'] = data['Line_of_Therapy_LOT'].astype(int)
            data = data[data['Line_of_Therapy_LOT'] >= 2]
            #latest_year_month = data['Year_Month'].max()
            latest_year_month= datetime.strptime(str(latest_year_month_value), "%Y%m")
            num_intervals = 5  # For example, generating 10 intervals
            end_dates = []
            for _ in range(num_intervals):
            #     .strftime('%B %Y')
                end_dates.append(latest_year_month)
                latest_year_month -= relativedelta(months=3)
            start_dates = [end_date - pd.DateOffset(months=11) for end_date in end_dates]

            adopter_results = pd.DataFrame()
            dabbler_results = pd.DataFrame()
            non_user_results = pd.DataFrame()

            for start_date, end_date in zip(start_dates, end_dates):
                period_data = data[(data['Year_Month'] >= start_date) & (data['Year_Month'] <= end_date)]

                total_prescriptions = period_data.groupby(['NPI_HCP_number']).agg({'Monthly_Adjusted_Count': 'sum'}).rename(columns={'Monthly_Adjusted_Count': 'total_prescriptions'})
                trodelvy_prescriptions = period_data[period_data['regimen_grouped_name'].str.contains('trodelvy', case=False)].groupby(['NPI_HCP_number']).agg({'Monthly_Adjusted_Count': 'sum'}).rename(columns={'Monthly_Adjusted_Count': 'trodelvy_prescriptions'})

                hcp_data = pd.merge(total_prescriptions, trodelvy_prescriptions, on=['NPI_HCP_number'], how='left').fillna(0)
                hcp_data['trodelvy_share'] = hcp_data['trodelvy_prescriptions'] / hcp_data['total_prescriptions']

                adopters = hcp_data[hcp_data['trodelvy_share'] >= 0.4]
                dabblers = hcp_data[(hcp_data['trodelvy_share'] > 0) & (hcp_data['trodelvy_share'] < 0.4)]
                non_users = hcp_data[hcp_data['trodelvy_prescriptions'] == 0]

                adopters_npi = adopters.index
                dabblers_npi = dabblers.index
                non_users_npi = non_users.index

                adopters_data = period_data[period_data['NPI_HCP_number'].isin(adopters_npi)]
                dabblers_data = period_data[period_data['NPI_HCP_number'].isin(dabblers_npi)]
                non_users_data = period_data[period_data['NPI_HCP_number'].isin(non_users_npi)]

                def calculate_drug_share(segment_data, end_date):
                    drug_monthly_counts = segment_data.groupby(['regimen_grouped_name']).agg({'Monthly_Adjusted_Count': 'sum'}).rename(columns={'Monthly_Adjusted_Count': 'drug_monthly_count'})
                    total_monthly_counts = segment_data['Monthly_Adjusted_Count'].sum()
                    drug_share = (drug_monthly_counts['drug_monthly_count'] / total_monthly_counts)
                    drug_share = drug_share.reset_index()
                    drug_share['Year_Month'] = end_date.strftime('%b-%y')
                    return drug_share

                adopters_share = calculate_drug_share(adopters_data, end_date)
                dabblers_share = calculate_drug_share(dabblers_data, end_date)
                non_users_share = calculate_drug_share(non_users_data, end_date)

                adopter_results = pd.concat([adopter_results, adopters_share])
                dabbler_results = pd.concat([dabbler_results, dabblers_share])
                non_user_results = pd.concat([non_user_results, non_users_share])


            adopter_pivot = adopter_results.pivot(index='Year_Month', columns='regimen_grouped_name', values='drug_monthly_count').fillna(0)
            dabbler_pivot = dabbler_results.pivot(index='Year_Month', columns='regimen_grouped_name', values='drug_monthly_count').fillna(0)
            non_user_pivot = non_user_results.pivot(index='Year_Month', columns='regimen_grouped_name', values='drug_monthly_count').fillna(0)

            return adopter_pivot.reset_index(), dabbler_pivot.reset_index(), non_user_pivot.reset_index()

        def calculate_metrics(df):
            df["TPC"] = df["HALAVEN BASED"] + df["XELODA BASED"] + \
                        df["VINORELBINE BASED"] + df["GEMZAR BASED"]
            df["Other"] = df["PLATINUM BASED"] + df["OTHER COMBINATION CHEMO"] + \
                          df["OTHER SINGLE CHEMO"] + df["PARP BASED"] + \
                          df["ANTHRACYCLINE BASED"]
            return df


        def process_data(df, calc_metrics):
            if calc_metrics:
                df = calculate_metrics(df)
            df['date_column'] = pd.to_datetime(df['Year_Month'], format='%b-%y')
            df_sorted = df.sort_values(by='date_column')
            df_sorted['date_column'] = df_sorted['date_column'].dt.strftime('%b-%y')
        #     df_sorted = df_sorted.reset_index()
            return df_sorted

        def select_columns(df, columns):
            return df[columns]




        adopter_results, dabbler_results, non_user_results = get_adopter_dabbler_shares(tnbc_nps_data_apr)
        print("adopter results",adopter_results.head())
        print("dabbler results",dabbler_results.head())
        print("non user results",non_user_results.head())
        adopter_df_sorted = process_data(adopter_results, True)
        adopter_df_sorted = select_columns(adopter_df_sorted, ["Year_Month", "Other", "TAXANE BASED","TPC","PD-1/PD-L1 BASED",  "ENHERTU BASED", "TRODELVY BASED"  ])
        dabbler_df_sorted = process_data(dabbler_results, True)
        dabbler_df_sorted = select_columns(dabbler_df_sorted, ["Year_Month", "Other", "TAXANE BASED","TPC","PD-1/PD-L1 BASED",  "ENHERTU BASED", "TRODELVY BASED"  ])
        non_user_df_sorted = process_data(non_user_results, True)
        non_user_df_sorted = select_columns(non_user_df_sorted, ["Year_Month", "Other", "TAXANE BASED","TPC","PD-1/PD-L1 BASED",  "ENHERTU BASED"])
        print("adopter_df_sorted",adopter_df_sorted.head())
        print("dabbler_df_sorted",dabbler_df_sorted.head())
        print("non_user_df_sorted",non_user_df_sorted.head())

        # end_date = pd.to_datetime(latest_year_month, format='%Y%m')
        end_date = latest_year_month_dt
        print("end_date!!!!!!", end_date)
        start_date_recent = end_date - pd.DateOffset(months=11)  # Start of recent 12 months
        end_date_past = start_date_recent - pd.DateOffset(days=1)  # End of past 12 months
        start_date_past = end_date_past - pd.DateOffset(months=11)  # Start of past 12 months
        tier_counts = {}
        # Function to process data
        def get_hcp_trend_data(orig_data, start_date, end_date, period_label, period):

            data = orig_data.copy()
            data['Year_Month'] = pd.to_datetime(data['Year_Month'], format='%Y%m')
            filtered_data = data[(data['Year_Month'] >= start_date) & (data['Year_Month'] <= end_date) & (data['Line_of_Therapy_LOT'].astype(int) >= 2)]

            # Calculate Monthly_Adjusted_Count for Trodelvy and total
            trodelvy_data = filtered_data[filtered_data['regimen_grouped_name'].str.contains('trodelvy', case=False)]
            trodelvy_counts = trodelvy_data.groupby('NPI_HCP_number')['Monthly_Adjusted_Count'].sum()
            total_counts = filtered_data.groupby('NPI_HCP_number')['Monthly_Adjusted_Count'].sum()

            # Classify each HCP
            classification = (trodelvy_counts / total_counts).fillna(0)
            classification = classification.map(lambda x: 'Adopter' if x >= 0.4 else ('Dabbler' if x > 0 else 'Non-user'))

            # Merge classification with tier data
            tier_data = filtered_data[['NPI_HCP_number', 'NPI_HCP_tier']].drop_duplicates().set_index('NPI_HCP_number')
            classification_with_tiers = classification.to_frame('Category').merge(tier_data, left_index=True, right_index=True)

            # Calculate the percentage of each category within each tier
            result = classification_with_tiers.groupby(['NPI_HCP_tier', 'Category']).size()
            total_by_tier = classification_with_tiers.groupby('NPI_HCP_tier').size()
            # Modify the keys to add a period and a space
            total_by_tier_dict = {
                f"{period} {key}": value for key, value in classification_with_tiers.groupby('NPI_HCP_tier').size().to_dict().items()
            }

            print( period +" "+ period_label,  "total_by_tier!!!!!!!", total_by_tier_dict)
            percentage_result =  result / total_by_tier.reindex(result.index, level=0)
        #     percentage_result = 100 * result / total_by_tier.reindex(result.index, level=0)

            return percentage_result.reset_index(name='Percentage').assign(Period=period_label).round(2), total_by_tier_dict

        # Apply processing to the data ranges
        past_result_df, tier_counts_past = get_hcp_trend_data(tnbc_nps_data_apr, start_date_past, end_date_past, f"{formatted_str_prev}", "past")
        # print("length!!!!!!!!", len(past_result_df))
        recent_result_df, tier_counts_present = get_hcp_trend_data(tnbc_nps_data_apr, start_date_recent, end_date, f"{formatted_str}", "present")

        # Concatenate results and pivot for formatting
        final_result_df = pd.concat([past_result_df, recent_result_df])
        formatted_result_df = final_result_df.pivot_table(index=['NPI_HCP_tier', 'Period'], columns='Category', values='Percentage', fill_value='0')
        formatted_result_df.reset_index(inplace=True)

        # Reorder for display according to your structure:
        formatted_result_df.sort_values(by=['NPI_HCP_tier', 'Period'], ascending=[True, True], inplace=True)

        # Print or return the formatted DataFrame
        formatted_result_df
        formatted_result_df_ = formatted_result_df[(formatted_result_df['NPI_HCP_tier'] == 'Tier 1') | (formatted_result_df['NPI_HCP_tier'] == 'Tier 2')
                                                  | (formatted_result_df['NPI_HCP_tier'] == 'Tier 3')]
        print("length formatted_result_df_!!!", len(formatted_result_df_))
        formatted_result_df_["NPI_Tier/Period"] = formatted_result_df_['NPI_HCP_tier'] +"\n ("+ formatted_result_df_['Period'] + ")"
        HCP_trend_segment = formatted_result_df_[['NPI_Tier/Period', 'Adopter','Dabbler', 'Non-user']]

        def calculate_source_business_percentages(orig_data):
            claims_data = orig_data.copy()
            claims_data['Year_Month'] = pd.to_datetime(claims_data['Year_Month'], format='%Y%m')

            # Get the latest 'Year_Month' from the data
        #     latest_year_month = claims_data['Year_Month'].max()
            latest_year_month = pd.to_datetime(f'{latest_year_month_value}', format='%Y%m')
            print("latest year month",latest_year_month.month)


            # Determine the current quarter of the latest 'Year_Month'
            current_quarter = (latest_year_month.month - 1) // 3 + 1
            current_year = latest_year_month.year

            # Calculate the quarter and year for the quarter before the latest
            previous_quarter = current_quarter - 1
            previous_year = current_year
            if previous_quarter == 0:
                previous_quarter = 4
                previous_year -= 1

            # Define quarters to consider
            quarters_to_consider = [(current_year, current_quarter), (previous_year, previous_quarter)]
            print("quarters_to_consider!!!!!!", quarters_to_consider)

            # Define a function to get the Year_Month values for a given quarter and year
            def get_quarter_year_months(year, quarter):
                month_start = (quarter - 1) * 3 + 1
                print("!!!!!!!!!", [pd.Timestamp(year=year, month=month, day=1) for month in range(month_start, month_start + 3)])
                return [pd.Timestamp(year=year, month=month, day=1) for month in range(month_start, month_start + 3)]

            # Prepare an empty DataFrame to collect results
            all_results = pd.DataFrame()

            # Iterate through the quarters to consider
            for year, quarter in quarters_to_consider:
                quarter_months = get_quarter_year_months(year, quarter)
                quarter_start = quarter_months[0]
                quarter_end = quarter_months[-1] + pd.offsets.MonthEnd()  # Adjusted to get the end of the month
                print("quarter_start!!!!!!!!", quarter_start, "quarter_end!!!!!!!", quarter_end)

                # Filter data within the quarter and for lines 2, 3, and 4L+
                for line in [2, 3, '4+']:


                    if line == '4+':
                        # When line is 4, include all lines >= 4
                        filtered_data = claims_data[(claims_data['Year_Month'] >= quarter_start) &
                                                    (claims_data['Year_Month'] <= quarter_end) &
                                                    (claims_data['Line_of_Therapy_LOT'].astype(int) >= 4)
                                                     & claims_data["regimen_grouped_name"].str.contains('trodelvy', case=False)
                                                    ]
                    else :
        #                 print("ragimen!!!!!!!", claims_data["regimen_grouped_name"].unique)
                        filtered_data = claims_data[
                                                    (claims_data['Year_Month'] >= quarter_start) &
                                                    (claims_data['Year_Month'] <= quarter_end) &
                                                    (claims_data['Line_of_Therapy_LOT'].astype(int) == line)
                                                     & claims_data["regimen_grouped_name"].str.contains('trodelvy', case=False)
                                                    ]
                        filtered_data.to_excel("line 2 sob data.xlsx")
                    # Group by 'source_of_business' and count occurrences
                    source_business_counts = filtered_data.groupby('source_of_business').size()
                    total_count = source_business_counts.sum()

                    # Calculate percentages
                    source_business_percentages = ((source_business_counts / total_count))

                    print("source of businesss index Husna",source_business_counts.index)
                    # Combine counts and percentages into a DataFrame
                    result = pd.DataFrame({
                        'Source of Business': source_business_counts.index,
                        'Percentage': source_business_percentages.values
                    })
                    result['Line'] = f'{line}L Q{quarter}\'{str(year)[-2:]}'
                    result['End Date'] = quarter_end
                    if line == '4+':
                        result['LOT'] = 4
                    else:
                        result['LOT'] = line

                    # Append the result to the collection DataFrame
                    all_results = pd.concat([all_results, result], ignore_index=True)

            # Pivot the DataFrame to match the desired output format
            final_results = all_results.pivot(index=['Line', 'End Date', 'LOT'], columns='Source of Business', values='Percentage').fillna(0)

            return final_results

        # Call the function with your claims_data DataFrame
        result = calculate_source_business_percentages(tnbc_nps_data_apr)

        if 'VINORELBINE' in result.columns and 'VINORELBINE BASED' not in result.columns:
            result = result.rename(columns={'VINORELBINE': 'VINORELBINE BASED'})
        else:
            print("Either 'VINORELBINE' column not found or 'VINORELBINE BASED' already present.")

        print("result!!!!!! Husna", result)
        # result = result.reset_index()

        result = result.reset_index()
        result['OTHERS'] = (
            result['ANTHRACYCLINE BASED'] +
            result['OTHER COMBINATION CHEMO'] +
            result['OTHER SINGLE CHEMO'] +
            result['PARP BASED'] +
            result['PLATINUM BASED']
        )
        print("colms!!!!!!!!", result.columns)
        source_of_business_ = result[['Line', 'End Date', 'LOT', "OTHERS",'TAXANE BASED', 'VINORELBINE BASED',
                                    'GEMZAR BASED', 'XELODA BASED','HALAVEN BASED', 'PD-1/PD-L1 BASED','ENHERTU BASED']]
        source_of_business_ = source_of_business_.reset_index()

        print("source_of_business!!!!!!! husna",source_of_business_)
        source_of_business_sorted = source_of_business_.sort_values(by=['LOT', 'End Date'], ascending=[False, False])
        source_of_business_final = source_of_business_sorted.drop(columns=['LOT', 'End Date'])

        new_df = pd.DataFrame(columns=source_of_business_final.columns)
        rows = []

        for i in range(len(source_of_business_final)):
            rows.append(source_of_business_final.iloc[i])
            if (i + 1) % 2 == 0:
                rows.append(pd.Series([None]*len(source_of_business_final.columns), index=source_of_business_final.columns))

        new_df = pd.concat([new_df, pd.DataFrame(rows)], ignore_index=True)


        # Use pd.concat to build the new DataFrame from the list of rows
        new_df = pd.concat([pd.DataFrame([row], columns=source_of_business_final.columns) for row in rows], ignore_index=True)



        source_of_business = new_df[['Line', "OTHERS",'TAXANE BASED', 'VINORELBINE BASED',
                                    'GEMZAR BASED', 'XELODA BASED','HALAVEN BASED', 'PD-1/PD-L1 BASED','ENHERTU BASED']]
        source_of_business


        print("source of business",source_of_business)

        def get_nps_share(data, line_of_therapy, num_of_months):
            # data['Year_Month'] = pd.to_datetime(data['Year_Month'], format='%Y%m')
            data['Line_of_Therapy_LOT'] = data['Line_of_Therapy_LOT'].astype(int)
            # data['Year_Month_dt'] = pd.to_datetime(data['Year_Month'], format='%Y%m')
            data['Year_Month_dt'] = pd.to_datetime(data['Year_Month'].astype(str), format='%Y%m')

            # Make sure start_date and end_date are datetime objects
            # If they are integers like 202301, convert them similarly:
            latest_year_month = pd.to_datetime(f'{latest_year_month_value}', format='%Y%m')
            end_dates = [latest_year_month - pd.DateOffset(months=i) for i in range(num_of_months)]
            # print("end dates list",end_dates)
            start_dates = [end_date - pd.DateOffset(months=3) for end_date in end_dates]

            print("start_dates!!!!!!", start_dates)
            print("end_dates!!!!!!", end_dates)

            # data['Year_Month'] = data['Year_Month'].astype(int)
        # #     latest_year_month = data['Year_Month'].max()
            all_results = pd.DataFrame()

            for start_date, end_date in zip(start_dates, end_dates):
                # start_date_dt = pd.to_datetime(str(start_date), format='%Y%m')
                # end_date_dt = pd.to_datetime(str(end_date), format='%Y%m')
                # period_data = data[(data['Year_Month'] >= start_date) & (data['Year_Month'] <= end_date)]

                # period_data = data[(data['Year_Month_dt'] >= start_date_dt) & (data['Year_Month_dt'] <= end_date_dt)]
                period_data = data[(data['Year_Month_dt'] >= start_date) & (data['Year_Month_dt'] <= end_date)]

                if line_of_therapy == '2+':
                    lot_data = period_data[period_data['Line_of_Therapy_LOT'] >= 2]
                elif line_of_therapy == '3+':
                    lot_data = period_data[period_data['Line_of_Therapy_LOT'] >= 3]
                elif line_of_therapy == '4+':
                    lot_data = period_data[period_data['Line_of_Therapy_LOT'] >= 4]
                else:
                    lot_data = period_data[period_data['Line_of_Therapy_LOT'] == int(line_of_therapy)]

                if lot_data.empty:
                    print(f"No data for line {line_of_therapy} in period {start_date} to {end_date}")
                    continue

                # Calculate the total for the entire period, not just monthly
                period_total = lot_data['Monthly_Adjusted_Count'].sum()
                period_count = lot_data['Monthly_Adjusted_Count'].count()
                if period_total == 0:
                    print(f"No total adjusted count for line {line_of_therapy} in period {start_date} to {end_date}")
                    continue

                # Group and calculate sums for each regimen within the period
                regimen_totals = lot_data.groupby('Regimen_Group')['Monthly_Adjusted_Count'].sum()
                percentage_shares = (regimen_totals / period_total)
                percentage_shares_df = percentage_shares.reset_index()
                percentage_shares_df.columns = ['Regimen_Group', 'Percentage']
                percentage_shares_df['Year_Month'] = end_date.strftime('%b-%y')
                percentage_shares_df['Line'] = line_of_therapy
                percentage_shares_df['Total_Count'] = period_count

                # Append to results
                all_results = pd.concat([all_results, percentage_shares_df], ignore_index=True)

            if all_results.empty:
                print("No data available for the given filters.")
                return all_results

            all_results = all_results.pivot_table(index=['Year_Month', 'Line', 'Total_Count'], columns='Regimen_Group', values='Percentage', fill_value=0).reset_index()

            # Flatten columns
            all_results.columns.name = None
            all_results = all_results.rename_axis(None, axis=1)

            # Sorting the results
            final_result = all_results.sort_values(by=['Year_Month', 'Line'])

            # Rounding the percentages
            percentage_columns = final_result.columns.difference(['Year_Month', 'Line'])
            final_result[percentage_columns] = final_result[percentage_columns]
            final_result = final_result.drop(columns= "Line")
            print("final_result!!!!!", final_result.columns)
        #     adc_df, tpc_df = '', ''

            adc_columns = ['ENHERTU BASED', 'TRODELVY BASED']
            tpc_columns = ['HALAVEN BASED', 'XELODA BASED', 'GEMZAR BASED','VINORELBINE BASED']

            adc_df = final_result[['Year_Month'] + [col for col in adc_columns if col in final_result.columns]]
            adc_df['ADC'] = (adc_df[adc_columns].sum(axis=1)*100).round(0).astype(int)
            adc_df = adc_df[['Year_Month', 'ADC']]

            tpc_df = final_result[['Year_Month'] + [col for col in tpc_columns if col in final_result.columns]]
            tpc_df['TPC'] = (tpc_df[tpc_columns].sum(axis=1)*100).round(0).astype(int)
            tpc_df = tpc_df[['Year_Month', 'TPC']]

            if line_of_therapy == "2+":
                final_result["OTHERS"] = final_result["PARP BASED"] + final_result["PD-1/PD-L1 BASED"]
                final_result['OTHER CHEMO'] = final_result["TAXANE BASED"] + final_result["PLATINUM BASED"] +final_result["OTHER COMBINATION CHEMO"] + final_result["OTHER SINGLE CHEMO"] +\
                                            final_result["ANTHRACYCLINE BASED"]
                final_result["TPC"] =   final_result["HALAVEN BASED"] + final_result["XELODA BASED"] +final_result["GEMZAR BASED"] + final_result["VINORELBINE BASED"]
                final_result = final_result[['Year_Month', 'OTHERS','OTHER CHEMO','TPC','ENHERTU BASED','TRODELVY BASED', 'Total_Count']]

            final_result = process_data(final_result, False)
            final_result = final_result.drop(columns=['date_column'])

            return final_result, adc_df, tpc_df

        def add_patient_counts(df_):
            df = df_.copy()
            df['Year_Month'] = df['Year_Month'] + " \n(" + df['Total_Count'].astype(str) + ")"
            df = df.drop(columns=['Total_Count'])
            return df

        nps_line_above_2_, adc_above_2, tpc_above_2 = get_nps_share(tnbc_nps_data_apr, line_of_therapy='2+', num_of_months = 5)
        nps_line_above_2 = add_patient_counts(nps_line_above_2_)
        print("nps_line_above 2 Husna",nps_line_above_2)

        forecast_data = {
            "Year_Month": ["Dec-23","Jan-24", "Feb-24", "Mar-24", "Apr-24", "May-24", "Jun-24"],
            "Unknown/L12": [0.04, 0.05, 0.06, 0.05, 0.05, 0.05, 0.05],
            "2L3": [0.30, 0.32, 0.33, 0.33, 0.33, 0.34, 0.35],
        #     "3L4": [0.25, 0.24, 0.25, 0.25, 0.26, 0.26, 0.27],
            "2L+5": [0.27, 0.29, 0.30, 0.30, 0.30, 0.31, 0.31],
            "3L+6": [0.24, 0.25, 0.25, 0.26, 0.26, 0.26, 0.27]

        }

        forecast_data_df = pd.DataFrame(forecast_data)
        trend_nps_line_1,adc_1, tpc_1 = get_nps_share(tnbc_nps_data_apr, line_of_therapy=1,num_of_months =22)

        trend_nps_line_2, trend_adc_2, trend_tpc_2 = get_nps_share(tnbc_nps_data_apr, line_of_therapy=2,num_of_months =22)
        # trend_nps_line_3 = get_nps_share(tnbc_nps_data_apr, line_of_therapy=3,num_of_months =22)
        trend_nps_line_above_2, trend_adc_above_2, trend_tpc_above_2 = get_nps_share(tnbc_nps_data_apr, line_of_therapy='2+', num_of_months =22)
        print("trend nps_line_above 2 Husna",trend_nps_line_above_2)
        trend_nps_line_above_3, trend_adc_above_3, trend_tpc_above_3 = get_nps_share(tnbc_nps_data_apr, line_of_therapy='3+', num_of_months =22)
        # trend_nps_line_above_4 = get_nps_share(tnbc_nps_data_apr, line_of_therapy='4+', num_of_months =22)

        nps_line_1_trend= trend_nps_line_1[['Year_Month', 'TRODELVY BASED']]
        nps_line_2_trend = trend_nps_line_2[['Year_Month', 'TRODELVY BASED']]
        # nps_line_3_trend = trend_nps_line_3[['Year_Month', 'TRODELVY BASED']]
        nps_line_above_2_trend = trend_nps_line_above_2[['Year_Month', 'TRODELVY BASED']]
        nps_line_above_3_trend = trend_nps_line_above_3[['Year_Month', 'TRODELVY BASED']]
        # nps_line_above_4_trend = trend_nps_line_above_4[['Year_Month', 'TRODELVY BASED']]

        nps_line_1_trend.rename(columns={'TRODELVY BASED': '1L'}, inplace=True)
        nps_line_2_trend.rename(columns={'TRODELVY BASED': '2L'}, inplace=True)
        # nps_line_3_trend.rename(columns={'TRODELVY BASED': '3L'}, inplace=True)
        nps_line_above_2_trend.rename(columns={'TRODELVY BASED': '2L+'}, inplace=True)
        print("nps_line_above 2 Husna trend",nps_line_above_2_trend)
        nps_line_above_3_trend.rename(columns={'TRODELVY BASED': '3L+'}, inplace=True)
        # nps_line_above_4_trend.rename(columns={'TRODELVY BASED': '4L+'}, inplace=True)

        final_merged_df = pd.merge(
            pd.merge(
                nps_line_1_trend,
                nps_line_2_trend,
                on='Year_Month'
            ),
            pd.merge(
                nps_line_above_2_trend,
                nps_line_above_3_trend,
                on='Year_Month'
            ),
            on='Year_Month'
        )
        nps_trends_in_all_lines = pd.merge(final_merged_df, forecast_data_df, on='Year_Month', how= "left")
        nps_trends_in_all_lines_final = process_data(nps_trends_in_all_lines, False)
        nps_trends_in_all_lines_final = nps_trends_in_all_lines_final.drop(columns= "date_column")

        last_row_values = nps_trends_in_all_lines_final.iloc[-1][['1L', '2L', '2L+', '3L+']]
        last_row_values_dict = last_row_values.to_dict()
        percentage_dict = {key: str(round(value * 100))+ "%" for key, value in last_row_values_dict.items()}

        req_columns = ['Year_Month', 'OTHER SINGLE CHEMO','OTHER COMBINATION CHEMO','ANTHRACYCLINE BASED','PARP BASED',
                          'PLATINUM BASED','TAXANE BASED','VINORELBINE BASED','GEMZAR BASED',
                          'XELODA BASED','HALAVEN BASED','PD-1/PD-L1 BASED','ENHERTU BASED','TRODELVY BASED', 'Total_Count']
        nps_line_2, adc_2l, tpc_2l = get_nps_share(tnbc_nps_data_apr, line_of_therapy=2, num_of_months = 8)
        nps_line_2_final = process_data(nps_line_2, False)
        nps_line_2_final = nps_line_2_final.drop(columns= "date_column")
        nps_line_2_final = nps_line_2_final[req_columns]
        nps_line_2_final = add_patient_counts(nps_line_2_final)

        adc_2l_final = process_data(adc_2l, False).drop(columns = "date_column")
        tpc_2l_final = process_data(tpc_2l, False).drop(columns = "date_column")

        nps_line_above_3, adc_above_3l, tpc_above_3l = get_nps_share(tnbc_nps_data_apr, line_of_therapy = "3+", num_of_months = 8)
        nps_line_above_3_final = process_data(nps_line_above_3, False)
        nps_line_above_3_final = nps_line_above_3_final.drop(columns= "date_column")
        nps_line_above_3_final = nps_line_above_3_final[req_columns]
        nps_line_above_3_final = add_patient_counts(nps_line_above_3_final)

        adc_above_3l = process_data(adc_above_3l, False).drop(columns = "date_column")
        tpc_above_3l = process_data(tpc_above_3l, False).drop(columns = "date_column")

        difference_2l = str(tpc_2l_final['TPC'].iloc[-1] - tpc_2l_final['TPC'].iloc[0]) + "%"
        difference_3l = str(tpc_above_3l['TPC'].iloc[-1] - tpc_above_3l['TPC'].iloc[0]) + "%"

        adc_tpc_merged_2l = pd.merge(adc_2l_final, tpc_2l_final,  on='Year_Month', how='inner')
        adc_tpc_merged_d_2l = adc_tpc_merged_2l.drop('Year_Month', axis=1)
        adc_tpc_merged_d_2l = adc_tpc_merged_d_2l.applymap(lambda x: f"{x}%" if pd.notnull(x) else x)
        adc_tpc_df_2l = adc_tpc_merged_d_2l.T

        adc_tpc_merged_d_3l = pd.merge(adc_above_3l, tpc_above_3l,  on='Year_Month', how='inner')
        adc_tpc_merged_d_3l = adc_tpc_merged_d_3l.drop('Year_Month', axis=1)
        adc_tpc_merged_d_3l = adc_tpc_merged_d_3l.applymap(lambda x: f"{x}%" if pd.notnull(x) else x)
        adc_tpc_df_3l = adc_tpc_merged_d_3l.T

        def get_nps_share_monthly(data, line_of_therapy):
            # data['Year_Month'] = pd.to_datetime(data['Year_Month'])
            #latest_year_month_value
            print("line of therapy retrived from get_nps_share_monthly",line_of_therapy)
            data=load_data_concurrently(latest_year_month_value,line_of_therapy=line_of_therapy)
            size_mb = data.memory_usage(deep=True).sum() / (1024 * 1024)

            print(f"DataFrame size for {line_of_therapy} tps data apr: {size_mb:.2f} MB")
            data['Year_Month_dt'] = pd.to_datetime(data['Year_Month'], format='%Y%m')
            data['Year_Month_dt'] = pd.to_datetime(data['Year_Month'].astype(str), format='%Y%m')
            data['Line_of_Therapy_LOT'] = data['Line_of_Therapy_LOT'].astype(int)
            unique_months = data['Year_Month'].sort_values(ascending=False).unique()
            all_results = pd.DataFrame()

            for month in unique_months[:8]:
                period_data = data[data['Year_Month'] == month]
                print("here is the month",month)
                month_dt = pd.to_datetime(str(month), format='%Y%m')  # or omit format if ISO date string

                print(month_dt.strftime('%b-%y'))

                if line_of_therapy == '2+':
                    lot_data = period_data[period_data['Line_of_Therapy_LOT'] >= 2]
                elif line_of_therapy == '3+':
                    lot_data = period_data[period_data['Line_of_Therapy_LOT'] >= 3]
                elif line_of_therapy == '4+':
                    lot_data = period_data[period_data['Line_of_Therapy_LOT'] >= 4]
                else:
                    lot_data = period_data[period_data['Line_of_Therapy_LOT'] == int(line_of_therapy)]

                if lot_data.empty:
                    print(f"No data for line {line_of_therapy} in {month_dt.strftime('%b-%y')}")
                    continue

                # Calculate the total adjusted count for the month
                period_total = lot_data['Monthly_Adjusted_Count'].sum()
                period_count = lot_data['Monthly_Adjusted_Count'].count()
                if period_total == 0:
                    print(f"No total adjusted count for line {line_of_therapy} in {month.strftime('%b-%y')}")
                    continue

                # Group and calculate sums for each regimen within the month
                regimen_totals = lot_data.groupby('Regimen_Group')['Monthly_Adjusted_Count'].sum()
                percentage_shares = (regimen_totals / period_total)
                percentage_shares_df = percentage_shares.reset_index()
                percentage_shares_df.columns = ['Regimen_Group', 'Percentage']
                percentage_shares_df['Year_Month'] = pd.to_datetime(str(month), format='%Y%m').strftime('%b-%y')
                # percentage_shares_df['Year_Month'] = pd.to_datetime(month).strftime('%b-%y')
        #         percentage_shares_df['Year_Month'] = month.strftime('%b-%y')
                percentage_shares_df['Line_of_Therapy_LOT'] = line_of_therapy
                percentage_shares_df['Total_Count'] = period_count

                # Append to results
                all_results = pd.concat([all_results, percentage_shares_df], ignore_index=True)


            if all_results.empty:
                print("No data available for the given filters.")
                return all_results

            all_results = all_results.pivot_table(index=['Year_Month', 'Line_of_Therapy_LOT', 'Total_Count'], columns='Regimen_Group', values='Percentage', fill_value=0).reset_index()

            all_results.columns.name = None
            all_results = all_results.rename_axis(None, axis=1)
            final_result = process_data(all_results, False)
            final_result = final_result.drop(columns=['Line_of_Therapy_LOT'])
            final_result['Year_Month'] = final_result['Year_Month'] + " \n(" + final_result['Total_Count'].astype(str) + ")"
            final_result = final_result.drop(columns=['Total_Count'])

            # Rounding the percentages
            percentage_columns = final_result.columns.difference(['Year_Month'])
            final_result[percentage_columns] = final_result[percentage_columns].round(2)

            adc_columns = ['ENHERTU BASED', 'TRODELVY BASED']
            tpc_columns = ['HALAVEN BASED', 'XELODA BASED', 'GEMZAR BASED','VINORELBINE BASED']

            tps_adc_df = final_result[['Year_Month'] + [col for col in adc_columns if col in final_result.columns]]
            tps_adc_df['ADC'] = (tps_adc_df[adc_columns].sum(axis=1)*100).round(0).astype(int)
            tps_adc_df = tps_adc_df[['Year_Month', 'ADC']]

            tps_tpc_df = final_result[['Year_Month'] + [col for col in tpc_columns if col in final_result.columns]]
            tps_tpc_df['TPC'] = (tps_tpc_df[tpc_columns].sum(axis=1)*100).round(0).astype(int)
            tps_tpc_df = tps_tpc_df[['Year_Month', 'TPC']]

            tps_column_order = ['Year_Month', "OTHER SINGLE CHEMO", "OTHER COMBINATION CHEMO", "ANTHRACYCLINE BASED", "PARP BASED", "PLATINUM BASED", "TAXANE BASED", "VINORELBINE BASED", "GEMZAR BASED", "XELODA BASED", "HALAVEN BASED", "PD-1/PD-L1 BASED", "ENHERTU BASED", "TRODELVY BASED"]
            final_result = final_result[tps_column_order]
            return final_result, tps_adc_df, tps_tpc_df

        tps_line_2, tps_adc_line_2, tps_tpc_line_2  = get_nps_share_monthly("tnbc_tps_data_apr", '2')
        tps_line_above_3, tps_adc_above_3, tps_tpc_above_3  = get_nps_share_monthly("tnbc_tps_data_apr", '3+')

        tps_difference_2l = str(tps_tpc_line_2['TPC'].iloc[-1] - tps_tpc_line_2['TPC'].iloc[0]) + "%"
        tps_difference_above_3l = str(tps_tpc_above_3['TPC'].iloc[-1] - tps_tpc_above_3['TPC'].iloc[0]) + "%"

        tps_adc_tpc_merged_2l_ = pd.merge(tps_adc_line_2, tps_tpc_line_2,  on='Year_Month', how='inner')
        tps_adc_tpc_merged_d_2l = tps_adc_tpc_merged_2l_.drop('Year_Month', axis=1)
        tps_adc_tpc_merged_d_2l = tps_adc_tpc_merged_d_2l.applymap(lambda x: f"{x}%" if pd.notnull(x) else x)
        tps_adc_tpc_df_2l = tps_adc_tpc_merged_d_2l.T

        tps_adc_tpc_merged_d_3l_ = pd.merge(tps_adc_above_3, tps_tpc_above_3,  on='Year_Month', how='inner')
        tps_adc_tpc_merged_d_3l = tps_adc_tpc_merged_d_3l_.drop('Year_Month', axis=1)
        tps_adc_tpc_merged_d_3l = tps_adc_tpc_merged_d_3l.applymap(lambda x: f"{x}%" if pd.notnull(x) else x)
        tps_adc_tpc_df_3l = tps_adc_tpc_merged_d_3l.T
        tps_adc_tpc_df_3l

        persistency_2l_forecast_data_ = {
            "Month_Num": ["M0", "M01", "M02", "M03", "M04", "M05", "M06", "M07", "M08", "M09", "M10", "M11", "M12", "M13", "M14", "M15", "M16", "M17", "M18", "M19", "M20", "M21", "M22", "M23",
                           "M24", "M25", "M26", "M27", "M28", "M29", "M30", "M31", "M32", "M33", "M34", "M35", "M36", "M37", "M38", "M39", "M40"],
            "Forecast": [1.0, 1.0, 0.82, 0.65, 0.51, 0.43, 0.36, 0.31, 0.28, 0.25, 0.23, 0.19, 0.16, 0.13, 0.11, 0.11, 0.1, 0.08, 0.07, 0.05, 0.04, 0.02, 0.0,"",
                         "", "", "", "", "", "", "", "", "", "", "", "", "", "", "", "", ""]
        }
        persistency_3l_forecast_data_ = {
            "Month_Num": ["M0", "M01", "M02", "M03", "M04", "M05", "M06", "M07", "M08", "M09", "M10", "M11", "M12", "M13", "M14", "M15", "M16", "M17", "M18", "M19", "M20", "M21", "M22", "M23",
                           "M24", "M25", "M26", "M27", "M28", "M29", "M30", "M31", "M32", "M33", "M34", "M35", "M36", "M37", "M38", "M39", "M40"],
            "Forecast": [1.0, 1.0, 0.84, 0.68, 0.52, 0.42, 0.35, 0.29, 0.23, 0.18, 0.14, 0.11, 0.11, 0.09, 0.08, 0.07, 0.06, 0.05, 0.04, 0.03, 0.02, 0.0, "",
                         "", "", "", "", "", "", "", "", "", "", "", "", "", "", "", "","",""]
        }
        persistency_above_4l_forecast_data_ = {
            "Month_Num": ["M0", "M01", "M02", "M03", "M04", "M05", "M06", "M07", "M08", "M09", "M10", "M11", "M12", "M13", "M14", "M15", "M16", "M17", "M18", "M19", "M20", "M21", "M22", "M23",
                           "M24", "M25", "M26", "M27", "M28", "M29", "M30", "M31", "M32", "M33", "M34", "M35", "M36", "M37", "M38", "M39", "M40"],
            "Forecast": [1.0, 1.0, 0.84, 0.65, 0.52, 0.45, 0.39, 0.32, 0.25, 0.21, 0.17, 0.13, 0.1, 0.07, 0.05, 0.04, 0.03, 0.02, 0.0, 0.0, "",
                         "", "", "", "", "", "", "", "", "", "", "", "", "", "", "", "", "", "","",""]
        }
        persistency_2l_forecast_data = pd.DataFrame(persistency_2l_forecast_data_)
        persistency_3l_forecast_data = pd.DataFrame(persistency_3l_forecast_data_)
        persistency_above_4l_forecast_data = pd.DataFrame(persistency_above_4l_forecast_data_)

        new_row = pd.DataFrame({ "Month_Num": ["M0"], "2020": [1],"2021": [1], "2022": [1],"2023": [1],"overall": [1]})

        def get_dot(df_perst):
            columns_to_check = df_perst.columns[1:-1]
            dot_dict = {}
            for col in columns_to_check:
                for i in range(1, len(df_perst)):
                    if df_perst[col][i-1] >= 0.50 > df_perst[col][i]:
                        month_start = int(df_perst['Month_Num'][i-1][1:])
                        month_end = int(df_perst['Month_Num'][i][1:])
                        dot_dict[col] = round(month_start * df_perst[col][i-1] + month_end * df_perst[col][i],1)
            dot_df = pd.DataFrame(list(dot_dict.items()), columns=['Year', 'dot'])
            return dot_df


        def get_persistency_data(lot):

            if lot == "4L+":
              start_time=time.time()
              data_query = f"""
                  SELECT *
                  FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`onc_persistent_data`
                  WHERE lot >= 4
              """

              perst_data = pd.DataFrame(self.cursor.execute(data_query).fetchall(),
                              columns=[elem[0] for elem in self.cursor.description])
              end_time=time.time()
              # perst_data = perst_data_old_batch[perst_data_old_batch['Line_of_Therapy_LOT'] >= 4]
            else:
                # perst_data = perst_data_old_batch[perst_data_old_batch['Line_of_Therapy_LOT'] == lot]
                start_time=time.time()
                data_query = f"""
                    SELECT *
                    FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`onc_persistent_data`
                    WHERE lot == {lot}
                """

                perst_data = pd.DataFrame(self.cursor.execute(data_query).fetchall(),
                                  columns=[elem[0] for elem in self.cursor.description])
                end_time=time.time()

            print("Columns in perst_data:", perst_data.columns)

            perst_data['Year_Month'] = pd.to_datetime(perst_data['first_claim_dt']).dt.strftime('%Y%m')
            print("Max first_claim_dt in perst_data:", perst_data['first_claim_dt'].max())

            sorted_df = perst_data.sort_values(by='first_claim_dt', ascending=False)
            max_date = sorted_df['first_claim_dt'].max()
            print("Max first_claim_dt:", max_date)

            perst_data['Year'] = pd.to_datetime(perst_data['first_claim_dt']).dt.strftime('%Y')
            perst_data_ = perst_data[perst_data['Year_Month'] <= max_date] # change here for date range

            print("here is the perst data",perst_data_)

            perst_data_['patients_count'] = (perst_data_['perst'] > 0).astype(int)
            perst_data_january = perst_data_[perst_data_['mnth_num'] == 1]
            yearly_patient_counts = perst_data_january.groupby('Year')['patients_count'].sum().reset_index()

            yearly_patient_counts.columns = ['Year', 'Total_Patients']
            total_sum = yearly_patient_counts['Total_Patients'].sum()
            overall_row = pd.DataFrame({'Year': ['Overall'], 'Total_Patients': [total_sum]})
            yearly_patient_counts = pd.concat([yearly_patient_counts, overall_row], ignore_index=True)

            average_perst = perst_data_.groupby(['Year', 'p_mnth'])['perst'].mean().reset_index()
            average_perst.columns = ['Year', 'Month_Num', 'Average_Perst']
            perst_by_yr = average_perst.pivot_table(index=['Month_Num'], columns=['Year'], values='Average_Perst').reset_index()

            perst_by_yr.columns.name = None
            perst_by_yr.columns = ['Month_Num'] + [f'{col}' for col in perst_by_yr.columns[1:]]

            print(perst_by_yr)
            overall_average_perst = perst_data_.groupby(['p_mnth'])['perst'].mean().reset_index()
            overall_average_perst.columns = [ 'Month_Num', 'Overall']
            all_perst_df = pd.merge(perst_by_yr, overall_average_perst, on='Month_Num', how='outer')

            perst_df_M0 = pd.concat([new_row, all_perst_df ], ignore_index=True)
            if lot == 2:
                forescat_data = persistency_2l_forecast_data
            elif lot == 3:
                forescat_data = persistency_3l_forecast_data
            elif lot == "4L+":
                forescat_data = persistency_above_4l_forecast_data

            final_perst_df  = pd.merge(forescat_data, perst_df_M0 , on='Month_Num', how='left')
            perst_column_order = ["Month_Num",  "2020", "2021", "2022", "2023","Overall", "Forecast"]
            final_perst_df = final_perst_df[perst_column_order]
            dot = get_dot(final_perst_df)
            yearly_patient_counts_dot  = pd.merge(yearly_patient_counts, dot , on='Year', how='outer')
            return final_perst_df, yearly_patient_counts_dot

        final_perst_df_2l, yearly_patient_counts_2l = get_persistency_data(2)
        final_perst_df_3l, yearly_patient_counts_3l = get_persistency_data(3)
        final_perst_df_above_4, yearly_patient_counts_above_4 = get_persistency_data('4L+')
        print("final persistent df 2l",final_perst_df_2l.head())

        from pptx import Presentation
        ppt_path = './Structured_Bot/TNBC_slides_report_automation_template.pptx'
        # ppt_path = 'mbc_refreshed_ppt_1.pptx'
        presentation = Presentation(ppt_path)


        for slide in presentation.slides:
            for shape in slide.shapes:
                final_df = []
                if shape.has_chart:
                    chart = shape.chart
                    chart_title = chart.chart_title.text_frame.text
                    print("chart_title!!!!", chart_title)
                    if "mTNBC New Patient Share (R3M) – 2L+" in chart_title:
                        
                        print("slide 1 sql chart 1 ",nps_line_above_2.head(10))
                        final_df = nps_line_above_2
                    elif "TRO mTNBC New Patient Share (R3M) – by LoT" in chart_title:
                        final_df = nps_trends_in_all_lines_final
                        print("slide 1 sql chart 2",nps_trends_in_all_lines_final.head(10))
                    elif "2L+ mTNBC Trended HCP Segment" in chart_title:
                        print("slide 2 sql chart 1",HCP_trend_segment.head(10))
                        final_df = HCP_trend_segment
                    elif chart_title == "Adopters" :
                        print("slide 2 sql chart 2",adopter_df_sorted.head(10))
                        final_df = adopter_df_sorted
                    elif chart_title == "Dabblers" :
                        print("slide 2 sql chart 3",dabbler_df_sorted.head(10))
                        final_df = dabbler_df_sorted
                    elif chart_title == "Non-users" :
                        print("slide 2 sql chart 4",non_user_df_sorted.head(10))
                        final_df = non_user_df_sorted


                    elif not chart_title:
                        # Try reading from paragraphs
                        text = ""
                        for paragraph in chart.chart_title.text_frame.paragraphs:
                            for run in paragraph.runs:
                                text += run.text
                        chart_title = text.strip()
                    # elif "Source of Business" in chart_title:
                        final_df = source_of_business
                        print("slide 5 sql chart 1",final_df.head(10))

                    # elif chart_title.contains("Source of Business") :
                    #     final_df = source_of_business
                    #     print("here is the final df of source of business",final_df)
                    elif "2L New Patients Shares by Regimen" in chart_title:
                        print("slide 3 sql chart 1",nps_line_2_final.head(10))
                        final_df = nps_line_2_final
                        
                    elif "3L+ New Patients Shares by Regimen" in chart_title:
                        print("slide 3 sql chart 2",nps_line_above_3_final.head(10))
                        final_df = nps_line_above_3_final
                    elif "2L Total Patients Shares by Regimen" in chart_title:
                         print("slide 4 sql chart 1",tps_line_2.head(10))
                         final_df = tps_line_2
                    elif "3L+ Total Patients Shares by Regimen" in chart_title:
                         print("slide 4 sql chart 2",tps_line_above_3.head(10))
                         final_df = tps_line_above_3
                    elif "TNBC 2L Persistency" in chart_title:
                        print("slide 6 sql chart 1",final_perst_df_2l.head(10))
                        final_df = final_perst_df_2l
                    elif "TNBC 3L Persistency" in chart_title:
                        print("slide 7 sql chart 1",final_perst_df_3l.head(10))
                        final_df = final_perst_df_3l
                    elif "TNBC 4L+ Persistency" in chart_title:
                        print("slide 8 sql chart 1",final_perst_df_above_4.head(10))
                        final_df = final_perst_df_above_4

                    if len(final_df) >0:
                        final_df.replace([np.nan, np.inf, -np.inf], '', inplace=True)
                        print("Chart found!", chart.chart_title.text_frame.text)
                        chart_data = CategoryChartData()
                        categories = final_df.iloc[:, 0].tolist()  # Assuming the first column is the category
                        chart_data.categories = categories
                        for col in final_df.columns[1:]:
                            chart_data.add_series(col, final_df[col].tolist())

                        chart.replace_data(chart_data)
                        print("replaced data!!")

                    print("at the end!!!")

        output_path = "claims_refreshed_ppt.pptx"
        presentation.save(output_path)

        def update_text(fill_text, paragraph):
            print("in update text!!!!!!!")
            paragraph.text = fill_text




        def fill_table_content(pptx_file):
            df_table = False
            start_row = 0
            prs = Presentation(pptx_file)
            for slide_number, slide in enumerate(prs.slides):
                print(f"Slide {slide_number + 1}")
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        flag_value = table.cell(0, 0).text
                        print("cell!!!", flag_value)
                        if flag_value == "adc_2l"  :
                            df_table = True
                            df = adc_tpc_df_2l
                        elif flag_value == "adc_3l":
                            df_table = True
                            df = adc_tpc_df_3l
                        elif "t_adc_2l" in flag_value:
                            df_table = True
                            df = tps_adc_tpc_df_2l
                        elif "t_adc_3l" in flag_value:
                            df_table = True
                            df = tps_adc_tpc_df_3l
                        elif "Period (2L)" in flag_value:
                            start_row = 1
                            df_table = True
                            df = yearly_patient_counts_2l
                        elif "Period (3L)" in flag_value:
                            start_row = 1
                            df_table = True
                            df = yearly_patient_counts_3l
                        elif "Period (4L+)" in flag_value:
                            start_row = 1
                            df_table = True
                            df = yearly_patient_counts_above_4

                        if df_table:
                            df_table = False
                            for i, row in enumerate(table.rows):
                                for j, cell in enumerate(row.cells):
                                    if i >= start_row:
                                        df_row_index = i - start_row
                                        if df_row_index < len(df) and j < len(df.columns):
                                            value = str(df.iloc[df_row_index, j])
                                            print("value!!!!!!!!!!", value)
                                            self.update_cell_text(cell, value)
                                        else:
                                            self.update_cell_text(cell, "")

                    if shape.has_text_frame:
                        for paragraph in shape.text_frame.paragraphs:
                            original_text = paragraph.text
                            print("Original Paragraph Text:", original_text)
        #                     {'1L': '5%', '2L': '31%', '2L+': '27%', '3L+': '22%'}
        #                     2L: 31%

                            if "1L: %" in original_text:
                                update_text('1L: '+ percentage_dict['1L'], paragraph)
                            elif "2L+: %" in original_text:
                                update_text('2L+: '+ percentage_dict['2L+'], paragraph)
                            elif "2L: %" in original_text:
                                update_text('2L: '+ percentage_dict['2L'], paragraph)
                            elif "3L+: %" in original_text:
                                update_text('3L+: '+ percentage_dict['3L+'], paragraph)

                            elif "PT1" in original_text:
                                update_text(str(tier_counts_past['past Tier 1']), paragraph)
                            elif "PT2" in original_text:
                                update_text(str(tier_counts_past['past Tier 2']), paragraph)
                            elif "PT3" in original_text:
                                update_text(str(tier_counts_past['past Tier 3']), paragraph)
                            elif "RT1" in original_text:
                                update_text(str(tier_counts_present['present Tier 1']), paragraph)
                            elif "RT2" in original_text:
                                update_text(str(tier_counts_present['present Tier 2']), paragraph)
                            elif "RT3" in original_text:
                                update_text(str(tier_counts_present['present Tier 3']), paragraph)
                            elif original_text == "tpc2"  :
                                update_text(difference_2l, paragraph)
                            elif original_text ==  "tpc3" :
                                update_text(difference_3l, paragraph)
                            elif original_text == "t_tpc2"  :
                                update_text(tps_difference_2l, paragraph)
                            elif original_text ==  "t_tpc3" :
                                update_text(tps_difference_above_3l, paragraph)

            prs.save("claims_refreshed_ppt.pptx")
            print("before the cursor closing in update onc claims")
            self.cursor.close()
            print("after the cursor closing in update onc claims")
            print("before the connection closing in update onc claims")
            self.connection.close()
            print("after the connection closing in update onc claims")

        # ppt_path = '/Users/kavyareddy/docker_scraping_2/agents_based_chatbots/TNBC slides/TNBC_slides_report_automation_test.pptx'
        fill_table_content(output_path)



class Update_Onc_Sales_Deck(Update_Deck):

    def __init__(self):
        super().__init__()
        self.connection = None
        self.cursor = None
        self.latest_weekend = ""

        self.latest_day_mon = ""
        self.latest_year = ""
        self.latest_month = ""

    def update_deck(self):
        print("in update deck onc sales")
        self.connect_to_databricks()
        print("Connection Established in update onc sales deck ")


        # tro_sales_data = tro_sales_data.compute()  # Convert Dask to Pandas
        # ddd_sales_data = ddd_sales_data.compute()  # Convert Dask to Pandas
        # tro_sales_data['drug_name']='trodelvy'
        start_time=time.time()
        sql_query = f"""
                SELECT * FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`867_data_drug-trodelvy_for_genie_filtered_dates`"""
        tro_sales_data= pd.DataFrame(self.cursor.execute(sql_query).fetchall(), columns=[elem[0] for elem in self.cursor.description])
        end_time=time.time()
        print("latency for getting the data",start_time-end_time)

        tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date']).dt.tz_localize(None)


        print("tro sales data sample",tro_sales_data.head())

                # Query to get the latest weekend for 867 data
        sorted_dates = pd.DataFrame(self.cursor.execute("""SELECT DISTINCT `week_end_date`
                    FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`867_data_drug-trodelvy_for_genie_filtered_dates`
                    ORDER BY `week_end_date` DESC""").fetchall(), columns=[elem[0] for elem in self.cursor.description])
        # Extract and format the latest weekend
        latest_weekend = pd.to_datetime(sorted_dates.iloc[0, 0]).strftime('%Y-%m-%d')
        self.latest_weekend = latest_weekend
        print("Latest weekend (867 data):", latest_weekend)

        # Query to get the latest weekend for DDD data
        start_time=time.time()
        latest_weekend_ddd = pd.DataFrame(self.cursor.execute("""SELECT
                    DATE_FORMAT(MAX(week_end_date), 'yyyy-MM-dd') AS latest_weekend_ddd
                FROM
                    `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`ddd_data_for_genie_filtered_dates`;""").fetchall(),
                columns=[elem[0] for elem in self.cursor.description])
        end_time=time.time()
        print("latency for getting the data ddd weekend date",start_time-end_time)
        # Convert latest_weekend_ddd to datetime
        latest_weekend_ddd = pd.to_datetime(latest_weekend_ddd.iloc[0, 0]).strftime('%Y-%m-%d')

        # Print the result
        print("Latest weekend for DDD:", latest_weekend_ddd)

        print("latest weekend we use",latest_weekend)
        # print("latest weekend for ddd",pd.to_datetime(ddd_sales_data['week_end_date']).max().date().strftime('%Y-%m-%d'))

        latest_day_mon= pd.to_datetime(latest_weekend).strftime('%m/%d')
        # latest_day_mon= pd.to_datetime(latest_weekend).strftime('%m/%d')

        self.latest_weekend = pd.to_datetime(self.latest_weekend)
        self.latest_day_mon = self.latest_weekend.strftime('%m/%d')
        self.latest_year = self.latest_weekend.year
        self.latest_month = self.latest_weekend.month

        latest_year = pd.to_datetime(latest_weekend).year
        latest_month = pd.to_datetime(latest_weekend).month
        latest_month, latest_year


        estimate_dates = [
            "2025-01-03", "2025-01-10", "2025-01-17", "2025-01-24", "2025-01-31", "2025-02-07", "2025-02-14",
            "2025-02-21", "2025-02-28", "2025-03-07", "2025-03-14", "2025-03-21", "2025-03-28", "2025-04-04",
            "2025-04-11", "2025-04-18", "2025-04-25", "2025-05-02", "2025-05-09", "2025-05-16", "2025-05-23",
            "2025-05-30", "2025-06-06", "2025-06-13", "2025-06-20", "2025-06-27", "2025-07-04", "2025-07-11",
            "2025-07-18", "2025-07-25", "2025-08-01", "2025-08-08", "2025-08-15", "2025-08-22", "2025-08-29",
            "2025-09-05", "2025-09-12", "2025-09-19", "2025-09-26", "2025-10-03", "2025-10-10", "2025-10-17",
            "2025-10-24", "2025-10-31", "2025-11-07", "2025-11-14", "2025-11-21", "2025-11-28", "2025-12-05",
            "2025-12-12", "2025-12-19", "2025-12-26"
        ]
        print("len!!!", len(estimate_dates))

        estimate_values = [
            7884, 9307, 8664, 8164, 9670, 8163, 7807, 7921, 8368, 8177, 7979, 8404, 8227, 8049,
            8107, 8068, 8025, 8436, 8488, 8588, 8616, 7510, 8001, 8022, 8083, 8289, 7535, 8917,
            8574, 8485, 8573, 8418, 8430, 8410, 8483, 7913, 8654, 8408, 8449, 8455, 8095, 8138,
            8053, 8172, 8439, 8489, 8498, 7379, 8177, 7796, 7925, 8526
        ]


        slide_1_chart_estimates = pd.DataFrame({
            'formatted_date': estimate_dates,
            '24 Budget Weekly Run Rate Estimate': estimate_values,
            'full_date': estimate_dates
        })
        # slide_1_chart_estimates['full_date'] = pd.to_datetime(slide_1_chart_estimates['formatted_date'] + '-2024', format='%d-%b-%Y')


        # In[ ]:


        def get_dates(current_date, weeks_before=None, weeks_after=None):
            print("current_date!!", current_date, weeks_before, weeks_after)
            given_date = datetime.strptime(current_date, "%Y-%m-%d")

            # Ensure weeks_before and weeks_after are valid integers
            if isinstance(weeks_before, str) and weeks_before.strip().isdigit():
                weeks_before = int(weeks_before)
            if isinstance(weeks_after, str) and weeks_after.strip().isdigit():
                weeks_after = int(weeks_after)

            # Handle weeks_before and weeks_after
            if weeks_before is not None and weeks_before != "":
                days_before = weeks_before * 7
                date_before = given_date - timedelta(days=days_before)
            elif weeks_after is not None and weeks_after != "":
                days_after = weeks_after * 7
                date_before = given_date + timedelta(days=days_after)
            else:
                raise ValueError("Either weeks_before or weeks_after must be provided and should be a valid number.")

            past_date = date_before.strftime("%Y-%m-%d")
            return past_date


        def get_quarter_start_date(latest_weekend):
                latest_weekend_= datetime.strptime(latest_weekend, "%Y-%m-%d")
                if latest_weekend_.month in [1, 2, 3]:
                    start_date  = datetime(latest_weekend_.year , 1, 1)
                elif latest_weekend_.month in [4, 5, 6]:
                    start_date  = datetime(latest_weekend_.year , 4, 1)
                elif latest_weekend_.month in [7, 8, 9]:
                    start_date  = datetime(latest_weekend_.year , 7, 1)
                else:
                    start_date  = datetime(latest_weekend_.year , 10, 1)
                return start_date.strftime("%Y-%m-%d"), latest_weekend_.strftime("%Y-%m-%d")

        def get_last_date_before_qtr(start_date, df, date_column):
                start_date = pd.to_datetime(start_date).date().strftime('%Y-%m-%d')
                filtered_df = df[df[date_column] < start_date]
                if not filtered_df.empty:
                    last_date_before_q4 = filtered_df[date_column].max()
                    return last_date_before_q4.strftime("%Y-%m-%d")
                else:
                    return None


        def get_date_range(base_date_str, weeks, date_format='%Y-%m-%d'):
            # Convert base_date_str to datetime object
            base_date = datetime.strptime(base_date_str, date_format)
            # Calculate start and end dates with time information
            end_date = base_date.replace(tzinfo=None) #remove tzinfo
            start_date = end_date - timedelta(weeks=weeks - 1)

            # Format start and end dates as strings with time information
            return start_date.strftime("%Y-%m-%d"), end_date.strftime("%Y-%m-%d")
        def get_previous_quarter_dates(quarter_start):
            quarter_start = pd.Timestamp(quarter_start)
            prev_quarters = {
                1: (pd.Timestamp(quarter_start.year - 1, 10, 1), pd.Timestamp(quarter_start.year - 1, 12, 31)),
                2: (pd.Timestamp(quarter_start.year, 1, 1), pd.Timestamp(quarter_start.year, 3, 31)),
                3: (pd.Timestamp(quarter_start.year, 4, 1), pd.Timestamp(quarter_start.year, 6, 30)),
                4: (pd.Timestamp(quarter_start.year, 7, 1), pd.Timestamp(quarter_start.year, 9, 30))
            }
            current_quarter = (quarter_start.month - 1) // 3 + 1
            prev_quarter_start, prev_quarter_end = prev_quarters[current_quarter]
            return prev_quarter_start.strftime("%Y-%m-%d"), prev_quarter_end.strftime("%Y-%m-%d")

        def sum_in_thousands(df, start_date, end_date):
            df['full_date'] = pd.to_datetime(df['full_date'])
            filtered_df = df[(df['full_date'] >= start_date) & (df['full_date'] <= end_date)]
            sum_values_in_thousands = filtered_df['24 Budget Weekly Run Rate Estimate'].sum() / 1000
            return sum_values_in_thousands

        def get_diff_in_weeks(latest_weekend_, start_date):
            difference_in_days = (pd.to_datetime(latest_weekend_) - pd.to_datetime(start_date)).days
            difference_in_weeks = difference_in_days / 7
            weeks_from_quarter_start = round(difference_in_weeks)
            return weeks_from_quarter_start

        drug_list=['gemzar', 'halaven', 'vinorelbine', 'xeloda','enhertu']

        p_start_date, p_end_date = get_date_range(get_dates(latest_weekend_ddd, 13), 13)

        sql_query = f'''SELECT *FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`ddd_data_for_genie_filtered_dates`
                                WHERE LOWER(drug_name) IN ({",".join([f"'{drug.lower()}'" for drug in drug_list])}) AND week_end_date >= '{p_start_date}'
                               '''
        ddd_sales_data = pd.DataFrame(self.cursor.execute(sql_query).fetchall(), columns=[elem[0] for elem in self.cursor.description])
        ddd_sales_data['week_end_date'] = pd.to_datetime(ddd_sales_data['week_end_date'], errors='coerce')
        ddd_sales_data = ddd_sales_data.rename(columns={'trx_cnt': 'qty_sold_pu'})
        print("ddd sales data sample",ddd_sales_data.head())

        # def calculate_growth(recent_sales, past_sales):
        #     # Avoid division by zero using np.where
        #     recent_sales = float(recent_sales)
        #     past_sales = float(past_sales)
        #     growth_percentage = np.where(past_sales != 0, ((recent_sales - past_sales) / past_sales) * 100, 0)
        #     return growth_percentage

        def calculate_growth(recent_sales, past_sales):
            # Avoid division by zero using simple conditional
            recent_sales = float(recent_sales)
            past_sales = float(past_sales)
            try:
                print("recent_sales!!!", recent_sales, "past_sales!!!", past_sales)
                print("recent_sales - past_sales!!!", recent_sales - past_sales)

                # Simple conditional to avoid division by zero
                if past_sales == 0:
                    growth_percentage = 0
                else:
                    growth_percentage = ((recent_sales - past_sales) / past_sales) * 100
                
                print("growth_percentage!!!", growth_percentage)
                return growth_percentage
            except Exception as e:
                print("error in calculate_growth", e)
                return 0




        def get_slide_1_data():
            # Validate sales_data
            print("latest weekend in get slide1+data",latest_weekend)
            if 'week_end_date' not in tro_sales_data.columns:
                raise KeyError("'week_end_date' is missing in sales_data!")
            tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date'], errors='coerce')

            # Filter data for the relevant date range
            filtered_data = tro_sales_data[
                (tro_sales_data['week_end_date'] >= '2024-01-05') &
                (tro_sales_data['week_end_date'] <= latest_weekend)
            ]
            if filtered_data.empty:
                raise ValueError("Filtered data is empty. Check the date range or data.")

            # Group by week and compute sales
            weekly_sales = filtered_data.groupby('week_end_date')['qty_sold_pu'].sum().reset_index()
            weekly_sales['formatted_date'] = weekly_sales['week_end_date'].dt.strftime('%d-%b')

            # Calculate R4W averages
            r4w_averages = []
            for current_date in filtered_data['week_end_date'].unique():
                if pd.isnull(current_date):
                    print(f"Skipping invalid date: {current_date}")
                    continue
                start_date = pd.to_datetime(current_date) - timedelta(weeks=3)
                end_date = pd.to_datetime(current_date)
                data_range = tro_sales_data[
                    (tro_sales_data['week_end_date'] >= start_date) &
                    (tro_sales_data['week_end_date'] <= end_date)
                ]
                if data_range.empty:
                    print(f"No data in range {start_date} to {end_date}")
                    continue
                average_sales = data_range['qty_sold_pu'].sum() / 4
                r4w_averages.append({'week_end_date': current_date, 'R4W Average': average_sales})

            # Validate R4W averages
            if not r4w_averages:
                raise ValueError("No data for R4W averages. Check your input data or logic.")
            r4w_averages_df = pd.DataFrame(r4w_averages)
            if 'week_end_date' not in r4w_averages_df.columns:
                raise KeyError("The 'week_end_date' column is missing from r4w_averages_df!")

            # Sort and process data
            r4w_averages_df['week_end_date'] = pd.to_datetime(r4w_averages_df['week_end_date'])
            r4w_averages_df_sorted = r4w_averages_df.sort_values(by='week_end_date')
            r4w_averages_df_sorted['formatted_date'] = r4w_averages_df_sorted['week_end_date'].dt.strftime('%d-%b')

            # Merge weekly sales with R4W averages
            result_df = pd.merge(weekly_sales, r4w_averages_df_sorted, on='formatted_date', how='inner')
            slide_1_chart_estimates['formatted_date'] = pd.to_datetime(slide_1_chart_estimates['formatted_date']).dt.strftime('%d-%b')
            result_df_ = pd.merge(slide_1_chart_estimates, result_df, on='formatted_date', how='left')

            # Final DataFrame
            final_df = result_df_[['formatted_date', 'qty_sold_pu', '24 Budget Weekly Run Rate Estimate', 'R4W Average']]
            final_df.rename(columns={'qty_sold_pu': 'Actuals'}, inplace=True)
            print("slide 1 final df is here",final_df.head())
            return final_df




        def get_slide_1_table_data():
            def get_weekly_average(start_date_= '', end_date_= '', weeks=''):
                if end_date_:
                    end_date = datetime.strptime(end_date_, '%Y-%m-%d').replace(tzinfo=pytz.UTC)
                    start_date = end_date - timedelta(weeks=weeks - 1)
                elif start_date_:
                    start_date = datetime.strptime(start_date_, '%Y-%m-%d').replace(tzinfo=pytz.UTC)
                    end_date = start_date + timedelta(weeks=weeks)
                print("start_date!!!!!", start_date, "end_date!!", end_date)
                start_date = start_date.replace(tzinfo=None)
                end_date = end_date.replace(tzinfo=None)
                mask = (tro_sales_data['week_end_date'] >= start_date) & (tro_sales_data['week_end_date'] <= end_date)
                recent_weeks_data = tro_sales_data.loc[mask, 'qty_sold_pu']
                total_sales = recent_weeks_data.sum()
                if weeks==0:
                    weekly_average=0
                else:
                    weekly_average = round((total_sales / weeks) / 1000,2)
                return weekly_average
                




            def calculate_total_vials(sales_data, start_date, end_date):
                filter_condition = (sales_data['week_end_date'] >= start_date) & (sales_data['week_end_date'] <= end_date)
                return sales_data.loc[filter_condition, 'qty_sold_pu'].sum() / 1000

            recent_avg_4w = get_weekly_average('', latest_weekend, 4)
            past_avg_4w = get_weekly_average('', get_dates(latest_weekend , 4,''), 4)

            recent_avg_8w = get_weekly_average('', latest_weekend, 8)
            past_avg_8w = get_weekly_average('', get_dates(latest_weekend , 8,''), 8)

            recent_avg_13w = get_weekly_average('', latest_weekend, 13)
            past_avg_13w = get_weekly_average('', get_dates(latest_weekend , 13,''), 13)

            recent_avg_52w = get_weekly_average('', latest_weekend, 52)
            past_avg_52w = get_weekly_average('', get_dates(latest_weekend , 52,''), 52)

            start_date, latest_weekend_ = get_quarter_start_date(latest_weekend)
            weeks_from_quarter_start = get_diff_in_weeks(latest_weekend_, start_date)

            recent_qtr = get_weekly_average('', latest_weekend, weeks_from_quarter_start)

            past_qtr_end_date = get_last_date_before_qtr(start_date, tro_sales_data, 'week_end_date')
            past_qtr_start_date, past_qtr_end_date_ = get_quarter_start_date(past_qtr_end_date)
            print("past_qtr_start_date!!!", past_qtr_start_date)
            past_qtr_end_date_new = get_dates(past_qtr_start_date , '',weeks_from_quarter_start)
            print("past_qtr_end_date_new!!", past_qtr_end_date_new, "weeks_from_quarter_start!!", weeks_from_quarter_start)
            past_qtr = get_weekly_average(past_qtr_start_date,'', weeks_from_quarter_start)

            growth_4w = calculate_growth(recent_avg_4w, past_avg_4w)
            print("growth_4w", growth_4w)
            growth_8w = calculate_growth(recent_avg_8w, past_avg_8w)
            print("growth_8w", growth_8w)
            growth_13w = calculate_growth(recent_avg_13w, past_avg_13w)
            growth_52w = calculate_growth(recent_avg_52w, past_avg_52w)
            print("recent_qtr!!!", recent_qtr, "past_qtr!!!", past_qtr)
            growth_qtr = calculate_growth(recent_qtr, past_qtr)
            # growth_qtr=growth_qtr[0]

            current_qtr_start, current_qtr_end = get_quarter_start_date(latest_weekend)
            prev_quarter_start, prev_quarter_end = get_previous_quarter_dates(current_qtr_start)
            ytd_est_avg = sum_in_thousands(slide_1_chart_estimates, f'{latest_year}-01-01', latest_weekend)
            current_qtr_est_avg = sum_in_thousands(slide_1_chart_estimates, current_qtr_start, current_qtr_end)
            prev_qtr_est_avg = sum_in_thousands(slide_1_chart_estimates, prev_quarter_start, prev_quarter_end)

            ytd_vials = calculate_total_vials(tro_sales_data, f'{latest_year}-01-01', latest_weekend)
            ytd_growth = calculate_growth(ytd_vials,  ytd_est_avg)
            prev_qtr_vials = calculate_total_vials(tro_sales_data, prev_quarter_start, prev_quarter_end)
            prev_qtr_growth = calculate_growth(prev_qtr_vials, prev_qtr_est_avg)

            current_qtr_vials = calculate_total_vials(tro_sales_data, current_qtr_start, current_qtr_end)
            curr_qtr_growth = calculate_growth(current_qtr_vials, current_qtr_est_avg)
            

            Q1_avg_2024 = get_weekly_average('', f'2024-03-29', 13)
            # def get_max_date_for_quarter(dates, start_month, end_month):
            #     quarter_dates = [date for date in dates if start_month <= date.month <= end_month and date.year == latest_year]
            #     print(f"Filtered dates for months {start_month}-{end_month}: {quarter_dates}")
            #     return max(quarter_dates) if quarter_dates else None
            def get_max_date_for_quarter(dates, start_month, end_month):
                # Convert 'dates' elements to datetime objects
                dates_dt = [pd.to_datetime(d) for d in dates['week_end_date']]
                quarter_dates = [date for date in dates_dt if start_month <= date.month <= end_month and date.year == latest_year]
                print(f"Filtered dates for months {start_month}-{end_month}: {quarter_dates}")
                return max(quarter_dates) if quarter_dates else None
            max_dates = {
            "Q1": get_max_date_for_quarter(sorted_dates, 1, 3),
            "Q2": get_max_date_for_quarter(sorted_dates, 4, 6),
            "Q3": get_max_date_for_quarter(sorted_dates, 7, 9),
            "Q4": get_max_date_for_quarter(sorted_dates, 10, 12),
            }

            if max_dates["Q2"] is None:
                print("No data available for Q2. Skipping calculation.")
                Q2_avg_2024 = 0
            else:
                Q2_avg_2024 = get_weekly_average('', max_dates["Q2"].strftime('%Y-%m-%d'), 13)


            final_output_dict = {
                "R4W": recent_avg_4w , "R8W": recent_avg_8w, "R13W": recent_avg_13w , "R52W": recent_avg_52w,
                "QTD": recent_qtr , "growth_4w": growth_4w , "growth_8w": growth_8w ,
                "growth_13w": growth_13w ,"growth_52w": growth_52w ,"growth_qtr": growth_qtr , "current_qtr_vials": current_qtr_vials,
            "prev_qtr_vials": prev_qtr_vials, "ytd_vials": ytd_vials, "curr_qtr_growth": curr_qtr_growth,
                "prev_qtr_growth": prev_qtr_growth, "ytd_growth": ytd_growth, "q1": Q1_avg_2024, "q2": Q2_avg_2024,
                "prev_qtr_est_avg": prev_qtr_est_avg, "ytd_est_avg": ytd_est_avg, "current_qtr_est_avg": current_qtr_est_avg
            }
            def round_value(val, decimals=2):
                if isinstance(val, pd.Series):
                    # Round the pandas Series
                    return val.round(decimals)
                elif isinstance(val, Decimal):
                    # Round Decimal (convert to float for rounding then back to Decimal if needed)
                    return Decimal(round(float(val), decimals))
                elif isinstance(val, (float, np.floating)):
                    # Round float or numpy float
                    return round(val, decimals)
                elif isinstance(val, np.ndarray):
                    # Check if single element array
                    if val.size == 1:
                        scalar = val.item()
                        return round(scalar, decimals)
                    else:
                        # For arrays with multiple elements, round all elements
                        return np.round(val, decimals)
                else:
                    # Leave other types unchanged
                    print("here in else")
                    return val

            # Keys you want to round
            keys_to_round = ["growth_4w", 'growth_qtr','curr_qtr_growth', 'prev_qtr_growth', 'ytd_growth',]

            for key in keys_to_round:
                if key in final_output_dict:
                     final_output_dict[key] = round_value(final_output_dict[key], 2)
                     print("here is the key and value",key,final_output_dict[key])

            return final_output_dict



        # In[ ]:


        slide_1_data = get_slide_1_data()
        slide_1_table_data = get_slide_1_table_data()
        slide_1_table_data["total_demand"] = f'''Total Demand Vials \n (Data through {latest_day_mon})'''
        print("slide 1 table data is here",slide_1_table_data)
        slide_1_table_data["avg_demand"] = f'''Weekly Avg Demand Vials \n (Data through {latest_day_mon})'''
        slide_1_table_data["qtr_avg_demand"] = f'''Weekly Avg Demand Vials \n (Actuals through {latest_day_mon})'''

        slide_1_table_data


        # In[ ]:


        def get_slide_2_data():
                def get_averages(filtered_data):
                    filtered_data['week_end_date'] = pd.to_datetime(filtered_data['week_end_date']).dt.tz_localize(None)

                    specified_date = pd.Timestamp(latest_weekend).tz_localize(None)
                    start_date_r13w = specified_date - timedelta(weeks=12)

                    end_date_p13w = start_date_r13w - timedelta(days=1)
                    start_date_p13w = end_date_p13w - timedelta(weeks=13)

                    end_date_pp13w = start_date_p13w - timedelta(days=1)
                    start_date_pp13w = end_date_pp13w - timedelta(weeks=13)

                    latest_week_qtr_start_dt, latest_weekend_  = get_quarter_start_date(latest_weekend)
                    start_date_ytd = pd.Timestamp(latest_week_qtr_start_dt).tz_localize(None)
                    end_date_ytd = pd.Timestamp(latest_weekend_).tz_localize(None)

                    r13w_data = filtered_data[(filtered_data['week_end_date'] >= start_date_r13w) & (filtered_data['week_end_date'] <= specified_date)]
                    r13w_total_qty = r13w_data['qty_sold_pu'].sum()
                    r13w_weekly_avg = r13w_total_qty / 13

                    p13w_data = filtered_data[(filtered_data['week_end_date'] >= start_date_p13w) & (filtered_data['week_end_date'] <= end_date_p13w)]
                    p13w_total_qty = p13w_data['qty_sold_pu'].sum()
                    p13w_weekly_avg = p13w_total_qty / 13
                    print("p13w_total_qty!!!", p13w_total_qty)

                    pp13w_data = filtered_data[(filtered_data['week_end_date'] >= start_date_pp13w) & (filtered_data['week_end_date'] <= end_date_pp13w)]
                    pp13w_total_qty = pp13w_data['qty_sold_pu'].sum()
                    pp13w_weekly_avg = pp13w_total_qty / 13

                    p13_growth = calculate_growth(p13w_weekly_avg, pp13w_weekly_avg)
                    r13_growth = calculate_growth(r13w_weekly_avg, p13w_weekly_avg)



                    ytd_data = filtered_data[(filtered_data['week_end_date'] >= start_date_ytd) & (filtered_data['week_end_date'] <= end_date_ytd)
            #                                    &(filtered_data['parent_name'] != "TEXAS ONCOLOGY")
                                            ]
                    ytd_ex_txs = filtered_data[(filtered_data['week_end_date'] >= start_date_ytd) & (filtered_data['week_end_date'] <= end_date_ytd)
                                            &(filtered_data['parent_name'] != "TEXAS ONCOLOGY")
                                            ]

                    ytd_total_qty = ytd_data['qty_sold_pu'].sum()
                    ytd_total_ex_txs = ytd_ex_txs['qty_sold_pu'].sum()

                    num_weeks_ytd = (end_date_ytd - start_date_ytd).days // 7 + 1
                    ytd_weekly_avg = ytd_total_qty / num_weeks_ytd
                    ytd_avg_ex_txs = ytd_total_ex_txs / num_weeks_ytd
                    return [r13w_weekly_avg, p13w_weekly_avg, pp13w_weekly_avg, ytd_avg_ex_txs, p13_growth, r13_growth]

                account_types = ['high','low', 'medium', 'academic', 'community','east','west','None']
                final_output_dict = {}
                final_growth_dict = {}
                for type_ in account_types:
                    print("type_!!!!!!", type_)
                    averages_dict = {}
                    growth_dict = {}
                    if type_ in ['academic', 'community']:
                        filtered_data = tro_sales_data[tro_sales_data['child_account_type'].str.lower() == type_]
                        filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize(None)
                    elif type_ in ['east','west']:
                        filtered_data = tro_sales_data[tro_sales_data['area_nm'].str.lower() == type_]
                        filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize(None)
                    elif type_ =='None':
                        filtered_data=tro_sales_data
                        filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize(None)
                    else:
                        filtered_data = tro_sales_data[tro_sales_data['linkedid_bc_segment'].str.lower() == type_]
                        filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize(None)
                    r13w, p13w, pp13w, ytd_ex_txs, p13_growth, r13_growth = get_averages(filtered_data)
                    averages_dict["PP13W"] = pp13w
                    averages_dict['P13W'] = p13w
                    averages_dict['R13W'] = r13w

                    growth_dict['P13W_growth'] = p13_growth
                    growth_dict['R13W_growth'] = r13_growth

                    # averages_dict["YTD-TX"] = ytd_ex_txs
                    df = pd.DataFrame(list(averages_dict.items()), columns=['time_period', type_])
                    final_output_dict[type_] = df

                    growth_df = pd.DataFrame(list(growth_dict.items()), columns=['time_period', "growth"])
                    final_growth_dict[type_] = growth_df
                return final_output_dict, final_growth_dict


            # In[ ]:


        slide_2_data,  slide_2_growth_data= get_slide_2_data()


        # In[ ]:


        def get_slide_3_data():
            def filter_data_by_date_range(sales_data, start_date, end_date):
                filtered_data = sales_data[(sales_data['week_end_date'] >= start_date) & (sales_data['week_end_date'] <= end_date)]
                return filtered_data

            def calculate_sales_metrics(filtered_data, col_name):
                total_sales_by_parent = filtered_data.groupby('parent_name')['qty_sold_pu'].sum().reset_index()
                total_weeks = filtered_data['week_end_date'].nunique()
                total_sales_by_parent['average_weekly_sales_'+ col_name] = total_sales_by_parent['qty_sold_pu'] / total_weeks
                total_sales = total_sales_by_parent['qty_sold_pu'].sum()
                total_sales_by_parent['sales_share_'+ col_name] = (total_sales_by_parent['qty_sold_pu'] / total_sales) * 100
                return total_sales_by_parent

            def calculate_growth_percentage(df, old_col, new_col, df_col):
              # Avoid division by zero by replacing zeros in the denominator with a small value
              # df[df_col] = (((df[new_col] - df[old_col]) * 100) / df[old_col].replace(0, 1e-10))
              # Convert 'old_col' to numeric, handling errors
              df[old_col] = pd.to_numeric(df[old_col], errors='coerce')
              # Convert 'new_col' to numeric, handling errors
              df[new_col] = pd.to_numeric(df[new_col], errors='coerce')
              # Avoid division by zero by replacing zeros in the denominator with a small value
              df[df_col] = (((df[new_col] - df[old_col]) * 100) / df[old_col].replace(0, 1e-10))
              return df

            def process_data(sales_data, start_date, end_date, col_name):
                filtered_data = filter_data_by_date_range(sales_data, start_date, end_date )
                return calculate_sales_metrics(filtered_data, col_name)

            def get_quarter_data(sales_data, start_date, end_date, column_name):
                quarter_data = filter_data_by_date_range(sales_data, start_date, end_date)
                quarter_sales = quarter_data.groupby('parent_name')['qty_sold_pu'].sum().reset_index()
                quarter_sales.rename(columns={'qty_sold_pu': column_name}, inplace=True)
                return quarter_sales

            def calculate_summary_row(df, label):
                sums = df.drop(columns=['parent_name']).sum()
                new_row_data = {'parent_name': label, **sums}
                return pd.DataFrame([new_row_data])

            def get_extra_rows(row_name, df):
                print("columns!!!!", df.columns)
                sum_columns = ['qty_sold_pu_current', 'average_weekly_sales_current','sales_share_current', 'qty_sold_pu_prev', 'average_weekly_sales_prev',
                            'sales_share_prev', 'qty_sold_pu_P13w', 'qty_sold_pu_R13w','average_weekly_sales_R13W_current', 'sales_share_R13W_current','average_weekly_sales_P13W_prev',
                            'qty_sold_pu_P_ytd', 'average_weekly_sales_P_qtd_current','sales_share_P13W_prev','sales_share_P_qtd_current', 'qty_sold_pu_R_ytd','average_weekly_sales_R_qtd_current', 'sales_share_R_qtd_current']
                growth_pairs = {"R12Mgrowth": ('qty_sold_pu_prev', 'qty_sold_pu_current'),
                            'R13Wgrowth': ('qty_sold_pu_P13w', 'qty_sold_pu_R13w'), "YTD_growth": ('qty_sold_pu_P_ytd', 'qty_sold_pu_R_ytd')
                            }
                new_row = {}
                for col in sum_columns:
                    new_row[col] = df[col].sum()
                for key, (from_col, to_col) in growth_pairs.items():
                    if df[from_col].sum() != 0:
                        new_row[key] = ((df[to_col].sum() / df[from_col].sum()) - 1) * 100
                    else:
                        new_row[key] = None
                new_row["parent_name"] = row_name
                new_df_ls.append(new_row)


            current_start_date, current_end_date = get_date_range(latest_weekend, 52)
            print( current_start_date, current_end_date)
            current_sales_by_parent = process_data(tro_sales_data, current_start_date, current_end_date, "current")
            current_sales_by_parent.rename(columns={'qty_sold_pu': 'qty_sold_pu_current'}, inplace=True)

            print("previous end date!!" , get_dates(latest_weekend , 52))
            prev_start_date, prev_end_date = get_date_range(get_dates(latest_weekend , 52), 52)
            prev_sales_by_parent = process_data(tro_sales_data, prev_start_date, prev_end_date, "prev")
            prev_sales_by_parent.rename(columns={'qty_sold_pu': 'qty_sold_pu_prev'}, inplace=True)
            print( "prev_start_date, prev_end_date!!",prev_start_date, prev_end_date)


            merged_sales = pd.merge(current_sales_by_parent, prev_sales_by_parent, on='parent_name', how= "outer").fillna(0)
            merged_sales.replace([float('inf'), -float('inf')], 0, inplace=True)
            top_10_parent_accounts = merged_sales.sort_values(by='qty_sold_pu_current', ascending=False)
            merged_sales_data = calculate_growth_percentage(top_10_parent_accounts, 'qty_sold_pu_prev', 'qty_sold_pu_current', "R12Mgrowth")

            latest_week_qtr_start_dt, latest_weekend_  = get_quarter_start_date(latest_weekend)
            R13_start_date_, R13_end_date_ = get_date_range(latest_weekend, 13)
            R13_sales_by_parent_ = process_data(tro_sales_data, R13_start_date_, R13_end_date_, "R13W_current")
            R13_sales_by_parent_.rename(columns={'qty_sold_pu': 'qty_sold_pu_R13w'}, inplace=True)

            past_13w_end_date = get_dates(latest_weekend , 13)
            print("past_13w_end_date!!!!!", past_13w_end_date)
            P13_start_date_, P13_end_date_ = get_date_range(past_13w_end_date, 13)
            P13_sales_by_parent_ = process_data(tro_sales_data, P13_start_date_, P13_end_date_, "P13W_prev")
            P13_sales_by_parent_.rename(columns={'qty_sold_pu': 'qty_sold_pu_P13w'}, inplace=True)

            w_13_merged_sales = pd.merge(P13_sales_by_parent_, R13_sales_by_parent_, on='parent_name', how= "outer").fillna(0)
            w_13_growth = calculate_growth_percentage(w_13_merged_sales, 'qty_sold_pu_P13w', 'qty_sold_pu_R13w', "R13Wgrowth")
            sub_df = pd.merge(merged_sales_data, w_13_merged_sales, on='parent_name', how= "outer").fillna(0)

            weeks_diff = get_diff_in_weeks( latest_weekend_, f'{latest_year}-01-01')
            print("weeks_diff!!!", weeks_diff)
            P_ytd_start_date = f'{latest_year-1}-01-01'
            previous_year = datetime.strptime(P_ytd_start_date, '%Y-%m-%d')
            date_after_weeks_diff = previous_year + timedelta(weeks= weeks_diff)
            P_ytd_end_date = date_after_weeks_diff.strftime('%Y-%m-%d')
            print("P_ytd_start_date, P_ytd_end_date!!!!!", P_ytd_start_date, P_ytd_end_date)

            P_ytd_sales= process_data(tro_sales_data, P_ytd_start_date, P_ytd_end_date, "P_qtd_current")
            print("P_ytd_sales!!!", P_ytd_sales)
            P_ytd_sales.rename(columns={'qty_sold_pu': 'qty_sold_pu_P_ytd'}, inplace=True)


            r_ytd_start_date, r_ytd_end_date = get_date_range(latest_weekend_, weeks_diff)
            R_ytd_sales= process_data(tro_sales_data, r_ytd_start_date, r_ytd_end_date, "R_qtd_current")
            R_ytd_sales.rename(columns={'qty_sold_pu': 'qty_sold_pu_R_ytd'}, inplace=True)

            Ytd_merged_sales = pd.merge(P_ytd_sales, R_ytd_sales, on='parent_name', how= "outer").fillna(0)
            Ytd_growth = calculate_growth_percentage(Ytd_merged_sales, 'qty_sold_pu_P_ytd', 'qty_sold_pu_R_ytd', "YTD_growth")

            final_df = pd.merge(sub_df, Ytd_growth, on='parent_name', how= "outer").fillna(0)
            final_df.replace([float('inf'), -float('inf')], 0, inplace=True)
            final_df = final_df.sort_values(by='qty_sold_pu_current', ascending=False)
            final_df = final_df.fillna(0)
            un_fil_top_15 = final_df.head(15)
            un_fil_top_10 = final_df.head(10)

            filtered_df = final_df[final_df['parent_name'].str.upper() != 'TEXAS ONCOLOGY']
            sorted_df = filtered_df.sort_values(by='qty_sold_pu_current', ascending=False)
            top_10 = sorted_df.head(10)
            top_15 = sorted_df.head(15)
            new_df_ls = []
            df_row = {'Top 10 - TOPA': top_10, 'Top 15 - TOPA': top_15, "Top 15 Parent Accounts":  un_fil_top_15,"Top 10 Parent Accounts":  un_fil_top_10,  "Nation - TOPA": sorted_df, "Nation": final_df}
        #     print("df_row!!", df_row)
            for key, value in df_row.items():
                get_extra_rows(key, value)

            calc_df =pd.DataFrame(new_df_ls)
            slide_df = pd.concat([final_df,calc_df], ignore_index=True)
            final_list = ['parent_name', 'sales_share_current', 'average_weekly_sales_current','R12Mgrowth',
            'average_weekly_sales_R13W_current', 'R13Wgrowth','YTD_growth']
            percent_col = ["sales_share_current", "R12Mgrowth", "R13Wgrowth", "YTD_growth"]
            new_final_df = pd.concat([slide_df.head(15), slide_df.tail(6) ])
            new_final_df = new_final_df[final_list]
            new_final_df.fillna(0, inplace=True)

            for column in new_final_df.columns:
                if pd.api.types.is_numeric_dtype(new_final_df[column]):
                    if column in percent_col:
                        new_final_df[column] = new_final_df[column].round(0).astype(int).astype(str) + '%'
                    else:
                        print("column!!!", column)
                        new_final_df[column] = new_final_df[column].round(0).astype(int).astype(str)

            columns_to_round = ['sales_share_current', 'average_weekly_sales_current', 'average_weekly_sales_R13W_current']

            for col in columns_to_round:
                new_final_df[col] = pd.to_numeric(new_final_df[col], errors='coerce')
            new_final_df[columns_to_round] = new_final_df[columns_to_round].applymap(np.ceil)
            return new_final_df






        slide_3_table_data = get_slide_3_data()


        # In[ ]:
        print("Before 4th Slide")
        def get_slide_4_5_data(account_type):
            # tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date'], errors='coerce')
            R13_start_date_, R13_end_date_ = get_date_range(latest_weekend, 13)
            print("R13_start_date_, R13_end_date_!!!", R13_start_date_, R13_end_date_)
            past_13w_end_date = get_dates(latest_weekend , 13)
            print("past_13w_end_date!!!!!", past_13w_end_date)
            P13_start_date_, P13_end_date_ = get_date_range(past_13w_end_date, 13)
        #     print("P13_start_date_, P13_end_date_, f'{latest_year}-01-01', latest_weekend!!!", P13_start_date_, P13_end_date_, f'{latest_year}-01-01', latest_weekend)
            account_mapping = {
                'academic': ('top_parent_account_type', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
                'community': ('top_parent_account_type', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
                'idn/hospital': ('top_parent_account_type', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
                'high': ('parent_bc_segment',  R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
                'medium': ('parent_bc_segment',  R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_)
            }
            column, start_date1, end_date1, start_date2, end_date2 = account_mapping[account_type]
            if tro_sales_data['week_end_date'].dt.tz is None:
                tro_sales_data['week_end_date'] = tro_sales_data['week_end_date'].dt.tz_localize('UTC')
            else:
                tro_sales_data['week_end_date'] = tro_sales_data['week_end_date'].dt.tz_convert('UTC')


            filtered_data = tro_sales_data[tro_sales_data[column].str.lower() == account_type]
            print("Filtered data in slide 4 and 5 ",filtered_data)
            # filtered_data = tro_sales_data[tro_sales_data[column].str.lower() == account_type]
            R13W = (pd.Timestamp(start_date1, tz='UTC'), pd.Timestamp(end_date1, tz='UTC'))
            Q124 = (pd.Timestamp(start_date2, tz='UTC'), pd.Timestamp(end_date2, tz='UTC'))

            print("Reached Before Calculate Average Sales")

            # Convert 'week_end_date' to datetime, if not already in datetime format
            filtered_data['week_end_date'] = pd.to_datetime(filtered_data['week_end_date'], errors='coerce')

            # Check if the datetime is timezone-aware
            if filtered_data['week_end_date'].dt.tz is None:
                filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize('UTC')
            else:
                filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_convert('UTC')


            def calculate_avg_sales(filtered_data, date_range):
                # Ensure 'week_end_date' is in datetime format
                filtered_data['week_end_date'] = pd.to_datetime(filtered_data['week_end_date'], errors='coerce')

                # Check and localize 'week_end_date' to UTC if it's not already localized
                if filtered_data['week_end_date'].dt.tz is None:
                    filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize('UTC')
                else:
                    filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_convert('UTC')

                # Filter data based on the range
                date_filtered_data = filtered_data[(filtered_data['week_end_date'] >= date_range[0]) & (filtered_data['week_end_date'] <= date_range[1])]

                # Calculate average sales
                avg_sales = date_filtered_data.groupby('parent_name')['qty_sold_pu'].sum() / date_filtered_data['week_end_date'].nunique()
                print("here is the avg sales",avg_sales.head())
                return avg_sales.reset_index().rename(columns={'qty_sold_pu': 'avg_sales'})

            
            # Assuming `filtered_data` is already defined
            R13w_avg_sales = calculate_avg_sales(filtered_data, R13W).rename(columns={'avg_sales': 'R13W_avg_sales'})
            Q124_TD_avg_sales = calculate_avg_sales(filtered_data, Q124).rename(columns={'avg_sales': 'P13W_avg_sales'})
            merged_data = pd.merge(R13w_avg_sales, Q124_TD_avg_sales, on='parent_name', how='outer')
            merged_data['absolute_volume'] = (merged_data['P13W_avg_sales'] - merged_data['R13W_avg_sales']).abs()
            merged_data['R13W_avg_sales'] = merged_data['R13W_avg_sales'].fillna(0)  # Replace NaN with 0 or another suitable value
            merged_data['P13W_avg_sales'] = merged_data['P13W_avg_sales'].fillna(0)
            # Round up the sales data using np.ceil
            merged_data['R13W_avg_sales'] = np.ceil(merged_data['R13W_avg_sales']).astype(float)
            merged_data['P13W_avg_sales'] = np.ceil(merged_data['P13W_avg_sales']).astype(float)
            # Fixing the rounding issue by converting columns to float before rounding

            # Calculate growth and round it up
            merged_data['sales_growth_percentage'] = ((merged_data['P13W_avg_sales'] - merged_data['R13W_avg_sales']) / merged_data['R13W_avg_sales'].replace(0, np.inf)).fillna(0) * 100
            merged_data['sales_growth_percentage'] = np.ceil(merged_data['sales_growth_percentage']).astype(int).astype(str) + '%'
            # merged_data['sales_growth_percentage'] = round(((merged_data['P13W_avg_sales'].astype(float) - merged_data['R13W_avg_sales'].astype(float)) / merged_data['R13W_avg_sales'].replace(0, np.inf).astype(float)) * 100, 2)

            merged_data = merged_data.sort_values(by='absolute_volume', ascending=False).head(10)
            merged_data = merged_data.drop(columns=['absolute_volume'])

            percent_col = ["sales_growth_percentage"]
            for column in merged_data.columns:
                if pd.api.types.is_numeric_dtype(merged_data[column]):
                    if column in percent_col:
                        merged_data[column] = merged_data[column].round(0).astype(int).astype(str) + '%'
                    else:
                        merged_data[column] = merged_data[column].round(0).astype(int).astype(str)
            print("merged data in slide 4 and 5 ",account_type ,merged_data)
            return merged_data



        high_acc_decliners = get_slide_4_5_data('high')
        med_acc_decliners = get_slide_4_5_data('medium')
        academic_acc_decliners = get_slide_4_5_data('academic')
        community_acc_decliners = get_slide_4_5_data('community')
        idn_hospital_decliners = get_slide_4_5_data('idn/hospital')
        idn_hospital_decliners


        def get_slide_6_7_east_west_data(region):
            print("region!!!!!!!", region)
            if region:
                filter_by = 'regn_nm'
            else:
                filter_by = 'area_nm'
            def get_extra_rows(row_name, df):
                new_row = {}
                for col in df.columns[1:]:
                    new_row[col] = df[col].sum()
                new_row[filter_by] = row_name
                df2 = pd.DataFrame([new_row])
                concatenated_df = pd.concat([df, df2], ignore_index=True)
                return concatenated_df

            def analyze_sales_growth_region(recent_start, recent_end, past_start, past_end, region_name, period, account_type = ''):

                if account_type and region:
                    recent_df = tro_sales_data[(tro_sales_data['week_end_date'] >= recent_start) & (tro_sales_data['week_end_date'] <= recent_end) &(tro_sales_data['area_nm'] == region) & (tro_sales_data['top_parent_account_type'].str.lower() == account_type)]
                    past_df = tro_sales_data[(tro_sales_data['week_end_date'] >= past_start) & (tro_sales_data['week_end_date'] <= past_end) & (tro_sales_data['area_nm'] == region) & (tro_sales_data['top_parent_account_type'].str.lower() == account_type)]
                elif region:
                    recent_df = tro_sales_data[(tro_sales_data['week_end_date'] >= recent_start) & (tro_sales_data['week_end_date'] <= recent_end) & (tro_sales_data['area_nm'] == region)]
                    past_df = tro_sales_data[(tro_sales_data['week_end_date'] >= past_start) & (tro_sales_data['week_end_date'] <= past_end) &(tro_sales_data['area_nm'] == region)]
                elif account_type and region == '':
                    recent_df = tro_sales_data[(tro_sales_data['week_end_date'] >= recent_start) & (tro_sales_data['week_end_date'] <= recent_end) & (tro_sales_data['top_parent_account_type'].str.lower() == account_type)]
                    past_df = tro_sales_data[(tro_sales_data['week_end_date'] >= past_start) & (tro_sales_data['week_end_date'] <= past_end)& (tro_sales_data['top_parent_account_type'].str.lower() == account_type)]
                else:
                    recent_df = tro_sales_data[(tro_sales_data['week_end_date'] >= recent_start) & (tro_sales_data['week_end_date'] <= recent_end) ]
                    past_df = tro_sales_data[(tro_sales_data['week_end_date'] >= past_start) & (tro_sales_data['week_end_date'] <= past_end)]

                recent_sales = recent_df.groupby(filter_by)['qty_sold_pu'].sum().reset_index()
                past_sales = past_df.groupby(filter_by)['qty_sold_pu'].sum().reset_index()
                merged_df = pd.merge(recent_sales, past_sales, on=filter_by, suffixes=(f'_recent_{period}', f'_past_{period}'))
                if region_name == '':
                    print("in region name empty!!!")
                    merged_df_= get_extra_rows("Nation",merged_df)
                else:
                    merged_df_= get_extra_rows(region_name,merged_df)
                merged_df_[f'growth_percentage_{period}'] = ((merged_df_[f'qty_sold_pu_recent_{period}'] - merged_df_[f'qty_sold_pu_past_{period}']) / merged_df_[f'qty_sold_pu_past_{period}']) * 100
        #         print(merged_df_.head())
                return merged_df_

            def analyze_and_store(recent_start, recent_end, past_start, past_end, region, name, account_type=None):
                df = analyze_sales_growth_region(recent_start, recent_end, past_start, past_end, region, name, account_type)
                print("df!!!!!!!!!", df)
                return df

            results = {}
            R4w_start_date, R4w_end_date = get_date_range(latest_weekend, 4)
            P4w_start_date, P4w_end_date = get_date_range(get_dates(latest_weekend , 4), 4)
            print("R4w_start_date, R4w_end_date!!!", R4w_start_date, R4w_end_date, "    P4w_start_date, P4w_end_date!!",     P4w_start_date, P4w_end_date)

            R13w_start_date, R13w_end_date = get_date_range(latest_weekend, 13)
            P13w_start_date, P13w_end_date = get_date_range(get_dates(latest_weekend , 13), 13)
            print("R13w_start_date, R13w_end_date!!!", R13w_start_date, R13w_end_date, "    P13w_start_date, P13w_end_date!!",     P13w_start_date, P13w_end_date)

            results['R4w_p4w'] = analyze_and_store(R4w_start_date, R4w_end_date, P4w_start_date, P4w_end_date, region, 'R4w_p4w')
            results['R13W_P13W'] = analyze_and_store(R13w_start_date, R13w_end_date, P13w_start_date, P13w_end_date, region, 'R13W_P13W')
#             results['R13W_Q1TD'] = analyze_and_store('2024-01-05', '2024-03-29', '2023-09-29', '2023-12-22', region, 'R13W_Q1TD')


            account_types = ['academic', 'community', 'idn/hospital']
            for account_type in account_types:
                results[f'{account_type}_4w'] = analyze_and_store(R4w_start_date, R4w_end_date, P4w_start_date, P4w_end_date, region, f'{account_type}_4w', account_type)
                results[f'{account_type}_13w'] = analyze_and_store(R13w_start_date, R13w_end_date, P13w_start_date, P13w_end_date, region, f'{account_type}_13w', account_type)
#                 results[f'{account_type}_13w_q1td'] = analyze_and_store('2024-01-05', '2024-03-29', '2023-09-29', '2023-12-22', region, f'{account_type}_13w_q1td', account_type)

            dfs = [results['R4w_p4w'],  results['R13W_P13W'],
                results['academic_4w'], results['academic_13w'],
                results['community_4w'], results['community_13w'],
                results['idn/hospital_4w'], results['idn/hospital_13w']]
#     results['R13W_Q1TD'],results['academic_13w_q1td'],results['community_13w_q1td'], results['idn/hospital_13w_q1td']

            merged_df = reduce(lambda left, right: pd.merge(left, right, on=filter_by), dfs)
            merged_df['Volume_Contribution_Academic'] = (merged_df['qty_sold_pu_recent_academic_13w'] / (merged_df['qty_sold_pu_recent_academic_13w']+merged_df['qty_sold_pu_recent_community_13w']+merged_df['qty_sold_pu_recent_idn/hospital_13w'])) * 100
            merged_df['Volume_Contribution_Community'] = (merged_df['qty_sold_pu_recent_community_13w'] / (merged_df['qty_sold_pu_recent_academic_13w']+merged_df['qty_sold_pu_recent_community_13w']+merged_df['qty_sold_pu_recent_idn/hospital_13w'])) * 100
            merged_df['Volume_Contribution_IDN'] = (merged_df['qty_sold_pu_recent_idn/hospital_13w'] / (merged_df['qty_sold_pu_recent_academic_13w']+merged_df['qty_sold_pu_recent_community_13w']+merged_df['qty_sold_pu_recent_idn/hospital_13w'])) * 100
            merged_df[[filter_by,'Volume_Contribution_Academic', 'Volume_Contribution_Community', 'Volume_Contribution_IDN' ]]
            merged_df = merged_df[[filter_by,'growth_percentage_R4w_p4w', 'growth_percentage_R13W_P13W',
                    'growth_percentage_academic_4w', 'growth_percentage_academic_13w',
                    'growth_percentage_community_4w', 'growth_percentage_community_13w',
                    'growth_percentage_idn/hospital_4w', 'growth_percentage_idn/hospital_13w',
                    'Volume_Contribution_Academic', 'Volume_Contribution_Community', 'Volume_Contribution_IDN']]
#             'growth_percentage_R13W_Q1TD',
#             merged_df.to_excel("merged_df_east.xlsx")
#             merged_df.to_excel("merged_df_west.xlsx")

            if region == "East":
                order = ['East', 'New England', 'New York', 'Atlantic Coastal', 'Blue Grass', 'Eastern Coastal', 'Atlanta', 'Southeast']
            elif region == "West":
                order = ['West', 'Pacific Northwest', 'SoCAL', 'Rocky Mountains', 'Gulf Plains', 'Texas', 'North Central', 'Great Lakes']
            else:
                order = ["Nation", 'East', 'West']

            merged_df = merged_df.set_index(filter_by).loc[order].reset_index()
            columns_to_round = ['growth_percentage_R4w_p4w','growth_percentage_R13W_P13W','growth_percentage_academic_4w','growth_percentage_academic_13w','growth_percentage_community_4w','growth_percentage_community_13w','growth_percentage_idn/hospital_4w','growth_percentage_idn/hospital_13w','Volume_Contribution_Academic','Volume_Contribution_Community','Volume_Contribution_IDN']
            for col in columns_to_round:
                merged_df[col] = pd.to_numeric(merged_df[col], errors='coerce')

            merged_df[columns_to_round] = merged_df[columns_to_round].applymap(np.ceil)

            print("Rounded specific columns in merged df:")
            print(merged_df.head())

            region_vials_growth = merged_df[[filter_by ,'growth_percentage_R13W_P13W','growth_percentage_R4w_p4w']]
            region_vials_growth_ = region_vials_growth.rename(columns={'growth_percentage_R13W_P13W': 'R13W', 'growth_percentage_R4w_p4w': 'R4W'})
            region_vials_growth_['R13W']= region_vials_growth_['R13W']/100
            region_vials_growth_['R4W']= region_vials_growth_['R4W']/100
            print("info of merged df: ",merged_df.info())
            # merged_df_rounded = merged_df.applymap(lambda x: str(round(x)) + "%" if isinstance(x, (int, float)) else x)
            merged_df_rounded = merged_df.copy()

            # Apply the transformation on numeric columns only
            for col in merged_df_rounded.select_dtypes(include=['float64', 'int64']).columns:
                merged_df_rounded[col] = merged_df_rounded[col].apply(lambda x: str(round(x)) + "%" if isinstance(x, (int, float)) else x)

            # Display the rounded DataFrame
            print(merged_df_rounded.head())

            print("merged_df_rounded: ",merged_df_rounded.head())

            # merged_df_rounded = merged_df.applymap(lambda x: str(round(x)) + "%" if isinstance(x, (int, float)) else x)
            def combine_columns(df, col1, col2, col3, new_col_name):
                df[new_col_name] = df[col1].astype(str)+ ' | ' + df[col2].astype(str)+  ' | ' + df[col3].astype(str)
                return df

            merged_df_rounded_ = combine_columns(merged_df_rounded, 'growth_percentage_academic_4w', 'growth_percentage_community_4w', 'growth_percentage_idn/hospital_4w', "tro_4w_combined")
            merged_df_rounded_ = combine_columns(merged_df_rounded_, 'growth_percentage_academic_13w', 'growth_percentage_community_13w', 'growth_percentage_idn/hospital_13w', "tro_13w_combined")
            merged_df_rounded_ = combine_columns(merged_df_rounded_, 'Volume_Contribution_Academic', 'Volume_Contribution_Community', 'Volume_Contribution_IDN', "volume_combined")
            final_merged_df = merged_df_rounded_[[filter_by,'growth_percentage_R4w_p4w', 'growth_percentage_R13W_P13W',
                            'tro_4w_combined',  'tro_13w_combined', 'volume_combined']]
            print(final_merged_df.head())

            def change_value(row_name, column_name, df):
                value = df[df[filter_by] == row_name][column_name].values[0]
                df.loc[df[filter_by] == row_name, column_name] = "Vials Growth " + str(value)
                return df
            if region :
                change_val_dict = {"growth_percentage_R4w_p4w": region, "growth_percentage_R13W_P13W": region } #"growth_percentage_R13W_Q1TD": region}
            else:
                change_val_dict = {"growth_percentage_R4w_p4w": "Nation", "growth_percentage_R13W_P13W": "Nation"}#, "growth_percentage_R13W_Q1TD": "Nation"}
            for column_name, row_name in change_val_dict.items():
                df_values_changed = change_value(row_name, column_name, final_merged_df)
            df_values_changed
            return region_vials_growth_, df_values_changed



        def get_regional_avg(region):
            if region:
                area_data = tro_sales_data[tro_sales_data['area_nm'].str.contains(region, na=False)]
                filter_by = "regn_nm"
            else:
                area_data = tro_sales_data
                filter_by = "area_nm"

            filtered_data = area_data[(area_data['week_end_date'] >= '2024-01-05') & (area_data['week_end_date'] <= latest_weekend)]
#             filtered_data.to_excel("east_waest_filtered_data.xlsx")
            try:
                weekly_total_sales = filtered_data.groupby(['week_end_date', filter_by])['qty_sold_pu'].sum().unstack().drop('Unknown', axis=1)
            except:
                weekly_total_sales = filtered_data.groupby(['week_end_date', filter_by])['qty_sold_pu'].sum().unstack()
            print(weekly_total_sales)
            weekly_total_sales['Regional Avg'] = weekly_total_sales.mean(axis=1)
            cols = ['Regional Avg'] + [col for col in weekly_total_sales.columns if col != 'Regional Avg']
            weekly_total_sales = weekly_total_sales[cols]
            weekly_total_sales.reset_index(inplace=True)
        #     print(weekly_total_sales)
            if region == "East":
                order = ['week_end_date', 'Regional Avg', 'New England','New York','Atlantic Coastal','Blue Grass','Eastern Coastal','Atlanta','Southeast']
            elif region == "West":
                order = ['week_end_date', 'Regional Avg', 'Pacific Northwest','SoCAL','Rocky Mountains','Gulf Plains','Texas','North Central','Great Lakes']
            else:
                weekly_total_sales = weekly_total_sales.rename(columns={'Regional Avg': 'Area Average'})
                order = ['week_end_date', 'Area Average', 'East', 'West']

            weekly_total_sales_east = weekly_total_sales[order]


            def convert_date_format(date_str):
                date_obj = datetime.fromisoformat(str(date_str).replace("Z", "+00:00"))
                return date_obj.strftime("%d-%b")

            weekly_total_sales_east['week_end_date'] = weekly_total_sales_east['week_end_date'].apply(convert_date_format)
            return weekly_total_sales_east
        print("here!!!! before west!!")
        region_vials_growth_west, west_area_sales = get_slide_6_7_east_west_data("West")
        print("here!!!! after west!!")
        region_vials_growth_east, east_area_sales = get_slide_6_7_east_west_data("East")
        print("here!!!! after west!!")

        east_regional_avg = get_regional_avg("East")
        print("east_regional_avg!!!!!!!!", east_regional_avg)
        west_regional_avg = get_regional_avg("West")
        nation_vials_growth, nation_area_sales = get_slide_6_7_east_west_data('')
        nation_area_avg = get_regional_avg("")





        def get_comp_vs_tro_data():
            def get_dates(current_date , weeks_before='', weeks_after=''):
                print("current_date!!", current_date, weeks_before, weeks_after )
                given_date = datetime.strptime(current_date, "%Y-%m-%d")
                if weeks_before:
                    days_before = weeks_before * 7
                    date_before = given_date - timedelta(days=days_before)
                if weeks_after:
                    days_before = weeks_after * 7
                    date_before = given_date + timedelta(days=days_before)
                past_date = date_before.strftime("%Y-%m-%d")
                return past_date

            def get_date_range(base_date_str, weeks, date_format='%Y-%m-%d'):
                # end_date = datetime.strptime(base_date_str, date_format).replace(tzinfo=pytz.UTC)
                # start_date = end_date - timedelta(weeks=weeks - 1)
                # return start_date.strftime("%Y-%m-%dT%H:%M:%S.%fZ"), end_date.strftime("%Y-%m-%dT%H:%M:%S.%fZ")
                # Convert base_date_str to datetime object
                base_date = datetime.strptime(base_date_str, date_format)
                # Calculate start and end dates with time information
                end_date = base_date.replace(tzinfo=None) #remove tzinfo
                start_date = end_date - timedelta(weeks=weeks - 1)

                # Format start and end dates as strings with time information
                # Return only the date part
                return start_date.strftime(date_format), end_date.strftime(date_format)

            def add_new_rows(df, parent_name):
                new_df_ls = []
                new_row = {}
                sum_columns = ['qty_sold_pu_r12m', 'qty_sold_pu_p12m']
                for col in sum_columns:
                    new_row[col] = df[col].sum()
                if parent_name:
                    col_name = "Nation - TOPA"
                else:
                    col_name = "Nation"
                new_row["top_parent_account_type"] = col_name
                new_df_ls.append(new_row)
                calc_df =pd.DataFrame(new_df_ls)
                df_new_rows = pd.concat([df,calc_df], ignore_index=True)
                return df_new_rows

            def get_growth(drug_list, end_date_, weeks, sub_group_by, col_name):
                if 'trodelvy' in drug_list:
                    whole_data=tro_sales_data
                else:
                    whole_data=ddd_sales_data
                def format_data(sub, total, nation, sub_group_by):
                    total[sub_group_by] = total['top_parent_account_type']
                    nation[sub_group_by] = nation['top_parent_account_type']
                    result = pd.concat([sub, total, nation.iloc[[-1]]]).reset_index(drop=True)
                    df_no_unknowns_null = result[(result['top_parent_account_type'] != 'Unknown') & (result[sub_group_by] != 'Unknown') & (result[sub_group_by] != None)]
                    df_no_unknowns_null =df_no_unknowns_null.dropna()
                    df_sorted = df_no_unknowns_null.sort_values(by=['top_parent_account_type', sub_group_by])
                    return df_sorted


                def calculate_growth(recent_sales, past_sales):
                    # Avoid division by zero using np.where
                    # Convert to numeric and handle potential errors
                    recent_sales = pd.to_numeric(recent_sales, errors='coerce').fillna(0)
                    past_sales = pd.to_numeric(past_sales, errors='coerce').fillna(0)
                    growth_percentage = np.where(past_sales != 0, ((recent_sales - past_sales) / past_sales) * 100, np.nan)
                    # Convert the result into a Pandas Series and then fill NaN values with 0
                    growth_percentage = pd.Series(growth_percentage).fillna(0).replace([np.inf, -np.inf], 0)
                    return growth_percentage

                def filter_data(whole_data, group_by, parent_name=''):
                    r_start_date, r_end_date  = get_date_range(latest_weekend_ddd, 13)
                    p_start_date, p_end_date = get_date_range(get_dates(latest_weekend_ddd , 13), 13)


                    print("r_start_date", r_start_date, "r_end_date", r_end_date, "p_start_date, p_end_date!!", p_start_date, p_end_date)
                    if parent_name:
                        trodelvy_data = whole_data[whole_data['drug_name'].str.lower().isin(drug_list) &(whole_data['parent_name'] != "TEXAS ONCOLOGY")]
                    else:
                        trodelvy_data = whole_data[whole_data['drug_name'].str.lower().isin(drug_list)]

                    r_start_date = pd.Timestamp(r_start_date).tz_localize(None)
                    r_end_date = pd.Timestamp(r_end_date).tz_localize(None)
                    p_start_date = pd.Timestamp(p_start_date).tz_localize(None)
                    p_end_date = pd.Timestamp(p_end_date).tz_localize(None)
                    trodelvy_data['week_end_date'] = pd.to_datetime(trodelvy_data['week_end_date']).dt.tz_localize(None)
                    r12m_data = trodelvy_data[(trodelvy_data['week_end_date'] >= r_start_date) & (trodelvy_data['week_end_date'] <= r_end_date)]
                    p12m_data = trodelvy_data[(trodelvy_data['week_end_date'] >= p_start_date) & (trodelvy_data['week_end_date'] <= p_end_date)]
                    r12m_sales = (r12m_data.groupby(group_by)['qty_sold_pu'].sum()/13).reset_index()

                    p12m_sales = (p12m_data.groupby(group_by)['qty_sold_pu'].sum()/13).reset_index()


                    print("recent 12 m sales",r12m_sales)
                    print("past 12 m sales",p12m_sales)


                    merged_sales_ = r12m_sales.merge(p12m_sales, on= group_by, suffixes=('_r12m', '_p12m'))
                    output_df = add_new_rows(merged_sales_, parent_name)
                    


                    growth_percentage = calculate_growth(output_df['qty_sold_pu_r12m'], output_df['qty_sold_pu_p12m'])

                    # growth_percentage = ((output_df['qty_sold_pu_r12m'] - output_df['qty_sold_pu_p12m']) /
                    #  output_df['qty_sold_pu_p12m']) * 100

                    # Convert growth_percentage to float to handle decimal.Decimal type
                    growth_percentage = growth_percentage.astype(float)

                    # Handle NaN, Inf, and -Inf values
                    growth_percentage = growth_percentage.fillna(0).replace([np.inf, -np.inf], 0)

                    # Apply rounding and formatting
                    output_df[f'growth_percentage_{col_name}'] = growth_percentage.round(0).astype(float).astype(str) + '%'

                    return output_df


                    return output_df
                sub_merged_sales_ = filter_data(whole_data, ['top_parent_account_type', sub_group_by], True)
                merged_sales = filter_data(whole_data, ['top_parent_account_type'], True)
                Nation_Sales = filter_data(whole_data, ['top_parent_account_type'], False)
                formatted_df = format_data(sub_merged_sales_, merged_sales, Nation_Sales, sub_group_by)
                return formatted_df[['top_parent_account_type', sub_group_by, f'growth_percentage_{col_name}']]

            def get_contribution(drug_list, end_date_, weeks, sub_group_by, col_name):
                if 'trodelvy' in drug_list:
                    whole_data=tro_sales_data
                else:
                    whole_data=ddd_sales_data
                def format_data(sub, total, nation, sub_group_by):
                    total[sub_group_by] = total['top_parent_account_type']
                    result = pd.concat([sub, total]).reset_index(drop=True)
                    df_no_unknowns_null = result[(result['top_parent_account_type'] != 'Unknown') & (result[sub_group_by] != 'Unknown') & (result[sub_group_by] != None)]
                    df_no_unknowns_null =df_no_unknowns_null.dropna()
                    df_sorted = df_no_unknowns_null.sort_values(by=['top_parent_account_type', sub_group_by])
                    return df_sorted

                def filter_data(group_by, parent_name= ''):
                    start_date, end_date = get_date_range(end_date_ , weeks, '%Y-%m-%d')
                    start_date = datetime.strptime(start_date, '%Y-%m-%d')
                    end_date = datetime.strptime(end_date, '%Y-%m-%d')
                    start_date = pd.to_datetime(start_date).tz_localize(None)
                    end_date = pd.to_datetime(end_date).tz_localize(None)
                    whole_data['week_end_date'] = pd.to_datetime(whole_data['week_end_date']).dt.tz_localize(None)

                    if parent_name:
#                         filtered_data = whole_data[(whole_data['drug_name'].str.lower().isin(drug_list)) & (whole_data['week_end_date'] >= start_date) & (whole_data['week_end_date'] <= end_date) ]
                        filtered_data = whole_data[(whole_data['drug_name'].str.lower().isin(drug_list)) & (whole_data['week_end_date'] >= start_date) & (whole_data['week_end_date'] <= end_date) &(whole_data['parent_name'] != "TEXAS ONCOLOGY")]
                    else:
                        filtered_data = whole_data[(whole_data['drug_name'].str.lower().isin(drug_list)) & (whole_data['week_end_date'] >= start_date) & (whole_data['week_end_date'] <= end_date)]
                    filtered_data = filtered_data[['top_parent_account_type',sub_group_by, 'qty_sold_pu']]
                    print(filtered_data.columns, "col!!!!!")
                    grouped_data = filtered_data.groupby(group_by)['qty_sold_pu'].sum().reset_index()
                    total_sales = grouped_data['qty_sold_pu'].sum()
                    grouped_data[f'contribution_{col_name}'] = (grouped_data['qty_sold_pu'] / total_sales)
                    return grouped_data

                out_put_df  = filter_data(['top_parent_account_type', sub_group_by], True).drop(columns=['qty_sold_pu'])
                out_put_df_2  = filter_data(['top_parent_account_type'], True).drop(columns=['qty_sold_pu'])
                Nation_output_df = filter_data(['top_parent_account_type'], False)
                formatted_df = format_data(out_put_df, out_put_df_2, Nation_output_df, sub_group_by)
                formatted_df[f'contribution_{col_name}'] = round(formatted_df[f'contribution_{col_name}']* 100, 0).astype(int).astype(str) + '%'
                return out_put_df, out_put_df_2, formatted_df[['top_parent_account_type', sub_group_by, f'contribution_{col_name}']]

            def get_growth_contrib_data(sub_group_by):
                trodelvy_growth_data = get_growth(["trodelvy"], latest_weekend_ddd, 52, sub_group_by, 'trodelvy')
                enhertu_growth_data = get_growth(["enhertu"], latest_weekend_ddd, 52, sub_group_by, 'enhertu')
                tpc_growth_data = get_growth(tpc, latest_weekend_ddd, 52, sub_group_by, 'tpc')

                growth_merged_df1_2 = pd.merge(trodelvy_growth_data, enhertu_growth_data, on=['top_parent_account_type', sub_group_by], how='outer')
                growth_merged_all = pd.merge(growth_merged_df1_2, tpc_growth_data, on=['top_parent_account_type', sub_group_by], how='outer')
                growth_merged_all['growth_combined'] = growth_merged_all.apply(lambda row: f"{row['growth_percentage_trodelvy']} | {row['growth_percentage_enhertu']} | {row['growth_percentage_tpc']}", axis=1)

                tro_sub_account_contr_data ,tro_account_contr_data, trodelvy_contrib_data = get_contribution(["trodelvy"], latest_weekend_ddd, 52, sub_group_by, 'trodelvy')
                enh_sub_account_contr_data ,enh_account_contr_data, enhertu_contrib_data = get_contribution(["enhertu"], latest_weekend_ddd, 52, sub_group_by, 'enhertu')
                tpc_sub_account_contr_data ,tpc_account_contr_data, tpc_contrib_data= get_contribution(tpc, latest_weekend_ddd, 52, sub_group_by, 'tpc')

                contrib_merged_df1_2 = pd.merge(trodelvy_contrib_data, enhertu_contrib_data, on=['top_parent_account_type', sub_group_by], how='outer')
                contrib_merged_all = pd.merge(contrib_merged_df1_2, tpc_contrib_data, on=['top_parent_account_type', sub_group_by], how='outer')
                contrib_merged_all['contrib_combined'] = contrib_merged_all.apply(lambda row: f"{row['contribution_trodelvy']} | {row['contribution_enhertu']} | {row['contribution_tpc']}", axis=1)
                contrib_merged_all
                if sub_group_by == 'parent_bc_segment':
                    final_df = growth_merged_all[['top_parent_account_type', 'parent_bc_segment', 'growth_combined']]
                else:
                    final_df = pd.merge(growth_merged_all, contrib_merged_all, on=['top_parent_account_type', 'kad_flag'], how='left')
                    final_df = final_df[['top_parent_account_type', 'kad_flag','growth_combined', 'contrib_combined']]

                return tro_sub_account_contr_data, tro_account_contr_data, final_df


            tpc = ['gemzar', 'halaven', 'vinorelbine', 'xeloda']
#             latest_weekend = '2024-05-03'

            tro_sub_account_contr_data_, tro_account_contr_data, final_df_account = get_growth_contrib_data('parent_bc_segment')
            final_df_account.iloc[8], final_df_account.iloc[9] = final_df_account.iloc[9].copy(), final_df_account.iloc[8].copy()

            tro_sub_kad_contr_data_, tro_account_kad_data, final_df_kad =   get_growth_contrib_data( 'kad_flag')

            tro_sub_account_contr_data = tro_sub_account_contr_data_[ (tro_sub_account_contr_data_['parent_bc_segment'] != 'Unknown') & (tro_sub_account_contr_data_['parent_bc_segment'] != None)]
            tro_sub_kad_contr_data = tro_sub_kad_contr_data_[ (tro_sub_kad_contr_data_['kad_flag'] != 'Unknown') & (tro_sub_kad_contr_data_['kad_flag'] != None)]
            final_df_kad = final_df_kad.replace({ 'Y': 'Yes', 'N': 'No'})
            final_df_kad = final_df_kad.fillna("-")

            tro_account_contr_data = tro_account_contr_data[['top_parent_account_type','contribution_trodelvy']]
            tro_account_contr_data['top_parent_account_type'] = tro_account_contr_data['top_parent_account_type'].replace({
                'Unknown': 'Other',
                'IDN/Other': 'IDN/Hospital'
            })
            tro_account_contr_data = tro_account_contr_data.groupby('top_parent_account_type').sum().reset_index()

            tro_sub_kad_contr_data['kad_flag'] = tro_sub_kad_contr_data['kad_flag'].replace({
                'Y': "KAD Targets",
                'N': "Non - KAD Targets"
            })

            tro_sub_account_contr_data.loc[tro_sub_account_contr_data['top_parent_account_type'] == 'Unknown', 'parent_bc_segment'] = 'Other'

            tro_sub_kad_contr_data.loc[tro_sub_kad_contr_data['top_parent_account_type'] == 'Unknown', 'kad_flag'] = 'Other'


            tro_sub_account_contr_data = tro_sub_account_contr_data.drop(columns=['top_parent_account_type'])
            tro_sub_kad_contr_data = tro_sub_kad_contr_data.drop(columns=['top_parent_account_type'])
#             tro_sub_account_contr_data = tro_sub_account_contr_data.groupby('parent_bc_segment').sum().reset_index()
            sum_other = tro_sub_account_contr_data[tro_sub_account_contr_data['parent_bc_segment'] == 'Other']['contribution_trodelvy'].sum()
            tro_sub_account_contr_data.loc[tro_sub_account_contr_data['parent_bc_segment'] == 'Other', 'contribution_trodelvy'] = sum_other
            tro_sub_account_contr_data = tro_sub_account_contr_data.drop_duplicates()

            return tro_sub_account_contr_data, tro_account_contr_data, tro_sub_kad_contr_data, final_df_account, final_df_kad

        tro_sub_account_contr_data, tro_account_contr_data, tro_sub_kad_contr_data, final_df_account, final_df_kad = get_comp_vs_tro_data()
#         tro_sub_account_contr_data, tro_account_contr_data, tro_sub_kad_contr_data, final_df_account, final_df_kad = "", "", "","",""





        global E4_st, E13_st, W13_st, W4_st, A4_st, A13_st, C13_st, C4_st, M4_st, M13_st, L13_st, L4_st, H13_st, H4_st ,overall13_st,overall4_st
        E4_st = E13_st = W13_st = W4_st = A4_st = A13_st = C13_st = C4_st = M4_st = M13_st = L13_st = L4_st = H13_st = H4_st = overall13_st=overall4_st=None

        def update_text(df, paragraph, growth_row, filter_):
            print("value!!!!!", df.loc[df['time_period'] == growth_row, filter_].values[0])
            val = int(round(df.loc[df['time_period'] == growth_row, filter_].values.item(), 0))
            paragraph.text = str(val)+"%"

        def update_placeholders(df, growth_row, filter_):
            if df.equals(acad_df):
                global A13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    A13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    A13_st="no growth"
                else:
                    A13_st="increased"
            elif df.equals(acad_df1):
                global A4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    A4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    A4_st="no growth"
                else:
                    A4_st="increased"
            elif df.equals(comm_df1):
                global C4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    C4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    C4_st="no growth"
                else:
                    C4_st="increased"
            elif df.equals(comm_df):
                global C13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    C13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    C13_st="no growth"
                else:
                    C13_st="increased"
            elif df.equals(area_df1):
                global E4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    E4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    E4_st="no growth"
                else:
                    E4_st="increased"
            elif df.equals(east_df):
                global E13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    E13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    E13_st="no growth"
                else:
                    E13_st="increased"
            elif df.equals(area_df2):
                global W4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    W4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    W4_st="no growth"
                else:
                    W4_st="increased"
            elif df.equals(west_df):
                global W13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    W13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    W13_st="no growth"
                else:
                    W13_st="increased"
            elif df.equals(low_df1):
                global L4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    L4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    L4_st="no growth"
                else:
                    L4_st="increased"
            elif df.equals(low_df):
                global L13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    L13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    L13_st="no growth"
                else:
                    L13_st="increased"
            elif df.equals(medium_df1):
                global M4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    M4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    M4_st="no growth"
                else:
                    M4_st="increased"
            elif df.equals(medium_df):
                global M13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    M13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    M13_st="no growth"
                else:
                    M13_st="increased"
            elif df.equals(high_df):
                global H13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    H13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    H13_st="no growth"
                else:
                    H13_st="increased"
            elif df.equals(high_df1):
                global H4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    H4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    H4_st="no growth"
                else:
                    H4_st="increased"
            elif df.equals(overall_df):
                global overall13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    overall13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    overall13_st="no growth"
                else:
                    overall13_st="increased"
            elif df.equals(overall_df1):
                global overall4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    overall4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    overall4_st="no growth"
                else:
                    overall4_st="increased"


            # return abs(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0)))
            return abs(int(round(df.loc[df['time_period'] == growth_row, filter_].iloc[0], 0)))


        def get_slide_12_data():
                def get_averages(filtered_data):
                    filtered_data['week_end_date'] = pd.to_datetime(filtered_data['week_end_date']).dt.tz_localize(None)

                    specified_date = pd.Timestamp(latest_weekend).tz_localize(None)
                    start_date_r4w = specified_date - timedelta(weeks=3)

                    end_date_p4w = start_date_r4w - timedelta(days=1)
                    start_date_p4w = end_date_p4w - timedelta(weeks=4)

                    latest_week_qtr_start_dt, latest_weekend_  = get_quarter_start_date(latest_weekend)
                    start_date_ytd = pd.Timestamp(latest_week_qtr_start_dt).tz_localize(None)
                    end_date_ytd = pd.Timestamp(latest_weekend_).tz_localize(None)

                    r4w_data = filtered_data[(filtered_data['week_end_date'] >= start_date_r4w) & (filtered_data['week_end_date'] <= specified_date)]
                    r4w_total_qty = r4w_data['qty_sold_pu'].sum()
                    r4w_weekly_avg = r4w_total_qty / 4

                    p4w_data = filtered_data[(filtered_data['week_end_date'] >= start_date_p4w) & (filtered_data['week_end_date'] <= end_date_p4w)]
                    p4w_total_qty = p4w_data['qty_sold_pu'].sum()
                    p4w_weekly_avg = p4w_total_qty / 4
                    print("p4w_total_qty!!!", p4w_total_qty)

                    r4w_growth = calculate_growth(r4w_weekly_avg, p4w_weekly_avg)



                    ytd_data = filtered_data[(filtered_data['week_end_date'] >= start_date_ytd) & (filtered_data['week_end_date'] <= end_date_ytd)
            #                                    &(filtered_data['parent_name'] != "TEXAS ONCOLOGY")
                                            ]
                    ytd_ex_txs = filtered_data[(filtered_data['week_end_date'] >= start_date_ytd) & (filtered_data['week_end_date'] <= end_date_ytd)
                                            &(filtered_data['parent_name'] != "TEXAS ONCOLOGY")
                                            ]

                    ytd_total_qty = ytd_data['qty_sold_pu'].sum()
                    ytd_total_ex_txs = ytd_ex_txs['qty_sold_pu'].sum()

                    num_weeks_ytd = (end_date_ytd - start_date_ytd).days // 7 + 1
                    ytd_weekly_avg = ytd_total_qty / num_weeks_ytd
                    ytd_avg_ex_txs = ytd_total_ex_txs / num_weeks_ytd
                    return [r4w_weekly_avg, p4w_weekly_avg,ytd_avg_ex_txs,r4w_growth]

                account_types = ['high','low', 'medium', 'academic', 'community','east','west','None']
                final_output_dict = {}
                final_growth_dict = {}
                for type_ in account_types:
                    print("type_!!!!!!", type_)
                    averages_dict = {}
                    growth_dict = {}
                    if type_ in ['academic', 'community']:
                        filtered_data = tro_sales_data[tro_sales_data['child_account_type'].str.lower() == type_].copy()
                        filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize(None)

                    elif type_ in ['east','west']:
                        filtered_data = tro_sales_data[tro_sales_data['area_nm'].str.lower() == type_].copy()
                        filtered_data['week_end_date'] = pd.to_datetime(filtered_data['week_end_date']).dt.tz_localize(None)

                    elif type_ in ["low","high","medium"]:
                        filtered_data = tro_sales_data[tro_sales_data['linkedid_bc_segment'].str.lower() == type_].copy()
                    else:
                        filtered_data=tro_sales_data
                    r4w, p4w, ytd_ex_txs, r4w_growth = get_averages(filtered_data)
                    averages_dict['P4W'] = p4w
                    averages_dict['R4W'] = r4w
                    growth_dict['R4W_growth'] = r4w_growth

                    # averages_dict["YTD-TX"] = ytd_ex_txs
                    df = pd.DataFrame(list(averages_dict.items()), columns=['time_period', type_])
                    final_output_dict[type_] = df

                    growth_df = pd.DataFrame(list(growth_dict.items()), columns=['time_period', "growth"])
                    final_growth_dict[type_] = growth_df
                return final_output_dict, final_growth_dict
        slide_12_data,  slide_12_growth_data= get_slide_12_data()
        def get_trend_r13w():
            # Get the latest available week_end_date
            tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date']).dt.strftime('%Y-%m-%dT%H:%M:%S.%fZ')
            tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date'])
            sorted_dates = tro_sales_data['week_end_date'].sort_values(ascending=False).unique()
            latest_weekend = sorted_dates[0].date().strftime('%Y-%m-%d')
            print(f"Latest week_end_date in the data: {latest_weekend}")

            # Calculate the start date for the most recent 13 weeks and previous 13 weeks
            latest_weekend = pd.to_datetime(latest_weekend)
            start_date = latest_weekend - timedelta(weeks=25)
            print(f"Recent 26-week period: {start_date} to {latest_weekend}")

            # Ensure the 'week_end_date' column has timezone removed
            tro_sales_data['week_end_date'] = tro_sales_data['week_end_date'].dt.tz_localize(None)

            # Filter sales data for the last 26 weeks dynamically
            tro_weekly_sales = (tro_sales_data[
                (tro_sales_data['week_end_date'] >= start_date) &
                (tro_sales_data['week_end_date'] <= latest_weekend)
            ]
            .groupby('week_end_date', as_index=False)['qty_sold_pu']
            .sum())
            tro_weekly_sales['week_end_date'] = tro_weekly_sales['week_end_date'].dt.strftime('%d-%b')
            return tro_weekly_sales

        r13w_trend=get_trend_r13w()

        acad_df = slide_2_growth_data["academic"]
        comm_df = slide_2_growth_data["community"]
        high_df = slide_2_growth_data["high"]
        medium_df = slide_2_growth_data["medium"]
        low_df = slide_2_growth_data["low"]
        east_df=slide_2_growth_data["east"]
        west_df=slide_2_growth_data["west"]
        overall_df=slide_2_growth_data["None"]


        acad_df1 = slide_12_growth_data["academic"]
        comm_df1 = slide_12_growth_data["community"]
        high_df1 = slide_12_growth_data["high"]
        medium_df1 = slide_12_growth_data["medium"]
        low_df1 = slide_12_growth_data["low"]
        area_df1=slide_12_growth_data["east"]
        area_df2=slide_12_growth_data["west"]
        overall_df1=slide_12_growth_data["None"]
        A13_g= update_placeholders(acad_df,'R13W_growth','growth')
        C13_g= update_placeholders(comm_df,'R13W_growth','growth')
        L13_g= update_placeholders(low_df,'R13W_growth','growth')
        M13_g= update_placeholders(medium_df,'R13W_growth','growth')
        H13_g= update_placeholders(high_df,'R13W_growth','growth')
        E13_g= update_placeholders(east_df,'R13W_growth','growth')
        W13_g= update_placeholders(west_df,'R13W_growth','growth')
        overall13_g=update_placeholders(overall_df,'R13W_growth','growth')
        #overall13_g = update_placeholders(overall_df, shape, 'R4W_growth', 'Value')

        A4_g= update_placeholders(acad_df1,'R4W_growth','growth')
        C4_g= update_placeholders(comm_df1,'R4W_growth','growth')
        L4_g= update_placeholders(low_df1,'R4W_growth','growth')
        M4_g= update_placeholders(medium_df1,'R4W_growth','growth')
        H4_g= update_placeholders(high_df1,'R4W_growth','growth')
        E4_g= update_placeholders(area_df1,'R4W_growth','growth')
        W4_g= update_placeholders(area_df2,'R4W_growth','growth')
        overall4_g=update_placeholders(overall_df1,'R4W_growth','growth')
        def get_text_r13():

            medium_segment_text_13 = f"Medium segment {M13_st} by {M13_g}%" if M13_st != 'no growth' else "Medium segment remained same"
            high_segment_text_13=f"High segment {H13_st} by {H13_g}%" if H13_st != 'no growth' else "High segment remained same"
            low_segment_text_13=f"Low segment {L13_st} by {L13_g}%" if L13_st != 'no growth' else "Low segment remained same"
            accounts_text_13=f"Academic {A13_st} by {A13_g}%" if A13_st != 'no growth' else "Academic remained same"
            community_text_13=f"Community {C13_st} by {C13_g}%" if C13_st != 'no growth' else "Community remained same"
            east_text_13=f"East {E13_st} by {E13_g}%" if E13_st != 'no growth' else "East remained same"
            west_text_13=f"West {W13_st} by {W13_g}%" if W13_st != 'no growth' else "West remained same"
            overall_text_13=f"Sales {overall13_st} by {overall13_g}%" if overall13_st != 'no growth' else "Sales remained same"


            # Now, use the variable 'medium_segment_text' in the formatted_text_r13w list
            formatted_text_r13w = [
                (f"R13W VS P13W {latest_day_mon}", True, 18),  # Bold title with font size 18
                (f"{overall_text_13} over the past 13 weeks (R13W) compared to the previous 13 weeks (P13W).\n", False, 14),
                ("Accounts: ", True, 14),
                (f"{accounts_text_13}, {community_text_13}", False, 14),
                ("Segments: ", True, 14),
                (f"{low_segment_text_13},{medium_segment_text_13} and {high_segment_text_13}", False, 14),
                ("Regions: ", True, 14),
                (f"{east_text_13}, {west_text_13}", False, 14)
            ]

            return formatted_text_r13w
        def get_text_r4():
            medium_segment_text_4 = f"Medium segment {M4_st} by {M4_g}%" if M4_st != 'no growth' else "Medium segment remained same"
            high_segment_text_4=f"High segment {H4_st} by {H4_g}%" if H4_st != 'no growth' else "High segment remained same"
            low_segment_text_4=f"Low segment {L4_st} by {L4_g}%" if L4_st != 'no growth' else "Low segment remained same"
            accounts_text_4=f"Academic {A4_st} by {A4_g}%" if A4_st != 'no growth' else "Academic remained same"
            community_text_4=f"Community {C4_st} by {C4_g}%" if C4_st != 'no growth' else "Community remained same"
            east_text_4=f"East {E4_st} by {E4_g}%" if E4_st != 'no growth' else "East remained same"
            west_text_4=f"West {W4_st} by {W4_g}%" if W4_st != 'no growth' else "West remained same"
            overall_text_4=f"Sales {overall4_st} by {overall4_g}%" if overall4_st != 'no growth' else "Sales remained same"

            formatted_text_r4w = [
                (f"R4W VS P4W {latest_day_mon}", True, 18),  # Bold title with font size 18
                (f"{overall_text_4} over the past 4 weeks (R4W) compared to the previous 4 weeks (P4W).\n", False, 14),
                ("Accounts: ", True, 14),
                (f"{accounts_text_4}, {community_text_4}", False, 14),
                ("Segments: ", True, 14),
                (f"{low_segment_text_4}, {medium_segment_text_4} and {high_segment_text_4}", False, 14),
                ("Regions: ", True, 14),
                (f"{east_text_4}, {west_text_4}", False, 14)

            ]
            return formatted_text_r4w

        # Function to update text with font size and bold styling
        def update_textbox(shape, formatted_text):
            text_frame = shape.text_frame
            text_frame.clear()
            for text, is_bold, font_size in formatted_text:
                para = text_frame.add_paragraph()
                run = para.add_run()
                run.text = text
                run.font.bold = is_bold
                run.font.size = Pt(font_size)  # Set font size
                para.space_after = Pt(0)


        from pptx import Presentation
        from pptx.chart.data import ChartData, CategoryChartData
        import os
        # script_dir = os.path.dirname(os.path.abspath(__file__))
        # ppt_path = os.path.join(script_dir, 'ppt_template_867_sales.pptx')
        # presentation = Presentation(ppt_path)

        ppt_path ='./Structured_Bot/ppt_template_867_sales.pptx'
        presentation = Presentation(ppt_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                final_df = []
                if shape.has_chart:
                    chart = shape.chart
                    chart_title = chart.chart_title.text_frame.text
                    print("chart title", chart_title)
                    if "Trodelvy Weekly Demand" in chart_title:
                        final_df = slide_1_data
                    elif "Academic account" in chart_title:
                        print("academic for slide 2")
                        res = slide_2_data["academic"]
                        if "TX" in chart_title:
                            print("inside TX")
                            final_df = res[res['time_period'] != 'YTD']
                            self.update_cell_text(chart.chart_title, f"Academic({latest_day_mon})-TX", "chart_title")
                        else:
                            print("inside else")
                            final_df = res[res['time_period'] != 'YTD-TX']
                            print(chart.chart_title)
                            print(f"Academic({latest_day_mon})")
                            print("before update cell text")
                            self.update_cell_text(chart.chart_title, f"Academic({latest_day_mon})", "chart_title")
                            print("after update cell text")
                        print("done with updating the cell text")
                    elif "High Potential" in chart_title:
                        res = slide_2_data["high"]
                        if "TX" in chart_title:
                            final_df = res[res['time_period'] != 'YTD']
                            self.update_cell_text(chart.chart_title, f"High Potential for before ({latest_day_mon})-TX", "chart_title")
                        else:
                            final_df = res[res['time_period'] != 'YTD-TX']
                            self.update_cell_text(chart.chart_title, f"High Potential({latest_day_mon})", "chart_title")
                    elif "Medium Potential" in chart_title:
                        res = slide_2_data["medium"]
                        if "TX" in chart_title:
                            final_df = res[res['time_period'] != 'YTD']
                            self.update_cell_text(chart.chart_title, f"Medium Potential({latest_day_mon})-TX", "chart_title")
                        else:
                            final_df = res[res['time_period'] != 'YTD-TX']
                            self.update_cell_text(chart.chart_title, f"Medium Potential({latest_day_mon})", "chart_title")

                    elif "Low Potential" in chart_title:
                        res = slide_2_data["low"]
                        if "TX" in chart_title:
                            final_df = res[res['time_period'] != 'YTD']
                            self.update_cell_text(chart.chart_title, f"Low Potential({latest_day_mon})-TX", "chart_title")
                        else:
                            final_df = res[res['time_period'] != 'YTD-TX']
                            self.update_cell_text(chart.chart_title, f"Low Potential({latest_day_mon})", "chart_title")
                    elif "Community account" in chart_title:
                        res = slide_2_data["community"]
                        if "TX" in chart_title:
                            final_df = res[res['time_period'] != 'YTD']
                            self.update_cell_text(chart.chart_title, f"Community({latest_day_mon})-TX", "chart_title")
                        else:
                            final_df = res[res['time_period'] != 'YTD-TX']
                            self.update_cell_text(chart.chart_title, f"Community({latest_day_mon})", "chart_title")
                    elif "R13W Trends" in chart_title:
                        final_df= r13w_trend
                        self.update_cell_text(chart.chart_title, f"Trodelvy Weekly Sales Trends({latest_day_mon})", "chart_title")
                    elif "Region Vials Growth" in chart_title and "East" in chart_title:
                        final_df = region_vials_growth_east
                        self.update_cell_text(chart.chart_title, f"Region Vials Growth % (East)w.e({latest_day_mon})", "chart_title")
                    elif "Region Vials Growth" in chart_title and "West" in chart_title:
                        final_df = region_vials_growth_west
                        self.update_cell_text(chart.chart_title, f"Region Vials Growth % (West)w.e({latest_day_mon})", "chart_title")
                    elif "Weekly Region Vials" in chart_title and "East" in chart_title:
                        self.final_df = east_regional_avg
                    elif "Weekly Region Vials" in chart_title and "West" in chart_title:
                        final_df = west_regional_avg
                    elif "Area Vials Growth" in chart_title:
                        final_df = nation_vials_growth
                        self.update_cell_text(chart.chart_title, f"Area Vials Growth % w.e({latest_day_mon})", "chart_title")
                    elif "Area Vials Sales" in chart_title:
                        final_df = nation_area_avg
                    elif chart_title == "segment contribution":
                        final_df = tro_sub_account_contr_data
                    elif chart_title == "Acc contribution ":
                        final_df = tro_account_contr_data
                    elif chart_title == "Kad contribution":
                        final_df = tro_sub_kad_contr_data
                    print("line 2738")
                    try:
                        final_df.replace('YTD-TX', 'YTD', inplace=True)
                    except:
                        print("line 2742")
                        pass
                    if len(final_df) >0:
                        final_df.replace([np.nan, np.inf, -np.inf], '', inplace=True)
                        print("Chart found!", chart.chart_title.text_frame.text)
                        chart_data = CategoryChartData()
                        categories = final_df.iloc[:, 0].tolist()  # Assuming the first column is the category
                        chart_data.categories = categories
                        for col in final_df.columns[1:]:
                            chart_data.add_series(col, final_df[col].tolist())
                        chart.replace_data(chart_data)
                        print("replaced data!!")
                else:
                    pass
                if shape.has_text_frame:
                    print("line 2757")
                    text = shape.text_frame.text
                    if "R13W vs P13W" in text:
                        update_textbox(shape, get_text_r13())
                    elif "R4W vs P4W" in text:
                        update_textbox(shape, get_text_r4())
                    for paragraph in shape.text_frame.paragraphs:
                        original_text = paragraph.text
                        print("Original Paragraph Text:", original_text)
                        if original_text == "aca PG":
                            update_text(acad_df, paragraph, 'P13W_growth', 'growth')
                        elif original_text == "aca RG":
                            update_text(acad_df, paragraph, 'R13W_growth', 'growth')
                        elif original_text == "com PG":
                            update_text(comm_df, paragraph, 'P13W_growth', 'growth')
                        elif original_text == "com RG":
                            update_text(comm_df, paragraph, 'R13W_growth', 'growth')
                        elif original_text == "high PG":
                            update_text(high_df, paragraph, 'P13W_growth', 'growth')
                        elif original_text == "high RG":
                            update_text(high_df, paragraph, 'R13W_growth', 'growth')
                        elif original_text == "med PG":
                            update_text(medium_df, paragraph, 'P13W_growth', 'growth')
                        elif original_text == "med RG":
                            update_text(medium_df, paragraph, 'R13W_growth', 'growth')
                        elif original_text == "low PG":
                            update_text(low_df, paragraph, 'P13W_growth', 'growth')
                        elif original_text == "low RG":
                            update_text(low_df, paragraph, 'R13W_growth', 'growth')
        presentation.save("867_sales_refreshed_ppt.pptx")


        from pptx import Presentation
        def print_table_content(pptx_file, slide_1_table_data, slide_3_data , high_acc_decliners, med_acc_decliners, academic_acc_decliners, community_acc_decliners ,idn_hospital_decliners, east_area_sales , west_area_sales, nation_area_sales):
            df_table = False
            start_row = 2
            prs = Presentation(pptx_file)
            for slide_number, slide in enumerate(prs.slides):
                print(f"Slide {slide_number + 1}")
                for shape in slide.shapes:
                    if not shape.has_table and not shape.has_chart:
                        continue
                    if shape.has_chart:
                        chart = shape.chart
                        print(chart.chart_title.text_frame.text)
                        if "Trodelvy Weekly Demand" in chart.chart_title.text_frame.text:
                            self.update_cell_text(chart.chart_title, f"Trodelvy Weekly Demand (k Vials) (WE {latest_day_mon})", "chart_title")

                    if shape.has_table:
                        table = shape.table
                        cell = table.cell(0, 0).text
                        print("cell!!!", cell)
                        flag_value = table.cell(0, 1).text

                        if flag_value == '':
                            flag_value = table.cell(0, 2).text
                        print("flag_value!!",flag_value)
                        if "% Contribution to Nation" in flag_value:
                            df = slide_3_data
                            df_table = True
                            start_row = 2
                            indices = [(0,2), (0,4),(0,5),(0,6),(0,7)]
                            keys = [f'R12M \nWeekly Sales \n(Ending {latest_day_mon})', f'R13W weekly avg\n(w.e. {latest_day_mon})', f'R13W Growth: (R13W w.e. {latest_day_mon} vs P13W)', f'YTD Sales Growth'  ]
                            for idx, value in zip(indices, keys):
                                self.update_cell_text(shape.table.cell(idx[0], idx[1]), value, "value")
                        elif "East" in flag_value or "West" in flag_value or "Nation" in flag_value:
                            print("in slide 6 & 10!!")
                            if "East" in flag_value:
                                df = east_area_sales
                            elif "West" in flag_value:
                                df = west_area_sales
                            elif "Nation" in flag_value:
                                df = nation_area_sales
                            start_row = 3
                            df_table = True
                        elif "high" in cell.lower():
                            df = high_acc_decliners
                            df_table = True
                            start_row = 2
                        elif "medium" in cell.lower():
                            df = med_acc_decliners
                            df_table = True
                            start_row = 2
                        elif "academic" in cell.lower():
                            df = academic_acc_decliners
                            df_table = True
                            start_row = 2
                        elif "community" in cell.lower():
                            df = community_acc_decliners
                            df_table = True
                            start_row = 2
                        elif "idn" in cell.lower():
                            print("in idn")
                            df = idn_hospital_decliners
                            df_table = True
                            start_row = 2
                        if "high" in cell.lower() or "medium" in cell.lower() or "academic" in cell.lower() or "community" in cell.lower() or "idn" in cell.lower():
                            indices = [(0,1), (0,2)]
                            keys = [f'R13W weekly avg (w.e. {latest_day_mon})', f'P13W weekly avg (w.e. {latest_day_mon})']
                            for idx, value in zip(indices, keys):
                                self.update_cell_text(shape.table.cell(idx[0], idx[1]), value, "value")

                        elif "Weekly Avg Demand Vials" in cell:
                            indices = [(0,0), (1, 1), (1, 2), (1, 3), (1, 4), (2, 1)]
                            keys = ['avg_demand', 'R4W', 'R8W', 'R13W', 'R52W', 'growth_4w', 'growth_8w', 'growth_13w', 'growth_qtr']
                            for idx, key in zip(indices, keys):
                                value = slide_1_table_data.get(key, '')
                                if key == "growth_4w":
                                    value = str(value) +"%"
                                self.update_cell_text(shape.table.cell(idx[0], idx[1]), value, "value")
                        elif "Total Demand Vials" in cell:
                            indices = [(0,0),(1, 1), (1, 3), (2, 1), (2, 3), (3,1), (3,3), (1,2),(2, 2), (3, 2), (2, 2), (2, 3), (2, 4)]
                            keys = ['total_demand', 'current_qtr_vials', 'curr_qtr_growth', 'ytd_vials', 'ytd_growth', 'prev_qtr_vials', 'prev_qtr_growth', 'current_qtr_est_avg', 'ytd_est_avg', 'prev_qtr_est_avg']
                            for idx, key in zip(indices, keys):
                                value = slide_1_table_data.get(key, '')
                                # if type(value) != str:
                                #     value = round(value, 1)
                                self.update_cell_text(shape.table.cell(idx[0], idx[1]), value , "value")

                        elif  "Quarterly Avg. Demand Vials" in cell:
                            indices = [(0,0), (1, 1), (1, 2)]
                            keys = ['qtr_avg_demand' ,'q1', 'q2']
                            for idx, key in zip(indices, keys):
                                value = slide_1_table_data.get(key, '')
                                # if type(value) != str:
                                #     value = round(value, 1)
                                self.update_cell_text(shape.table.cell(idx[0], idx[1]), value , "value")

                        elif "BC Segment type" in flag_value:
                            df = final_df_account
                            df_table = True
                            start_row = 1
                        elif "KAD Targets" in flag_value:
                            df_table = True
                            df = final_df_kad

                        if df_table:
                            for i, row in enumerate(table.rows):
                                for j, cell in enumerate(row.cells):
                                    if i >= start_row:
                                        df_row_index = i - start_row
                                        if df_row_index < len(df) and j < len(df.columns):
                                            value = str(df.iloc[df_row_index, j])
                                            self.update_cell_text(cell, value)
                                        else:
                                            self.update_cell_text(cell, "")

                    print("-" * 20)
            print("saving file!!!")
            prs.save("867_sales_refreshed_ppt.pptx")
            print("saved file!!!")

        print_table_content('867_sales_refreshed_ppt.pptx', slide_1_table_data, slide_3_table_data, high_acc_decliners, med_acc_decliners, academic_acc_decliners, community_acc_decliners, idn_hospital_decliners, east_area_sales, west_area_sales, nation_area_sales)
        print("before closing cursor in update sales oncology slides")
        self.cursor.close()
        print("after closing cursor in update sales oncology slides")
        print("before closing connection in update sales oncology slides")
        self.connection.close()
        print("after closing connection in update sales oncology slides")
