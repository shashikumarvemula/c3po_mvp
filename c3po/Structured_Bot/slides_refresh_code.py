#!/usr/bin/env python
# coding: utf-8

# In[ ]:

import numpy as np
import time
import zipfile
import xml.etree.ElementTree as ET
import openpyxl
import os
import tempfile
import shutil
import pandas as pd
import matplotlib.pyplot as plt
from datetime import datetime, timedelta
import pytz
from functools import reduce
import dask.dataframe as dd
from pptx.util import Pt
from databricks import sql
import os
import pandas as pd
from datetime import datetime, timedelta
import pytz
from decimal import Decimal

class UpdateSalesDeck:
    def __init__(self):
        # Databricks SQL Warehouse configuration
        # Read Databricks credentials from environment variables
        self.DATABRICKS_SERVER_HOSTNAME = os.getenv("DATABRICKS_SERVER_HOSTNAME")
        self.DATABRICKS_HTTP_PATH = os.getenv("DATABRICKS_HTTP_PATH")
        self.DATABRICKS_TOKEN = os.getenv('DATABRICKS_TOKEN')
        self.connection = None
        self.cursor = None
        self.latest_weekend = ""

        self.latest_day_mon = ""
        self.latest_year = ""
        self.latest_month = ""


    def connect_to_databricks(self):
        try:
            from databricks import sql
            import os
            self.connection = sql.connect(
                server_hostname=self.DATABRICKS_SERVER_HOSTNAME,
                http_path=self.DATABRICKS_HTTP_PATH,
                access_token=self.DATABRICKS_TOKEN,
                _tls_no_verify=True
            )
            self.cursor = self.connection.cursor()
            print("Databricks Connected Successfully")
        except Exception as e:
            print(f"Database connection error: {str(e)}")
            raise

    def update_deck(self):
        self.connect_to_databricks()
        print("Connection Established")


        # tro_sales_data = tro_sales_data.compute()  # Convert Dask to Pandas
        # ddd_sales_data = ddd_sales_data.compute()  # Convert Dask to Pandas
        # tro_sales_data['drug_name']='trodelvy'
        start_time=time.time()
        sql_query = f"""
                SELECT * FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`867_data_drug-trodelvy_for_genie_filtered_dates`"""
        tro_sales_data= pd.DataFrame(self.cursor.execute(sql_query).fetchall(), columns=[elem[0] for elem in self.cursor.description])
        end_time=time.time()
        print("latency for getting the data",start_time-end_time)

        tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date']).dt.tz_localize(None)


        print("tro sales data sample",tro_sales_data.head())

                # Query to get the latest weekend for 867 data
        sorted_dates = pd.DataFrame(self.cursor.execute("""SELECT DISTINCT `week_end_date`
                    FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`867_data_drug-trodelvy_for_genie_filtered_dates`
                    ORDER BY `week_end_date` DESC""").fetchall(), columns=[elem[0] for elem in self.cursor.description])
        # Extract and format the latest weekend
        latest_weekend = pd.to_datetime(sorted_dates.iloc[0, 0]).strftime('%Y-%m-%d')
        self.latest_weekend = latest_weekend
        print("Latest weekend (867 data):", latest_weekend)

        # Query to get the latest weekend for DDD data
        start_time=time.time()
        latest_weekend_ddd = pd.DataFrame(self.cursor.execute("""SELECT
                    DATE_FORMAT(MAX(week_end_date), 'yyyy-MM-dd') AS latest_weekend_ddd
                FROM
                    `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`ddd_data_for_genie_filtered_dates`;""").fetchall(),
                columns=[elem[0] for elem in self.cursor.description])
        end_time=time.time()
        print("latency for getting the data ddd weekend date",start_time-end_time)
        # Convert latest_weekend_ddd to datetime
        latest_weekend_ddd = pd.to_datetime(latest_weekend_ddd.iloc[0, 0]).strftime('%Y-%m-%d')

        # Print the result
        print("Latest weekend for DDD:", latest_weekend_ddd)

        print("latest weekend we use",latest_weekend)
        # print("latest weekend for ddd",pd.to_datetime(ddd_sales_data['week_end_date']).max().date().strftime('%Y-%m-%d'))

        latest_day_mon= pd.to_datetime(latest_weekend).strftime('%m/%d')
        # latest_day_mon= pd.to_datetime(latest_weekend).strftime('%m/%d')

        self.latest_weekend = pd.to_datetime(self.latest_weekend)
        self.latest_day_mon = self.latest_weekend.strftime('%m/%d')
        self.latest_year = self.latest_weekend.year
        self.latest_month = self.latest_weekend.month

        latest_year = pd.to_datetime(latest_weekend).year
        latest_month = pd.to_datetime(latest_weekend).month
        latest_month, latest_year


        # ddd_sales_data = ddd_sales_data.rename(columns={'trx_cnt': 'qty_sold_pu'})
        # tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date']).dt.strftime('%Y-%m-%dT%H:%M:%S.%fZ')
        # tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date'])
        # sorted_dates = tro_sales_data['week_end_date'].sort_values(ascending=False).unique()
        # latest_weekend = sorted_dates[0].date().strftime('%Y-%m-%d')
        # print("latest weekend",latest_weekend)
        # latest_weekend_ddd=pd.to_datetime(ddd_sales_data['week_end_date']).max().date().strftime('%Y-%m-%d')
        # print("latest weekend we use",latest_weekend)
        # print("latest weekend for ddd",pd.to_datetime(ddd_sales_data['week_end_date']).max().date().strftime('%Y-%m-%d'))

        # latest_day_mon= pd.to_datetime(latest_weekend).strftime('%m/%d')
        # latest_day_mon= pd.to_datetime(latest_weekend).strftime('%m/%d')

        # latest_year = pd.to_datetime(latest_weekend).year
        # latest_month = pd.to_datetime(latest_weekend).month
        # latest_month, latest_year
        estimate_dates = [
            "2025-01-03", "2025-01-10", "2025-01-17", "2025-01-24", "2025-01-31", "2025-02-07", "2025-02-14",
            "2025-02-21", "2025-02-28", "2025-03-07", "2025-03-14", "2025-03-21", "2025-03-28", "2025-04-04",
            "2025-04-11", "2025-04-18", "2025-04-25", "2025-05-02", "2025-05-09", "2025-05-16", "2025-05-23",
            "2025-05-30", "2025-06-06", "2025-06-13", "2025-06-20", "2025-06-27", "2025-07-04", "2025-07-11",
            "2025-07-18", "2025-07-25", "2025-08-01", "2025-08-08", "2025-08-15", "2025-08-22", "2025-08-29",
            "2025-09-05", "2025-09-12", "2025-09-19", "2025-09-26", "2025-10-03", "2025-10-10", "2025-10-17",
            "2025-10-24", "2025-10-31", "2025-11-07", "2025-11-14", "2025-11-21", "2025-11-28", "2025-12-05",
            "2025-12-12", "2025-12-19", "2025-12-26"
        ]
        print("len!!!", len(estimate_dates))

        estimate_values = [
            7884, 9307, 8664, 8164, 9670, 8163, 7807, 7921, 8368, 8177, 7979, 8404, 8227, 8049,
            8107, 8068, 8025, 8436, 8488, 8588, 8616, 7510, 8001, 8022, 8083, 8289, 7535, 8917,
            8574, 8485, 8573, 8418, 8430, 8410, 8483, 7913, 8654, 8408, 8449, 8455, 8095, 8138,
            8053, 8172, 8439, 8489, 8498, 7379, 8177, 7796, 7925, 8526
        ]


        slide_1_chart_estimates = pd.DataFrame({
            'formatted_date': estimate_dates,
            '24 Budget Weekly Run Rate Estimate': estimate_values,
            'full_date': estimate_dates
        })
        # slide_1_chart_estimates['full_date'] = pd.to_datetime(slide_1_chart_estimates['formatted_date'] + '-2024', format='%d-%b-%Y')


        # In[ ]:


        def get_dates(current_date, weeks_before=None, weeks_after=None):
            print("current_date!!", current_date, weeks_before, weeks_after)
            given_date = datetime.strptime(current_date, "%Y-%m-%d")

            # Ensure weeks_before and weeks_after are valid integers
            if isinstance(weeks_before, str) and weeks_before.strip().isdigit():
                weeks_before = int(weeks_before)
            if isinstance(weeks_after, str) and weeks_after.strip().isdigit():
                weeks_after = int(weeks_after)

            # Handle weeks_before and weeks_after
            if weeks_before is not None and weeks_before != "":
                days_before = weeks_before * 7
                date_before = given_date - timedelta(days=days_before)
            elif weeks_after is not None and weeks_after != "":
                days_after = weeks_after * 7
                date_before = given_date + timedelta(days=days_after)
            else:
                raise ValueError("Either weeks_before or weeks_after must be provided and should be a valid number.")

            past_date = date_before.strftime("%Y-%m-%d")
            return past_date


        def get_quarter_start_date(latest_weekend):
                latest_weekend_= datetime.strptime(latest_weekend, "%Y-%m-%d")
                if latest_weekend_.month in [1, 2, 3]:
                    start_date  = datetime(latest_weekend_.year , 1, 1)
                elif latest_weekend_.month in [4, 5, 6]:
                    start_date  = datetime(latest_weekend_.year , 4, 1)
                elif latest_weekend_.month in [7, 8, 9]:
                    start_date  = datetime(latest_weekend_.year , 7, 1)
                else:
                    start_date  = datetime(latest_weekend_.year , 10, 1)
                return start_date.strftime("%Y-%m-%d"), latest_weekend_.strftime("%Y-%m-%d")

        def get_last_date_before_qtr(start_date, df, date_column):
                start_date = pd.to_datetime(start_date).date().strftime('%Y-%m-%d')
                filtered_df = df[df[date_column] < start_date]
                if not filtered_df.empty:
                    last_date_before_q4 = filtered_df[date_column].max()
                    return last_date_before_q4.strftime("%Y-%m-%d")
                else:
                    return None

        # def get_date_range(base_date_str, weeks, date_format='%Y-%m-%d'):
        #     end_date = datetime.strptime(base_date_str, date_format).replace(tzinfo=pytz.UTC)
        #     start_date = end_date - timedelta(weeks=weeks - 1)
        #     return start_date.strftime("%Y-%m-%d"), end_date.strftime("%Y-%m-%d")
        # def get_date_range(base_date_str, weeks, date_format='%Y-%m-%d'):
        #     # Convert base_date_str to datetime object
        #     base_date = datetime.strptime(base_date_str, date_format)

        #     # Calculate start and end dates with time information
        #     end_date = base_date.replace(tzinfo=pytz.UTC)
        #     start_date = end_date - timedelta(weeks=weeks - 1)

        #     # Format start and end dates as strings with time information
        #     return start_date.strftime("%Y-%m-%dT%H:%M:%S.%fZ"), end_date.strftime("%Y-%m-%dT%H:%M:%S.%fZ")

        def get_date_range(base_date_str, weeks, date_format='%Y-%m-%d'):
            # Convert base_date_str to datetime object
            base_date = datetime.strptime(base_date_str, date_format)
            # Calculate start and end dates with time information
            end_date = base_date.replace(tzinfo=None) #remove tzinfo
            start_date = end_date - timedelta(weeks=weeks - 1)

            # Format start and end dates as strings with time information
            return start_date.strftime("%Y-%m-%d"), end_date.strftime("%Y-%m-%d")
        def get_previous_quarter_dates(quarter_start):
            quarter_start = pd.Timestamp(quarter_start)
            prev_quarters = {
                1: (pd.Timestamp(quarter_start.year - 1, 10, 1), pd.Timestamp(quarter_start.year - 1, 12, 31)),
                2: (pd.Timestamp(quarter_start.year, 1, 1), pd.Timestamp(quarter_start.year, 3, 31)),
                3: (pd.Timestamp(quarter_start.year, 4, 1), pd.Timestamp(quarter_start.year, 6, 30)),
                4: (pd.Timestamp(quarter_start.year, 7, 1), pd.Timestamp(quarter_start.year, 9, 30))
            }
            current_quarter = (quarter_start.month - 1) // 3 + 1
            prev_quarter_start, prev_quarter_end = prev_quarters[current_quarter]
            return prev_quarter_start.strftime("%Y-%m-%d"), prev_quarter_end.strftime("%Y-%m-%d")

        def sum_in_thousands(df, start_date, end_date):
            df['full_date'] = pd.to_datetime(df['full_date'])
            filtered_df = df[(df['full_date'] >= start_date) & (df['full_date'] <= end_date)]
            sum_values_in_thousands = filtered_df['24 Budget Weekly Run Rate Estimate'].sum() / 1000
            return sum_values_in_thousands

        def get_diff_in_weeks(latest_weekend_, start_date):
            difference_in_days = (pd.to_datetime(latest_weekend_) - pd.to_datetime(start_date)).days
            difference_in_weeks = difference_in_days / 7
            weeks_from_quarter_start = round(difference_in_weeks)
            return weeks_from_quarter_start

        drug_list=['gemzar', 'halaven', 'vinorelbine', 'xeloda','enhertu']

        p_start_date, p_end_date = get_date_range(get_dates(latest_weekend_ddd, 13), 13)

        sql_query = f'''SELECT *FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`ddd_data_for_genie_filtered_dates`
                                WHERE LOWER(drug_name) IN ({",".join([f"'{drug.lower()}'" for drug in drug_list])}) AND week_end_date >= '{p_start_date}'
                               '''
        ddd_sales_data = pd.DataFrame(self.cursor.execute(sql_query).fetchall(), columns=[elem[0] for elem in self.cursor.description])
        ddd_sales_data['week_end_date'] = pd.to_datetime(ddd_sales_data['week_end_date'], errors='coerce')
        ddd_sales_data = ddd_sales_data.rename(columns={'trx_cnt': 'qty_sold_pu'})
        print("ddd sales data sample",ddd_sales_data.head())
        # def calculate_growth(recent_sales, past_sales):
        #     if past_sales == 0:
        #         return float('inf')
        #     growth_percentage = ((recent_sales - past_sales) / past_sales) * 100
        #     return growth_percentage
        # def calculate_growth(recent_sales, past_sales):
        #     if past_sales == 0:
        #         return float('inf')
        #     # Convert both recent_sales and past_sales to float before performing calculations
        #     growth_percentage = ((float(recent_sales) - float(past_sales)) / float(past_sales)) * 100
        #     return growth_percentage

        # def calculate_growth(recent_sales, past_sales):
        #     if past_sales == 0:
        #         return float('inf')
        #     # Ensure both inputs are floats before calculation
        #     recent_sales = float(recent_sales)
        #     past_sales = float(past_sales)
        #     growth_percentage = ((recent_sales - past_sales) / past_sales) * 100
        #     return growth_percentage
        def calculate_growth(recent_sales, past_sales):
            # Avoid division by zero using np.where
            recent_sales = float(recent_sales)
            past_sales = float(past_sales)
            growth_percentage = np.where(past_sales != 0, ((recent_sales - past_sales) / past_sales) * 100, 0)
            return growth_percentage

        def calculate_growth(recent_sales, past_sales):
            # Avoid division by zero using np.where
            recent_sales = float(recent_sales)
            past_sales = float(past_sales)
            growth_percentage = np.where(past_sales != 0, ((recent_sales - past_sales) / past_sales) * 100, np.nan)
            # Convert the result into a Pandas Series and then fill NaN values
            growth_percentage = pd.Series(growth_percentage).fillna(0).replace([np.inf, -np.inf], 0)
            return growth_percentage






        # In[ ]:


        def get_slide_1_data():
            # Validate sales_data
            print("latest weekend in get slide1+data",latest_weekend)
            if 'week_end_date' not in tro_sales_data.columns:
                raise KeyError("'week_end_date' is missing in sales_data!")
            tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date'], errors='coerce')

            # Filter data for the relevant date range
            filtered_data = tro_sales_data[
                (tro_sales_data['week_end_date'] >= '2024-01-05') &
                (tro_sales_data['week_end_date'] <= latest_weekend)
            ]
            if filtered_data.empty:
                raise ValueError("Filtered data is empty. Check the date range or data.")

            # Group by week and compute sales
            weekly_sales = filtered_data.groupby('week_end_date')['qty_sold_pu'].sum().reset_index()
            weekly_sales['formatted_date'] = weekly_sales['week_end_date'].dt.strftime('%d-%b')

            # Calculate R4W averages
            r4w_averages = []
            for current_date in filtered_data['week_end_date'].unique():
                if pd.isnull(current_date):
                    print(f"Skipping invalid date: {current_date}")
                    continue
                start_date = pd.to_datetime(current_date) - timedelta(weeks=3)
                end_date = pd.to_datetime(current_date)
                data_range = tro_sales_data[
                    (tro_sales_data['week_end_date'] >= start_date) &
                    (tro_sales_data['week_end_date'] <= end_date)
                ]
                if data_range.empty:
                    print(f"No data in range {start_date} to {end_date}")
                    continue
                average_sales = data_range['qty_sold_pu'].sum() / 4
                r4w_averages.append({'week_end_date': current_date, 'R4W Average': average_sales})

            # Validate R4W averages
            if not r4w_averages:
                raise ValueError("No data for R4W averages. Check your input data or logic.")
            r4w_averages_df = pd.DataFrame(r4w_averages)
            if 'week_end_date' not in r4w_averages_df.columns:
                raise KeyError("The 'week_end_date' column is missing from r4w_averages_df!")

            # Sort and process data
            r4w_averages_df['week_end_date'] = pd.to_datetime(r4w_averages_df['week_end_date'])
            r4w_averages_df_sorted = r4w_averages_df.sort_values(by='week_end_date')
            r4w_averages_df_sorted['formatted_date'] = r4w_averages_df_sorted['week_end_date'].dt.strftime('%d-%b')

            # Merge weekly sales with R4W averages
            result_df = pd.merge(weekly_sales, r4w_averages_df_sorted, on='formatted_date', how='inner')
            slide_1_chart_estimates['formatted_date'] = pd.to_datetime(slide_1_chart_estimates['formatted_date']).dt.strftime('%d-%b')
            result_df_ = pd.merge(slide_1_chart_estimates, result_df, on='formatted_date', how='left')

            # Final DataFrame
            final_df = result_df_[['formatted_date', 'qty_sold_pu', '24 Budget Weekly Run Rate Estimate', 'R4W Average']]
            final_df.rename(columns={'qty_sold_pu': 'Actuals'}, inplace=True)
            print("slide 1 final df is here",final_df.head())
            return final_df




        def get_slide_1_table_data():
            def get_weekly_average(start_date_= '', end_date_= '', weeks=''):
                if end_date_:
                    end_date = datetime.strptime(end_date_, '%Y-%m-%d').replace(tzinfo=pytz.UTC)
                    start_date = end_date - timedelta(weeks=weeks - 1)
                elif start_date_:
                    start_date = datetime.strptime(start_date_, '%Y-%m-%d').replace(tzinfo=pytz.UTC)
                    end_date = start_date + timedelta(weeks=weeks)
                print("start_date!!!!!", start_date, "end_date!!", end_date)
                start_date = start_date.replace(tzinfo=None)
                end_date = end_date.replace(tzinfo=None)
                mask = (tro_sales_data['week_end_date'] >= start_date) & (tro_sales_data['week_end_date'] <= end_date)
                recent_weeks_data = tro_sales_data.loc[mask, 'qty_sold_pu']
                total_sales = recent_weeks_data.sum()
                weekly_average = round((total_sales / weeks) / 1000,2)
                return weekly_average



            def calculate_total_vials(sales_data, start_date, end_date):
                filter_condition = (sales_data['week_end_date'] >= start_date) & (sales_data['week_end_date'] <= end_date)
                return sales_data.loc[filter_condition, 'qty_sold_pu'].sum() / 1000

            recent_avg_4w = get_weekly_average('', latest_weekend, 4)
            past_avg_4w = get_weekly_average('', get_dates(latest_weekend , 4,''), 4)

            recent_avg_8w = get_weekly_average('', latest_weekend, 8)
            past_avg_8w = get_weekly_average('', get_dates(latest_weekend , 8,''), 8)

            recent_avg_13w = get_weekly_average('', latest_weekend, 13)
            past_avg_13w = get_weekly_average('', get_dates(latest_weekend , 13,''), 13)

            recent_avg_52w = get_weekly_average('', latest_weekend, 52)
            past_avg_52w = get_weekly_average('', get_dates(latest_weekend , 52,''), 52)

            start_date, latest_weekend_ = get_quarter_start_date(latest_weekend)
            weeks_from_quarter_start = get_diff_in_weeks(latest_weekend_, start_date)

            recent_qtr = get_weekly_average('', latest_weekend, weeks_from_quarter_start)

            past_qtr_end_date = get_last_date_before_qtr(start_date, tro_sales_data, 'week_end_date')
            past_qtr_start_date, past_qtr_end_date_ = get_quarter_start_date(past_qtr_end_date)
            print("past_qtr_start_date!!!", past_qtr_start_date)
            past_qtr_end_date_new = get_dates(past_qtr_start_date , '',weeks_from_quarter_start)
            print("past_qtr_end_date_new!!", past_qtr_end_date_new, "weeks_from_quarter_start!!", weeks_from_quarter_start)
            past_qtr = get_weekly_average(past_qtr_start_date,'', weeks_from_quarter_start)

            growth_4w = calculate_growth(recent_avg_4w, past_avg_4w)
            print("before modifying",growth_4w)
            growth_4w=growth_4w[0]
            print("after modifying",growth_4w)
            growth_8w = calculate_growth(recent_avg_8w, past_avg_8w)
            print("before modifying",growth_8w)
            growth_8w=growth_8w[0]
            print("after modifying",growth_8w)
            growth_13w = calculate_growth(recent_avg_13w, past_avg_13w)
            growth_13w = growth_13w[0]
            growth_52w = calculate_growth(recent_avg_52w, past_avg_52w)
            growth_52w = growth_52w[0]
            print("recent_qtr!!!", recent_qtr, "past_qtr!!!", past_qtr)
            growth_qtr = calculate_growth(recent_qtr, past_qtr)
            growth_qtr=growth_qtr[0]

            current_qtr_start, current_qtr_end = get_quarter_start_date(latest_weekend)
            prev_quarter_start, prev_quarter_end = get_previous_quarter_dates(current_qtr_start)
            ytd_est_avg = sum_in_thousands(slide_1_chart_estimates, f'{latest_year}-01-01', latest_weekend)
            current_qtr_est_avg = sum_in_thousands(slide_1_chart_estimates, current_qtr_start, current_qtr_end)
            prev_qtr_est_avg = sum_in_thousands(slide_1_chart_estimates, prev_quarter_start, prev_quarter_end)

            ytd_vials = calculate_total_vials(tro_sales_data, f'{latest_year}-01-01', latest_weekend)
            ytd_growth = calculate_growth(ytd_vials,  ytd_est_avg)
            ytd_growth=ytd_growth[0]
            prev_qtr_vials = calculate_total_vials(tro_sales_data, prev_quarter_start, prev_quarter_end)
            prev_qtr_growth = calculate_growth(prev_qtr_vials, prev_qtr_est_avg )
            prev_qtr_growth=prev_qtr_growth[0]

            current_qtr_vials = calculate_total_vials(tro_sales_data, current_qtr_start, current_qtr_end)
            curr_qtr_growth = calculate_growth(current_qtr_vials, current_qtr_est_avg)
            curr_qtr_growth=curr_qtr_growth[0]
            

            Q1_avg_2024 = get_weekly_average('', f'2024-03-29', 13)
            # def get_max_date_for_quarter(dates, start_month, end_month):
            #     quarter_dates = [date for date in dates if start_month <= date.month <= end_month and date.year == latest_year]
            #     print(f"Filtered dates for months {start_month}-{end_month}: {quarter_dates}")
            #     return max(quarter_dates) if quarter_dates else None
            def get_max_date_for_quarter(dates, start_month, end_month):
                # Convert 'dates' elements to datetime objects
                dates_dt = [pd.to_datetime(d) for d in dates['week_end_date']]
                quarter_dates = [date for date in dates_dt if start_month <= date.month <= end_month and date.year == latest_year]
                print(f"Filtered dates for months {start_month}-{end_month}: {quarter_dates}")
                return max(quarter_dates) if quarter_dates else None
            max_dates = {
            "Q1": get_max_date_for_quarter(sorted_dates, 1, 3),
            "Q2": get_max_date_for_quarter(sorted_dates, 4, 6),
            "Q3": get_max_date_for_quarter(sorted_dates, 7, 9),
            "Q4": get_max_date_for_quarter(sorted_dates, 10, 12),
            }

            if max_dates["Q2"] is None:
                print("No data available for Q2. Skipping calculation.")
                Q2_avg_2024 = 0
            else:
                Q2_avg_2024 = get_weekly_average('', max_dates["Q2"].strftime('%Y-%m-%d'), 13)


            final_output_dict = {
                "R4W": recent_avg_4w , "R8W": recent_avg_8w, "R13W": recent_avg_13w , "R52W": recent_avg_52w,
                "QTD": recent_qtr , "growth_4w": growth_4w , "growth_8w": growth_8w ,
                "growth_13w": growth_13w ,"growth_52w": growth_52w ,"growth_qtr": growth_qtr , "current_qtr_vials": current_qtr_vials,
            "prev_qtr_vials": prev_qtr_vials, "ytd_vials": ytd_vials, "curr_qtr_growth": curr_qtr_growth,
                "prev_qtr_growth": prev_qtr_growth, "ytd_growth": ytd_growth, "q1": Q1_avg_2024, "q2": Q2_avg_2024,
                "prev_qtr_est_avg": prev_qtr_est_avg, "ytd_est_avg": ytd_est_avg, "current_qtr_est_avg": current_qtr_est_avg
            }
            def round_value(val, decimals=2):
                if isinstance(val, pd.Series):
                    # Round the pandas Series
                    return val.round(decimals)
                elif isinstance(val, Decimal):
                    # Round Decimal (convert to float for rounding then back to Decimal if needed)
                    return Decimal(round(float(val), decimals))
                elif isinstance(val, (float, np.floating)):
                    # Round float or numpy float
                    return round(val, decimals)
                elif isinstance(val, np.ndarray):
                    # Check if single element array
                    if val.size == 1:
                        scalar = val.item()
                        return round(scalar, decimals)
                    else:
                        # For arrays with multiple elements, round all elements
                        return np.round(val, decimals)
                else:
                    # Leave other types unchanged
                    print("here in else")
                    return val

            # Keys you want to round
            keys_to_round = ["growth_4w", 'growth_qtr','curr_qtr_growth', 'prev_qtr_growth', 'ytd_growth',]

            for key in keys_to_round:
                if key in final_output_dict:
                     final_output_dict[key] = round_value(final_output_dict[key], 2)
                     print("here is the key and value",key,final_output_dict[key])

            return final_output_dict



        # In[ ]:


        slide_1_data = get_slide_1_data()
        slide_1_table_data = get_slide_1_table_data()
        slide_1_table_data["total_demand"] = f'''Total Demand Vials \n (Data through {latest_day_mon})'''
        print("slide 1 table data is here",slide_1_table_data)
        slide_1_table_data["avg_demand"] = f'''Weekly Avg Demand Vials \n (Data through {latest_day_mon})'''
        slide_1_table_data["qtr_avg_demand"] = f'''Weekly Avg Demand Vials \n (Actuals through {latest_day_mon})'''

        slide_1_table_data


        # In[ ]:


        def get_slide_2_data():
                def get_averages(filtered_data):
                    filtered_data['week_end_date'] = pd.to_datetime(filtered_data['week_end_date']).dt.tz_localize(None)

                    specified_date = pd.Timestamp(latest_weekend).tz_localize(None)
                    start_date_r13w = specified_date - timedelta(weeks=12)

                    end_date_p13w = start_date_r13w - timedelta(days=1)
                    start_date_p13w = end_date_p13w - timedelta(weeks=13)

                    end_date_pp13w = start_date_p13w - timedelta(days=1)
                    start_date_pp13w = end_date_pp13w - timedelta(weeks=13)

                    latest_week_qtr_start_dt, latest_weekend_  = get_quarter_start_date(latest_weekend)
                    start_date_ytd = pd.Timestamp(latest_week_qtr_start_dt).tz_localize(None)
                    end_date_ytd = pd.Timestamp(latest_weekend_).tz_localize(None)

                    r13w_data = filtered_data[(filtered_data['week_end_date'] >= start_date_r13w) & (filtered_data['week_end_date'] <= specified_date)]
                    r13w_total_qty = r13w_data['qty_sold_pu'].sum()
                    r13w_weekly_avg = r13w_total_qty / 13

                    p13w_data = filtered_data[(filtered_data['week_end_date'] >= start_date_p13w) & (filtered_data['week_end_date'] <= end_date_p13w)]
                    p13w_total_qty = p13w_data['qty_sold_pu'].sum()
                    p13w_weekly_avg = p13w_total_qty / 13
                    print("p13w_total_qty!!!", p13w_total_qty)

                    pp13w_data = filtered_data[(filtered_data['week_end_date'] >= start_date_pp13w) & (filtered_data['week_end_date'] <= end_date_pp13w)]
                    pp13w_total_qty = pp13w_data['qty_sold_pu'].sum()
                    pp13w_weekly_avg = pp13w_total_qty / 13

                    p13_growth = calculate_growth(p13w_weekly_avg, pp13w_weekly_avg)
                    r13_growth = calculate_growth(r13w_weekly_avg, p13w_weekly_avg)



                    ytd_data = filtered_data[(filtered_data['week_end_date'] >= start_date_ytd) & (filtered_data['week_end_date'] <= end_date_ytd)
            #                                    &(filtered_data['parent_name'] != "TEXAS ONCOLOGY")
                                            ]
                    ytd_ex_txs = filtered_data[(filtered_data['week_end_date'] >= start_date_ytd) & (filtered_data['week_end_date'] <= end_date_ytd)
                                            &(filtered_data['parent_name'] != "TEXAS ONCOLOGY")
                                            ]

                    ytd_total_qty = ytd_data['qty_sold_pu'].sum()
                    ytd_total_ex_txs = ytd_ex_txs['qty_sold_pu'].sum()

                    num_weeks_ytd = (end_date_ytd - start_date_ytd).days // 7 + 1
                    ytd_weekly_avg = ytd_total_qty / num_weeks_ytd
                    ytd_avg_ex_txs = ytd_total_ex_txs / num_weeks_ytd
                    return [r13w_weekly_avg, p13w_weekly_avg, pp13w_weekly_avg, ytd_avg_ex_txs, p13_growth, r13_growth]

                account_types = ['high','low', 'medium', 'academic', 'community','east','west','None']
                final_output_dict = {}
                final_growth_dict = {}
                for type_ in account_types:
                    print("type_!!!!!!", type_)
                    averages_dict = {}
                    growth_dict = {}
                    if type_ in ['academic', 'community']:
                        filtered_data = tro_sales_data[tro_sales_data['child_account_type'].str.lower() == type_]
                        filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize(None)
                    elif type_ in ['east','west']:
                        filtered_data = tro_sales_data[tro_sales_data['area_nm'].str.lower() == type_]
                        filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize(None)
                    elif type_ =='None':
                        filtered_data=tro_sales_data
                        filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize(None)
                    else:
                        filtered_data = tro_sales_data[tro_sales_data['linkedid_bc_segment'].str.lower() == type_]
                        filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize(None)
                    r13w, p13w, pp13w, ytd_ex_txs, p13_growth, r13_growth = get_averages(filtered_data)
                    averages_dict["PP13W"] = pp13w
                    averages_dict['P13W'] = p13w
                    averages_dict['R13W'] = r13w

                    growth_dict['P13W_growth'] = p13_growth
                    growth_dict['R13W_growth'] = r13_growth

                    # averages_dict["YTD-TX"] = ytd_ex_txs
                    df = pd.DataFrame(list(averages_dict.items()), columns=['time_period', type_])
                    final_output_dict[type_] = df

                    growth_df = pd.DataFrame(list(growth_dict.items()), columns=['time_period', "growth"])
                    final_growth_dict[type_] = growth_df
                return final_output_dict, final_growth_dict


            # In[ ]:


        slide_2_data,  slide_2_growth_data= get_slide_2_data()


        # In[ ]:


        def get_slide_3_data():
            def filter_data_by_date_range(sales_data, start_date, end_date):
                filtered_data = sales_data[(sales_data['week_end_date'] >= start_date) & (sales_data['week_end_date'] <= end_date)]
                return filtered_data

            def calculate_sales_metrics(filtered_data, col_name):
                total_sales_by_parent = filtered_data.groupby('parent_name')['qty_sold_pu'].sum().reset_index()
                total_weeks = filtered_data['week_end_date'].nunique()
                total_sales_by_parent['average_weekly_sales_'+ col_name] = total_sales_by_parent['qty_sold_pu'] / total_weeks
                total_sales = total_sales_by_parent['qty_sold_pu'].sum()
                total_sales_by_parent['sales_share_'+ col_name] = (total_sales_by_parent['qty_sold_pu'] / total_sales) * 100
                return total_sales_by_parent

            # def calculate_growth_percentage(df , old_col, new_col, df_col):
            #     df[df_col] = (((df[ new_col] - df[ old_col])*100) / df[old_col])
            #     return df

            # def calculate_growth_percentage(df, old_col, new_col, df_col):
            #     # Avoid division by zero by replacing zeros in the denominator with a small value
            #     df[df_col] = (((df[new_col] - df[old_col]) * 100) / df[old_col].replace(0, 0.0001))
            #     # Alternatively, you can use np.where to handle zero values:
            #     # df[df_col] = np.where(df[old_col] != 0, (((df[new_col] - df[old_col]) * 100) / df[old_col]), 0)
            #     return df

            # def calculate_growth_percentage(df, old_col, new_col, df_col):
            #     # Convert both columns to float to ensure compatibility
            #     df[old_col] = df[old_col].astype(float)
            #     df[new_col] = df[new_col].astype(float)
            #     # Avoid division by zero by replacing zeros in the denominator with a small value
            #     df[df_col] = (((df[new_col] - df[old_col]) * 100) / df[old_col].replace(0, 0.0001))
            #     # Alternatively, you can use np.where to handle zero values:
            #     # df[df_col] = np.where(df[old_col] != 0, (((df[new_col] - df[old_col]) * 100) / df[old_col]), 0)
            #     return df
    #         def calculate_growth_percentage(df, old_col, new_col, df_col):
    # # Convert both columns to float to ensure compatibility
    #             df[old_col] = df[old_col].astype(float)
    #             df[new_col] = df[new_col].astype(float)
    #             # Avoid division by zero using np.where
    #             df[df_col] = np.where(df[old_col] != 0, (((df[new_col] - df[old_col]) * 100) / df[old_col]), 0)
    #             return df
            def calculate_growth_percentage(df, old_col, new_col, df_col):
              # Avoid division by zero by replacing zeros in the denominator with a small value
              # df[df_col] = (((df[new_col] - df[old_col]) * 100) / df[old_col].replace(0, 1e-10))
              # Convert 'old_col' to numeric, handling errors
              df[old_col] = pd.to_numeric(df[old_col], errors='coerce')
              # Convert 'new_col' to numeric, handling errors
              df[new_col] = pd.to_numeric(df[new_col], errors='coerce')
              # Avoid division by zero by replacing zeros in the denominator with a small value
              df[df_col] = (((df[new_col] - df[old_col]) * 100) / df[old_col].replace(0, 1e-10))
              return df

            def process_data(sales_data, start_date, end_date, col_name):
                filtered_data = filter_data_by_date_range(sales_data, start_date, end_date )
                return calculate_sales_metrics(filtered_data, col_name)

            def get_quarter_data(sales_data, start_date, end_date, column_name):
                quarter_data = filter_data_by_date_range(sales_data, start_date, end_date)
                quarter_sales = quarter_data.groupby('parent_name')['qty_sold_pu'].sum().reset_index()
                quarter_sales.rename(columns={'qty_sold_pu': column_name}, inplace=True)
                return quarter_sales

            def calculate_summary_row(df, label):
                sums = df.drop(columns=['parent_name']).sum()
                new_row_data = {'parent_name': label, **sums}
                return pd.DataFrame([new_row_data])

            def get_extra_rows(row_name, df):
                print("columns!!!!", df.columns)
                sum_columns = ['qty_sold_pu_current', 'average_weekly_sales_current','sales_share_current', 'qty_sold_pu_prev', 'average_weekly_sales_prev',
                            'sales_share_prev', 'qty_sold_pu_P13w', 'qty_sold_pu_R13w','average_weekly_sales_R13W_current', 'sales_share_R13W_current','average_weekly_sales_P13W_prev',
                            'qty_sold_pu_P_ytd', 'average_weekly_sales_P_qtd_current','sales_share_P13W_prev','sales_share_P_qtd_current', 'qty_sold_pu_R_ytd','average_weekly_sales_R_qtd_current', 'sales_share_R_qtd_current']
                growth_pairs = {"R12Mgrowth": ('qty_sold_pu_prev', 'qty_sold_pu_current'),
                            'R13Wgrowth': ('qty_sold_pu_P13w', 'qty_sold_pu_R13w'), "YTD_growth": ('qty_sold_pu_P_ytd', 'qty_sold_pu_R_ytd')
                            }
                new_row = {}
                for col in sum_columns:
                    new_row[col] = df[col].sum()
                for key, (from_col, to_col) in growth_pairs.items():
                    if df[from_col].sum() != 0:
                        new_row[key] = ((df[to_col].sum() / df[from_col].sum()) - 1) * 100
                    else:
                        new_row[key] = None
                new_row["parent_name"] = row_name
                new_df_ls.append(new_row)


            current_start_date, current_end_date = get_date_range(latest_weekend, 52)
            print( current_start_date, current_end_date)
            current_sales_by_parent = process_data(tro_sales_data, current_start_date, current_end_date, "current")
            current_sales_by_parent.rename(columns={'qty_sold_pu': 'qty_sold_pu_current'}, inplace=True)

            print("previous end date!!" , get_dates(latest_weekend , 52))
            prev_start_date, prev_end_date = get_date_range(get_dates(latest_weekend , 52), 52)
            prev_sales_by_parent = process_data(tro_sales_data, prev_start_date, prev_end_date, "prev")
            prev_sales_by_parent.rename(columns={'qty_sold_pu': 'qty_sold_pu_prev'}, inplace=True)
            print( "prev_start_date, prev_end_date!!",prev_start_date, prev_end_date)


            merged_sales = pd.merge(current_sales_by_parent, prev_sales_by_parent, on='parent_name', how= "outer").fillna(0)
            merged_sales.replace([float('inf'), -float('inf')], 0, inplace=True)
            top_10_parent_accounts = merged_sales.sort_values(by='qty_sold_pu_current', ascending=False)
            merged_sales_data = calculate_growth_percentage(top_10_parent_accounts, 'qty_sold_pu_prev', 'qty_sold_pu_current', "R12Mgrowth")

            latest_week_qtr_start_dt, latest_weekend_  = get_quarter_start_date(latest_weekend)
            R13_start_date_, R13_end_date_ = get_date_range(latest_weekend, 13)
            R13_sales_by_parent_ = process_data(tro_sales_data, R13_start_date_, R13_end_date_, "R13W_current")
            R13_sales_by_parent_.rename(columns={'qty_sold_pu': 'qty_sold_pu_R13w'}, inplace=True)

            past_13w_end_date = get_dates(latest_weekend , 13)
            print("past_13w_end_date!!!!!", past_13w_end_date)
            P13_start_date_, P13_end_date_ = get_date_range(past_13w_end_date, 13)
            P13_sales_by_parent_ = process_data(tro_sales_data, P13_start_date_, P13_end_date_, "P13W_prev")
            P13_sales_by_parent_.rename(columns={'qty_sold_pu': 'qty_sold_pu_P13w'}, inplace=True)

            w_13_merged_sales = pd.merge(P13_sales_by_parent_, R13_sales_by_parent_, on='parent_name', how= "outer").fillna(0)
            w_13_growth = calculate_growth_percentage(w_13_merged_sales, 'qty_sold_pu_P13w', 'qty_sold_pu_R13w', "R13Wgrowth")
            sub_df = pd.merge(merged_sales_data, w_13_merged_sales, on='parent_name', how= "outer").fillna(0)

            weeks_diff = get_diff_in_weeks( latest_weekend_, f'{latest_year}-01-01')
            print("weeks_diff!!!", weeks_diff)
            P_ytd_start_date = f'{latest_year-1}-01-01'
            previous_year = datetime.strptime(P_ytd_start_date, '%Y-%m-%d')
            date_after_weeks_diff = previous_year + timedelta(weeks= weeks_diff)
            P_ytd_end_date = date_after_weeks_diff.strftime('%Y-%m-%d')
            print("P_ytd_start_date, P_ytd_end_date!!!!!", P_ytd_start_date, P_ytd_end_date)

            P_ytd_sales= process_data(tro_sales_data, P_ytd_start_date, P_ytd_end_date, "P_qtd_current")
            print("P_ytd_sales!!!", P_ytd_sales)
            P_ytd_sales.rename(columns={'qty_sold_pu': 'qty_sold_pu_P_ytd'}, inplace=True)


            r_ytd_start_date, r_ytd_end_date = get_date_range(latest_weekend_, weeks_diff)
            R_ytd_sales= process_data(tro_sales_data, r_ytd_start_date, r_ytd_end_date, "R_qtd_current")
            R_ytd_sales.rename(columns={'qty_sold_pu': 'qty_sold_pu_R_ytd'}, inplace=True)

            Ytd_merged_sales = pd.merge(P_ytd_sales, R_ytd_sales, on='parent_name', how= "outer").fillna(0)
            Ytd_growth = calculate_growth_percentage(Ytd_merged_sales, 'qty_sold_pu_P_ytd', 'qty_sold_pu_R_ytd', "YTD_growth")

            final_df = pd.merge(sub_df, Ytd_growth, on='parent_name', how= "outer").fillna(0)
            final_df.replace([float('inf'), -float('inf')], 0, inplace=True)
            final_df = final_df.sort_values(by='qty_sold_pu_current', ascending=False)
            final_df = final_df.fillna(0)
            un_fil_top_15 = final_df.head(15)
            un_fil_top_10 = final_df.head(10)

            filtered_df = final_df[final_df['parent_name'].str.upper() != 'TEXAS ONCOLOGY']
            sorted_df = filtered_df.sort_values(by='qty_sold_pu_current', ascending=False)
            top_10 = sorted_df.head(10)
            top_15 = sorted_df.head(15)
            new_df_ls = []
            df_row = {'Top 10 - TOPA': top_10, 'Top 15 - TOPA': top_15, "Top 15 Parent Accounts":  un_fil_top_15,"Top 10 Parent Accounts":  un_fil_top_10,  "Nation - TOPA": sorted_df, "Nation": final_df}
        #     print("df_row!!", df_row)
            for key, value in df_row.items():
                get_extra_rows(key, value)

            calc_df =pd.DataFrame(new_df_ls)
            slide_df = pd.concat([final_df,calc_df], ignore_index=True)
            final_list = ['parent_name', 'sales_share_current', 'average_weekly_sales_current','R12Mgrowth',
            'average_weekly_sales_R13W_current', 'R13Wgrowth','YTD_growth']
            percent_col = ["sales_share_current", "R12Mgrowth", "R13Wgrowth", "YTD_growth"]
            new_final_df = pd.concat([slide_df.head(15), slide_df.tail(6) ])
            new_final_df = new_final_df[final_list]
            new_final_df.fillna(0, inplace=True)

            for column in new_final_df.columns:
                if pd.api.types.is_numeric_dtype(new_final_df[column]):
                    if column in percent_col:
                        new_final_df[column] = new_final_df[column].round(0).astype(int).astype(str) + '%'
                    else:
                        print("column!!!", column)
                        new_final_df[column] = new_final_df[column].round(0).astype(int).astype(str)

            columns_to_round = ['sales_share_current', 'average_weekly_sales_current', 'average_weekly_sales_R13W_current']

            for col in columns_to_round:
                new_final_df[col] = pd.to_numeric(new_final_df[col], errors='coerce')
            new_final_df[columns_to_round] = new_final_df[columns_to_round].applymap(np.ceil)
            return new_final_df





        # In[ ]:


        slide_3_table_data = get_slide_3_data()


        # In[ ]:
        print("Before 4th Slide")
        def get_slide_4_5_data(account_type):
            # tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date'], errors='coerce')
            R13_start_date_, R13_end_date_ = get_date_range(latest_weekend, 13)
            print("R13_start_date_, R13_end_date_!!!", R13_start_date_, R13_end_date_)
            past_13w_end_date = get_dates(latest_weekend , 13)
            print("past_13w_end_date!!!!!", past_13w_end_date)
            P13_start_date_, P13_end_date_ = get_date_range(past_13w_end_date, 13)
        #     print("P13_start_date_, P13_end_date_, f'{latest_year}-01-01', latest_weekend!!!", P13_start_date_, P13_end_date_, f'{latest_year}-01-01', latest_weekend)
            account_mapping = {
                'academic': ('top_parent_account_type', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
                'community': ('top_parent_account_type', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
                'idn/hospital': ('top_parent_account_type', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
                'high': ('parent_bc_segment',  R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
                'medium': ('parent_bc_segment',  R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_)
            }
            column, start_date1, end_date1, start_date2, end_date2 = account_mapping[account_type]
            if tro_sales_data['week_end_date'].dt.tz is None:
                tro_sales_data['week_end_date'] = tro_sales_data['week_end_date'].dt.tz_localize('UTC')
            else:
                tro_sales_data['week_end_date'] = tro_sales_data['week_end_date'].dt.tz_convert('UTC')

            # sql_query = f'''
            #     SELECT *
            #     FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`867_data_drug-trodelvy_for_genie_filtered_dates`'''
            # tro_sales_data = pd.DataFrame(self.cursor.execute(sql_query).fetchall(), columns=[elem[0] for elem in self.cursor.description])
            # tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date'], errors='coerce')
            # if filtered_data['week_end_date'].dt.tz is None:
            #     filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize('UTC')
            # else:
            #     filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_convert('UTC')
            # if tro_sales_data['week_end_date'].dt.tz is None:
            #     tro_sales_data['week_end_date'] = tro_sales_data['week_end_date'].dt.tz_localize('UTC')
            # else:
            #     tro_sales_data['week_end_date'] = tro_sales_data['week_end_date'].dt.tz_convert('UTC')
            filtered_data = tro_sales_data[tro_sales_data[column].str.lower() == account_type]
            print("Filtered data in slide 4 and 5 ",filtered_data)
            # filtered_data = tro_sales_data[tro_sales_data[column].str.lower() == account_type]
            R13W = (pd.Timestamp(start_date1, tz='UTC'), pd.Timestamp(end_date1, tz='UTC'))
            Q124 = (pd.Timestamp(start_date2, tz='UTC'), pd.Timestamp(end_date2, tz='UTC'))

            print("Reached Before Calculate Average Sales")

            # Convert 'week_end_date' to datetime, if not already in datetime format
            filtered_data['week_end_date'] = pd.to_datetime(filtered_data['week_end_date'], errors='coerce')

            # Check if the datetime is timezone-aware
            if filtered_data['week_end_date'].dt.tz is None:
                filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize('UTC')
            else:
                filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_convert('UTC')


            def calculate_avg_sales(filtered_data, date_range):
                # Ensure 'week_end_date' is in datetime format
                filtered_data['week_end_date'] = pd.to_datetime(filtered_data['week_end_date'], errors='coerce')

                # Check and localize 'week_end_date' to UTC if it's not already localized
                if filtered_data['week_end_date'].dt.tz is None:
                    filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize('UTC')
                else:
                    filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_convert('UTC')

                # Filter data based on the range
                date_filtered_data = filtered_data[(filtered_data['week_end_date'] >= date_range[0]) & (filtered_data['week_end_date'] <= date_range[1])]

                # Calculate average sales
                avg_sales = date_filtered_data.groupby('parent_name')['qty_sold_pu'].sum() / date_filtered_data['week_end_date'].nunique()
                print("here is the avg sales",avg_sales.head())
                return avg_sales.reset_index().rename(columns={'qty_sold_pu': 'avg_sales'})

            
            # Assuming `filtered_data` is already defined
            R13w_avg_sales = calculate_avg_sales(filtered_data, R13W).rename(columns={'avg_sales': 'R13W_avg_sales'})
            Q124_TD_avg_sales = calculate_avg_sales(filtered_data, Q124).rename(columns={'avg_sales': 'P13W_avg_sales'})
            merged_data = pd.merge(R13w_avg_sales, Q124_TD_avg_sales, on='parent_name', how='outer')
            merged_data['absolute_volume'] = (merged_data['P13W_avg_sales'] - merged_data['R13W_avg_sales']).abs()
            merged_data['R13W_avg_sales'] = merged_data['R13W_avg_sales'].fillna(0)  # Replace NaN with 0 or another suitable value
            merged_data['P13W_avg_sales'] = merged_data['P13W_avg_sales'].fillna(0)
            # Round up the sales data using np.ceil
            merged_data['R13W_avg_sales'] = np.ceil(merged_data['R13W_avg_sales']).astype(float)
            merged_data['P13W_avg_sales'] = np.ceil(merged_data['P13W_avg_sales']).astype(float)
            # Fixing the rounding issue by converting columns to float before rounding

            # Calculate growth and round it up
            merged_data['sales_growth_percentage'] = ((merged_data['P13W_avg_sales'] - merged_data['R13W_avg_sales']) / merged_data['R13W_avg_sales'].replace(0, np.inf)).fillna(0) * 100
            merged_data['sales_growth_percentage'] = np.ceil(merged_data['sales_growth_percentage']).astype(int).astype(str) + '%'
            # merged_data['sales_growth_percentage'] = round(((merged_data['P13W_avg_sales'].astype(float) - merged_data['R13W_avg_sales'].astype(float)) / merged_data['R13W_avg_sales'].replace(0, np.inf).astype(float)) * 100, 2)

            merged_data = merged_data.sort_values(by='absolute_volume', ascending=False).head(10)
            merged_data = merged_data.drop(columns=['absolute_volume'])

            percent_col = ["sales_growth_percentage"]
            for column in merged_data.columns:
                if pd.api.types.is_numeric_dtype(merged_data[column]):
                    if column in percent_col:
                        merged_data[column] = merged_data[column].round(0).astype(int).astype(str) + '%'
                    else:
                        merged_data[column] = merged_data[column].round(0).astype(int).astype(str)
            print("merged data in slide 4 and 5 ",account_type ,merged_data)
            return merged_data



        # def get_slide_4_5_data(account_type):
        #     tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date'], errors='coerce')
        #     R13_start_date_, R13_end_date_ = get_date_range(latest_weekend, 13)
        #     print("R13_start_date_, R13_end_date_!!!", R13_start_date_, R13_end_date_)
        #     past_13w_end_date = get_dates(latest_weekend , 13)
        #     print("past_13w_end_date!!!!!", past_13w_end_date)
        #     P13_start_date_, P13_end_date_ = get_date_range(past_13w_end_date, 13)
        # #     print("P13_start_date_, P13_end_date_, f'{latest_year}-01-01', latest_weekend!!!", P13_start_date_, P13_end_date_, f'{latest_year}-01-01', latest_weekend)
        #     account_mapping = {
        #         'academic': ('top_parent_account_type', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
        #         'community': ('top_parent_account_type', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
        #         'high': ('linkedid_bc_segment', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
        #         'idn/hospital': ('child_account_type', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
        #         'low': ('linkedid_bc_segment', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
        #         'medium': ('linkedid_bc_segment', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),

        #     }


        #     def get_weekly_average_sales(sales_data, start_date, end_date, group_by_column):
        #         # Filter sales data within the specified date range
        #         filtered_data = sales_data[(sales_data['week_end_date'] >= start_date) & (sales_data['week_end_date'] <= end_date)]

        #         # Group by specified column and calculate total sales
        #         total_sales_by_parent = filtered_data.groupby(group_by_column)['qty_sold_pu'].sum().reset_index()

        #         # Calculate the number of weeks within the date range
        #         total_weeks = filtered_data['week_end_date'].nunique()

        #         # Calculate average weekly sales
        #         total_sales_by_parent['avg_sales'] = total_sales_by_parent['qty_sold_pu'] / total_weeks

        #         # Rename columns for clarity
        #         total_sales_by_parent.rename(columns={'qty_sold_pu': 'total_sales', group_by_column: 'parent_name'}, inplace=True)

        #         return total_sales_by_parent

        #     group_by_column, R13_start_date, R13_end_date, P13_start_date, P13_end_date = account_mapping[account_type]
        #     R13w_avg_sales = get_weekly_average_sales(tro_sales_data, R13_start_date, R13_end_date, group_by_column)
        #     R13w_avg_sales.rename(columns={'avg_sales': 'R13W_avg_sales'}, inplace=True)
        #     # print("R13w_avg_sales!!!", R13w_avg_sales)
        #     P13w_avg_sales = get_weekly_average_sales(tro_sales_data, P13_start_date, P13_end_date, group_by_column)
        #     P13w_avg_sales.rename(columns={'avg_sales': 'P13W_avg_sales'}, inplace=True)

        #     # Q124_TD_avg_sales = get_weekly_average_sales(tro_sales_data, f'{latest_year}-01-01', latest_weekend, group_by_column)
        #     # Q124_TD_avg_sales.rename(columns={'avg_sales': 'Q124_TD_avg_sales'}, inplace=True)

        #     merged_data = pd.merge(R13w_avg_sales, P13w_avg_sales, on='parent_name', how='outer')
        #     # print("merged_data!!!", merged_data)
        #     # merged_data = pd.merge(round(R13w_avg_sales), round(Q124_TD_avg_sales), on='parent_name', how='outer')

        #     merged_data['absolute_volume'] = (merged_data['P13W_avg_sales'] - merged_data['R13W_avg_sales']).abs()
        #     # Convert 'R13W_avg_sales' to float
        #     merged_data['R13W_avg_sales'] = merged_data['R13W_avg_sales'].astype(float)
        #     # Now calculate the percentage

        #     merged_data['P13W_avg_sales'] = merged_data['P13W_avg_sales'].astype(float)
        #     merged_data['R13W_avg_sales'] = merged_data['R13W_avg_sales'].astype(float)
        #     merged_data['sales_growth_percentage'] = round(((merged_data['P13W_avg_sales'] - merged_data['R13W_avg_sales']) / merged_data['R13W_avg_sales'].replace(0, np.inf)) * 100)
        #     merged_data = merged_data.sort_values(by='absolute_volume', ascending=False).head(10)
        #     merged_data = merged_data.drop(columns=['absolute_volume'])
        #     # merged_data = merged_data.sort_values(by='sales_growth_percentage', ascending=True)

        #     return merged_data
        # def get_slide_4_5_data(account_type):
        #     tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date'], errors='coerce')
        #     R13_start_date_, R13_end_date_ = get_date_range(latest_weekend, 13)
        #     print("R13_start_date_, R13_end_date_!!!", R13_start_date_, R13_end_date_)
        #     past_13w_end_date = get_dates(latest_weekend , 13)
        #     print("past_13w_end_date!!!!!", past_13w_end_date)
        #     P13_start_date_, P13_end_date_ = get_date_range(past_13w_end_date, 13)
        # #     print("P13_start_date_, P13_end_date_, f'{latest_year}-01-01', latest_weekend!!!", P13_start_date_, P13_end_date_, f'{latest_year}-01-01', latest_weekend)
        #     account_mapping = {
        #         'academic': ('top_parent_account_type', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
        #         'community': ('top_parent_account_type', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
        #         'idn/hospital': ('top_parent_account_type', R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
        #         'high': ('parent_bc_segment',  R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_),
        #         'medium': ('parent_bc_segment',  R13_start_date_, R13_end_date_, P13_start_date_, P13_end_date_)

        #     }
        #     column, start_date1, end_date1, start_date2, end_date2 = account_mapping[account_type]
        #     if tro_sales_data['week_end_date'].dt.tz is None:
        #         tro_sales_data['week_end_date'] = tro_sales_data['week_end_date'].dt.tz_localize('UTC')
        #     else:
        #         tro_sales_data['week_end_date'] = tro_sales_data['week_end_date'].dt.tz_convert('UTC')
        #     filtered_data = tro_sales_data[tro_sales_data[column].str.lower() == account_type]
        #     R13W = (pd.Timestamp(start_date1, tz='UTC'), pd.Timestamp(end_date1, tz='UTC'))
        #     Q124 = (pd.Timestamp(start_date2, tz='UTC'), pd.Timestamp(end_date2, tz='UTC'))

        #     def calculate_avg_sales(filtered_data, date_range):
        #         date_filtered_data = filtered_data[(filtered_data['week_end_date'] >= date_range[0]) & (filtered_data['week_end_date'] <= date_range[1])]
        #         avg_sales = date_filtered_data.groupby('parent_name')['qty_sold_pu'].sum() / date_filtered_data['week_end_date'].nunique()
        #         return avg_sales.reset_index().rename(columns={'qty_sold_pu': 'avg_sales'})

        #     R13w_avg_sales = calculate_avg_sales(filtered_data,  R13W).rename(columns={'avg_sales': 'R13W_avg_sales'})
        #     Q124_TD_avg_sales = calculate_avg_sales(filtered_data, Q124).rename(columns={'avg_sales': 'P13W_avg_sales'})
        #     merged_data = pd.merge(round(R13w_avg_sales), round(Q124_TD_avg_sales), on='parent_name', how='outer')
        #     merged_data['absolute_volume'] = (merged_data['P13W_avg_sales'] - merged_data['R13W_avg_sales']).abs()
        #     merged_data['sales_growth_percentage'] = round(((merged_data['P13W_avg_sales'] - merged_data['R13W_avg_sales']) / merged_data['R13W_avg_sales'].replace(0, np.inf)) * 100)
        #     merged_data = merged_data.sort_values(by='absolute_volume', ascending=False).head(10)
        #     merged_data = merged_data.drop(columns=['absolute_volume'])
        #     percent_col = ["sales_growth_percentage"]
        #     for column in merged_data.columns:
        #         if pd.api.types.is_numeric_dtype(merged_data[column]):
        #             if column in percent_col:
        #                 merged_data[column] = merged_data[column].round(0).astype(int).astype(str) + '%'
        #             else:
        #                 print("column!!!", column)
        #                 merged_data[column] = merged_data[column].round(0).astype(int).astype(str)
        #     return merged_data


        # In[ ]:


        high_acc_decliners = get_slide_4_5_data('high')
        med_acc_decliners = get_slide_4_5_data('medium')
        academic_acc_decliners = get_slide_4_5_data('academic')
        community_acc_decliners = get_slide_4_5_data('community')
        idn_hospital_decliners = get_slide_4_5_data('idn/hospital')
        idn_hospital_decliners


        def get_slide_6_7_east_west_data(region):
            print("region!!!!!!!", region)
            if region:
                filter_by = 'regn_nm'
            else:
                filter_by = 'area_nm'
            def get_extra_rows(row_name, df):
                new_row = {}
                for col in df.columns[1:]:
                    new_row[col] = df[col].sum()
                new_row[filter_by] = row_name
                df2 = pd.DataFrame([new_row])
                concatenated_df = pd.concat([df, df2], ignore_index=True)
                return concatenated_df

            def analyze_sales_growth_region(recent_start, recent_end, past_start, past_end, region_name, period, account_type = ''):

                if account_type and region:
                    recent_df = tro_sales_data[(tro_sales_data['week_end_date'] >= recent_start) & (tro_sales_data['week_end_date'] <= recent_end) &(tro_sales_data['area_nm'] == region) & (tro_sales_data['top_parent_account_type'].str.lower() == account_type)]
                    past_df = tro_sales_data[(tro_sales_data['week_end_date'] >= past_start) & (tro_sales_data['week_end_date'] <= past_end) & (tro_sales_data['area_nm'] == region) & (tro_sales_data['top_parent_account_type'].str.lower() == account_type)]
                elif region:
                    recent_df = tro_sales_data[(tro_sales_data['week_end_date'] >= recent_start) & (tro_sales_data['week_end_date'] <= recent_end) & (tro_sales_data['area_nm'] == region)]
                    past_df = tro_sales_data[(tro_sales_data['week_end_date'] >= past_start) & (tro_sales_data['week_end_date'] <= past_end) &(tro_sales_data['area_nm'] == region)]
                elif account_type and region == '':
                    recent_df = tro_sales_data[(tro_sales_data['week_end_date'] >= recent_start) & (tro_sales_data['week_end_date'] <= recent_end) & (tro_sales_data['top_parent_account_type'].str.lower() == account_type)]
                    past_df = tro_sales_data[(tro_sales_data['week_end_date'] >= past_start) & (tro_sales_data['week_end_date'] <= past_end)& (tro_sales_data['top_parent_account_type'].str.lower() == account_type)]
                else:
                    recent_df = tro_sales_data[(tro_sales_data['week_end_date'] >= recent_start) & (tro_sales_data['week_end_date'] <= recent_end) ]
                    past_df = tro_sales_data[(tro_sales_data['week_end_date'] >= past_start) & (tro_sales_data['week_end_date'] <= past_end)]

                recent_sales = recent_df.groupby(filter_by)['qty_sold_pu'].sum().reset_index()
                past_sales = past_df.groupby(filter_by)['qty_sold_pu'].sum().reset_index()
                merged_df = pd.merge(recent_sales, past_sales, on=filter_by, suffixes=(f'_recent_{period}', f'_past_{period}'))
                if region_name == '':
                    print("in region name empty!!!")
                    merged_df_= get_extra_rows("Nation",merged_df)
                else:
                    merged_df_= get_extra_rows(region_name,merged_df)
                merged_df_[f'growth_percentage_{period}'] = ((merged_df_[f'qty_sold_pu_recent_{period}'] - merged_df_[f'qty_sold_pu_past_{period}']) / merged_df_[f'qty_sold_pu_past_{period}']) * 100
        #         print(merged_df_.head())
                return merged_df_

            def analyze_and_store(recent_start, recent_end, past_start, past_end, region, name, account_type=None):
                df = analyze_sales_growth_region(recent_start, recent_end, past_start, past_end, region, name, account_type)
                print("df!!!!!!!!!", df)
                return df

            results = {}
            R4w_start_date, R4w_end_date = get_date_range(latest_weekend, 4)
            P4w_start_date, P4w_end_date = get_date_range(get_dates(latest_weekend , 4), 4)
            print("R4w_start_date, R4w_end_date!!!", R4w_start_date, R4w_end_date, "    P4w_start_date, P4w_end_date!!",     P4w_start_date, P4w_end_date)

            R13w_start_date, R13w_end_date = get_date_range(latest_weekend, 13)
            P13w_start_date, P13w_end_date = get_date_range(get_dates(latest_weekend , 13), 13)
            print("R13w_start_date, R13w_end_date!!!", R13w_start_date, R13w_end_date, "    P13w_start_date, P13w_end_date!!",     P13w_start_date, P13w_end_date)

            results['R4w_p4w'] = analyze_and_store(R4w_start_date, R4w_end_date, P4w_start_date, P4w_end_date, region, 'R4w_p4w')
            results['R13W_P13W'] = analyze_and_store(R13w_start_date, R13w_end_date, P13w_start_date, P13w_end_date, region, 'R13W_P13W')
#             results['R13W_Q1TD'] = analyze_and_store('2024-01-05', '2024-03-29', '2023-09-29', '2023-12-22', region, 'R13W_Q1TD')


            account_types = ['academic', 'community', 'idn/hospital']
            for account_type in account_types:
                results[f'{account_type}_4w'] = analyze_and_store(R4w_start_date, R4w_end_date, P4w_start_date, P4w_end_date, region, f'{account_type}_4w', account_type)
                results[f'{account_type}_13w'] = analyze_and_store(R13w_start_date, R13w_end_date, P13w_start_date, P13w_end_date, region, f'{account_type}_13w', account_type)
#                 results[f'{account_type}_13w_q1td'] = analyze_and_store('2024-01-05', '2024-03-29', '2023-09-29', '2023-12-22', region, f'{account_type}_13w_q1td', account_type)

            dfs = [results['R4w_p4w'],  results['R13W_P13W'],
                results['academic_4w'], results['academic_13w'],
                results['community_4w'], results['community_13w'],
                results['idn/hospital_4w'], results['idn/hospital_13w']]
#     results['R13W_Q1TD'],results['academic_13w_q1td'],results['community_13w_q1td'], results['idn/hospital_13w_q1td']

            merged_df = reduce(lambda left, right: pd.merge(left, right, on=filter_by), dfs)
            merged_df['Volume_Contribution_Academic'] = (merged_df['qty_sold_pu_recent_academic_13w'] / (merged_df['qty_sold_pu_recent_academic_13w']+merged_df['qty_sold_pu_recent_community_13w']+merged_df['qty_sold_pu_recent_idn/hospital_13w'])) * 100
            merged_df['Volume_Contribution_Community'] = (merged_df['qty_sold_pu_recent_community_13w'] / (merged_df['qty_sold_pu_recent_academic_13w']+merged_df['qty_sold_pu_recent_community_13w']+merged_df['qty_sold_pu_recent_idn/hospital_13w'])) * 100
            merged_df['Volume_Contribution_IDN'] = (merged_df['qty_sold_pu_recent_idn/hospital_13w'] / (merged_df['qty_sold_pu_recent_academic_13w']+merged_df['qty_sold_pu_recent_community_13w']+merged_df['qty_sold_pu_recent_idn/hospital_13w'])) * 100
            merged_df[[filter_by,'Volume_Contribution_Academic', 'Volume_Contribution_Community', 'Volume_Contribution_IDN' ]]
            merged_df = merged_df[[filter_by,'growth_percentage_R4w_p4w', 'growth_percentage_R13W_P13W',
                    'growth_percentage_academic_4w', 'growth_percentage_academic_13w',
                    'growth_percentage_community_4w', 'growth_percentage_community_13w',
                    'growth_percentage_idn/hospital_4w', 'growth_percentage_idn/hospital_13w',
                    'Volume_Contribution_Academic', 'Volume_Contribution_Community', 'Volume_Contribution_IDN']]
#             'growth_percentage_R13W_Q1TD',
#             merged_df.to_excel("merged_df_east.xlsx")
#             merged_df.to_excel("merged_df_west.xlsx")

            if region == "East":
                order = ['East', 'New England', 'New York', 'Atlantic Coastal', 'Blue Grass', 'Eastern Coastal', 'Atlanta', 'Southeast']
            elif region == "West":
                order = ['West', 'Pacific Northwest', 'SoCAL', 'Rocky Mountains', 'Gulf Plains', 'Texas', 'North Central', 'Great Lakes']
            else:
                order = ["Nation", 'East', 'West']

            merged_df = merged_df.set_index(filter_by).loc[order].reset_index()
            columns_to_round = ['growth_percentage_R4w_p4w','growth_percentage_R13W_P13W','growth_percentage_academic_4w','growth_percentage_academic_13w','growth_percentage_community_4w','growth_percentage_community_13w','growth_percentage_idn/hospital_4w','growth_percentage_idn/hospital_13w','Volume_Contribution_Academic','Volume_Contribution_Community','Volume_Contribution_IDN']
            for col in columns_to_round:
                merged_df[col] = pd.to_numeric(merged_df[col], errors='coerce')

            merged_df[columns_to_round] = merged_df[columns_to_round].applymap(np.ceil)

            print("Rounded specific columns in merged df:")
            print(merged_df.head())

            region_vials_growth = merged_df[[filter_by ,'growth_percentage_R13W_P13W','growth_percentage_R4w_p4w']]
            region_vials_growth_ = region_vials_growth.rename(columns={'growth_percentage_R13W_P13W': 'R13W', 'growth_percentage_R4w_p4w': 'R4W'})
            region_vials_growth_['R13W']= region_vials_growth_['R13W']/100
            region_vials_growth_['R4W']= region_vials_growth_['R4W']/100
            print("info of merged df: ",merged_df.info())
            # merged_df_rounded = merged_df.applymap(lambda x: str(round(x)) + "%" if isinstance(x, (int, float)) else x)
            merged_df_rounded = merged_df.copy()

            # Apply the transformation on numeric columns only
            for col in merged_df_rounded.select_dtypes(include=['float64', 'int64']).columns:
                merged_df_rounded[col] = merged_df_rounded[col].apply(lambda x: str(round(x)) + "%" if isinstance(x, (int, float)) else x)

            # Display the rounded DataFrame
            print(merged_df_rounded.head())

            print("merged_df_rounded: ",merged_df_rounded.head())

            # merged_df_rounded = merged_df.applymap(lambda x: str(round(x)) + "%" if isinstance(x, (int, float)) else x)
            def combine_columns(df, col1, col2, col3, new_col_name):
                df[new_col_name] = df[col1].astype(str)+ ' | ' + df[col2].astype(str)+  ' | ' + df[col3].astype(str)
                return df

            merged_df_rounded_ = combine_columns(merged_df_rounded, 'growth_percentage_academic_4w', 'growth_percentage_community_4w', 'growth_percentage_idn/hospital_4w', "tro_4w_combined")
            merged_df_rounded_ = combine_columns(merged_df_rounded_, 'growth_percentage_academic_13w', 'growth_percentage_community_13w', 'growth_percentage_idn/hospital_13w', "tro_13w_combined")
            merged_df_rounded_ = combine_columns(merged_df_rounded_, 'Volume_Contribution_Academic', 'Volume_Contribution_Community', 'Volume_Contribution_IDN', "volume_combined")
            final_merged_df = merged_df_rounded_[[filter_by,'growth_percentage_R4w_p4w', 'growth_percentage_R13W_P13W',
                            'tro_4w_combined',  'tro_13w_combined', 'volume_combined']]
            print(final_merged_df.head())

            def change_value(row_name, column_name, df):
                value = df[df[filter_by] == row_name][column_name].values[0]
                df.loc[df[filter_by] == row_name, column_name] = "Vials Growth " + str(value)
                return df
            if region :
                change_val_dict = {"growth_percentage_R4w_p4w": region, "growth_percentage_R13W_P13W": region } #"growth_percentage_R13W_Q1TD": region}
            else:
                change_val_dict = {"growth_percentage_R4w_p4w": "Nation", "growth_percentage_R13W_P13W": "Nation"}#, "growth_percentage_R13W_Q1TD": "Nation"}
            for column_name, row_name in change_val_dict.items():
                df_values_changed = change_value(row_name, column_name, final_merged_df)
            df_values_changed
            return region_vials_growth_, df_values_changed



        def get_regional_avg(region):
            if region:
                area_data = tro_sales_data[tro_sales_data['area_nm'].str.contains(region, na=False)]
                filter_by = "regn_nm"
            else:
                area_data = tro_sales_data
                filter_by = "area_nm"

            filtered_data = area_data[(area_data['week_end_date'] >= '2024-01-05') & (area_data['week_end_date'] <= latest_weekend)]
#             filtered_data.to_excel("east_waest_filtered_data.xlsx")
            try:
                weekly_total_sales = filtered_data.groupby(['week_end_date', filter_by])['qty_sold_pu'].sum().unstack().drop('Unknown', axis=1)
            except:
                weekly_total_sales = filtered_data.groupby(['week_end_date', filter_by])['qty_sold_pu'].sum().unstack()
            print(weekly_total_sales)
            weekly_total_sales['Regional Avg'] = weekly_total_sales.mean(axis=1)
            cols = ['Regional Avg'] + [col for col in weekly_total_sales.columns if col != 'Regional Avg']
            weekly_total_sales = weekly_total_sales[cols]
            weekly_total_sales.reset_index(inplace=True)
        #     print(weekly_total_sales)
            if region == "East":
                order = ['week_end_date', 'Regional Avg', 'New England','New York','Atlantic Coastal','Blue Grass','Eastern Coastal','Atlanta','Southeast']
            elif region == "West":
                order = ['week_end_date', 'Regional Avg', 'Pacific Northwest','SoCAL','Rocky Mountains','Gulf Plains','Texas','North Central','Great Lakes']
            else:
                weekly_total_sales = weekly_total_sales.rename(columns={'Regional Avg': 'Area Average'})
                order = ['week_end_date', 'Area Average', 'East', 'West']

            weekly_total_sales_east = weekly_total_sales[order]


            def convert_date_format(date_str):
                date_obj = datetime.fromisoformat(str(date_str).replace("Z", "+00:00"))
                return date_obj.strftime("%d-%b")

            weekly_total_sales_east['week_end_date'] = weekly_total_sales_east['week_end_date'].apply(convert_date_format)
            return weekly_total_sales_east
        print("here!!!! before west!!")
        region_vials_growth_west, west_area_sales = get_slide_6_7_east_west_data("West")
        print("here!!!! after west!!")
        region_vials_growth_east, east_area_sales = get_slide_6_7_east_west_data("East")
        print("here!!!! after west!!")

        east_regional_avg = get_regional_avg("East")
        print("east_regional_avg!!!!!!!!", east_regional_avg)
        west_regional_avg = get_regional_avg("West")
        nation_vials_growth, nation_area_sales = get_slide_6_7_east_west_data('')
        nation_area_avg = get_regional_avg("")





        def get_comp_vs_tro_data():
            def get_dates(current_date , weeks_before='', weeks_after=''):
                print("current_date!!", current_date, weeks_before, weeks_after )
                given_date = datetime.strptime(current_date, "%Y-%m-%d")
                if weeks_before:
                    days_before = weeks_before * 7
                    date_before = given_date - timedelta(days=days_before)
                if weeks_after:
                    days_before = weeks_after * 7
                    date_before = given_date + timedelta(days=days_before)
                past_date = date_before.strftime("%Y-%m-%d")
                return past_date

            def get_date_range(base_date_str, weeks, date_format='%Y-%m-%d'):
                # end_date = datetime.strptime(base_date_str, date_format).replace(tzinfo=pytz.UTC)
                # start_date = end_date - timedelta(weeks=weeks - 1)
                # return start_date.strftime("%Y-%m-%dT%H:%M:%S.%fZ"), end_date.strftime("%Y-%m-%dT%H:%M:%S.%fZ")
                # Convert base_date_str to datetime object
                base_date = datetime.strptime(base_date_str, date_format)
                # Calculate start and end dates with time information
                end_date = base_date.replace(tzinfo=None) #remove tzinfo
                start_date = end_date - timedelta(weeks=weeks - 1)

                # Format start and end dates as strings with time information
                # Return only the date part
                return start_date.strftime(date_format), end_date.strftime(date_format)

            def add_new_rows(df, parent_name):
                new_df_ls = []
                new_row = {}
                sum_columns = ['qty_sold_pu_r12m', 'qty_sold_pu_p12m']
                for col in sum_columns:
                    new_row[col] = df[col].sum()
                if parent_name:
                    col_name = "Nation - TOPA"
                else:
                    col_name = "Nation"
                new_row["top_parent_account_type"] = col_name
                new_df_ls.append(new_row)
                calc_df =pd.DataFrame(new_df_ls)
                df_new_rows = pd.concat([df,calc_df], ignore_index=True)
                return df_new_rows

            def get_growth(drug_list, end_date_, weeks, sub_group_by, col_name):
                if 'trodelvy' in drug_list:
                    whole_data=tro_sales_data
                else:
                    whole_data=ddd_sales_data
                def format_data(sub, total, nation, sub_group_by):
                    total[sub_group_by] = total['top_parent_account_type']
                    nation[sub_group_by] = nation['top_parent_account_type']
                    result = pd.concat([sub, total, nation.iloc[[-1]]]).reset_index(drop=True)
                    df_no_unknowns_null = result[(result['top_parent_account_type'] != 'Unknown') & (result[sub_group_by] != 'Unknown') & (result[sub_group_by] != None)]
                    df_no_unknowns_null =df_no_unknowns_null.dropna()
                    df_sorted = df_no_unknowns_null.sort_values(by=['top_parent_account_type', sub_group_by])
                    return df_sorted
                # def calculate_growth(recent_sales, past_sales):
                #         # Avoid division by zero using np.where
                #         recent_sales = float(recent_sales)
                #         past_sales = float(past_sales)
                #         growth_percentage = np.where(past_sales != 0, ((recent_sales - past_sales) / past_sales) * 100, np.nan)
                #         # Convert the result into a Pandas Series and then fill NaN values
                #         growth_percentage = pd.Series(growth_percentage).fillna(0).replace([np.inf, -np.inf], 0)
                #         return growth_percentage

                def calculate_growth(recent_sales, past_sales):
                    # Avoid division by zero using np.where
                    # Convert to numeric and handle potential errors
                    recent_sales = pd.to_numeric(recent_sales, errors='coerce').fillna(0)
                    past_sales = pd.to_numeric(past_sales, errors='coerce').fillna(0)
                    growth_percentage = np.where(past_sales != 0, ((recent_sales - past_sales) / past_sales) * 100, np.nan)
                    # Convert the result into a Pandas Series and then fill NaN values with 0
                    growth_percentage = pd.Series(growth_percentage).fillna(0).replace([np.inf, -np.inf], 0)
                    return growth_percentage

                def filter_data(whole_data, group_by, parent_name=''):
                    r_start_date, r_end_date  = get_date_range(latest_weekend_ddd, 13)
                    p_start_date, p_end_date = get_date_range(get_dates(latest_weekend_ddd , 13), 13)


                    print("r_start_date", r_start_date, "r_end_date", r_end_date, "p_start_date, p_end_date!!", p_start_date, p_end_date)
                    if parent_name:
                        trodelvy_data = whole_data[whole_data['drug_name'].str.lower().isin(drug_list) &(whole_data['parent_name'] != "TEXAS ONCOLOGY")]
                    else:
                        trodelvy_data = whole_data[whole_data['drug_name'].str.lower().isin(drug_list)]
                    # r_start_date = datetime.strptime(r_start_date, '%Y-%m-%d').date()
                    # r_end_date = datetime.strptime(r_end_date, '%Y-%m-%d').date()
                    # p_start_date = datetime.strptime(p_start_date, '%Y-%m-%d').date()
                    # p_end_date = datetime.strptime(p_end_date, '%Y-%m-%d').date()
                    # r_start_date = datetime.strptime(r_start_date, '%Y-%m-%d')  # Remove .date()
                    # r_end_date = datetime.strptime(r_end_date, '%Y-%m-%d')    # Remove .date()
                    # p_start_date = datetime.strptime(p_start_date, '%Y-%m-%d')  # Remove .date()
                    # p_end_date = datetime.strptime(p_end_date, '%Y-%m-%d')    # Remove .date()

                    # r_start_date = pd.Timestamp(r_start_date, tz=pytz.UTC)
                    # r_end_date = pd.Timestamp(r_end_date, tz=pytz.UTC)
                    # p_start_date = pd.Timestamp(p_start_date, tz=pytz.UTC)
                    # p_end_date = pd.Timestamp(p_end_date, tz=pytz.UTC)
                    r_start_date = pd.Timestamp(r_start_date).tz_localize(None)
                    r_end_date = pd.Timestamp(r_end_date).tz_localize(None)
                    p_start_date = pd.Timestamp(p_start_date).tz_localize(None)
                    p_end_date = pd.Timestamp(p_end_date).tz_localize(None)
                    trodelvy_data['week_end_date'] = pd.to_datetime(trodelvy_data['week_end_date']).dt.tz_localize(None)
                    r12m_data = trodelvy_data[(trodelvy_data['week_end_date'] >= r_start_date) & (trodelvy_data['week_end_date'] <= r_end_date)]
                    p12m_data = trodelvy_data[(trodelvy_data['week_end_date'] >= p_start_date) & (trodelvy_data['week_end_date'] <= p_end_date)]
                    r12m_sales = (r12m_data.groupby(group_by)['qty_sold_pu'].sum()/13).reset_index()

                    p12m_sales = (p12m_data.groupby(group_by)['qty_sold_pu'].sum()/13).reset_index()


                    print("recent 12 m sales",r12m_sales)
                    print("past 12 m sales",p12m_sales)


                    merged_sales_ = r12m_sales.merge(p12m_sales, on= group_by, suffixes=('_r12m', '_p12m'))
                    output_df = add_new_rows(merged_sales_, parent_name)
                    # growth_percentage = np.where(output_df['qty_sold_pu_p12m'] != 0,
                    #              ((output_df['qty_sold_pu_r12m'] - output_df['qty_sold_pu_p12m']) / output_df['qty_sold_pu_p12m']) * 100,
                    #              0)
                    # # growth_percentage = ((output_df['qty_sold_pu_r12m'] - output_df['qty_sold_pu_p12m']) / output_df['qty_sold_pu_p12m']) * 100
                    # growth_percentage = growth_percentage.fillna(0).replace([np.inf, -np.inf], 0)
                    # output_df[f'growth_percentage_{col_name}'] = growth_percentage.round(0).astype(int).astype(str) + '%'
                    # Compute growth_percentage as a pandas Series
                    # growth_percentage = pd.Series(((output_df['qty_sold_pu_r12m'] - output_df['qty_sold_pu_p12m']) / output_df['qty_sold_pu_p12m']) * 100)

                    # # Fill NaNs and replace infinities
                    # growth_percentage = growth_percentage.fillna(0).replace([np.inf, -np.inf], 0)

                    # # Add it to the DataFrame
                    # output_df[f'growth_percentage_{col_name}'] = growth_percentage.round(0).astype(int).astype(str) + '%'
                    # growth_percentage = ((output_df['qty_sold_pu_r12m'] - output_df['qty_sold_pu_p12m']) /
                    #  output_df['qty_sold_pu_p12m']) * 100






                    # output_df['qty_sold_pu_r12m'] = pd.to_numeric(output_df['qty_sold_pu_r12m'], errors='coerce')
                    # output_df['qty_sold_pu_p12m'] = pd.to_numeric(output_df['qty_sold_pu_p12m'], errors='coerce')
                    # growth_percentage= np.where(
                    #     output_df['qty_sold_pu_p12m'] == 0,
                    #     0,  # Safe value for division by zero (or `NaN`)
                    #     ((output_df['qty_sold_pu_r12m'] - output_df['qty_sold_pu_p12m']) / output_df['qty_sold_pu_p12m']) * 100
                    # )

                    # print("output df in slide 10,11 Husna",output_df.head())


                    # # growth_percentage = ((output_df['qty_sold_pu_r12m'] - output_df['qty_sold_pu_p12m']) / output_df['qty_sold_pu_p12m']) * 100
                    # # growth_percentage = growth_percentage.fillna(0).replace([np.inf, -np.inf], 0)
                    # # output_df[f'growth_percentage_{col_name}'] = growth_percentage.round(0).astype(int).astype(str) + '%'
                    # output_df[f'growth_percentage_{col_name}'] = growth_percentage.astype(float).round(0).astype(int).astype(str) + '%'
                    # print("output df after calc growth",output_df.head(),output_df.columns)



                    growth_percentage = calculate_growth(output_df['qty_sold_pu_r12m'], output_df['qty_sold_pu_p12m'])

                    # growth_percentage = ((output_df['qty_sold_pu_r12m'] - output_df['qty_sold_pu_p12m']) /
                    #  output_df['qty_sold_pu_p12m']) * 100

                    # Convert growth_percentage to float to handle decimal.Decimal type
                    growth_percentage = growth_percentage.astype(float)

                    # Handle NaN, Inf, and -Inf values
                    growth_percentage = growth_percentage.fillna(0).replace([np.inf, -np.inf], 0)

                    # Apply rounding and formatting
                    output_df[f'growth_percentage_{col_name}'] = growth_percentage.round(0).astype(float).astype(str) + '%'

                    return output_df
                    # growth_percentage = np.where(
                    #       output_df['qty_sold_pu_p12m'] != 0,
                    #       ((output_df['qty_sold_pu_r12m'] - output_df['qty_sold_pu_p12m']) / output_df['qty_sold_pu_p12m']) * 100,
                    #       0  # Replace with 0 or another suitable value when the denominator is 0
                    #   )

                    # # Convert growth_percentage to float to handle decimal.Decimal type
                    # growth_percentage = growth_percentage.astype(float)

                    # Handle NaN, Inf, and -Inf values
                    # growth_percentage = growth_percentage.fillna(0).replace([np.inf, -np.inf], 0)

                    # Apply rounding and formatting
                    # output_df[f'growth_percentage_{col_name}'] = growth_percentage.round(0).astype(int).astype(str) + '%'

                    return output_df
                sub_merged_sales_ = filter_data(whole_data, ['top_parent_account_type', sub_group_by], True)
                merged_sales = filter_data(whole_data, ['top_parent_account_type'], True)
                Nation_Sales = filter_data(whole_data, ['top_parent_account_type'], False)
                formatted_df = format_data(sub_merged_sales_, merged_sales, Nation_Sales, sub_group_by)
                return formatted_df[['top_parent_account_type', sub_group_by, f'growth_percentage_{col_name}']]

            def get_contribution(drug_list, end_date_, weeks, sub_group_by, col_name):
                if 'trodelvy' in drug_list:
                    whole_data=tro_sales_data
                else:
                    whole_data=ddd_sales_data
                def format_data(sub, total, nation, sub_group_by):
                    total[sub_group_by] = total['top_parent_account_type']
                    result = pd.concat([sub, total]).reset_index(drop=True)
                    df_no_unknowns_null = result[(result['top_parent_account_type'] != 'Unknown') & (result[sub_group_by] != 'Unknown') & (result[sub_group_by] != None)]
                    df_no_unknowns_null =df_no_unknowns_null.dropna()
                    df_sorted = df_no_unknowns_null.sort_values(by=['top_parent_account_type', sub_group_by])
                    return df_sorted

                def filter_data(group_by, parent_name= ''):
                    start_date, end_date = get_date_range(end_date_ , weeks, '%Y-%m-%d')
                    start_date = datetime.strptime(start_date, '%Y-%m-%d')
                    end_date = datetime.strptime(end_date, '%Y-%m-%d')
                    start_date = pd.to_datetime(start_date).tz_localize(None)
                    end_date = pd.to_datetime(end_date).tz_localize(None)
                    whole_data['week_end_date'] = pd.to_datetime(whole_data['week_end_date']).dt.tz_localize(None)

                    if parent_name:
#                         filtered_data = whole_data[(whole_data['drug_name'].str.lower().isin(drug_list)) & (whole_data['week_end_date'] >= start_date) & (whole_data['week_end_date'] <= end_date) ]
                        filtered_data = whole_data[(whole_data['drug_name'].str.lower().isin(drug_list)) & (whole_data['week_end_date'] >= start_date) & (whole_data['week_end_date'] <= end_date) &(whole_data['parent_name'] != "TEXAS ONCOLOGY")]
                    else:
                        filtered_data = whole_data[(whole_data['drug_name'].str.lower().isin(drug_list)) & (whole_data['week_end_date'] >= start_date) & (whole_data['week_end_date'] <= end_date)]
                    filtered_data = filtered_data[['top_parent_account_type',sub_group_by, 'qty_sold_pu']]
                    print(filtered_data.columns, "col!!!!!")
                    grouped_data = filtered_data.groupby(group_by)['qty_sold_pu'].sum().reset_index()
                    total_sales = grouped_data['qty_sold_pu'].sum()
                    grouped_data[f'contribution_{col_name}'] = (grouped_data['qty_sold_pu'] / total_sales)
                    return grouped_data

                out_put_df  = filter_data(['top_parent_account_type', sub_group_by], True).drop(columns=['qty_sold_pu'])
                out_put_df_2  = filter_data(['top_parent_account_type'], True).drop(columns=['qty_sold_pu'])
                Nation_output_df = filter_data(['top_parent_account_type'], False)
                formatted_df = format_data(out_put_df, out_put_df_2, Nation_output_df, sub_group_by)
                formatted_df[f'contribution_{col_name}'] = round(formatted_df[f'contribution_{col_name}']* 100, 0).astype(int).astype(str) + '%'
                return out_put_df, out_put_df_2, formatted_df[['top_parent_account_type', sub_group_by, f'contribution_{col_name}']]

            def get_growth_contrib_data(sub_group_by):
                trodelvy_growth_data = get_growth(["trodelvy"], latest_weekend_ddd, 52, sub_group_by, 'trodelvy')
                enhertu_growth_data = get_growth(["enhertu"], latest_weekend_ddd, 52, sub_group_by, 'enhertu')
                tpc_growth_data = get_growth(tpc, latest_weekend_ddd, 52, sub_group_by, 'tpc')

                growth_merged_df1_2 = pd.merge(trodelvy_growth_data, enhertu_growth_data, on=['top_parent_account_type', sub_group_by], how='outer')
                growth_merged_all = pd.merge(growth_merged_df1_2, tpc_growth_data, on=['top_parent_account_type', sub_group_by], how='outer')
                growth_merged_all['growth_combined'] = growth_merged_all.apply(lambda row: f"{row['growth_percentage_trodelvy']} | {row['growth_percentage_enhertu']} | {row['growth_percentage_tpc']}", axis=1)

                tro_sub_account_contr_data ,tro_account_contr_data, trodelvy_contrib_data = get_contribution(["trodelvy"], latest_weekend_ddd, 52, sub_group_by, 'trodelvy')
                enh_sub_account_contr_data ,enh_account_contr_data, enhertu_contrib_data = get_contribution(["enhertu"], latest_weekend_ddd, 52, sub_group_by, 'enhertu')
                tpc_sub_account_contr_data ,tpc_account_contr_data, tpc_contrib_data= get_contribution(tpc, latest_weekend_ddd, 52, sub_group_by, 'tpc')

                contrib_merged_df1_2 = pd.merge(trodelvy_contrib_data, enhertu_contrib_data, on=['top_parent_account_type', sub_group_by], how='outer')
                contrib_merged_all = pd.merge(contrib_merged_df1_2, tpc_contrib_data, on=['top_parent_account_type', sub_group_by], how='outer')
                contrib_merged_all['contrib_combined'] = contrib_merged_all.apply(lambda row: f"{row['contribution_trodelvy']} | {row['contribution_enhertu']} | {row['contribution_tpc']}", axis=1)
                contrib_merged_all
                if sub_group_by == 'parent_bc_segment':
                    final_df = growth_merged_all[['top_parent_account_type', 'parent_bc_segment', 'growth_combined']]
                else:
                    final_df = pd.merge(growth_merged_all, contrib_merged_all, on=['top_parent_account_type', 'kad_flag'], how='left')
                    final_df = final_df[['top_parent_account_type', 'kad_flag','growth_combined', 'contrib_combined']]

                return tro_sub_account_contr_data, tro_account_contr_data, final_df


            tpc = ['gemzar', 'halaven', 'vinorelbine', 'xeloda']
#             latest_weekend = '2024-05-03'

            tro_sub_account_contr_data_, tro_account_contr_data, final_df_account = get_growth_contrib_data('parent_bc_segment')
            final_df_account.iloc[8], final_df_account.iloc[9] = final_df_account.iloc[9].copy(), final_df_account.iloc[8].copy()

            tro_sub_kad_contr_data_, tro_account_kad_data, final_df_kad =   get_growth_contrib_data( 'kad_flag')

            tro_sub_account_contr_data = tro_sub_account_contr_data_[ (tro_sub_account_contr_data_['parent_bc_segment'] != 'Unknown') & (tro_sub_account_contr_data_['parent_bc_segment'] != None)]
            tro_sub_kad_contr_data = tro_sub_kad_contr_data_[ (tro_sub_kad_contr_data_['kad_flag'] != 'Unknown') & (tro_sub_kad_contr_data_['kad_flag'] != None)]
            final_df_kad = final_df_kad.replace({ 'Y': 'Yes', 'N': 'No'})
            final_df_kad = final_df_kad.fillna("-")

            tro_account_contr_data = tro_account_contr_data[['top_parent_account_type','contribution_trodelvy']]
            tro_account_contr_data['top_parent_account_type'] = tro_account_contr_data['top_parent_account_type'].replace({
                'Unknown': 'Other',
                'IDN/Other': 'IDN/Hospital'
            })
            tro_account_contr_data = tro_account_contr_data.groupby('top_parent_account_type').sum().reset_index()

            tro_sub_kad_contr_data['kad_flag'] = tro_sub_kad_contr_data['kad_flag'].replace({
                'Y': "KAD Targets",
                'N': "Non - KAD Targets"
            })

            tro_sub_account_contr_data.loc[tro_sub_account_contr_data['top_parent_account_type'] == 'Unknown', 'parent_bc_segment'] = 'Other'

            tro_sub_kad_contr_data.loc[tro_sub_kad_contr_data['top_parent_account_type'] == 'Unknown', 'kad_flag'] = 'Other'


            tro_sub_account_contr_data = tro_sub_account_contr_data.drop(columns=['top_parent_account_type'])
            tro_sub_kad_contr_data = tro_sub_kad_contr_data.drop(columns=['top_parent_account_type'])
#             tro_sub_account_contr_data = tro_sub_account_contr_data.groupby('parent_bc_segment').sum().reset_index()
            sum_other = tro_sub_account_contr_data[tro_sub_account_contr_data['parent_bc_segment'] == 'Other']['contribution_trodelvy'].sum()
            tro_sub_account_contr_data.loc[tro_sub_account_contr_data['parent_bc_segment'] == 'Other', 'contribution_trodelvy'] = sum_other
            tro_sub_account_contr_data = tro_sub_account_contr_data.drop_duplicates()

            return tro_sub_account_contr_data, tro_account_contr_data, tro_sub_kad_contr_data, final_df_account, final_df_kad

        tro_sub_account_contr_data, tro_account_contr_data, tro_sub_kad_contr_data, final_df_account, final_df_kad = get_comp_vs_tro_data()
#         tro_sub_account_contr_data, tro_account_contr_data, tro_sub_kad_contr_data, final_df_account, final_df_kad = "", "", "","",""



#         def get_comp_vs_tro_data():
#             # Issue 2: Inconsistent date handling
#             def get_dates(current_date, weeks_before=None, weeks_after=None):
#                 # Fix: Ensure current_date is always datetime
#                 given_date = pd.to_datetime(current_date)

#                 if weeks_before and str(weeks_before).strip().isdigit():
#                     weeks_before = int(weeks_before)
#                     date_before = given_date - timedelta(days=weeks_before * 7)
#                 elif weeks_after and str(weeks_after).strip().isdigit():
#                     weeks_after = int(weeks_after)
#                     date_before = given_date + timedelta(days=weeks_after * 7)
#                 else:
#                     raise ValueError("Either weeks_before or weeks_after must be provided")

#                 return date_before.strftime("%Y-%m-%d")

#             def get_date_range(base_date_str, weeks, date_format='%Y-%m-%d'):
#                 end_date = datetime.strptime(base_date_str, date_format).replace(tzinfo=pytz.UTC)
#                 start_date = end_date - timedelta(weeks=weeks - 1)
#                 return start_date.strftime("%Y-%m-%d"), end_date.strftime("%Y-%m-%d")

#             def add_new_rows(df, parent_name):
#                 new_df_ls = []
#                 new_row = {}
#                 sum_columns = ['qty_sold_pu_r12m', 'qty_sold_pu_p12m']
#                 for col in sum_columns:
#                     new_row[col] = df[col].sum()
#                 if parent_name:
#                     col_name = "Nation - TOPA"
#                 else:
#                     col_name = "Nation"
#                 new_row["top_parent_account_type"] = col_name
#                 new_df_ls.append(new_row)
#                 calc_df =pd.DataFrame(new_df_ls)
#                 df_new_rows = pd.concat([df,calc_df], ignore_index=True)
#                 return df_new_rows

#             def get_growth(drug_list, end_date_, weeks, sub_group_by, col_name):
#                 def format_data(sub, total, nation, sub_group_by):
#                     total[sub_group_by] = total['top_parent_account_type']
#                     nation[sub_group_by] = nation['top_parent_account_type']
#                     result = pd.concat([sub, total, nation.iloc[[-1]]]).reset_index(drop=True)
#                     df_no_unknowns_null = result[(result['top_parent_account_type'] != 'Unknown') &
#                                                  (result[sub_group_by] != 'Unknown') &
#                                                  (result[sub_group_by] != None)]
#                     df_no_unknowns_null = df_no_unknowns_null.dropna()
#                     df_sorted = df_no_unknowns_null.sort_values(by=['top_parent_account_type', sub_group_by])
#                     return df_sorted

#                 def filter_data(group_by, parent_name=''):
#                     # Get the start and end dates
#                     r_start_date, r_end_date = get_date_range(latest_weekend_ddd, 13)
#                     p_start_date, p_end_date = get_date_range(get_dates(latest_weekend_ddd, 13), 13)

#                     print("r_start_date", r_start_date, "r_end_date", r_end_date, "p_start_date, p_end_date!!", p_start_date, p_end_date)

#                     # Ensure dates are datetime objects
#                     r_start_date = pd.to_datetime(r_start_date)
#                     r_end_date = pd.to_datetime(r_end_date)
#                     p_start_date = pd.to_datetime(p_start_date)
#                     p_end_date = pd.to_datetime(p_end_date)

#                     # Determine the SQL query based on the presence of 'trodelvy' in drug_list
#                     if 'trodelvy' in drug_list:

#                         if parent_name:
#                             filtered_data = tro_sales_data[(tro_sales_data['drug_name'].str.lower().isin(drug_list))&(tro_sales_data['parent_name'] != "TEXAS ONCOLOGY")]
#                             # Query for Trodelvy with parent name filter
#                             # sql_query = f'''
#                             #     SELECT *
#                             #     FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`867_data_drug-trodelvy_for_genie_filtered_dates`
#                             #     WHERE LOWER(drug_name) IN ({",".join([f"'{drug.lower()}'" for drug in drug_list])})
#                             #     AND parent_name != "TEXAS ONCOLOGY"
#                             # '''
#                             trodelvy_data = tro_sales_data[tro_sales_data['drug_name'].str.lower().isin(drug_list) &(tro_sales_data['parent_name'] != "TEXAS ONCOLOGY")]
#                         else:
#                             # Query for Trodelvy without parent name filter
#                             # sql_query = f'''
#                             #     SELECT *
#                             #     FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`867_data_drug-trodelvy_for_genie_filtered_dates`
#                             #     WHERE LOWER(drug_name) IN ({",".join([f"'{drug.lower()}'" for drug in drug_list])})
#                             # '''
#                             trodelvy_data = tro_sales_data[tro_sales_data['drug_name'].str.lower().isin(drug_list)]
#                     else:
#                         if parent_name:
#                             # Query for non-Trodelvy drugs with parent name filter
#                             sql_query = f'''
#                                 SELECT *
#                                 FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`ddd_data_for_genie_filtered_dates`
#                                 WHERE LOWER(drug_name) IN ({",".join([f"'{drug.lower()}'" for drug in drug_list])})
#                                 AND parent_name != "TEXAS ONCOLOGY"
#                             '''
#                             trodelvy_data = pd.DataFrame(self.cursor.execute(sql_query).fetchall(), columns=[elem[0] for elem in self.cursor.description])
#                         else:
#                             # Query for non-Trodelvy drugs without parent name filter
#                             sql_query = f'''
#                                 SELECT *
#                                 FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`ddd_data_for_genie_filtered_dates`
#                                 WHERE LOWER(drug_name) IN ({",".join([f"'{drug.lower()}'" for drug in drug_list])})
#                             '''
#                             trodelvy_data = pd.DataFrame(self.cursor.execute(sql_query).fetchall(), columns=[elem[0] for elem in self.cursor.description])
#                     # Execute the query and store the data in a DataFrame


#                     # Rename 'trx_cnt' to 'qty_sold_pu'
#                     if 'trx_cnt' in trodelvy_data.columns:
#                         trodelvy_data = trodelvy_data.rename(columns={'trx_cnt': 'qty_sold_pu'})

#                     # Check if the column 'qty_sold_pu' exists
#                     if 'qty_sold_pu' not in trodelvy_data.columns:
#                         print(f"Error: 'qty_sold_pu' column not found in the query result.")
#                         print("Available columns:", trodelvy_data.columns)
#                         return pd.DataFrame()  # Return empty DataFrame if column is not found

#                     # Filter the data based on date range
#                     r12m_data = trodelvy_data[(pd.to_datetime(trodelvy_data['week_end_date']) >= r_start_date) &
#                                                (pd.to_datetime(trodelvy_data['week_end_date']) <= r_end_date)]
#                     p12m_data = trodelvy_data[(pd.to_datetime(trodelvy_data['week_end_date']) >= p_start_date) &
#                                                (pd.to_datetime(trodelvy_data['week_end_date']) <= p_end_date)]

#                     # Calculate sales for r12m and p12m
#                     r12m_sales = (r12m_data.groupby(group_by)['qty_sold_pu'].sum() / 13).reset_index()
#                     p12m_sales = (p12m_data.groupby(group_by)['qty_sold_pu'].sum() / 13).reset_index()

#                     # Merge the r12m and p12m sales data
#                     merged_sales_ = r12m_sales.merge(p12m_sales, on=group_by, suffixes=('_r12m', '_p12m'))

#                     # Add new rows if needed and calculate growth percentage
#                     output_df = add_new_rows(merged_sales_, parent_name)
#                     growth_percentage = ((output_df['qty_sold_pu_r12m'] - output_df['qty_sold_pu_p12m']) /
#                      output_df['qty_sold_pu_p12m']) * 100

#                     # Convert growth_percentage to float to handle decimal.Decimal type
#                     growth_percentage = growth_percentage.astype(float)

#                     # Handle NaN, Inf, and -Inf values
#                     growth_percentage = growth_percentage.fillna(0).replace([np.inf, -np.inf], 0)

#                     # Apply rounding and formatting
#                     output_df[f'growth_percentage_{col_name}'] = growth_percentage.round(0).astype(int).astype(str) + '%'


#                     return output_df



#                 # Assuming you want to format and return the data after filter_data
#                 sub_merged_sales_ = filter_data(['top_parent_account_type', sub_group_by], True)
#                 merged_sales = filter_data(['top_parent_account_type'], True)
#                 Nation_Sales = filter_data(['top_parent_account_type'], False)

#                 formatted_df = format_data(sub_merged_sales_, merged_sales, Nation_Sales, sub_group_by)
#                 return formatted_df[['top_parent_account_type', sub_group_by, f'growth_percentage_{col_name}']]


#             def get_contribution(drug_list, end_date_, weeks, sub_group_by, col_name):
#                 def format_data(sub, total, nation, sub_group_by):
#                     total[sub_group_by] = total['top_parent_account_type']
#                     result = pd.concat([sub, total]).reset_index(drop=True)
#                     df_no_unknowns_null = result[(result['top_parent_account_type'] != 'Unknown') &
#                                                  (result[sub_group_by] != 'Unknown') &
#                                                  (result[sub_group_by] != None)]
#                     df_no_unknowns_null = df_no_unknowns_null.dropna()
#                     df_sorted = df_no_unknowns_null.sort_values(by=['top_parent_account_type', sub_group_by])
#                     return df_sorted

#                 def filter_data(group_by, parent_name=''):
#                     # Get the start and end dates
#                     start_date, end_date = get_date_range(end_date_, weeks, '%Y-%m-%d')
#                     print("start_date in slide 10,11, end_date!!!", start_date, end_date)

#                     # Query the database based on drug list and parent name
#                     if 'trodelvy' in drug_list:
#                         if parent_name:
#                             filtered_data = tro_sales_data[(tro_sales_data['drug_name'].str.lower().isin(drug_list)) & (tro_sales_data['week_end_date'] >= start_date) & (tro_sales_data['week_end_date'] <= end_date) &(tro_sales_data['parent_name'] != "TEXAS ONCOLOGY")]
#                         else:
#                             filtered_data = tro_sales_data[(tro_sales_data['drug_name'].str.lower().isin(drug_list)) & (tro_sales_data['week_end_date'] >= start_date) & (tro_sales_data['week_end_date'] <= end_date)]
#                     else:
#                         if parent_name:
#                             sql_query = f'''
#                                 SELECT *
#                                 FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`ddd_data_for_genie_filtered_dates`
#                                 WHERE LOWER(drug_name) IN ({",".join([f"'{drug.lower()}'" for drug in drug_list])})
#                                 AND week_end_date >= '{start_date}'
#                                 AND week_end_date <= '{end_date}'
#                                 AND parent_name != "TEXAS ONCOLOGY"
#                             '''
#                             filtered_data = pd.DataFrame(self.cursor.execute(sql_query).fetchall(), columns=[elem[0] for elem in self.cursor.description])
#                         else:
#                             sql_query = f'''
#                                 SELECT *
#                                 FROM `commercial-us-apo-onc-iia-dev`.`setuserv_adhocs`.`ddd_data_for_genie_filtered_dates`
#                                 WHERE LOWER(drug_name) IN ({",".join([f"'{drug.lower()}'" for drug in drug_list])})
#                                 AND week_end_date >= '{start_date}'
#                                 AND week_end_date <= '{end_date}'
#                             '''
#                             filtered_data = pd.DataFrame(self.cursor.execute(sql_query).fetchall(), columns=[elem[0] for elem in self.cursor.description])
#                     # Execute the query and load the data into a DataFrame
#                     print(f"Filtered Data columns before trx_cnt conversion: {filtered_data}")

#                     # If no data is returned, return empty DataFrame
#                     if filtered_data.empty:
#                         return pd.DataFrame()

#                     # Rename 'txn_cnt' to 'qty_sold_pu' if it exists
#                     if 'trx_cnt' in filtered_data.columns:
#                         filtered_data = filtered_data.rename(columns={'trx_cnt': 'qty_sold_pu'})

#                     filtered_data = filtered_data[['top_parent_account_type',sub_group_by, 'qty_sold_pu']]
#                     print(f"Filtered Data columns After trx_cnt conversion: {filtered_data}")
#                     print(filtered_data.head(), "col!!!!!")

#                     grouped_data = filtered_data.groupby(group_by)['qty_sold_pu'].sum().reset_index()
#                     print("grouped data for contribution",grouped_data.head())
#                     total_sales = grouped_data['qty_sold_pu'].sum()
#                     grouped_data[f'contribution_{col_name}'] = (grouped_data['qty_sold_pu'] / total_sales)

#                     # Return filtered data
#                     # return filtered_data
#                     return grouped_data

#                 out_put_df  = filter_data(['top_parent_account_type', sub_group_by], True).drop(columns=['qty_sold_pu'])
#                 out_put_df_2  = filter_data(['top_parent_account_type'], True).drop(columns=['qty_sold_pu'])
#                 Nation_output_df = filter_data(['top_parent_account_type'], False)
#                 formatted_df = format_data(out_put_df, out_put_df_2, Nation_output_df, sub_group_by)
#                 formatted_df[f'contribution_{col_name}'] = round(formatted_df[f'contribution_{col_name}']* 100, 0).astype(int).astype(str) + '%'
#                 return out_put_df, out_put_df_2, formatted_df[['top_parent_account_type', sub_group_by, f'contribution_{col_name}']]
#                 # Get filtered data
#                 # filtered_data = filter_data(sub_group_by)

#                 # # If no data is returned, return empty DataFrames
#                 # if filtered_data.empty:
#                 #     return pd.DataFrame(), pd.DataFrame(), pd.DataFrame()

#                 # # Group by the sub_group_by column and calculate the sum of 'qty_sold_pu'
#                 # grouped_data = filtered_data.groupby(sub_group_by).agg({'qty_sold_pu': 'sum'}).reset_index()

#                 # # Format data (add the totals, remove unknowns, and sort the data)
#                 # print(f'Filtered Data Columns before format data{filtered_data}')
#                 # formatted_data = format_data(filtered_data, grouped_data, grouped_data, sub_group_by)

#                 # # Calculate the total sales
#                 # total_sales = grouped_data['qty_sold_pu'].sum()

#                 # # Calculate the contribution for each group
#                 # grouped_data[f'contribution_{col_name}'] = (grouped_data['qty_sold_pu'] / total_sales)

#                 # # Return the processed data as a tuple
#                 # print(f"filtered_data: {filtered_data.head()}")
#                 # print(f"grouped data: {grouped_data.head()}")
#                 # print(f"formatted_data: {formatted_data.head()}")
#                 # return filtered_data, grouped_data, formatted_data

#             print("Before Growth Contribution Data")
#             def get_growth_contrib_data(sub_group_by):
#                 print("Before Growth Contribution Data")

#                 trodelvy_growth_data = get_growth(["trodelvy"], latest_weekend_ddd, 52, sub_group_by, 'trodelvy')
#                 print(f"trodelvy_growth_data: {trodelvy_growth_data}")

#                 enhertu_growth_data = get_growth(["enhertu"], latest_weekend_ddd, 52, sub_group_by, 'enhertu')
#                 print(f"enhertu_growth_data: {enhertu_growth_data}")

#                 tpc_growth_data = get_growth(tpc, latest_weekend_ddd, 52, sub_group_by, 'tpc')
#                 print(f"tpc_growth_data: {tpc_growth_data}")

#                 print("After Growth Data Calculation")

#                 growth_merged_df1_2 = pd.merge(trodelvy_growth_data, enhertu_growth_data, on=['top_parent_account_type', sub_group_by], how='outer')
#                 growth_merged_all = pd.merge(growth_merged_df1_2, tpc_growth_data, on=['top_parent_account_type', sub_group_by], how='outer')
#                 growth_merged_all['growth_combined'] = growth_merged_all.apply(lambda row: f"{row['growth_percentage_trodelvy']} | {row['growth_percentage_enhertu']} | {row['growth_percentage_tpc']}", axis=1)

#                 # Ensure we receive valid data from get_contribution
#                 tro_sub_account_contr_data, tro_account_contr_data, trodelvy_contrib_data = get_contribution(["trodelvy"], latest_weekend_ddd, 52, sub_group_by, 'trodelvy')
#                 # if tro_sub_account_contr_data is None:
#                 #     return None  # Handle the case if no data is returned

#                 enh_sub_account_contr_data, enh_account_contr_data, enhertu_contrib_data = get_contribution(["enhertu"], latest_weekend_ddd, 52, sub_group_by, 'enhertu')
#                 # if enh_sub_account_contr_data is None:
#                 #     return None

#                 tpc_sub_account_contr_data, tpc_account_contr_data, tpc_contrib_data = get_contribution(tpc, latest_weekend_ddd, 52, sub_group_by, 'tpc')
#                 # if tpc_sub_account_contr_data is None:
#                 #     return None



#                 contrib_merged_df1_2 = pd.merge(trodelvy_contrib_data, enhertu_contrib_data, on=['top_parent_account_type', sub_group_by], how='outer')
#                 contrib_merged_all = pd.merge(contrib_merged_df1_2, tpc_contrib_data, on=['top_parent_account_type', sub_group_by], how='outer')
#                 contrib_merged_all['contrib_combined'] = contrib_merged_all.apply(lambda row: f"{row['contribution_trodelvy']} | {row['contribution_enhertu']} | {row['contribution_tpc']}", axis=1)

#                 if sub_group_by == 'parent_bc_segment':
#                     final_df = growth_merged_all[['top_parent_account_type', 'parent_bc_segment', 'growth_combined']]
#                 else:
#                     final_df = pd.merge(growth_merged_all, contrib_merged_all, on=['top_parent_account_type', 'kad_flag'], how='left')
#                     final_df = final_df[['top_parent_account_type', 'kad_flag', 'growth_combined', 'contrib_combined']]

#                 print("After Final df calculation in get Growth")
#                 return tro_sub_account_contr_data, tro_account_contr_data, final_df
#                 print(f"tro_sub_account_contr_data_: {tro_sub_account_contr_data_}")
#                 print(f"tro_account_contr_data: {tro_account_contr_data}")
#                 print(f"final_df_account: {final_df}")

#             '*********************************************'



#             tpc = ['gemzar', 'halaven', 'vinorelbine', 'xeloda']
# #             latest_weekend = '2024-05-03'

#             tro_sub_account_contr_data_, tro_account_contr_data, final_df_account = get_growth_contrib_data('parent_bc_segment')
#             print(f"tro_sub_account_contr_data_ at calling: {tro_sub_account_contr_data_}")
#             print(f"tro_account_contr_data at calling: {tro_account_contr_data}")
#             print(f"final_df_account at calling: {final_df_account}")

#             final_df_account.iloc[8], final_df_account.iloc[9] = final_df_account.iloc[9].copy(), final_df_account.iloc[8].copy()

#             tro_sub_kad_contr_data_, tro_account_kad_data, final_df_kad =   get_growth_contrib_data( 'kad_flag')

#             tro_sub_account_contr_data = tro_sub_account_contr_data_[ (tro_sub_account_contr_data_['parent_bc_segment'] != 'Unknown') & (tro_sub_account_contr_data_['parent_bc_segment'] != None)]
#             tro_sub_kad_contr_data = tro_sub_kad_contr_data_[ (tro_sub_kad_contr_data_['kad_flag'] != 'Unknown') & (tro_sub_kad_contr_data_['kad_flag'] != None)]
#             final_df_kad = final_df_kad.replace({ 'Y': 'Yes', 'N': 'No'})
#             final_df_kad = final_df_kad.fillna("-")

#             tro_account_contr_data = tro_account_contr_data[['top_parent_account_type','contribution_trodelvy']]
#             tro_account_contr_data['top_parent_account_type'] = tro_account_contr_data['top_parent_account_type'].replace({
#                 'Unknown': 'Other',
#                 'IDN/Other': 'IDN/Hospital'
#             })
#             tro_account_contr_data = tro_account_contr_data.groupby('top_parent_account_type').sum().reset_index()

#             tro_sub_kad_contr_data['kad_flag'] = tro_sub_kad_contr_data['kad_flag'].replace({
#                 'Y': "KAD Targets",
#                 'N': "Non - KAD Targets"
#             })

#             tro_sub_account_contr_data.loc[tro_sub_account_contr_data['top_parent_account_type'] == 'Unknown', 'parent_bc_segment'] = 'Other'

#             tro_sub_kad_contr_data.loc[tro_sub_kad_contr_data['top_parent_account_type'] == 'Unknown', 'kad_flag'] = 'Other'


#             tro_sub_account_contr_data = tro_sub_account_contr_data.drop(columns=['top_parent_account_type'])
#             tro_sub_kad_contr_data = tro_sub_kad_contr_data.drop(columns=['top_parent_account_type'])
# #             tro_sub_account_contr_data = tro_sub_account_contr_data.groupby('parent_bc_segment').sum().reset_index()
#             sum_other = tro_sub_account_contr_data[tro_sub_account_contr_data['parent_bc_segment'] == 'Other']['contribution_trodelvy'].sum()
#             tro_sub_account_contr_data.loc[tro_sub_account_contr_data['parent_bc_segment'] == 'Other', 'contribution_trodelvy'] = sum_other
#             tro_sub_account_contr_data = tro_sub_account_contr_data.drop_duplicates()

#             return tro_sub_account_contr_data, tro_account_contr_data, tro_sub_kad_contr_data, final_df_account, final_df_kad

#         tro_sub_account_contr_data, tro_account_contr_data, tro_sub_kad_contr_data, final_df_account, final_df_kad = get_comp_vs_tro_data()
# #         tro_sub_account_contr_data, tro_account_contr_data, tro_sub_kad_contr_data, final_df_account, final_df_kad = "", "", "","",""



        global E4_st, E13_st, W13_st, W4_st, A4_st, A13_st, C13_st, C4_st, M4_st, M13_st, L13_st, L4_st, H13_st, H4_st ,overall13_st,overall4_st
        E4_st = E13_st = W13_st = W4_st = A4_st = A13_st = C13_st = C4_st = M4_st = M13_st = L13_st = L4_st = H13_st = H4_st = overall13_st=overall4_st=None

        def update_text(df, paragraph, growth_row, filter_):
            print("value!!!!!", df.loc[df['time_period'] == growth_row, filter_].values[0])
            val = int(round(df.loc[df['time_period'] == growth_row, filter_].values.item(), 0))
            paragraph.text = str(val)+"%"

        def update_placeholders(df, growth_row, filter_):
            if df.equals(acad_df):
                global A13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    A13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    A13_st="no growth"
                else:
                    A13_st="increased"
            elif df.equals(acad_df1):
                global A4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    A4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    A4_st="no growth"
                else:
                    A4_st="increased"
            elif df.equals(comm_df1):
                global C4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    C4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    C4_st="no growth"
                else:
                    C4_st="increased"
            elif df.equals(comm_df):
                global C13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    C13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    C13_st="no growth"
                else:
                    C13_st="increased"
            elif df.equals(area_df1):
                global E4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    E4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    E4_st="no growth"
                else:
                    E4_st="increased"
            elif df.equals(east_df):
                global E13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    E13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    E13_st="no growth"
                else:
                    E13_st="increased"
            elif df.equals(area_df2):
                global W4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    W4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    W4_st="no growth"
                else:
                    W4_st="increased"
            elif df.equals(west_df):
                global W13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    W13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    W13_st="no growth"
                else:
                    W13_st="increased"
            elif df.equals(low_df1):
                global L4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    L4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    L4_st="no growth"
                else:
                    L4_st="increased"
            elif df.equals(low_df):
                global L13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    L13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    L13_st="no growth"
                else:
                    L13_st="increased"
            elif df.equals(medium_df1):
                global M4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    M4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    M4_st="no growth"
                else:
                    M4_st="increased"
            elif df.equals(medium_df):
                global M13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    M13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    M13_st="no growth"
                else:
                    M13_st="increased"
            elif df.equals(high_df):
                global H13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    H13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    H13_st="no growth"
                else:
                    H13_st="increased"
            elif df.equals(high_df1):
                global H4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    H4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    H4_st="no growth"
                else:
                    H4_st="increased"
            elif df.equals(overall_df):
                global overall13_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    overall13_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    overall13_st="no growth"
                else:
                    overall13_st="increased"
            elif df.equals(overall_df1):
                global overall4_st
                if(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<0):
                    overall4_st="declined"
                elif(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))>=0 and int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0))<1):
                    overall4_st="no growth"
                else:
                    overall4_st="increased"


            # return abs(int(round(df.loc[df['time_period'] == growth_row, filter_].values[0],0)))
            return abs(int(round(df.loc[df['time_period'] == growth_row, filter_].iloc[0], 0)))
        def update_cell_text(cell, value, type_= ''):
            if type_ == "chart_title_":
                print("in update chart title")
                cell.text_frame.text = value
            else:
                text_frame = cell.text_frame
                if text_frame.paragraphs:
                    paragraph = text_frame.paragraphs[0]
                    if paragraph.runs:
                        para_len = len(paragraph.runs)
                        if para_len > 1:
                            print("len> 1", para_len)
                            for run_idx, run in enumerate(paragraph.runs):
                                run.text = ''
                        print("type of value!!!!!", type(value))
                        if "." not in str(value):
                            try:
                                value = int(value)
                                value = f"{value:,}"
                            except:
                                pass
                        paragraph.runs[0].text = str(value)

        def get_slide_12_data():
                def get_averages(filtered_data):
                    filtered_data['week_end_date'] = pd.to_datetime(filtered_data['week_end_date']).dt.tz_localize(None)

                    specified_date = pd.Timestamp(latest_weekend).tz_localize(None)
                    start_date_r4w = specified_date - timedelta(weeks=3)

                    end_date_p4w = start_date_r4w - timedelta(days=1)
                    start_date_p4w = end_date_p4w - timedelta(weeks=4)

                    latest_week_qtr_start_dt, latest_weekend_  = get_quarter_start_date(latest_weekend)
                    start_date_ytd = pd.Timestamp(latest_week_qtr_start_dt).tz_localize(None)
                    end_date_ytd = pd.Timestamp(latest_weekend_).tz_localize(None)

                    r4w_data = filtered_data[(filtered_data['week_end_date'] >= start_date_r4w) & (filtered_data['week_end_date'] <= specified_date)]
                    r4w_total_qty = r4w_data['qty_sold_pu'].sum()
                    r4w_weekly_avg = r4w_total_qty / 4

                    p4w_data = filtered_data[(filtered_data['week_end_date'] >= start_date_p4w) & (filtered_data['week_end_date'] <= end_date_p4w)]
                    p4w_total_qty = p4w_data['qty_sold_pu'].sum()
                    p4w_weekly_avg = p4w_total_qty / 4
                    print("p4w_total_qty!!!", p4w_total_qty)

                    r4w_growth = calculate_growth(r4w_weekly_avg, p4w_weekly_avg)



                    ytd_data = filtered_data[(filtered_data['week_end_date'] >= start_date_ytd) & (filtered_data['week_end_date'] <= end_date_ytd)
            #                                    &(filtered_data['parent_name'] != "TEXAS ONCOLOGY")
                                            ]
                    ytd_ex_txs = filtered_data[(filtered_data['week_end_date'] >= start_date_ytd) & (filtered_data['week_end_date'] <= end_date_ytd)
                                            &(filtered_data['parent_name'] != "TEXAS ONCOLOGY")
                                            ]

                    ytd_total_qty = ytd_data['qty_sold_pu'].sum()
                    ytd_total_ex_txs = ytd_ex_txs['qty_sold_pu'].sum()

                    num_weeks_ytd = (end_date_ytd - start_date_ytd).days // 7 + 1
                    ytd_weekly_avg = ytd_total_qty / num_weeks_ytd
                    ytd_avg_ex_txs = ytd_total_ex_txs / num_weeks_ytd
                    return [r4w_weekly_avg, p4w_weekly_avg,ytd_avg_ex_txs,r4w_growth]

                account_types = ['high','low', 'medium', 'academic', 'community','east','west','None']
                final_output_dict = {}
                final_growth_dict = {}
                for type_ in account_types:
                    print("type_!!!!!!", type_)
                    averages_dict = {}
                    growth_dict = {}
                    if type_ in ['academic', 'community']:
                        filtered_data = tro_sales_data[tro_sales_data['child_account_type'].str.lower() == type_].copy()
                        filtered_data['week_end_date'] = filtered_data['week_end_date'].dt.tz_localize(None)

                    elif type_ in ['east','west']:
                        filtered_data = tro_sales_data[tro_sales_data['area_nm'].str.lower() == type_].copy()
                        filtered_data['week_end_date'] = pd.to_datetime(filtered_data['week_end_date']).dt.tz_localize(None)

                    elif type_ in ["low","high","medium"]:
                        filtered_data = tro_sales_data[tro_sales_data['linkedid_bc_segment'].str.lower() == type_].copy()
                    else:
                        filtered_data=tro_sales_data
                    r4w, p4w, ytd_ex_txs, r4w_growth = get_averages(filtered_data)
                    averages_dict['P4W'] = p4w
                    averages_dict['R4W'] = r4w
                    growth_dict['R4W_growth'] = r4w_growth

                    # averages_dict["YTD-TX"] = ytd_ex_txs
                    df = pd.DataFrame(list(averages_dict.items()), columns=['time_period', type_])
                    final_output_dict[type_] = df

                    growth_df = pd.DataFrame(list(growth_dict.items()), columns=['time_period', "growth"])
                    final_growth_dict[type_] = growth_df
                return final_output_dict, final_growth_dict
        slide_12_data,  slide_12_growth_data= get_slide_12_data()
        def get_trend_r13w():
            # Get the latest available week_end_date
            tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date']).dt.strftime('%Y-%m-%dT%H:%M:%S.%fZ')
            tro_sales_data['week_end_date'] = pd.to_datetime(tro_sales_data['week_end_date'])
            sorted_dates = tro_sales_data['week_end_date'].sort_values(ascending=False).unique()
            latest_weekend = sorted_dates[0].date().strftime('%Y-%m-%d')
            print(f"Latest week_end_date in the data: {latest_weekend}")

            # Calculate the start date for the most recent 13 weeks and previous 13 weeks
            latest_weekend = pd.to_datetime(latest_weekend)
            start_date = latest_weekend - timedelta(weeks=25)
            print(f"Recent 26-week period: {start_date} to {latest_weekend}")

            # Ensure the 'week_end_date' column has timezone removed
            tro_sales_data['week_end_date'] = tro_sales_data['week_end_date'].dt.tz_localize(None)

            # Filter sales data for the last 26 weeks dynamically
            tro_weekly_sales = (tro_sales_data[
                (tro_sales_data['week_end_date'] >= start_date) &
                (tro_sales_data['week_end_date'] <= latest_weekend)
            ]
            .groupby('week_end_date', as_index=False)['qty_sold_pu']
            .sum())
            tro_weekly_sales['week_end_date'] = tro_weekly_sales['week_end_date'].dt.strftime('%d-%b')
            return tro_weekly_sales

        r13w_trend=get_trend_r13w()

        acad_df = slide_2_growth_data["academic"]
        comm_df = slide_2_growth_data["community"]
        high_df = slide_2_growth_data["high"]
        medium_df = slide_2_growth_data["medium"]
        low_df = slide_2_growth_data["low"]
        east_df=slide_2_growth_data["east"]
        west_df=slide_2_growth_data["west"]
        overall_df=slide_2_growth_data["None"]


        acad_df1 = slide_12_growth_data["academic"]
        comm_df1 = slide_12_growth_data["community"]
        high_df1 = slide_12_growth_data["high"]
        medium_df1 = slide_12_growth_data["medium"]
        low_df1 = slide_12_growth_data["low"]
        area_df1=slide_12_growth_data["east"]
        area_df2=slide_12_growth_data["west"]
        overall_df1=slide_12_growth_data["None"]
        A13_g= update_placeholders(acad_df,'R13W_growth','growth')
        C13_g= update_placeholders(comm_df,'R13W_growth','growth')
        L13_g= update_placeholders(low_df,'R13W_growth','growth')
        M13_g= update_placeholders(medium_df,'R13W_growth','growth')
        H13_g= update_placeholders(high_df,'R13W_growth','growth')
        E13_g= update_placeholders(east_df,'R13W_growth','growth')
        W13_g= update_placeholders(west_df,'R13W_growth','growth')
        overall13_g=update_placeholders(overall_df,'R13W_growth','growth')
        #overall13_g = update_placeholders(overall_df, shape, 'R4W_growth', 'Value')

        A4_g= update_placeholders(acad_df1,'R4W_growth','growth')
        C4_g= update_placeholders(comm_df1,'R4W_growth','growth')
        L4_g= update_placeholders(low_df1,'R4W_growth','growth')
        M4_g= update_placeholders(medium_df1,'R4W_growth','growth')
        H4_g= update_placeholders(high_df1,'R4W_growth','growth')
        E4_g= update_placeholders(area_df1,'R4W_growth','growth')
        W4_g= update_placeholders(area_df2,'R4W_growth','growth')
        overall4_g=update_placeholders(overall_df1,'R4W_growth','growth')
        def get_text_r13():

            medium_segment_text_13 = f"Medium segment {M13_st} by {M13_g}%" if M13_st != 'no growth' else "Medium segment remained same"
            high_segment_text_13=f"High segment {H13_st} by {H13_g}%" if H13_st != 'no growth' else "High segment remained same"
            low_segment_text_13=f"Low segment {L13_st} by {L13_g}%" if L13_st != 'no growth' else "Low segment remained same"
            accounts_text_13=f"Academic {A13_st} by {A13_g}%" if A13_st != 'no growth' else "Academic remained same"
            community_text_13=f"Community {C13_st} by {C13_g}%" if C13_st != 'no growth' else "Community remained same"
            east_text_13=f"East {E13_st} by {E13_g}%" if E13_st != 'no growth' else "East remained same"
            west_text_13=f"West {W13_st} by {W13_g}%" if W13_st != 'no growth' else "West remained same"
            overall_text_13=f"Sales {overall13_st} by {overall13_g}%" if overall13_st != 'no growth' else "Sales remained same"


            # Now, use the variable 'medium_segment_text' in the formatted_text_r13w list
            formatted_text_r13w = [
                (f"R13W VS P13W {latest_day_mon}", True, 18),  # Bold title with font size 18
                (f"{overall_text_13} over the past 13 weeks (R13W) compared to the previous 13 weeks (P13W).\n", False, 14),
                ("Accounts: ", True, 14),
                (f"{accounts_text_13}, {community_text_13}", False, 14),
                ("Segments: ", True, 14),
                (f"{low_segment_text_13},{medium_segment_text_13} and {high_segment_text_13}", False, 14),
                ("Regions: ", True, 14),
                (f"{east_text_13}, {west_text_13}", False, 14)
            ]

            return formatted_text_r13w
        def get_text_r4():
            medium_segment_text_4 = f"Medium segment {M4_st} by {M4_g}%" if M4_st != 'no growth' else "Medium segment remained same"
            high_segment_text_4=f"High segment {H4_st} by {H4_g}%" if H4_st != 'no growth' else "High segment remained same"
            low_segment_text_4=f"Low segment {L4_st} by {L4_g}%" if L4_st != 'no growth' else "Low segment remained same"
            accounts_text_4=f"Academic {A4_st} by {A4_g}%" if A4_st != 'no growth' else "Academic remained same"
            community_text_4=f"Community {C4_st} by {C4_g}%" if C4_st != 'no growth' else "Community remained same"
            east_text_4=f"East {E4_st} by {E4_g}%" if E4_st != 'no growth' else "East remained same"
            west_text_4=f"West {W4_st} by {W4_g}%" if W4_st != 'no growth' else "West remained same"
            overall_text_4=f"Sales {overall4_st} by {overall4_g}%" if overall4_st != 'no growth' else "Sales remained same"

            formatted_text_r4w = [
                (f"R4W VS P4W {latest_day_mon}", True, 18),  # Bold title with font size 18
                (f"{overall_text_4} over the past 4 weeks (R4W) compared to the previous 4 weeks (P4W).\n", False, 14),
                ("Accounts: ", True, 14),
                (f"{accounts_text_4}, {community_text_4}", False, 14),
                ("Segments: ", True, 14),
                (f"{low_segment_text_4}, {medium_segment_text_4} and {high_segment_text_4}", False, 14),
                ("Regions: ", True, 14),
                (f"{east_text_4}, {west_text_4}", False, 14)

            ]
            return formatted_text_r4w

        # Function to update text with font size and bold styling
        def update_textbox(shape, formatted_text):
            text_frame = shape.text_frame
            text_frame.clear()
            for text, is_bold, font_size in formatted_text:
                para = text_frame.add_paragraph()
                run = para.add_run()
                run.text = text
                run.font.bold = is_bold
                run.font.size = Pt(font_size)  # Set font size
                para.space_after = Pt(0)


        from pptx import Presentation
        from pptx.chart.data import ChartData, CategoryChartData
        import os
        # script_dir = os.path.dirname(os.path.abspath(__file__))
        # ppt_path = os.path.join(script_dir, 'ppt_template_867_sales.pptx')
        # presentation = Presentation(ppt_path)

        ppt_path ='./Structured_Bot/ppt_template_867_sales.pptx'
        presentation = Presentation(ppt_path)
        for slide in presentation.slides:
            for shape in slide.shapes:
                final_df = []
                if shape.has_chart:
                    chart = shape.chart
                    chart_title = chart.chart_title.text_frame.text
                    if "Trodelvy Weekly Demand" in chart_title:
                        final_df = slide_1_data
                    elif "Academic account" in chart_title:
                        print("academic for slide 2")
                        res = slide_2_data["academic"]
                        if "TX" in chart_title:
                            final_df = res[res['time_period'] != 'YTD']
                            update_cell_text(chart.chart_title, f"Academic({latest_day_mon})-TX", "chart_title")
                        else:
                            final_df = res[res['time_period'] != 'YTD-TX']
                            update_cell_text(chart.chart_title, f"Academic({latest_day_mon})", "chart_title")
                    elif "High Potential" in chart_title:
                        res = slide_2_data["high"]
                        if "TX" in chart_title:
                            final_df = res[res['time_period'] != 'YTD']
                            update_cell_text(chart.chart_title, f"High Potential for before ({latest_day_mon})-TX", "chart_title")
                        else:
                            final_df = res[res['time_period'] != 'YTD-TX']
                            update_cell_text(chart.chart_title, f"High Potential({latest_day_mon})", "chart_title")
                    elif "Medium Potential" in chart_title:
                        res = slide_2_data["medium"]
                        if "TX" in chart_title:
                            final_df = res[res['time_period'] != 'YTD']
                            update_cell_text(chart.chart_title, f"Medium Potential({latest_day_mon})-TX", "chart_title")
                        else:
                            final_df = res[res['time_period'] != 'YTD-TX']
                            update_cell_text(chart.chart_title, f"Medium Potential({latest_day_mon})", "chart_title")

                    elif "Low Potential" in chart_title:
                        res = slide_2_data["low"]
                        if "TX" in chart_title:
                            final_df = res[res['time_period'] != 'YTD']
                            update_cell_text(chart.chart_title, f"Low Potential({latest_day_mon})-TX", "chart_title")
                        else:
                            final_df = res[res['time_period'] != 'YTD-TX']
                            update_cell_text(chart.chart_title, f"Low Potential({latest_day_mon})", "chart_title")
                    elif "Community account" in chart_title:
                        res = slide_2_data["community"]
                        if "TX" in chart_title:
                            final_df = res[res['time_period'] != 'YTD']
                            update_cell_text(chart.chart_title, f"Community({latest_day_mon})-TX", "chart_title")
                        else:
                            final_df = res[res['time_period'] != 'YTD-TX']
                            update_cell_text(chart.chart_title, f"Community({latest_day_mon})", "chart_title")
                    elif "R13W Trends" in chart_title:
                        final_df= r13w_trend
                        update_cell_text(chart.chart_title, f"Trodelvy Weekly Sales Trends({latest_day_mon})", "chart_title")
                    elif "Region Vials Growth" in chart_title and "East" in chart_title:
                        final_df = region_vials_growth_east
                        update_cell_text(chart.chart_title, f"Region Vials Growth % (East)w.e({latest_day_mon})", "chart_title")
                    elif "Region Vials Growth" in chart_title and "West" in chart_title:
                        final_df = region_vials_growth_west
                        update_cell_text(chart.chart_title, f"Region Vials Growth % (West)w.e({latest_day_mon})", "chart_title")
                    elif "Weekly Region Vials" in chart_title and "East" in chart_title:
                        final_df = east_regional_avg
                    elif "Weekly Region Vials" in chart_title and "West" in chart_title:
                        final_df = west_regional_avg
                    elif "Area Vials Growth" in chart_title:
                        final_df = nation_vials_growth
                        update_cell_text(chart.chart_title, f"Area Vials Growth % w.e({latest_day_mon})", "chart_title")
                    elif "Area Vials Sales" in chart_title:
                        final_df = nation_area_avg
                    elif chart_title == "segment contribution":
                        final_df = tro_sub_account_contr_data
                    elif chart_title == "Acc contribution ":
                        final_df = tro_account_contr_data
                    elif chart_title == "Kad contribution":
                        final_df = tro_sub_kad_contr_data
                    try:
                        final_df.replace('YTD-TX', 'YTD', inplace=True)
                    except:
                        pass
                    if len(final_df) >0:
                        final_df.replace([np.nan, np.inf, -np.inf], '', inplace=True)
                        print("Chart found!", chart.chart_title.text_frame.text)
                        chart_data = CategoryChartData()
                        categories = final_df.iloc[:, 0].tolist()  # Assuming the first column is the category
                        chart_data.categories = categories
                        for col in final_df.columns[1:]:
                            chart_data.add_series(col, final_df[col].tolist())
                        chart.replace_data(chart_data)
                        print("replaced data!!")
                else:
                    pass
                if shape.has_text_frame:
                    text = shape.text_frame.text
                    if "R13W vs P13W" in text:
                        update_textbox(shape, get_text_r13())
                    elif "R4W vs P4W" in text:
                        update_textbox(shape, get_text_r4())
                    for paragraph in shape.text_frame.paragraphs:
                        original_text = paragraph.text
                        print("Original Paragraph Text:", original_text)
                        if original_text == "aca PG":
                            update_text(acad_df, paragraph, 'P13W_growth', 'growth')
                        elif original_text == "aca RG":
                            update_text(acad_df, paragraph, 'R13W_growth', 'growth')
                        elif original_text == "com PG":
                            update_text(comm_df, paragraph, 'P13W_growth', 'growth')
                        elif original_text == "com RG":
                            update_text(comm_df, paragraph, 'R13W_growth', 'growth')
                        elif original_text == "high PG":
                            update_text(high_df, paragraph, 'P13W_growth', 'growth')
                        elif original_text == "high RG":
                            update_text(high_df, paragraph, 'R13W_growth', 'growth')
                        elif original_text == "med PG":
                            update_text(medium_df, paragraph, 'P13W_growth', 'growth')
                        elif original_text == "med RG":
                            update_text(medium_df, paragraph, 'R13W_growth', 'growth')
                        elif original_text == "low PG":
                            update_text(low_df, paragraph, 'P13W_growth', 'growth')
                        elif original_text == "low RG":
                            update_text(low_df, paragraph, 'R13W_growth', 'growth')
        presentation.save("867_sales_refreshed_ppt.pptx")

        # In[ ]:

        from pptx import Presentation
        def print_table_content(pptx_file, slide_1_table_data, slide_3_data , high_acc_decliners, med_acc_decliners, academic_acc_decliners, community_acc_decliners ,idn_hospital_decliners, east_area_sales , west_area_sales, nation_area_sales):
            df_table = False
            start_row = 2
            prs = Presentation(pptx_file)
            for slide_number, slide in enumerate(prs.slides):
                print(f"Slide {slide_number + 1}")
                for shape in slide.shapes:
                    if not shape.has_table and not shape.has_chart:
                        continue
                    if shape.has_chart:
                        chart = shape.chart
                        print(chart.chart_title.text_frame.text)
                        if "Trodelvy Weekly Demand" in chart.chart_title.text_frame.text:
                            update_cell_text(chart.chart_title, f"Trodelvy Weekly Demand (k Vials) (WE {latest_day_mon})", "chart_title")

                    if shape.has_table:
                        table = shape.table
                        cell = table.cell(0, 0).text
                        print("cell!!!", cell)
                        flag_value = table.cell(0, 1).text

                        if flag_value == '':
                            flag_value = table.cell(0, 2).text
                        print("flag_value!!",flag_value)
                        if "% Contribution to Nation" in flag_value:
                            df = slide_3_data
                            df_table = True
                            start_row = 2
                            indices = [(0,2), (0,4),(0,5),(0,6),(0,7)]
                            keys = [f'R12M \nWeekly Sales \n(Ending {latest_day_mon})', f'R13W weekly avg\n(w.e. {latest_day_mon})', f'R13W Growth: (R13W w.e. {latest_day_mon} vs P13W)', f'YTD Sales Growth'  ]
                            for idx, value in zip(indices, keys):
                                update_cell_text(shape.table.cell(idx[0], idx[1]), value, "value")
                        elif "East" in flag_value or "West" in flag_value or "Nation" in flag_value:
                            print("in slide 6 & 10!!")
                            if "East" in flag_value:
                                df = east_area_sales
                            elif "West" in flag_value:
                                df = west_area_sales
                            elif "Nation" in flag_value:
                                df = nation_area_sales
                            start_row = 3
                            df_table = True
                        elif "high" in cell.lower():
                            df = high_acc_decliners
                            df_table = True
                            start_row = 2
                        elif "medium" in cell.lower():
                            df = med_acc_decliners
                            df_table = True
                            start_row = 2
                        elif "academic" in cell.lower():
                            df = academic_acc_decliners
                            df_table = True
                            start_row = 2
                        elif "community" in cell.lower():
                            df = community_acc_decliners
                            df_table = True
                            start_row = 2
                        elif "idn" in cell.lower():
                            print("in idn")
                            df = idn_hospital_decliners
                            df_table = True
                            start_row = 2
                        if "high" in cell.lower() or "medium" in cell.lower() or "academic" in cell.lower() or "community" in cell.lower() or "idn" in cell.lower():
                            indices = [(0,1), (0,2)]
                            keys = [f'R13W weekly avg (w.e. {latest_day_mon})', f'P13W weekly avg (w.e. {latest_day_mon})']
                            for idx, value in zip(indices, keys):
                                update_cell_text(shape.table.cell(idx[0], idx[1]), value, "value")

                        elif "Weekly Avg Demand Vials" in cell:
                            indices = [(0,0), (1, 1), (1, 2), (1, 3), (1, 4), (2, 1)]
                            keys = ['avg_demand', 'R4W', 'R8W', 'R13W', 'R52W', 'growth_4w', 'growth_8w', 'growth_13w', 'growth_qtr']
                            for idx, key in zip(indices, keys):
                                value = slide_1_table_data.get(key, '')
                                if key == "growth_4w":
                                    value = str(value) +"%"
                                update_cell_text(shape.table.cell(idx[0], idx[1]), value, "value")
                        elif "Total Demand Vials" in cell:
                            indices = [(0,0),(1, 1), (1, 3), (2, 1), (2, 3), (3,1), (3,3), (1,2),(2, 2), (3, 2), (2, 2), (2, 3), (2, 4)]
                            keys = ['total_demand', 'current_qtr_vials', 'curr_qtr_growth', 'ytd_vials', 'ytd_growth', 'prev_qtr_vials', 'prev_qtr_growth', 'current_qtr_est_avg', 'ytd_est_avg', 'prev_qtr_est_avg']
                            for idx, key in zip(indices, keys):
                                value = slide_1_table_data.get(key, '')
                                # if type(value) != str:
                                #     value = round(value, 1)
                                update_cell_text(shape.table.cell(idx[0], idx[1]), value , "value")

                        elif  "Quarterly Avg. Demand Vials" in cell:
                            indices = [(0,0), (1, 1), (1, 2)]
                            keys = ['qtr_avg_demand' ,'q1', 'q2']
                            for idx, key in zip(indices, keys):
                                value = slide_1_table_data.get(key, '')
                                # if type(value) != str:
                                #     value = round(value, 1)
                                update_cell_text(shape.table.cell(idx[0], idx[1]), value , "value")

                        elif "BC Segment type" in flag_value:
                            df = final_df_account
                            df_table = True
                            start_row = 1
                        elif "KAD Targets" in flag_value:
                            df_table = True
                            df = final_df_kad

                        if df_table:
                            for i, row in enumerate(table.rows):
                                for j, cell in enumerate(row.cells):
                                    if i >= start_row:
                                        df_row_index = i - start_row
                                        if df_row_index < len(df) and j < len(df.columns):
                                            value = str(df.iloc[df_row_index, j])
                                            update_cell_text(cell, value)
                                        else:
                                            update_cell_text(cell, "")

                    print("-" * 20)
            print("saving file!!!")
            prs.save("867_sales_refreshed_ppt.pptx")
            print("saved file!!!")

        print_table_content('867_sales_refreshed_ppt.pptx', slide_1_table_data, slide_3_table_data, high_acc_decliners, med_acc_decliners, academic_acc_decliners, community_acc_decliners, idn_hospital_decliners, east_area_sales, west_area_sales, nation_area_sales)

